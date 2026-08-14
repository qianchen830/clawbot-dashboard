# -*- coding: utf-8 -*-
"""
金蝶品牌色配置
所有PPT生成器统一使用此配置
"""

from pptx.dml.color import RGBColor

# ── 金蝶品牌色方案 ──
KINGDEE_COLORS = {
    # 主色调
    'primary': RGBColor(0x1F, 0x4E, 0x79),       # 深蓝 #1F4E79
    'kingdee_blue': RGBColor(0x00, 0x6C, 0xB8),  # 金蝶蓝 #006CB8
    'secondary': RGBColor(0x2E, 0x75, 0xB6),     # 辅色蓝 #2E75B6
    'light_blue': RGBColor(0x00, 0x8E, 0xD4),    # 亮蓝 #008ED4
    
    # 强调色
    'accent': RGBColor(0xED, 0x7D, 0x31),        # 强调橙 #ED7D31
    'warning': RGBColor(0xFF, 0xC0, 0x00),       # 警告金 #FFC000
    'success': RGBColor(0x70, 0xAD, 0x47),       # 成功绿 #70AD47
    'error': RGBColor(0xC0, 0x00, 0x00),         # 错误红 #C00000
    
    # 文字色
    'text_primary': RGBColor(0x33, 0x33, 0x33),  # 正文深灰 #333333
    'text_secondary': RGBColor(0x66, 0x66, 0x66),# 二级灰 #666666
    'text_light': RGBColor(0x99, 0x99, 0x99),    # 三级灰 #999999
    
    # 背景色
    'white': RGBColor(0xFF, 0xFF, 0xFF),         # 白色
    'bg_light': RGBColor(0xF2, 0xF6, 0xFA),      # 浅蓝背景 #F2F6FA
    'bg_card': RGBColor(0xFA, 0xFB, 0xFC),       # 卡片背景 #FAFBFC
    
    # 边框色
    'border': RGBColor(0xD0, 0xD8, 0xE0),        # 边框灰 #D0D8E0
    
    # 渐变色
    'gradient_start': RGBColor(0x1A, 0x23, 0x7E),# 渐变起点 #1A237E
    'gradient_end': RGBColor(0x2E, 0x75, 0xB6),  # 渐变终点 #2E75B6
}

# 简写别名
C = KINGDEE_COLORS

# 字体配置
FONT_NAME = '微软雅黑'
FONT_SIZE_TITLE = Pt(28)
FONT_SIZE_SUBTITLE = Pt(20)
FONT_SIZE_BODY = Pt(14)
FONT_SIZE_SMALL = Pt(12)

# 页面尺寸（16:9）
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def get_layout(prs, layout_name):
    """根据名称获取布局"""
    layout_map = {
        '封面': 0,
        '目录': 1,
        '内页': 2,
        '封底': 3,
        '空白': 5,
    }
    
    # 尝试按名称匹配
    for i, layout in enumerate(prs.slide_layouts):
        if layout_name in layout.name:
            return layout
    
    # 默认返回内页布局
    if len(prs.slide_layouts) > 2:
        return prs.slide_layouts[2]
    return prs.slide_layouts[0]

def add_title_bar(slide, title_text, colors=None):
    """添加统一的标题栏"""
    if colors is None:
        colors = C
    
    # 标题栏背景
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(0.75)
    
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt
    
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = colors['primary']
    title_bg.line.fill.background()
    
    # 标题文本
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.name = '微软雅黑'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    return title_bg

def add_footer(slide, colors=None):
    """添加统一的页脚"""
    if colors is None:
        colors = C
    
    from pptx.util import Pt
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.0), Inches(3), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = colors['kingdee_blue']
    line.line.fill.background()
    
    # 公司信息
    footer_box = slide.shapes.add_textbox(Inches(11), Inches(7.0), Inches(2), Inches(0.3))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "金蝶软件（中国）有限公司"
    p.font.name = '微软雅黑'
    p.font.size = Pt(9)
    p.font.color.rgb = colors['text_light']
    p.alignment = PP_ALIGN.RIGHT
    
    return line, footer_box

def add_cover_slide(prs, title, subtitle="", company="", date="", colors=None):
    """添加封面页"""
    if colors is None:
        colors = C
    
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    
    layout = get_layout(prs, '封面')
    slide = prs.slides.add_slide(layout)
    
    # 渐变背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['gradient_start']
    bg.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = '微软雅黑'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = '微软雅黑'
        p.font.size = Pt(20)
        p.font.color.rgb = colors['white']
    
    # 公司信息
    if company:
        comp_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.25), Inches(6), Inches(0.5))
        tf = comp_box.text_frame
        p = tf.paragraphs[0]
        p.text = company
        p.font.name = '微软雅黑'
        p.font.size = Pt(14)
        p.font.color.rgb = colors['white']
    
    # 日期
    if date:
        date_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.75), Inches(6), Inches(0.4))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date
        p.font.name = '微软雅黑'
        p.font.size = Pt(12)
        p.font.color.rgb = colors['white']
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9), Inches(7.0), Inches(3.4), Inches(0.17)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = colors['accent']
    line.line.fill.background()
    
    return slide
