#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT高级模板库 - v2.0
提供更专业的PPT幻灯片模板，包括KPI卡片、时间线、表格等
新增：渐变背景、进度条、图标装饰
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 扩展配色方案
COLORS_EXTENDED = {
    'primary': RGBColor(0, 102, 153),
    'primary_dark': RGBColor(0, 76, 115),
    'primary_light': RGBColor(0, 153, 204),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'accent_green': RGBColor(46, 139, 87),
    'accent_red': RGBColor(204, 51, 51),
    'accent_blue': RGBColor(51, 102, 204),
    'text': RGBColor(51, 51, 51),
    'text_light': RGBColor(102, 102, 102),
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(240, 240, 240),
    'light_blue': RGBColor(230, 243, 255),
    'light_green': RGBColor(230, 255, 240),
    'light_orange': RGBColor(255, 243, 224),
    'light_red': RGBColor(255, 230, 230),
    'card_bg': RGBColor(248, 250, 252),
    'border': RGBColor(220, 220, 220),
    'gradient_start': RGBColor(0, 102, 153),
    'gradient_end': RGBColor(0, 153, 204),
    'gold': RGBColor(212, 175, 55),
    'silver': RGBColor(169, 169, 169),
    'bronze': RGBColor(205, 127, 50),
}


def add_kpi_card_slide(prs, title, kpi_data):
    """
    添加KPI指标卡片页
    kpi_data: [{"label": "指标名", "value": "数值", "change": "+30%", "color": "green/blue/red"}]
    最多6个KPI卡片，分两行排列
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    # KPI卡片布局
    card_count = min(len(kpi_data), 6)
    cols = 3 if card_count > 3 else card_count
    rows = 2 if card_count > 3 else 1
    
    card_width = Inches(3.8)
    card_height = Inches(2.5)
    h_gap = Inches(0.5)
    v_gap = Inches(0.4)
    start_x = (Inches(13.333) - (cols * card_width + (cols - 1) * h_gap)) / 2
    start_y = Inches(1.6)
    
    for i, kpi in enumerate(kpi_data[:card_count]):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_width + h_gap)
        y = start_y + row * (card_height + v_gap)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS_EXTENDED['card_bg']
        card.line.color.rgb = COLORS_EXTENDED['border']
        card.line.width = Pt(1)
        
        # 左侧色条
        color_key = kpi.get('color', 'blue')
        color_map = {
            'green': COLORS_EXTENDED['accent_green'],
            'blue': COLORS_EXTENDED['accent_blue'],
            'red': COLORS_EXTENDED['accent_red'],
            'orange': COLORS_EXTENDED['accent'],
            'primary': COLORS_EXTENDED['primary']
        }
        bar_color = color_map.get(color_key, COLORS_EXTENDED['primary'])
        
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()
        
        # 指标标签
        label_tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), card_width - Inches(0.5), Inches(0.5))
        p = label_tb.text_frame.paragraphs[0]
        p.text = kpi.get('label', '')
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS_EXTENDED['text_light']
        
        # 指标数值
        value_tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), card_width - Inches(0.5), Inches(0.8))
        p = value_tb.text_frame.paragraphs[0]
        p.text = kpi.get('value', '')
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['text']
        
        # 变化值
        if kpi.get('change'):
            change_tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.8), card_width - Inches(0.5), Inches(0.4))
            p = change_tb.text_frame.paragraphs[0]
            p.text = kpi.get('change', '')
            p.font.size = Pt(14)
            p.font.bold = True
            is_positive = '+' in kpi['change'] or '↑' in kpi['change'] or '↓70' in kpi['change'] or '缩短' in kpi['change'] or '降低' in kpi['change']
            p.font.color.rgb = COLORS_EXTENDED['accent_green'] if is_positive else COLORS_EXTENDED['accent_red']
    
    return slide


def add_table_slide(prs, title, headers, rows, highlight_cols=None):
    """
    添加表格页
    headers: ["列1", "列2", "列3"]
    rows: [["值1", "值2", "值3"], ...]
    highlight_cols: [0, 2]  # 需要高亮的列索引
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    # 创建表格
    col_count = len(headers)
    row_count = len(rows) + 1  # +1 for header
    table_width = Inches(12)
    table_height = Inches(min(5.5, 0.6 + 0.5 * row_count))
    table_left = Inches(0.667)
    table_top = Inches(1.5)
    
    table = slide.shapes.add_table(row_count, col_count, table_left, table_top, table_width, table_height).table
    
    # 设置列宽
    col_width = table_width // col_count
    for i in range(col_count):
        table.columns[i].width = col_width
    
    # 填充表头
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS_EXTENDED['primary']
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS_EXTENDED['white']
            paragraph.alignment = PP_ALIGN.CENTER
    
    # 填充数据
    for i, row_data in enumerate(rows):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_value)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS_EXTENDED['card_bg']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS_EXTENDED['white']
            
            if highlight_cols and j in highlight_cols:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS_EXTENDED['primary']
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = COLORS_EXTENDED['text']
                    paragraph.alignment = PP_ALIGN.CENTER
    
    return slide


def add_timeline_slide(prs, title, phases):
    """
    添加时间线/阶段页
    phases: [{"name": "阶段1", "duration": "2周", "items": ["任务1", "任务2"]}]
    最多5个阶段
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    # 时间线横线
    line_y = Inches(2.2)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), line_y, Inches(12.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    line.line.fill.background()
    
    # 阶段节点
    phase_count = min(len(phases), 5)
    phase_width = Inches(12.333) / phase_count
    
    # 阶段颜色
    phase_colors = [
        COLORS_EXTENDED['primary'],
        COLORS_EXTENDED['primary_light'],
        COLORS_EXTENDED['accent'],
        COLORS_EXTENDED['accent_green'],
        COLORS_EXTENDED['accent_blue']
    ]
    
    for i, phase in enumerate(phases[:phase_count]):
        cx = Inches(0.5) + phase_width * (i + 0.5)
        
        # 圆形节点
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.2), line_y - Inches(0.2), Inches(0.4), Inches(0.4))
        node.fill.solid()
        node.fill.fore_color.rgb = phase_colors[i % len(phase_colors)]
        node.line.fill.background()
        
        # 阶段名称
        name_tb = slide.shapes.add_textbox(cx - Inches(1.5), line_y + Inches(0.4), Inches(3), Inches(0.5))
        p = name_tb.text_frame.paragraphs[0]
        p.text = phase.get('name', f'阶段{i+1}')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 时长
        if phase.get('duration'):
            dur_tb = slide.shapes.add_textbox(cx - Inches(1.5), line_y + Inches(0.8), Inches(3), Inches(0.3))
            p = dur_tb.text_frame.paragraphs[0]
            p.text = phase.get('duration')
            p.font.size = Pt(12)
            p.font.color.rgb = phase_colors[i % len(phase_colors)]
            p.alignment = PP_ALIGN.CENTER
        
        # 任务列表
        items = phase.get('items', [])
        items_tb = slide.shapes.add_textbox(cx - Inches(1.5), line_y + Inches(1.2), Inches(3), Inches(3.5))
        tf = items_tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS_EXTENDED['text_light']
            p.space_before = Pt(4)
    
    return slide


def add_three_column_slide(prs, title, col1, col2, col3):
    """
    添加三栏内容页
    col: {"title": "标题", "items": ["项1", "项2"], "icon": "🏭"}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    cols = [col1, col2, col3]
    col_width = Inches(3.8)
    h_gap = Inches(0.5)
    start_x = (Inches(13.333) - (3 * col_width + 2 * h_gap)) / 2
    
    for i, col in enumerate(cols):
        x = start_x + i * (col_width + h_gap)
        y = Inches(1.5)
        
        # 列卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS_EXTENDED['card_bg']
        card.line.color.rgb = COLORS_EXTENDED['border']
        card.line.width = Pt(1)
        
        # 列标题
        col_title = col.get('title', '')
        col_icon = col.get('icon', '')
        
        title_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), col_width - Inches(0.4), Inches(0.6))
        p = title_tb.text_frame.paragraphs[0]
        p.text = f"{col_icon} {col_title}" if col_icon else col_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['primary']
        
        # 分隔线
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.8), col_width - Inches(0.4), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = COLORS_EXTENDED['primary_light']
        sep.line.fill.background()
        
        # 内容
        items = col.get('items', [])
        items_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0), col_width - Inches(0.4), Inches(4.2))
        tf = items_tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            if item.startswith('###'):
                p.text = item.replace('###', '').strip()
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = COLORS_EXTENDED['primary']
                p.space_before = Pt(10)
            elif item.startswith('-'):
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(13)
                p.font.color.rgb = COLORS_EXTENDED['text']
                p.space_before = Pt(4)
            elif item.strip():
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = COLORS_EXTENDED['text']
                p.space_before = Pt(4)
    
    return slide


def add_comparison_slide(prs, title, before_title, before_items, after_title, after_items):
    """
    添加前后对比页（带对比箭头）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    # 左侧（Before）
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.3), Inches(5.5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS_EXTENDED['light_red']
    left_card.line.color.rgb = COLORS_EXTENDED['accent_red']
    left_card.line.width = Pt(2)
    
    left_title_tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.7), Inches(0.6))
    p = left_title_tb.text_frame.paragraphs[0]
    p.text = f"❌ {before_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['accent_red']
    
    left_content_tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.7), Inches(4.2))
    tf = left_content_tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(before_items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.space_before = Pt(8)
    
    # 中间箭头
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(3.8), Inches(1.1), Inches(0.8))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS_EXTENDED['accent']
    arrow.line.fill.background()
    
    # 右侧（After）
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(5.3), Inches(5.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLORS_EXTENDED['light_green']
    right_card.line.color.rgb = COLORS_EXTENDED['accent_green']
    right_card.line.width = Pt(2)
    
    right_title_tb = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(0.6))
    p = right_title_tb.text_frame.paragraphs[0]
    p.text = f"✅ {after_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['accent_green']
    
    right_content_tb = slide.shapes.add_textbox(Inches(7.8), Inches(2.5), Inches(4.7), Inches(4.2))
    tf = right_content_tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(after_items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.space_before = Pt(8)
    
    return slide


def add_icon_card_slide(prs, title, cards):
    """
    添加图标卡片页（用于展示模块/功能）
    cards: [{"icon": "💰", "title": "标题", "desc": "描述"}, ...]
    最多6个卡片
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    card_count = min(len(cards), 6)
    cols = 3
    rows = (card_count + 2) // 3
    
    card_width = Inches(3.8)
    card_height = Inches(2.5)
    h_gap = Inches(0.5)
    v_gap = Inches(0.4)
    start_x = (Inches(13.333) - (cols * card_width + (cols - 1) * h_gap)) / 2
    start_y = Inches(1.6)
    
    for i, card in enumerate(cards[:card_count]):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_width + h_gap)
        y = start_y + row * (card_height + v_gap)
        
        # 卡片背景
        card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLORS_EXTENDED['white']
        card_shape.line.color.rgb = COLORS_EXTENDED['border']
        card_shape.line.width = Pt(1)
        
        # 顶部色条
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, Inches(0.05))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLORS_EXTENDED['primary']
        top_bar.line.fill.background()
        
        # 图标
        icon_tb = slide.shapes.add_textbox(x, y + Inches(0.2), card_width, Inches(0.6))
        p = icon_tb.text_frame.paragraphs[0]
        p.text = card.get('icon', '📌')
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), card_width - Inches(0.4), Inches(0.5))
        p = title_tb.text_frame.paragraphs[0]
        p.text = card.get('title', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.4), card_width - Inches(0.4), Inches(1.0))
        p = desc_tb.text_frame.paragraphs[0]
        p.text = card.get('desc', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS_EXTENDED['text_light']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_big_number_slide(prs, title, big_numbers):
    """
    添加大数字展示页（用于展示关键业绩指标）
    big_numbers: [{"value": "50%", "label": "效率提升", "desc": "财务结账时间缩短", "color": "green"}]
    最多4个大数字
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    # 大数字布局
    num_count = min(len(big_numbers), 4)
    num_width = Inches(2.8)
    h_gap = Inches(0.5)
    total_width = num_count * num_width + (num_count - 1) * h_gap
    start_x = (Inches(13.333) - total_width) / 2
    start_y = Inches(2.0)
    
    color_map = {
        'green': COLORS_EXTENDED['accent_green'],
        'blue': COLORS_EXTENDED['accent_blue'],
        'orange': COLORS_EXTENDED['accent'],
        'red': COLORS_EXTENDED['accent_red'],
        'primary': COLORS_EXTENDED['primary']
    }
    
    for i, num in enumerate(big_numbers[:num_count]):
        x = start_x + i * (num_width + h_gap)
        y = start_y
        
        # 数字背景圆
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), y, Inches(2), Inches(2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color_map.get(num.get('color', 'blue'), COLORS_EXTENDED['primary'])
        circle.line.fill.background()
        
        # 大数字
        value_tb = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(0.5), Inches(2), Inches(1.0))
        p = value_tb.text_frame.paragraphs[0]
        p.text = num.get('value', '0%')
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_tb = slide.shapes.add_textbox(x, y + Inches(2.3), num_width, Inches(0.5))
        p = label_tb.text_frame.paragraphs[0]
        p.text = num.get('label', '')
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_tb = slide.shapes.add_textbox(x, y + Inches(2.8), num_width, Inches(0.5))
        p = desc_tb.text_frame.paragraphs[0]
        p.text = num.get('desc', '')
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS_EXTENDED['text_light']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_process_flow_slide(prs, title, steps):
    """
    添加流程图页（用于展示业务流程/实施步骤）
    steps: [{"name": "步骤1", "desc": "描述"}]
    最多6个步骤
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    step_count = min(len(steps), 6)
    step_width = Inches(1.8)
    arrow_width = Inches(0.3)
    total_width = step_count * step_width + (step_count - 1) * arrow_width
    start_x = (Inches(13.333) - total_width) / 2
    start_y = Inches(2.5)
    
    colors = [
        COLORS_EXTENDED['primary'],
        COLORS_EXTENDED['primary_light'],
        COLORS_EXTENDED['accent'],
        COLORS_EXTENDED['accent_green'],
        COLORS_EXTENDED['accent_blue'],
        COLORS_EXTENDED['secondary']
    ]
    
    for i, step in enumerate(steps[:step_count]):
        x = start_x + i * (step_width + arrow_width)
        
        # 步骤圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), start_y, Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i % len(colors)]
        circle.line.fill.background()
        
        # 步骤数字
        num_tb = slide.shapes.add_textbox(x + Inches(0.4), start_y + Inches(0.2), Inches(1.0), Inches(0.6))
        p = num_tb.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 步骤名称
        name_tb = slide.shapes.add_textbox(x, start_y + Inches(1.3), step_width, Inches(0.6))
        p = name_tb.text_frame.paragraphs[0]
        p.text = step.get('name', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 步骤描述
        desc_tb = slide.shapes.add_textbox(x, start_y + Inches(1.9), step_width, Inches(1.5))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step.get('desc', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS_EXTENDED['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 箭头（最后一个不需要）
        if i < step_count - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_width + Inches(0.05), start_y + Inches(0.3), Inches(0.2), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS_EXTENDED['light_gray']
            arrow.line.fill.background()
    
    return slide


def add_quote_slide(prs, quote, author, title=None):
    """
    添加引用页（用于客户评价、领导寄语等）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS_EXTENDED['light_blue']
    bg.line.fill.background()
    
    # 左引号
    left_quote = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(1), Inches(1))
    p = left_quote.text_frame.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = COLORS_EXTENDED['primary_light']
    
    # 引用内容
    quote_tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2.5))
    tf = quote_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = COLORS_EXTENDED['text']
    p.alignment = PP_ALIGN.CENTER
    
    # 作者
    author_tb = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.333), Inches(0.5))
    p = author_tb.text_frame.paragraphs[0]
    p.text = f"— {author}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS_EXTENDED['primary']
    p.alignment = PP_ALIGN.RIGHT
    
    # 标题（如果有）
    if title:
        title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        p = title_tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['primary']
    
    return slide


def add_stats_grid_slide(prs, title, stats):
    """
    添加统计网格页（用于展示多个统计数据）
    stats: [{"value": "100+", "label": "客户案例", "icon": "🏢"}]
    最多12个统计项，排列成3x4网格
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS_EXTENDED['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS_EXTENDED['white']
    
    stat_count = min(len(stats), 12)
    cols = 4
    rows = (stat_count + cols - 1) // cols
    
    cell_width = Inches(2.8)
    cell_height = Inches(2.5)
    h_gap = Inches(0.3)
    v_gap = Inches(0.3)
    total_width = cols * cell_width + (cols - 1) * h_gap
    start_x = (Inches(13.333) - total_width) / 2
    start_y = Inches(1.5)
    
    for i, stat in enumerate(stats[:stat_count]):
        row = i // cols
        col = i % cols
        x = start_x + col * (cell_width + h_gap)
        y = start_y + row * (cell_height + v_gap)
        
        # 统计项背景
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_width, cell_height)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS_EXTENDED['white']
        cell.line.color.rgb = COLORS_EXTENDED['border']
        cell.line.width = Pt(1)
        
        # 图标
        icon_tb = slide.shapes.add_textbox(x, y + Inches(0.2), cell_width, Inches(0.5))
        p = icon_tb.text_frame.paragraphs[0]
        p.text = stat.get('icon', '📊')
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        
        # 数值
        value_tb = slide.shapes.add_textbox(x, y + Inches(0.8), cell_width, Inches(0.8))
        p = value_tb.text_frame.paragraphs[0]
        p.text = stat.get('value', '')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS_EXTENDED['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_tb = slide.shapes.add_textbox(x, y + Inches(1.6), cell_width, Inches(0.6))
        p = label_tb.text_frame.paragraphs[0]
        p.text = stat.get('label', '')
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS_EXTENDED['text_light']
        p.alignment = PP_ALIGN.CENTER
    
    return slide
