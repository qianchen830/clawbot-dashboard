# -*- coding: utf-8 -*-
"""
金蝶PPT样式组件库 V1.0
提供统一的视觉元素生成方法

配色方案：
- 主蓝色: #0886EC
- 次蓝色: #0070C0  
- 强调色: #FF7401 (橙色)
- 深蓝: #003F56
- 浅蓝背景: #CCDDEA
- 紫色辅助: #B073FC
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import copy

# 金蝶官方配色
KINGDEE_COLORS = {
    'primary': RGBColor(0x08, 0x86, 0xEC),      # 主蓝色 #0886EC
    'secondary': RGBColor(0x00, 0x70, 0xC0),    # 次蓝色 #0070C0
    'accent': RGBColor(0xFF, 0x74, 0x01),       # 橙色强调 #FF7401
    'dark': RGBColor(0x00, 0x3F, 0x56),         # 深蓝 #003F56
    'light': RGBColor(0xCC, 0xDD, 0xEA),        # 浅蓝背景 #CCDDEA
    'text_dark': RGBColor(0x00, 0x00, 0x00),    # 深色文字 #000000
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),   # 浅色文字 #FFFFFF
    'purple': RGBColor(0xB0, 0x73, 0xFC),       # 紫色辅助 #B073FC
    'light_blue': RGBColor(0x23, 0x86, 0xEE),   # 浅蓝 #2386EE
}


class KingdeeStyleGenerator:
    """金蝶PPT样式生成器"""
    
    def __init__(self, slide):
        self.slide = slide
        self.colors = KINGDEE_COLORS
    
    # ==================== 基础形状 ====================
    
    def add_rectangle(self, left, top, width, height, fill_color, text="", 
                      font_size=11, font_color=None, corner_radius=0):
        """添加矩形"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()  # 无边框
        
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or self.colors['text_dark']
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.name = "微软雅黑"
        
        # 圆角
        if corner_radius > 0:
            try:
                shape.adjustments[0] = corner_radius
            except:
                pass
        
        return shape
    
    def add_rounded_rectangle(self, left, top, width, height, fill_color, 
                              text="", font_size=11, font_color=None):
        """添加圆角矩形"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or self.colors['text_dark']
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "微软雅黑"
        
        return shape
    
    # ==================== 卡片组件 ====================
    
    def add_card(self, left, top, width, height, title, content_lines=None,
                 title_color=None, bg_color=None):
        """添加卡片组件
        
        样式：标题栏(蓝色底白字) + 内容区(浅色背景)
        """
        title_height = 0.4
        title_color = title_color or self.colors['primary']
        bg_color = bg_color or self.colors['light']
        
        # 标题栏
        title_shape = self.add_rounded_rectangle(
            left, top, width, title_height,
            title_color, title, 12, self.colors['text_light']
        )
        
        # 内容区
        if content_lines:
            content_shape = self.add_rounded_rectangle(
                left, top + title_height, width, height - title_height,
                bg_color, "", 11, self.colors['text_dark']
            )
            tf = content_shape.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(11)
                p.font.color.rgb = self.colors['text_dark']
                p.font.name = "微软雅黑"
        
        return title_shape
    
    def add_kpi_card(self, left, top, width, height, value, unit, label,
                     value_color=None, bg_color=None):
        """添加KPI展示卡片
        
        样式：大数字 + 单位 + 底部标签
        """
        value_color = value_color or self.colors['primary']
        bg_color = bg_color or self.colors['light']
        
        # 背景卡片
        card = self.add_rounded_rectangle(
            left, top, width, height,
            bg_color, "", 11, self.colors['text_dark']
        )
        
        # 大数字
        value_box = self.slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.15),
            Inches(width - 0.2), Inches(0.5)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        
        # 数字部分
        run1 = p.add_run()
        run1.text = str(value)
        run1.font.size = Pt(24)
        run1.font.bold = True
        run1.font.color.rgb = value_color
        run1.font.name = "微软雅黑"
        
        # 单位部分
        run2 = p.add_run()
        run2.text = unit
        run2.font.size = Pt(12)
        run2.font.color.rgb = self.colors['text_dark']
        run2.font.name = "微软雅黑"
        
        # 底部标签
        label_box = self.slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + height - 0.35),
            Inches(width - 0.2), Inches(0.3)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['text_dark']
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
        
        return card
    
    def add_kpi_row(self, left, top, kpis, card_width=1.5, card_height=1.2, gap=0.2):
        """添加KPI卡片行
        
        Args:
            kpis: [(value, unit, label), ...]
        """
        shapes = []
        for i, (value, unit, label) in enumerate(kpis):
            x = left + i * (card_width + gap)
            shape = self.add_kpi_card(x, top, card_width, card_height, 
                                       value, unit, label)
            shapes.append(shape)
        return shapes
    
    # ==================== 功能模块色块 ====================
    
    def add_function_block(self, left, top, width, height, title, items,
                           accent=False):
        """添加功能模块色块
        
        Args:
            title: 模块标题
            items: 功能项列表
            accent: 是否使用强调色(橙色)
        """
        bg_color = self.colors['accent'] if accent else self.colors['primary']
        
        # 标题栏
        title_shape = self.add_rounded_rectangle(
            left, top, width, 0.35,
            bg_color, title, 11, self.colors['text_light']
        )
        
        # 功能项
        if items:
            item_height = min(0.35, (height - 0.35) / len(items))
            for i, item in enumerate(items[:int((height-0.35)/0.35)]):
                item_shape = self.add_rounded_rectangle(
                    left, top + 0.35 + i * item_height,
                    width, item_height,
                    self.colors['light'], item, 10, self.colors['text_dark']
                )
        
        return title_shape
    
    def add_function_blocks_row(self, left, top, blocks, width=1.3, height=1.8, gap=0.1):
        """添加功能模块行
        
        Args:
            blocks: [(title, items, accent), ...]
        """
        shapes = []
        for i, (title, items, accent) in enumerate(blocks):
            x = left + i * (width + gap)
            shape = self.add_function_block(x, top, width, height, 
                                            title, items, accent)
            shapes.append(shape)
        return shapes
    
    # ==================== 时间线组件 ====================
    
    def add_timeline(self, left, top, width, milestones):
        """添加时间线
        
        Args:
            milestones: [(date, title), ...]
        """
        shapes = []
        total = len(milestones)
        
        # 渐变底条
        bar = self.add_rectangle(
            left, top + 0.5, width, 0.1,
            self.colors['light'], "", 11, self.colors['text_dark']
        )
        shapes.append(bar)
        
        # 节点
        for i, (date, title) in enumerate(milestones):
            x = left + (width / (total - 1)) * i if total > 1 else left
            
            # 圆形节点
            node = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.15), Inches(top + 0.4),
                Inches(0.3), Inches(0.3)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = self.colors['primary']
            node.line.fill.background()
            shapes.append(node)
            
            # 编号
            num_box = self.slide.shapes.add_textbox(
                Inches(x - 0.1), Inches(top + 0.43),
                Inches(0.2), Inches(0.25)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['text_light']
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
            
            # 日期
            date_box = self.slide.shapes.add_textbox(
                Inches(x - 0.4), Inches(top + 0.8),
                Inches(0.8), Inches(0.3)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['text_dark']
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
            
            # 标题
            title_box = self.slide.shapes.add_textbox(
                Inches(x - 0.5), Inches(top + 1.1),
                Inches(1), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['text_dark']
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
        
        return shapes
    
    # ==================== 架构图组件 ====================
    
    def add_architecture_layer(self, left, top, width, height, title, 
                               modules, layer_color=None):
        """添加架构层
        
        Args:
            title: 层名称(如"前台"、"中台"、"后台")
            modules: 模块列表
        """
        layer_color = layer_color or self.colors['dark']
        
        # 左侧标签
        label = self.add_rounded_rectangle(
            left, top, 0.8, height,
            layer_color, title, 12, self.colors['text_light']
        )
        
        # 模块卡片
        module_width = (width - 1) / len(modules) if modules else 1
        shapes = [label]
        
        for i, module in enumerate(modules):
            x = left + 0.9 + i * module_width
            card = self.add_rounded_rectangle(
                x, top + 0.1, module_width - 0.1, height - 0.2,
                self.colors['primary'], module, 11, self.colors['text_light']
            )
            shapes.append(card)
        
        return shapes
    
    # ==================== 流程图组件 ====================
    
    def add_process_flow(self, left, top, steps, step_width=1.5, step_height=0.6, gap=0.3):
        """添加横向流程图
        
        Args:
            steps: [step1, step2, ...]
        """
        shapes = []
        arrow_width = 0.2
        
        for i, step in enumerate(steps):
            x = left + i * (step_width + gap + arrow_width)
            
            # 步骤框
            box = self.add_rounded_rectangle(
                x, top, step_width, step_height,
                self.colors['primary'], step, 12, self.colors['text_light']
            )
            shapes.append(box)
            
            # 箭头(最后一个不加)
            if i < len(steps) - 1:
                arrow = self.slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(x + step_width + 0.05), Inches(top + step_height/2 - 0.1),
                    Inches(arrow_width), Inches(0.2)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['accent']
                arrow.line.fill.background()
                shapes.append(arrow)
        
        return shapes
    
    # ==================== 表格美化 ====================
    
    def add_styled_table(self, left, top, rows, cols, data, header_color=None):
        """添加美化表格"""
        from pptx.table import Table
        
        # 简化版：用矩形模拟
        cell_width = 2.0
        cell_height = 0.4
        header_color = header_color or self.colors['primary']
        
        shapes = []
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_text in enumerate(row_data):
                x = left + col_idx * cell_width
                y = top + row_idx * cell_height
                
                # 背景色
                if row_idx == 0:
                    bg = header_color
                    fc = self.colors['text_light']
                else:
                    bg = self.colors['light'] if row_idx % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                    fc = self.colors['text_dark']
                
                cell = self.add_rectangle(
                    x, y, cell_width, cell_height,
                    bg, str(cell_text), 10, fc
                )
                shapes.append(cell)
        
        return shapes


# ==================== 辅助函数 ====================

def hex_to_rgb(hex_color):
    """十六进制颜色转RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), 
                   int(hex_color[2:4], 16), 
                   int(hex_color[4:6], 16))


def create_slide_with_title(slide, title, subtitle=None):
    """创建带标题的幻灯片"""
    # 使用母版的标题占位符
    if slide.slide_layout.name == '3_白色内页':
        for shape in slide.shapes:
            if shape.has_text_frame and shape.placeholder_format:
                if shape.placeholder_format.type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = title
                    break
    else:
        # 手动添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = KINGDEE_COLORS['dark']
        p.font.name = "微软雅黑"


if __name__ == "__main__":
    print("金蝶PPT样式组件库 V1.0")
    print("配色方案：")
    for name, color in KINGDEE_COLORS.items():
        print(f"  {name}: #{color}")
