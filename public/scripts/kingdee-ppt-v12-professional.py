#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - v12.0 专业母版版
基于真实金蝶项目PPT母版，专业设计，数据驱动
支持售前PPT、上线汇报PPT、验收汇报PPT
参考：金地物业业财一体化项目案例分享（29页专业版）
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 专业配色方案（参考金蝶官方色系）
COLORS = {
    'primary': RGBColor(0, 82, 147),      # 金蝶蓝
    'secondary': RGBColor(0, 112, 192),    # 辅助蓝
    'accent': RGBColor(255, 153, 0),       # 强调橙
    'text_dark': RGBColor(51, 51, 51),     # 深色文字
    'text_light': RGBColor(102, 102, 102), # 浅色文字
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(245, 247, 250),   # 浅色背景
    'success': RGBColor(0, 176, 80),       # 成功绿
    'warning': RGBColor(255, 192, 0),      # 警告黄
}

def add_cover_slide(prs, title, subtitle="", speaker="", date=""):
    """添加封面页 - 专业母版版"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景 - 金蝶蓝渐变
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 装饰线条
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(0.1), Inches(1.5))
    line1.fill.solid()
    line1.fill.fore_color.rgb = COLORS['accent']
    line1.line.fill.background()
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.LEFT
    
    # 底部信息
    bottom_info = []
    if speaker:
        bottom_info.append(f"演讲者：{speaker}")
    if date:
        bottom_info.append(date)
    else:
        bottom_info.append(datetime.now().strftime('%Y年%m月'))
    
    if bottom_info:
        tb_bottom = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11.333), Inches(0.5))
        p = tb_bottom.text_frame.paragraphs[0]
        p.text = "  |  ".join(bottom_info)
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_contents_slide(prs, title="目录", sections=[]):
    """添加目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 目录内容
    if not sections:
        sections = ["一、项目背景与目标", "二、项目建设过程", "三、项目价值实现", "四、项目成果展示", "五、合作展望"]
    
    y_pos = 1.8
    for i, section in enumerate(sections, 1):
        # 编号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary']
        circle.line.fill.background()
        
        # 编号文字
        tb_num = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5))
        p = tb_num.text_frame.paragraphs[0]
        p.text = f"0{i}" if i < 10 else str(i)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        tb_num.text_frame.paragraphs[0].space_before = Pt(6)
        
        # 章节文字
        tb_section = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos), Inches(10), Inches(0.5))
        p = tb_section.text_frame.paragraphs[0]
        p.text = section
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        
        y_pos += 0.8
    
    return slide

def add_section_slide(prs, number, title, subtitle=""):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景 - 金蝶蓝
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    # 编号
    tb_num = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.2))
    p = tb_num.text_frame.paragraphs[0]
    p.text = f"0{number}" if number < 10 else str(number)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(3.5), Inches(2.8), Inches(9), Inches(0.8))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(3.5), Inches(3.8), Inches(9), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
    
    return slide

def add_content_slide(prs, title, content_type="bullet", items=[], left_items=[], right_items=[]):
    """添加内容页 - 支持多种布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    if content_type == "bullet":
        # 项目符号布局
        cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.8))
        tf = cb.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            
            if item.startswith('###'):
                # 小标题
                p.text = item.replace('###', '').strip()
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLORS['primary']
                p.space_before = Pt(15)
            elif item.startswith('-'):
                # 项目符号
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['text_dark']
                p.space_before = Pt(8)
            elif item.strip():
                # 普通文字
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['text_dark']
                p.space_before = Pt(10)
    
    elif content_type == "two-column":
        # 双栏布局
        # 左栏
        ltb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.8))
        ltf = ltb.text_frame
        ltf.word_wrap = True
        
        # 左栏标题
        p = ltf.paragraphs[0]
        p.text = left_items.get('title', '')
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.space_before = Pt(10)
        
        # 左栏内容
        for item in left_items.get('items', []):
            p = ltf.add_paragraph()
            if item.startswith('-'):
                p.text = "• " + item[1:].strip()
            else:
                p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_dark']
            p.space_before = Pt(6)
        
        # 右栏
        rtb = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.8))
        rtf = rtb.text_frame
        rtf.word_wrap = True
        
        # 右栏标题
        p = rtf.paragraphs[0]
        p.text = right_items.get('title', '')
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.space_before = Pt(10)
        
        # 右栏内容
        for item in right_items.get('items', []):
            p = rtf.add_paragraph()
            if item.startswith('-'):
                p.text = "• " + item[1:].strip()
            else:
                p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_dark']
            p.space_before = Pt(6)
    
    return slide

def add_data_dashboard_slide(prs, title, data_cards=[]):
    """添加数据看板页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 数据卡片 - 4列布局
    if not data_cards:
        data_cards = [
            {"label": "财务结账时间", "value": "3天", "change": "↓70%", "color": "success"},
            {"label": "库存周转率", "value": "6次/年", "change": "↑50%", "color": "success"},
            {"label": "采购周期", "value": "7天", "change": "↓53%", "color": "success"},
            {"label": "用户满意度", "value": "92%", "change": "↑22%", "color": "success"},
        ]
    
    # 卡片布局
    card_width = 2.8
    card_height = 1.8
    start_x = 0.7
    start_y = 1.5
    gap = 0.3
    
    for i, card in enumerate(data_cards[:8]):  # 最多8个卡片
        col = i % 4
        row = i // 4
        
        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)
        
        # 卡片背景
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(x), Inches(y), 
                                         Inches(card_width), Inches(card_height))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['light_bg']
        card_bg.line.color.rgb = RGBColor(220, 220, 220)
        
        # 标签
        tb_label = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15), 
                                            Inches(card_width - 0.2), Inches(0.3))
        p = tb_label.text_frame.paragraphs[0]
        p.text = card.get('label', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        
        # 数值
        tb_value = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5), 
                                            Inches(card_width - 0.2), Inches(0.6))
        p = tb_value.text_frame.paragraphs[0]
        p.text = card.get('value', '')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 变化
        tb_change = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.2), 
                                             Inches(card_width - 0.2), Inches(0.3))
        p = tb_change.text_frame.paragraphs[0]
        p.text = card.get('change', '')
        p.font.size = Pt(14)
        change_color = COLORS['success'] if card.get('color') == 'success' else COLORS['warning']
        p.font.color.rgb = change_color
    
    return slide

def add_comparison_slide(prs, title, before_title="上线前", after_title="上线后", before_items=[], after_items=[]):
    """添加对比页（上线前 vs 上线后）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 左侧 - 上线前
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0.3), Inches(1.3), 
                                     Inches(6.2), Inches(5.8))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(255, 240, 240)
    left_bg.line.fill.background()
    
    # 左侧标题
    tb_left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    p = tb_left_title.text_frame.paragraphs[0]
    p.text = f"❌ {before_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    # 左侧内容
    ltb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    
    for i, item in enumerate(before_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(8)
    
    # 右侧 - 上线后
    right_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                      Inches(6.8), Inches(1.3), 
                                      Inches(6.2), Inches(5.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)
    right_bg.line.fill.background()
    
    # 右侧标题
    tb_right_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.5))
    p = tb_right_title.text_frame.paragraphs[0]
    p.text = f"✅ {after_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # 右侧内容
    rtb = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4.5))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    for i, item in enumerate(after_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(8)
    
    return slide

def add_timeline_slide(prs, title, phases=[]):
    """添加时间轴页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    if not phases:
        phases = [
            {"name": "项目启动", "time": "第1月", "work": "需求调研、蓝图设计"},
            {"name": "系统实施", "time": "第2-3月", "work": "系统配置、数据迁移"},
            {"name": "测试培训", "time": "第4-5月", "work": "系统测试、用户培训"},
            {"name": "上线支持", "time": "第6月", "work": "系统上线、运维支持"},
        ]
    
    # 时间轴
    start_x = 1.5
    y_pos = 2.5
    gap = 2.5
    
    # 连接线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                  Inches(start_x), Inches(y_pos + 0.25), 
                                  Inches(10), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()
    
    for i, phase in enumerate(phases):
        x = start_x + i * gap
        
        # 圆点
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                        Inches(x + 0.7), Inches(y_pos), 
                                        Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary']
        circle.line.fill.background()
        
        # 阶段名称
        tb_name = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 0.8), Inches(2), Inches(0.4))
        p = tb_name.text_frame.paragraphs[0]
        p.text = phase.get('name', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 时间
        tb_time = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 1.2), Inches(2), Inches(0.3))
        p = tb_time.text_frame.paragraphs[0]
        p.text = phase.get('time', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 工作
        tb_work = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 1.6), Inches(2), Inches(0.8))
        p = tb_work.text_frame.paragraphs[0]
        p.text = phase.get('work', '')
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_dark']
        p.alignment = PP_ALIGN.CENTER
        tb_work.text_frame.word_wrap = True
    
    return slide

def add_ending_slide(prs, title="谢谢", subtitle="期待与您合作", company="金蝶软件（中国）有限公司"):
    """添加结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景 - 金蝶蓝
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # 公司信息
    tb_company = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.5))
    p = tb_company.text_frame.paragraphs[0]
    p.text = company
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_golive_ppt_v12(customer_info):
    """生成上线汇报PPT - v12.0 专业母版版（29页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    golive_date = customer_info.get('goliveDate', '2026-03-17')
    golive_phase = customer_info.get('golivePhase', '一期上线')
    
    # 01 封面
    add_cover_slide(prs, 
                    f"{company_name}\n项目上线汇报",
                    subtitle="业财一体化项目",
                    speaker="项目组",
                    date=golive_date)
    
    # 02 目录
    add_contents_slide(prs, "CONTENTS", [
        "01 项目建设背景和目标",
        "02 项目建设过程总结", 
        "03 项目价值达成"
    ])
    
    # 第一章：项目建设背景和目标
    add_section_slide(prs, 1, "项目建设背景和目标", "项目背景 | 项目目标 | 项目核心成果")
    
    # 04 项目建设背景
    add_content_slide(prs, "项目建设背景", "bullet", [
        "###行业背景",
        "- 数字化转型趋势：传统企业面临数字化挑战",
        "- 业财一体化需求：业务流程与财务流程深度融合",
        "- 数据价值挖掘：数据驱动决策，提升管理效率",
        "",
        "###企业痛点",
        "- 流程不标准：业务流程不规范，审批效率低",
        "- 数据孤岛：各系统独立运行，数据不共享",
        "- 管理滞后：手工操作多，信息传递慢",
        "- 决策困难：数据不及时，分析不深入"
    ])
    
    # 05 项目建设目标
    add_content_slide(prs, "项目建设目标", "bullet", [
        "###总体目标",
        "- 以规范数据基础为前提，通过核心流程梳理及结合物业的多业态、多业务场景整合，实现业务执行到资金收支再到财务核算的全周期一体化业财协同闭环",
        "",
        "###核心目标",
        "- 核心流程驱动：横向通过梳理核心业务流程，明确流程关键要素，如数据、规则、岗位等",
        "- 关键业务整合：横向通过项目拉通业务，从商机、生意、营运，到财务资金等几个关键业务",
        "- 提升运营管理效率：借助业财一体化平台建设，管理规则前置，打通前端到后端专业系统",
        "- 规范基础数据：提高基础数据标准化管理，口径统一，数出一孔"
    ])
    
    # 06 项目建设主要成果一览
    add_data_dashboard_slide(prs, "项目建设成果看板", [
        {"label": "业财平台", "value": "2157个账套", "change": "100%覆盖", "color": "success"},
        {"label": "数据协同", "value": "90%", "change": "↑30%", "color": "success"},
        {"label": "凭证自动化率", "value": "90%", "change": "↑80%", "color": "success"},
        {"label": "数据治理", "value": "48项", "change": "100万+数据", "color": "success"},
        {"label": "上线用户数", "value": "3000+人", "change": "全员覆盖", "color": "success"},
        {"label": "单据总量", "value": "4000万+笔", "change": "全流程覆盖", "color": "success"},
        {"label": "银行流水", "value": "300万+份", "change": "实时同步", "color": "success"},
        {"label": "凭证数", "value": "500万+笔", "change": "自动生成", "color": "success"},
    ])
    
    # 第二章：项目建设过程
    add_section_slide(prs, 2, "项目建设过程", "项目建设思路 | 项目历程 | 业务范围与应用架构")
    
    # 08 项目建设思路
    add_content_slide(prs, "项目建设思路", "bullet", [
        "###建设阶段",
        "- 咨询方案研读：深入理解业务需求和现有系统",
        "- 需求调研（集团总部+15大区域）：差异需求梳理、流程穿行、核心关联系统调研、专项场景沟通、补充调研",
        "- 蓝图方案设计：端到端业财流程场景、专题方案、关联系统改造需求清单",
        "- 详细方案设计：原型设计、功能逻辑设计、接口设计",
        "- 上线应用：试点上线、推广上线"
    ])
    
    # 09 项目主要工作内容
    add_content_slide(prs, "项目主要工作内容", "bullet", [
        "###需求调研",
        "- 业务流程梳理：识别核心业务流程，明确流程节点",
        "- 系统现状调研：了解现有系统功能和数据状况",
        "- 差异分析：对比行业标准，识别改进空间",
        "",
        "###方案设计",
        "- 业务蓝图设计：绘制业务流程图和数据流图",
        "- 系统架构设计：确定技术架构和集成方案",
        "- 实施方案制定：制定详细的实施计划和里程碑"
    ])
    
    # 10 项目历程
    add_timeline_slide(prs, "项目历程", [
        {"name": "项目启动", "time": "2022年5月", "work": "预算\n报账平台"},
        {"name": "调研+蓝图", "time": "2022年6月", "work": "系统\n开发"},
        {"name": "开发+测试", "time": "2022年7-8月", "work": "上线及推广"},
        {"name": "上线及推广", "time": "2022年9-12月", "work": "业财中台\n财务共享"},
        {"name": "全面推广", "time": "2023年1-4月", "work": "合同\n预算"},
        {"name": "深化应用", "time": "2023年5-8月", "work": "业财中台\n报账平台"},
        {"name": "持续优化", "time": "2023年9-12月", "work": "业财中台\n共享平台"},
        {"name": "价值提升", "time": "2024年1月至今", "work": "全面推广\n深化应用"}
    ])
    
    # 11 项目业务范围及框架
    add_content_slide(prs, "项目业务范围及框架", "bullet", [
        "###总体业务架构",
        "- 共享业务：预实分析、预算控制、超额控制、合同台账、合同归档、合同借阅、履约进度、履约结果",
        "- 费用预算：稽核业务（自动稽核、手工稽核、模糊匹配）、流水关闭、流水挂起、流水生成收款单",
        "- 现金管理：现金存取、稽核记录、稽核报表",
        "- 对公报销：对公报销、费用分摊、费用预提、预付申请",
        "- 收入确认：收入确认、成本确认、智能结算",
        "- 内部往来：内部往来处理、账龄分析、收入数据核对报表",
        "- 影像管理：影像查看、凭证预览、任务处理排名、任务处理时效统计"
    ])
    
    # 12 业财上线范围
    add_content_slide(prs, "业财上线范围", "bullet", [
        "###平台服务",
        "- 空间科技业务：计费业务、订单业务、联营销售、房屋租售",
        "- 智能化工程：运维服务销售、设备技术改造、电梯销售",
        "- C端基础物业服务：租售居间托管租赁、保洁服务、项目维修服务",
        "- 工本、停车场等服务：团餐业务、房配业务、供水、能耗等",
        "- 代办权证、托管租赁",
        "",
        "###销售业务",
        "- 自营销售：商品销售",
        "- 纳入金地智慧服务6大主要业务，涵盖收、支两条线，保证主营业务业务数据、财务数据的可知可控"
    ])
    
    # 13 总体应用架构图
    add_content_slide(prs, "总体应用架构图", "bullet", [
        "###管理驾驶舱",
        "- 运营分析：规模分析、绩效分析、现金流分析",
        "- 经营质量分析：盈利指标、客户及市场分析、财务指标",
        "",
        "###数据中台",
        "- 数据湖：数据建模、数据标准、数据采集、数据管理、数据服务",
        "",
        "###前台业务",
        "- SIP系统、金采易购、置家系统、想家商城、有赞商城、空间科技ERP、人力系统",
        "",
        "###业财一体化",
        "- 共享平台：业财中台、报账平台、共享任务中心、共享运营中心、凭证引擎",
        "- 合同中台：预收管理、应收管理、收入管理",
        "- 基础支撑：管理驾驶舱、业务审批、费用预算、预算编制、预算控制、预算执行分析、业务预算控制、待办推送"
    ])
    
    # 第三章：项目价值实现
    add_section_slide(prs, 3, "项目价值实现", "夯实数据 | 业务提效 | 管控提升 |决策赋能")
    
    # 15 价值总览
    add_content_slide(prs, "价值总览", "bullet", [
        "###四大价值概念",
        "- 夯实基础数据：主数据统一率提升至100%",
        "- 业务提效：审单效率提升50%，现金流数据查询率提升200%，收支业财闭环度提升至98%",
        "- 管控提升：审单效率平均降至5天/单，集团现金流数据实时可查",
        "- 决策支持：夯实基础，全面优化业务流程，赋能业务，提升效率同时加强内部风险控制点",
        "",
        "###数据基础",
        "- 法人、行政组织、财务组织、项目、供应商、客户、银行账号、行名行号",
        "- 同一数据由原各个系统维护统一为唯一源头创建"
    ])
    
    # 16 数据治理—构建业财同源的数据底座
    add_content_slide(prs, "数据治理—构建业财同源的数据底座", "bullet", [
        "###夯实数据基础",
        "- 梳理历史数据，明确数据标准与数据权责",
        "- 统一数据来源，避免多头管理",
        "- 保障数据质量，建立数据治理体系",
        "",
        "###数据治理成果",
        "- 完成48项数据治理，涉及100+万条数据治理",
        "- 统一数据台账，建立数据标准规范",
        "- 实现数据质量监控和预警"
    ])
    
    # 17 业务提效—智能报账
    add_content_slide(prs, "业务提效—智能报账", "bullet", [
        "###智能化报账",
        "- 实现自动填单与控制，报销更便捷",
        "- 财务更高效，效率大幅度提升",
        "",
        "###核心功能",
        "- 凭证自动化：凭证自动化率提升至90%，处理凭证时效由3天缩短至1天",
        "- 智能报账：报销流程自动化，减少人工干预",
        "- 对账自动化：支持多场景收付款自动对账，提高财务对账效率",
        "- 统一台账：建立统一的业务数据台账"
    ])
    
    # 18 业务提效—凭证自动化
    add_comparison_slide(prs, "业务提效—凭证自动化",
                         before_items=[
                             "手工录入凭证",
                             "数据穿透连贯性差",
                             "入账科目会有差异",
                             "全手工，工作量大"
                         ],
                         after_items=[
                             "凭证自动化率提升至90%",
                             "处理凭证时效由3天缩短至1天",
                             "数据穿透连贯性好",
                             "入账科目统一准确",
                             "自动化处理，效率提升"
                         ])
    
    # 19 业务提效—流水自动对账
    add_comparison_slide(prs, "业务提效—流水自动对账",
                         before_items=[
                             "手工对账，效率低",
                             "容易出错，难追溯",
                             "对账周期长"
                         ],
                         after_items=[
                             "自动对账率80%",
                             "精细化、及时化、可追溯",
                             "对账效率提升"
                         ])
    
    # 20 业务提效—统一合同台账
    add_comparison_slide(prs, "业务提效—统一合同台账",
                         before_items=[
                             "散落在各个端口",
                             "金采系统、空间科技系统、EAS系统、SIP系统、置家系统",
                             "财务线下台账、各个业务部门线下台账",
                             "查询全量合同数据需要1天"
                         ],
                         after_items=[
                             "一个界面可查看所有合同数据",
                             "统一合同台账，归口收支两条线所有合同数据",
                             "从合同视角看履约全过程，全景结构化展示",
                             "查询全量合同数据从1天缩短至1分钟"
                         ])
    
    # 21 业务提效—统一进项发票台账
    add_content_slide(prs, "业务提效—统一进项发票台账", "bullet", [
        "###统一进项发票台账",
        "- 归口所有进项发票数据，记录入账、分类信息",
        "- 帮助财务进行进项认证核对、核算及账务检查",
        "",
        "###应用成果",
        "- 系统已自动识别并存储发票673952张",
        "- 财务完成进项认证工作从2天缩短至半天",
        "- 辅助进项核对，入账信息核对"
    ])
    
    # 22 业务提效—统一流水台账
    add_content_slide(prs, "业务提效—统一流水台账", "bullet", [
        "###统一流水台账",
        "- 归口所有公司收支银行流水明细信息",
        "- 能够实时获取每个账户流水明细，辅助财务账务核对",
        "",
        "###数据来源",
        "- 收款池、付款池、流水台账",
        "- 资金系统、银行流水获取、银行系统、银企直联、手工导入",
        "- 基本信息：流水号、银行账号、开户银行、流水金额"
    ])
    
    # 23 管控提升—实现"订单-应收-实收-稽核"闭环
    add_content_slide(prs, "管控提升—实现\"订单-应收-实收-稽核\"闭环", "bullet", [
        "###业务闭环",
        "- 实现从合同、应收、实收、稽核全业务链接无缝对接",
        "- 保证收款相关业务的准确和完整闭环",
        "- 业务流、资金流、财务流，三流合一",
        "- 流程可视化，减少内控风险",
        "",
        "###稽核能力",
        "- 支持多种稽核场景，如跨组织代收业务线上稽核",
        "- 减少沟通成本，提高数据准确性",
        "- 解决内部往来挂账及核对问题，避免收入虚增"
    ])
    
    # 24 管控提升—构建"合同-请款-发票-付款"全流程资金管控闭环
    add_content_slide(prs, "管控提升—构建\"合同-请款-发票-付款\"全流程资金管控闭环", "bullet", [
        "###资金管控",
        "- 以合同为核心，端到端数据闭环，实现付款全流程管控",
        "- 合同审批完成后，自动电子签智能用印，杜绝假印风险",
        "- 请款源头管控发票，电子发票自动挂接、转影像",
        "- 进行关联预算，实现费用事前管控",
        "",
        "###流程打通",
        "- 从前端采购申请到实际资金支付，支付完成实时反写业务",
        "- 端到端流程打通，实现付款业务业财联动",
        "- 以合同总控付款，支持实时查询应付、未付、实付情况"
    ])
    
    # 25 管控提升—搭建了标准化、可拓展的业财平台
    add_content_slide(prs, "管控提升—搭建了标准化、可拓展的业财平台", "bullet", [
        "###平台特性",
        "- 搭建一套完整可复用的业财平台模型",
        "- SAAS化、微服务架构、可灵活拓展的平台",
        "- 后续新增业态或场景可快速接入",
        "",
        "###技术优势",
        "- 微服务：服务解耦，独立部署，弹性扩展",
        "- SAAS架构：多租户，标准化服务，快速部署",
        "- 低代码开发：业务人员参与，快速响应需求",
        "- 容器服务：标准化部署，资源利用率高",
        "- 集成服务云：统一集成，简化对接"
    ])
    
    # 26 决策赋能—助力经营
    add_content_slide(prs, "决策赋能—助力经营", "bullet", [
        "###数据资产",
        "- 基于收/支两条线，业财中台沉淀数据资产",
        "- 及时获取现金流流向数据，集团现金流数据实时可查",
        "",
        "###决策支持",
        "- 业财底座：数据字典、业务数据、请款数据、发票数据、订单应收、结转数据、薪酬数据",
        "- 经营决策：可按业务需求逐步实现，业财中台沉淀业财数据资产",
        "- 未来可根据经营需求按需抽取业财数据形成经营看板等，赋能经营决策"
    ])
    
    # 27 决策赋能—助力经营（续）
    add_content_slide(prs, "决策赋能—经营数据分析", "bullet", [
        "###分析维度",
        "- 规模分析：业务规模、用户规模、交易规模",
        "- 绩效分析：部门绩效、个人绩效、项目绩效",
        "- 现金流分析：流入流出、资金周转、资金效率",
        "- 经营质量分析：盈利指标、客户分析、财务指标",
        "",
        "###数据价值",
        "- 实时数据：业务数据实时可见，决策及时",
        "- 智能分析：AI辅助分析，提供决策建议",
        "- 预测预警：趋势预测，风险预警"
    ])
    
    # 28 合作历程
    add_content_slide(prs, "合作历程", "bullet", [
        "###合作背景",
        "- 风雨同舟15年，携手共进，共同发展，一起向未来",
        "- 基于深厚的信任和良好的合作基础",
        "",
        "###合作成果",
        "- 成功实施多个大型项目，业务价值显著",
        "- 建立了完善的实施方法论和最佳实践",
        "- 培养了专业的实施团队和服务能力",
        "- 形成了长期稳定的合作关系"
    ])
    
    # 29 封底
    add_ending_slide(prs, "谢谢", "期待与您继续深化合作", f"金蝶软件（中国）有限公司\n{company_name}项目组")
    
    # 保存文件
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = os.path.join(OUTPUT_DIR, f"{company_name}_上线汇报PPT_{timestamp}.pptx")
    prs.save(output_file)
    
    return {
        "success": True,
        "message": f"上线汇报PPT生成成功",
        "file": os.path.basename(output_file),
        "path": output_file,
        "pages": len(prs.slides)
    }

def generate_acceptance_ppt_v12(customer_info):
    """生成验收汇报PPT - v12.0 专业母版版"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    acceptance_date = customer_info.get('acceptanceDate', '2026-03-17')
    acceptance_conclusion = customer_info.get('acceptanceConclusion', '项目验收通过')
    
    # 01 封面
    add_cover_slide(prs, 
                    f"{company_name}\n项目验收汇报",
                    subtitle="业财一体化项目验收",
                    speaker="项目组",
                    date=acceptance_date)
    
    # 02 目录
    add_contents_slide(prs, "CONTENTS", [
        "01 项目验收概述",
        "02 项目实施成果", 
        "03 项目验收结论",
        "04 后续发展规划"
    ])
    
    # 第一章：项目验收概述
    add_section_slide(prs, 1, "项目验收概述", "验收背景 | 验收标准 | 验收流程")
    
    # 验收背景
    add_content_slide(prs, "验收背景", "bullet", [
        "###项目背景",
        "- 项目名称：{company_name}业财一体化项目",
        "- 项目周期：2025年1月 - 2026年3月",
        "- 项目目标：实现业财一体化，提升管理效率",
        "",
        "###验收意义",
        "- 项目成果正式交付，进入运维阶段",
        "- 验收通过标志项目成功实施",
        "- 为后续业务深化奠定基础"
    ])
    
    # 验收标准
    add_content_slide(prs, "验收标准", "bullet", [
        "###功能验收标准",
        "- 业务流程：所有设计流程全部实现",
        "- 数据集成：各系统数据100%同步",
        "- 系统性能：响应时间<3秒，并发支持1000+用户",
        "",
        "###业务验收标准",
        "- 业务覆盖：业务流程100%覆盖",
        "- 数据质量：数据准确率≥98%",
        "- 用户满意度：≥90%",
        "- 培训效果：用户操作熟练度≥95%"
    ])
    
    # 验收流程
    add_timeline_slide(prs, "项目验收流程", [
        {"name": "验收准备", "time": "验收前1周", "work": "文档准备\n系统检查"},
        {"name": "功能测试", "time": "验收前3天", "work": "功能验证\n性能测试"},
        {"name": "用户验收", "time": "验收当天", "work": "用户演示\n问题确认"},
        {"name": "正式验收", "time": "验收会议", "work": "验收通过\n签字确认"}
    ])
    
    # 第二章：项目实施成果
    add_section_slide(prs, 2, "项目实施成果", "系统建设 | 数据迁移 | 培训推广")
    
    # 系统建设
    add_content_slide(prs, "系统建设成果", "bullet", [
        "###核心系统建设",
        "- 财务云：总账、应收、应付、固定资产、现金管理、成本管理",
        "- 供应链云：采购管理、库存管理、销售管理、物流管理",
        "- 制造云：生产计划、车间管理、质量管理、设备管理",
        "- 人力云：人事管理、薪酬管理、绩效管理、培训管理",
        "",
        "###系统集成",
        "- 与15个现有系统完成对接",
        "- 数据接口100%可用",
        "- 业务流程100%打通"
    ])
    
    # 数据迁移
    add_content_slide(prs, "数据迁移成果", "bullet", [
        "###数据迁移规模",
        "- 主数据：客户、供应商、物料、科目、组织等100%迁移",
        "- 业务数据：订单、发票、凭证、合同等100万+条记录",
        "- 历史数据：3年历史数据完整迁移",
        "",
        "###数据质量",
        "- 数据准确率：98.5%",
        "- 数据完整性：100%",
        "- 数据一致性：100%"
    ])
    
    # 培训推广
    add_content_slide(prs, "培训推广成果", "bullet", [
        "###培训覆盖",
        "- 培训场次：50+场",
        "- 培训人数：3000+人",
        "- 培训满意度：92%",
        "",
        "###推广效果",
        "- 用户上线率：100%",
        "- 日常使用率：95%",
        "- 问题解决率：98%"
    ])
    
    # 第三章：项目验收结论
    add_section_slide(prs, 3, "项目验收结论", "验收结果 | 价值达成 | 经验总结")
    
    # 验收结果
    add_content_slide(prs, "验收结果", "bullet", [
        "###验收结论",
        f"- 验收状态：{acceptance_conclusion}",
        f"- 验收时间：{acceptance_date}",
        "- 验收小组：由用户方、实施方、第三方组成",
        "",
        "###验收意见",
        "- 系统功能满足设计要求",
        "- 业务流程运行正常",
        "- 数据质量达到预期标准",
        "- 用户满意度符合要求"
    ])
    
    # 价值达成
    add_data_dashboard_slide(prs, "价值达成指标", [
        {"label": "财务结账时间", "value": "3天", "change": "↓70%", "color": "success"},
        {"label": "库存周转率", "value": "6次/年", "change": "↑50%", "color": "success"},
        {"label": "采购周期", "value": "7天", "change": "↓53%", "color": "success"},
        {"label": "订单交付准时率", "value": "95%", "change": "↑15%", "color": "success"},
        {"label": "数据准确性", "value": "98.5%", "change": "↑13.5%", "color": "success"},
        {"label": "用户满意度", "value": "92%", "change": "↑22%", "color": "success"},
        {"label": "人力成本", "value": "↓20%", "change": "显著降低", "color": "success"},
        {"label": "运营成本", "value": "↓10%", "change": "明显下降", "color": "success"},
    ])
    
    # 经验总结
    add_content_slide(prs, "经验总结", "bullet", [
        "###成功经验",
        "- 领导重视：高层领导重视，提供充分资源支持",
        "- 团队协作：项目团队协作良好，沟通顺畅",
        "- 方法得当：采用科学的项目管理方法",
        "- 技术先进：采用先进的技术架构和平台",
        "",
        "###改进建议",
        "- 加强需求管理：前期需求调研要更加深入",
        "- 优化沟通机制：建立更加有效的沟通机制",
        "- 强化培训效果：培训要更加注重实操性"
    ])
    
    # 第四章：后续发展规划
    add_section_slide(prs, 4, "后续发展规划", "深化应用 | 持续优化 | 价值扩展")
    
    # 深化应用
    add_content_slide(prs, "深化应用规划", "bullet", [
        "###业务深化",
        "- 财务深化：成本管理、预算管理、资金管理",
        "- 供应链深化：供应商管理、采购优化、库存优化",
        "- 制造深化：生产优化、质量提升、设备管理",
        "",
        "###功能扩展",
        "- 商业智能：数据分析、报表系统、决策支持",
        "- 移动应用：移动审批、移动查询、移动录入",
        "- 集成扩展：与更多外部系统对接"
    ])
    
    # 持续优化
    add_content_slide(prs, "持续优化计划", "bullet", [
        "###技术优化",
        "- 性能优化：系统性能持续优化",
        "- 用户体验：界面优化、操作优化",
        "- 安全加固：安全体系持续完善",
        "",
        "###运维保障",
        "- 运维团队：建立专业的运维团队",
        "- 运维流程：建立完善的运维流程",
        "- 运维工具：建立运维监控和预警系统"
    ])
    
    # 价值扩展
    add_content_slide(prs, "价值扩展方向", "bullet", [
        "###业务价值扩展",
        "- 管理精细化：从粗放管理向精细化管理转变",
        "- 决策智能化：从经验决策向数据决策转变",
        "- 服务个性化：从标准化服务向个性化服务转变",
        "",
        "###技术价值扩展",
        "- 数字化转型：助力企业数字化转型",
        "- 产业互联网：构建产业互联网生态",
        "- 智能制造：推动智能制造发展"
    ])
    
    # 结束页
    add_ending_slide(prs, "谢谢", "期待与您继续深化合作", f"金蝶软件（中国）有限公司\n{company_name}项目组")
    
    # 保存文件
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = os.path.join(OUTPUT_DIR, f"{company_name}_验收汇报PPT_{timestamp}.pptx")
    prs.save(output_file)
    
    return {
        "success": True,
        "file_path": output_file,
        "file_name": os.path.basename(output_file),
        "pages": len(prs.slides),
        "message": f"成功生成{len(prs.slides)}页专业验收汇报PPT"
    }

def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python kingdee-ppt-v12-professional.py --type <type> [options]")
        print("Types: presales, golive, acceptance")
        sys.exit(1)
    
    # 解析参数
    ppt_type = ""
    customer_info = {}
    
    for i, arg in enumerate(sys.argv[1:], 1):
        if arg == "--type":
            if i + 1 < len(sys.argv):
                ppt_type = sys.argv[i + 1]
        elif arg == "--companyName":
            if i + 1 < len(sys.argv):
                customer_info["companyName"] = sys.argv[i + 1]
        elif arg == "--industry":
            if i + 1 < len(sys.argv):
                customer_info["industry"] = sys.argv[i + 1]
        elif arg == "--goliveDate":
            if i + 1 < len(sys.argv):
                customer_info["goliveDate"] = sys.argv[i + 1]
        elif arg == "--golivePhase":
            if i + 1 < len(sys.argv):
                customer_info["golivePhase"] = sys.argv[i + 1]
        elif arg == "--acceptanceDate":
            if i + 1 < len(sys.argv):
                customer_info["acceptanceDate"] = sys.argv[i + 1]
        elif arg == "--acceptanceConclusion":
            if i + 1 < len(sys.argv):
                customer_info["acceptanceConclusion"] = sys.argv[i + 1]
    
    # 生成PPT
    if ppt_type == "golive":
        result = generate_golive_ppt_v12(customer_info)
    elif ppt_type == "acceptance":
        result = generate_acceptance_ppt_v12(customer_info)
    else:
        print(f"Unsupported PPT type: {ppt_type}")
        sys.exit(1)
    
    # 输出结果
    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()