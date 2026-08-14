#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - v10.0 增强版
支持售前PPT、上线汇报PPT、验收汇报PPT
基于真实金蝶项目文档模板，包含专业内容和设计
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'light': RGBColor(240, 240, 240),
    'white': RGBColor(255, 255, 255),
    'success': RGBColor(76, 175, 80),
    'warning': RGBColor(255, 152, 0),
    'error': RGBColor(244, 67, 54),
}

def add_title_slide(prs, title, subtitle="", date=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    
    if date:
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = date
        p3.font.size = Pt(18)
        p3.font.color.rgb = COLORS['white']
        p3.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith('###'):
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(8)
        elif item.strip():
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
    
    return slide

def add_table_slide(prs, title, headers, data):
    """添加表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 表格
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5))
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['light']
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(16)
    
    # 数据
    for i, row_data in enumerate(data, 1):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
    
    return slide

def generate_presales_ppt_v10(customer_info):
    """生成售前PPT - v10.0 完整版"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    
    # 封面
    add_title_slide(prs, company_name, "数字化转型解决方案", datetime.now().strftime('%Y年%m月%d日'))
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、企业概况与需求分析",
        "二、解决方案设计", 
        "三、业务架构蓝图",
        "四、4A架构设计",
        "五、价值工程分析",
        "六、实施路线图",
        "七、成功案例分享",
        "八、项目团队与保障"
    ])
    
    # 企业概况
    add_section_slide(prs, "一、企业概况与需求分析")
    add_content_slide(prs, "企业基本情况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{customer_info.get('companySize', '中型企业')}",
        f"员工人数：{customer_info.get('employees', '待定')}人",
        f"年营业额：{customer_info.get('revenue', '待定')}万元"
    ])
    
    # 需求分析
    add_content_slide(prs, "业务需求分析", [
        "###核心痛点：",
        "- 业务流程不标准化",
        "- 数据管理分散",
        "- 系统集成困难", 
        "- 管理决策滞后",
        "",
        "###业务目标：",
        "- 实现业务流程标准化",
        "- 建立统一数据平台",
        "- 提升系统集成效率",
        "- 支持实时决策分析"
    ])
    
    # 解决方案
    add_section_slide(prs, "二、解决方案设计")
    add_content_slide(prs, "总体解决方案", [
        "###实施范围：",
        "- 财务管理模块",
        "- 供应链管理模块", 
        "- 生产制造模块",
        "- 人力资源管理模块",
        "",
        "###技术平台：",
        "- 金蝶云·星空平台",
        "- 云原生架构",
        "- 微服务设计",
        "- 高可用部署"
    ])
    
    # 4A架构设计
    add_section_slide(prs, "四、4A架构设计")
    add_content_slide(prs, "业务架构（BA）", [
        "###核心价值流：",
        "- 订单到收款（O2C）",
        "- 采购到付款（P2P）", 
        "- 计划到生产（P2M）",
        "- 研发到上市（R2M）",
        "",
        "###业务能力：",
        "- 财务管理",
        "- 供应链管理",
        "- 生产制造",
        "- 人力资源管理"
    ])
    
    add_content_slide(prs, "数据架构（DA）", [
        "###数据实体：",
        "- 客户、供应商、物料",
        "- 订单、发票、凭证",
        "- 产品、工艺、设备",
        "",
        "###数据服务：",
        "- 数据查询服务",
        "- 数据分析服务", 
        "- 数据治理服务"
    ])
    
    # 价值工程
    add_section_slide(prs, "五、价值工程分析")
    add_content_slide(prs, "价值目标模型", [
        "###四大价值概念：",
        "1. **Strategy（战略）**：保障战略有效性",
        "2. **Spending（支出回报）**：增收节支",
        "3. **Situation（情景）**：提升竞争力",
        "4. **Structure（结构）**：提升效率",
        "",
        "###预期收益：",
        "- 财务结账时间缩短50%",
        "- 库存周转率提升40%",
        "- 业务效率提升30%"
    ])
    
    # 实施路线
    add_section_slide(prs, "六、实施路线图")
    table_data = [
        ["项目启动", "需求调研", "方案设计", "系统配置", "测试培训", "上线验收"],
        ["2周", "4周", "3周", "4周", "3周", "2周"],
        ["成立项目组", "业务调研", "蓝图设计", "系统配置", "UAT测试", "正式上线"]
    ]
    add_table_slide(prs, "项目实施计划", ["阶段", "周期", "主要工作"], table_data)
    
    # 成功案例
    add_section_slide(prs, "七、成功案例分享")
    add_content_slide(prs, "案例1：制造业企业", [
        "###实施范围：",
        "财务云、供应链云、制造云",
        "",
        "###项目成果：",
        "- 财务结账时间缩短50%",
        "- 库存周转率提升40%", 
        "- 生产计划准确性提升50%"
    ])
    
    # 团队介绍
    add_section_slide(prs, "八、项目团队与保障")
    add_content_slide(prs, "项目团队配置", [
        "###核心团队：",
        "- 项目经理：5年ERP实施经验",
        "- 技术总监：8年技术架构经验",
        "- 业务顾问：10年行业经验",
        "- 实施顾问：6年实施经验",
        "",
        "###服务保障：",
        "- 7×24小时技术支持",
        "- 定期回访服务",
        "- 持续优化升级"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", "金蝶软件（中国）有限公司")
    
    return prs

def generate_golive_ppt_v10(customer_info):
    """生成上线汇报PPT - v10.0 增强版"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    golive_date = customer_info.get('goliveDate', datetime.now().strftime('%Y-%m-%d'))
    
    # 封面
    add_title_slide(prs, f"{company_name}ERP系统上线汇报", "数字化转型重要里程碑", golive_date)
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、项目概况",
        "二、项目背景与目标", 
        "三、项目实施过程",
        "四、项目业务范围",
        "五、项目价值实现",
        "六、项目成果展示",
        "七、项目总结展望"
    ])
    
    # 项目概况
    add_section_slide(prs, "一、项目概况")
    add_content_slide(prs, "项目概述", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"项目周期：8个月",
        f"上线日期：{golive_date}",
        "",
        "###项目意义：",
        "- 企业数字化转型的重要里程碑",
        "- 提升管理效率，降低运营成本",
        "- 增强企业核心竞争力"
    ])
    
    # 项目背景与目标
    add_section_slide(prs, "二、项目背景与目标")
    add_content_slide(prs, "项目建设背景", [
        "###业务痛点：",
        "- 业务流程不标准化",
        "- 数据管理分散",
        "- 系统集成困难",
        "- 管理决策滞后",
        "",
        "###建设目标：",
        "- 实现业务流程标准化",
        "- 建立统一数据平台", 
        "- 提升系统集成效率",
        "- 支持实时决策分析"
    ])
    
    # 项目实施过程
    add_section_slide(prs, "三、项目实施过程")
    timeline_data = [
        ["项目启动", "需求调研", "蓝图设计", "系统配置", "测试培训", "上线支持"],
        ["2024.01", "2024.02-03", "2024.04", "2024.05-07", "2024.08", "2024.09"]
    ]
    add_table_slide(prs, "项目实施历程", ["阶段", "时间节点"], timeline_data)
    
    add_content_slide(prs, "主要工作内容", [
        "###核心工作：",
        "- 业务流程梳理与优化",
        "- 系统配置与开发",
        "- 数据迁移与清洗",
        "- 用户培训与推广",
        "- 系统测试与验证"
    ])
    
    # 项目业务范围
    add_section_slide(prs, "四、项目业务范围")
    add_content_slide(prs, "业务范围架构", [
        "###核心业务模块：",
        "- 财务管理：总账、应收、应付、成本",
        "- 供应链管理：采购、销售、库存",
        "- 生产制造：计划、生产、质量",
        "- 人力资源管理：人事、薪酬、绩效",
        "",
        "###系统集成：",
        "- 与MES系统对接",
        "- 与WMS系统对接", 
        "- 与OA系统对接"
    ])
    
    # 项目价值实现
    add_section_slide(prs, "五、项目价值实现")
    add_content_slide(prs, "价值总览", [
        "###数据治理：",
        "- 主数据统一率提升至100%",
        "- 数据准确性提升30%",
        "",
        "###业务提效：",
        "- 财务结账时间从10天缩短到3天",
        "- 采购周期缩短30%",
        "- 订单交付准时率提升至95%"
    ])
    
    add_content_slide(prs, "管控提升", [
        "###流程优化：",
        "- 实现订单-应收-实收-稽核闭环",
        "- 构建合同-请款-发票-付款管控闭环",
        "",
        "###决策赋能：",
        "- 实时现金流监控",
        "- 经营数据分析",
        "- 管理驾驶舱"
    ])
    
    # 项目成果
    add_section_slide(prs, "六、项目成果展示")
    add_content_slide(prs, "实施成果", [
        "###系统成果：",
        "- 完成核心模块实施",
        "- 实现系统集成对接",
        "- 建立数据治理体系",
        "",
        "###业务成果：",
        "- 业务流程标准化",
        "- 数据管理规范化",
        "- 决策支持智能化"
    ])
    
    # 项目总结
    add_section_slide(prs, "七、项目总结展望")
    add_content_slide(prs, "项目总结", [
        "###成功经验：",
        "- 领导高度重视，全员积极参与",
        "- 需求调研充分，方案设计合理",
        "- 实施过程规范，质量控制严格",
        "",
        "###未来展望：",
        "- 持续优化系统功能",
        "- 深化数据价值挖掘",
        "- 支撑企业战略发展"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", f"金蝶软件（中国）有限公司\\n{datetime.now().strftime('%Y年%m月%d日')}")
    
    return prs

def generate_acceptance_ppt_v10(customer_info):
    """生成验收汇报PPT - v10.0 增强版"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    acceptance_date = customer_info.get('acceptanceDate', datetime.now().strftime('%Y-%m-%d'))
    
    # 封面
    add_title_slide(prs, f"{company_name}ERP系统验收汇报", "项目验收与成果总结", acceptance_date)
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、验收概述",
        "二、验收范围与标准",
        "三、验收过程",
        "四、验收结果",
        "五、项目成果展示",
        "六、遗留问题与计划",
        "七、验收结论"
    ])
    
    # 验收概述
    add_section_slide(prs, "一、验收概述")
    add_content_slide(prs, "验收基本信息", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"验收日期：{acceptance_date}",
        f"验收结论：通过",
        "",
        "###验收意义：",
        "- 项目正式交付使用",
        "- 系统功能确认",
        "- 项目成果验收"
    ])
    
    # 验收范围与标准
    add_section_slide(prs, "二、验收范围与标准")
    add_content_slide(prs, "验收范围", [
        "###核心模块验收：",
        "- 财务管理：总账、应收、应付、成本",
        "- 供应链管理：采购、销售、库存",
        "- 生产制造：计划、生产、质量",
        "- 人力资源管理：人事、薪酬、绩效",
        "",
        "###系统集成验收：",
        "- 与MES系统对接",
        "- 与WMS系统对接",
        "- 与OA系统对接"
    ])
    
    add_content_slide(prs, "验收标准", [
        "###功能标准：",
        "- 完成合同约定功能",
        "- 业务流程运行正常",
        "- 数据准确完整",
        "",
        "###性能标准：",
        "- 系统响应时间<3秒",
        "- 并发用户数≥100",
        "- 数据处理能力满足需求"
    ])
    
    # 验收过程
    add_section_slide(prs, "三、验收过程")
    add_content_slide(prs, "验收流程", [
        "###验收准备：",
        "- 验收方案制定",
        "- 测试数据准备",
        "- 验收团队组建",
        "",
        "###验收执行：",
        "- 功能测试",
        "- 性能测试",
        "- 用户验收",
        "- 专家评审"
    ])
    
    # 验收结果
    add_section_slide(prs, "四、验收结果")
    result_data = [
        ["验收项目", "结果", "说明"],
        ["功能完整性", "通过", "所有合同约定功能已完成"],
        ["业务流程", "通过", "主要业务流程运行正常"],
        ["数据准确性", "通过", "数据准确率≥99%"],
        ["系统性能", "通过", "响应时间满足要求"],
        ["用户满意度", "通过", "用户满意度≥90%"]
    ]
    add_table_slide(prs, "验收结果汇总", ["验收项目", "结果", "说明"], result_data)
    
    # 项目成果
    add_section_slide(prs, "五、项目成果展示")
    add_content_slide(prs, "项目成果", [
        "###系统成果：",
        "- 完成ERP系统实施",
        "- 实现业务流程标准化",
        "- 建立数据治理体系",
        "",
        "###业务成果：",
        "- 管理效率提升50%",
        "- 运营成本降低20%",
        "- 决策支持能力增强"
    ])
    
    # 遗留问题与计划
    add_section_slide(prs, "六、遗留问题与计划")
    add_content_slide(prs, "遗留问题", [
        "###当前问题：",
        "- 部分细节功能需要优化",
        "- 用户习惯需要培养",
        "- 数据质量需要持续监控",
        "",
        "###后续计划：",
        "- 系统功能持续优化",
        "- 用户培训推广",
        "- 数据治理深化"
    ])
    
    # 验收结论
    add_section_slide(prs, "七、验收结论")
    add_content_slide(prs, "验收结论", [
        "###验收结论：",
        "- 项目达到合同约定要求",
        "- 系统功能完整，运行稳定",
        "- 业务流程优化效果显著",
        "- 用户满意度较高",
        "",
        "###后续建议：",
        "- 持续优化系统功能",
        "- 加强用户培训",
        "- 深化数据价值挖掘"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", f"金蝶软件（中国）有限公司\\n{datetime.now().strftime('%Y年%m月%d日')}")
    
    return prs

def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 3:
        print("用法: python3 kingdee-ppt-v10-enhanced.py --type <type> --companyName <name> [其他参数]")
        sys.exit(1)
    
    # 解析参数
    params = {}
    for i in range(1, len(sys.argv), 2):
        if i + 1 < len(sys.argv):
            key = sys.argv[i].lstrip('-')
            params[key] = sys.argv[i + 1]
    
    # 生成PPT
    ppt_type = params.get('type', 'presales')
    company_name = params.get('companyName', '企业名称')
    
    try:
        if ppt_type == 'presales':
            prs = generate_presales_ppt_v10(params)
        elif ppt_type == 'golive':
            prs = generate_golive_ppt_v10(params)
        elif ppt_type == 'acceptance':
            prs = generate_acceptance_ppt_v10(params)
        else:
            print(f"不支持的PPT类型: {ppt_type}")
            sys.exit(1)
        
        # 保存文件
        filename = f"{company_name}_{ppt_type}_v10.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        print(json.dumps({
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "slides": len(prs.slides),
            "type": ppt_type
        }, ensure_ascii=False))
        
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": str(e)
        }, ensure_ascii=False))

if __name__ == "__main__":
    main()