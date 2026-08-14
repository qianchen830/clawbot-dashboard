#!/usr/bin/env python3
"""
验证生成的PPT
"""
from pptx import Presentation

def verify_ppt(ppt_path):
    """验证PPT结构"""
    print(f"验证PPT: {ppt_path}")
    print("=" * 80)
    
    prs = Presentation(ppt_path)
    
    print(f"\n总页数: {len(prs.slides)}")
    print(f"幻灯片尺寸: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    
    # 统计布局使用
    layout_stats = {}
    for slide in prs.slides:
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        layout_stats[layout_name] = layout_stats.get(layout_name, 0) + 1
    
    print("\n布局使用统计:")
    for layout, count in sorted(layout_stats.items(), key=lambda x: x[1], reverse=True):
        print(f"  {layout}: {count}页")
    
    # 显示前20页的标题
    print("\n前20页标题:")
    for idx in range(min(20, len(prs.slides))):
        slide = prs.slides[idx]
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title = text[:60]
                    break
        print(f"  [{idx+1:3d}] {slide.slide_layout.name:15s} - {title}")
    
    return len(prs.slides)

if __name__ == "__main__":
    ppt_path = "/home/openclaw/.openclaw/workspace/output/kingdee-presales-ppt-v18.pptx"
    verify_ppt(ppt_path)
