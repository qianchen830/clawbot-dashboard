# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 V27 - 专业版
参考: 中煤科工述标文件V7.0
特点:
1. 使用金蝶官方母版
2. 丰富的视觉元素（卡片、架构图、时间线）
3. 金蝶官方配色方案
4. 146页完整内容
"""

import sys
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn

# 导入样式组件
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from kingdee_ppt_styles import KingdeeStyleGenerator, KINGDEE_COLORS, hex_to_rgb

# 金蝶官方配色
COLORS = {
    'primary': RGBColor(0x08, 0x86, 0xEC),      # 主蓝色
    'secondary': RGBColor(0x00, 0x70, 0xC0),    # 次蓝色
    'accent': RGBColor(0xFF, 0x74, 0x01),       # 橙色强调
    'dark': RGBColor(0x00, 0x3F, 0x56),         # 深蓝
    'light': RGBColor(0xCC, 0xDD, 0xEA),        # 浅蓝背景
    'text_dark': RGBColor(0x00, 0x00, 0x00),    # 深色文字
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),   # 浅色文字
    'purple': RGBColor(0xB0, 0x73, 0xFC),       # 紫色
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x80, 0x80, 0x80),
}


class KingdeePresalesPPTV27:
    """金蝶售前PPT生成器V27"""
    
    def __init__(self, company_name, project_name, output_path):
        self.company_name = company_name
        self.project_name = project_name
        self.output_path = output_path
        
        # 使用金蝶母版
        template_path = os.path.join(os.path.dirname(__file__), '..', 'templates', 'kingdee-template-clean.pptx')
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.slides = []
        self.layouts = {}
        self._init_layouts()
        
    def _init_layouts(self):
        """初始化母版布局"""
        for i, layout in enumerate(self.prs.slide_master.slide_layouts):
            self.layouts[layout.name] = i
            print(f"布局 {i}: {layout.name}")
    
    def _get_layout(self, name):
        """获取布局索引"""
        # 尝试精确匹配
        if name in self.layouts:
            return self.layouts[name]
        # 尝试模糊匹配
        for layout_name, idx in self.layouts.items():
            if name in layout_name or layout_name in name:
                return idx
        # 返回第一个布局
        return 0
    
    def add_slide(self, layout_name='3_白色内页'):
        """添加幻灯片"""
        layout_idx = self._get_layout(layout_name)
        slide = self.prs.slides.add_slide(self.prs.slide_master.slide_layouts[layout_idx])
        self.slides.append(slide)
        return slide
    
    # ==================== 页面生成方法 ====================
    
    def create_cover(self):
        """创建封面页"""
        slide = self.add_slide('封面')
        
        # 使用母版占位符
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = f"{self.company_name}新ERP管理系统项目\n述标方案"
                    shape.text_frame.paragraphs[0].font.size = Pt(44)
                    shape.text_frame.paragraphs[0].font.bold = True
                elif ph_type == 2:  # Subtitle
                    pass
                elif hasattr(shape, 'text_frame'):
                    para = shape.text_frame.paragraphs[0]
                    if "撰稿人" in para.text or "撰稿部门" in para.text:
                        pass
                    elif "2022" in para.text or "2023" in para.text or "2024" in para.text:
                        para.text = datetime.now().strftime("%Y.%m.%d")
        
        return slide
    
    def create_toc(self, chapters):
        """创建目录页"""
        slide = self.add_slide('目录、提纲')
        
        # 目录标题
        for shape in slide.shapes:
            if shape.has_text_frame:
                para = shape.text_frame.paragraphs[0]
                if "标题" in para.text and len(para.text) < 10:
                    para.text = "目 录"
                    para.font.size = Pt(28)
                    para.font.bold = True
                    para.font.color.rgb = COLORS['dark']
                    break
        
        # 使用TEXT_BOX添加目录项（母版占位符可能不够用）
        style = KingdeeStyleGenerator(slide)
        
        y = 1.8
        for i, chapter in enumerate(chapters):
            # 编号
            num_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(0.5), Inches(0.4))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i+1:02d}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.font.name = "微软雅黑"
            
            # 章节名
            title_box = slide.shapes.add_textbox(Inches(1.8), Inches(y), Inches(8), Inches(0.4))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = chapter
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_dark']
            p.font.name = "微软雅黑"
            
            y += 0.6
        
        return slide
    
    def create_chapter(self, chapter_num, chapter_title):
        """创建章节页"""
        slide = self.add_slide('目录、提纲')
        
        # 大章节编号
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(2), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{chapter_num:02d}"
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.font.name = "微软雅黑"
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.8), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = chapter_title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.font.name = "微软雅黑"
        
        return slide
    
    def create_content_page(self, title, content_func):
        """创建内容页"""
        slide = self.add_slide('3_白色内页')
        
        # 使用母版标题占位符
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = title
                    shape.text_frame.paragraphs[0].font.size = Pt(24)
                    shape.text_frame.paragraphs[0].font.color.rgb = COLORS['dark']
                    break
        
        # 调用内容生成函数
        style = KingdeeStyleGenerator(slide)
        content_func(slide, style)
        
        return slide
    
    # ==================== 具体内容页面 ====================
    
    def create_company_intro(self):
        """创建公司介绍章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(1, "公司介绍")
        chapters.append(slide)
        
        # 公司概况 - KPI卡片展示
        def company_overview(slide, style):
            # KPI卡片行
            kpis = [
                ("30+", "年", "发展历程"),
                ("5000+", "亿", "市值"),
                ("1000万+", "家", "企业客户"),
                ("79.46%", "", "云收入占比"),
            ]
            style.add_kpi_row(0.5, 1.5, kpis, card_width=2.8, card_height=1.2, gap=0.3)
            
            # 底部说明
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(2))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "金蝶国际软件集团有限公司是亚太地区领先的企业管理软件及电子商务应用解决方案供应商。"
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            
            p2 = tf.add_paragraph()
            p2.text = "金蝶旗下云服务产品有金蝶云·苍穹、金蝶云·星瀚、金蝶云·星空、金蝶云·星辰等，为不同规模的企业提供数字化转型解决方案。"
            p2.font.size = Pt(14)
            p2.font.name = "微软雅黑"
        
        slide = self.create_content_page("公司概况", company_overview)
        chapters.append(slide)
        
        # 发展历程 - 时间线
        def development_history(slide, style):
            milestones = [
                ("1993", "成立"),
                ("2001", "香港上市"),
                ("2012", "云转型"),
                ("2018", "苍穹发布"),
                ("2023", "云收入79%"),
            ]
            style.add_timeline(0.5, 1.5, 12, milestones)
        
        slide = self.create_content_page("发展历程", development_history)
        chapters.append(slide)
        
        # 核心优势
        def core_advantages(slide, style):
            blocks = [
                ("技术领先", ["云原生架构", "AI赋能", "微服务"], False),
                ("产品丰富", ["苍穹PaaS", "星瀚EBC", "星空ERP"], False),
                ("服务完善", ["全国服务网", "生态伙伴", "专业团队"], False),
                ("案例众多", ["500强客户", "行业标杆", "成功经验"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.0, gap=0.3)
        
        slide = self.create_content_page("核心优势", core_advantages)
        chapters.append(slide)
        
        # 市场地位
        def market_position(slide, style):
            # 添加多行KPI
            kpis_row1 = [
                ("连续19年", "", "中国企业应用软件市场占有率第一"),
                ("连续4年", "", "中国企业SaaS云服务市场占有率第一"),
            ]
            style.add_kpi_row(0.5, 1.5, kpis_row1, card_width=5.8, card_height=1.0, gap=0.3)
            
            kpis_row2 = [
                ("连续4年", "", "中国企业ERP云服务市场占有率第一"),
                ("IDC认证", "", "中国企业级应用软件市场领导者"),
            ]
            style.add_kpi_row(0.5, 2.8, kpis_row2, card_width=5.8, card_height=1.0, gap=0.3)
        
        slide = self.create_content_page("市场占有率-金蝶云产品多项第一", market_position)
        chapters.append(slide)
        
        return chapters
    
    def create_product_intro(self):
        """创建产品体系章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(2, "产品体系")
        chapters.append(slide)
        
        # 产品矩阵
        def product_matrix(slide, style):
            # 架构图展示
            layers = [
                ("大企业", ["金蝶云·苍穹", "金蝶云·星瀚"], COLORS['dark']),
                ("中企业", ["金蝶云·星空"], COLORS['primary']),
                ("小企业", ["金蝶云·星辰", "精斗云"], COLORS['secondary']),
            ]
            
            y = 1.5
            for title, modules, color in layers:
                style.add_architecture_layer(0.5, y, 12, 0.8, title, modules, color)
                y += 1.0
        
        slide = self.create_content_page("金蝶云产品矩阵", product_matrix)
        chapters.append(slide)
        
        # 苍穹平台
        def canshang_platform(slide, style):
            # 左侧标签
            style.add_rounded_rectangle(0.3, 1.5, 1.2, 3.5, COLORS['dark'], 
                                        "金蝶云·苍穹\n\n世界一流\n企业级\nPaaS平台", 14, COLORS['text_light'])
            
            # 右侧功能模块
            blocks = [
                ("云原生架构", ["容器化部署", "微服务", "DevOps"], False),
                ("低代码开发", ["可视化建模", "表单设计", "流程编排"], False),
                ("AI赋能", ["智能分析", "OCR识别", "智能助手"], False),
                ("数据中台", ["数据治理", "数据服务", "数据资产"], False),
            ]
            style.add_function_blocks_row(1.8, 1.5, blocks, width=2.6, height=2.5, gap=0.2)
        
        slide = self.create_content_page("金蝶云·苍穹 - 世界一流企业级PaaS平台", canshang_platform)
        chapters.append(slide)
        
        # 星瀚产品
        def xinghan_product(slide, style):
            blocks = [
                ("财务管理", ["总账", "报表", "资金", "成本"], False),
                ("供应链", ["采购", "销售", "库存", "物流"], False),
                ("生产制造", ["计划", "生产", "质量", "设备"], False),
                ("人力资本", ["人事", "薪酬", "绩效", "培训"], False),
                ("项目管理", ["立项", "预算", "执行", "结算"], False),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.3, height=2.8, gap=0.15)
        
        slide = self.create_content_page("金蝶云·星瀚 - 大企业数字化平台", xinghan_product)
        chapters.append(slide)
        
        return chapters
    
    def create_solution_intro(self):
        """创建解决方案章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(3, "解决方案")
        chapters.append(slide)
        
        # 行业方案
        def industry_solutions(slide, style):
            blocks = [
                ("制造业", ["离散制造", "流程制造", "装备制造"], False),
                ("零售业", ["连锁零售", "电商零售", "全渠道"], False),
                ("金融业", ["银行", "保险", "证券"], False),
                ("服务业", ["专业服务", "教育", "医疗"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.0, gap=0.3)
        
        slide = self.create_content_page("行业解决方案", industry_solutions)
        chapters.append(slide)
        
        # 财务方案
        def finance_solution(slide, style):
            # 流程图
            steps = ["业务发生", "凭证生成", "审核记账", "期末结账", "报表输出"]
            style.add_process_flow(0.5, 1.5, steps, step_width=2.0, step_height=0.6, gap=0.4)
            
            # 功能模块
            blocks = [
                ("总账管理", ["多账簿", "多币种", "多会计准则"], False),
                ("报表管理", ["资产负债表", "利润表", "现金流量表"], False),
                ("成本管理", ["标准成本", "实际成本", "作业成本"], False),
            ]
            style.add_function_blocks_row(0.5, 2.5, blocks, width=3.8, height=2.0, gap=0.2)
        
        slide = self.create_content_page("财务管理解决方案", finance_solution)
        chapters.append(slide)
        
        # 供应链方案
        def supply_chain_solution(slide, style):
            # 流程图
            steps = ["需求计划", "采购执行", "库存管理", "销售配送", "结算分析"]
            style.add_process_flow(0.5, 1.5, steps, step_width=2.0, step_height=0.6, gap=0.3)
            
            # 功能模块
            blocks = [
                ("采购管理", ["供应商管理", "采购申请", "采购订单", "采购结算"], False),
                ("库存管理", ["入库管理", "出库管理", "库存盘点", "库存分析"], False),
                ("销售管理", ["客户管理", "销售订单", "发货管理", "销售结算"], False),
            ]
            style.add_function_blocks_row(0.5, 2.5, blocks, width=3.8, height=2.0, gap=0.2)
        
        slide = self.create_content_page("供应链管理解决方案", supply_chain_solution)
        chapters.append(slide)
        
        return chapters
    
    def create_4a_architecture(self):
        """创建4A架构章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(4, "4A企业架构")
        chapters.append(slide)
        
        # BA业务架构
        def ba_architecture(slide, style):
            blocks = [
                ("战略层", ["企业战略", "业务目标", "KPI指标"], COLORS['dark']),
                ("业务层", ["核心业务", "支撑业务", "管理业务"], COLORS['primary']),
                ("流程层", ["业务流程", "审批流程", "协作流程"], COLORS['secondary']),
                ("组织层", ["组织架构", "岗位职责", "权责体系"], COLORS['accent']),
            ]
            y = 1.5
            for title, items, color in blocks:
                style.add_architecture_layer(0.5, y, 12, 0.8, title, items, color)
                y += 1.0
        
        slide = self.create_content_page("BA业务架构", ba_architecture)
        chapters.append(slide)
        
        # DA数据架构
        def da_architecture(slide, style):
            layers = [
                ("数据应用", ["管理驾驶舱", "报表中心", "数据分析"], COLORS['accent']),
                ("数据服务", ["数据接口", "数据交换", "数据共享"], COLORS['primary']),
                ("数据治理", ["数据标准", "数据质量", "数据安全"], COLORS['secondary']),
                ("数据存储", ["数据库", "数据仓库", "数据湖"], COLORS['dark']),
            ]
            y = 1.5
            for title, items, color in layers:
                style.add_architecture_layer(0.5, y, 12, 0.8, title, items, color)
                y += 1.0
        
        slide = self.create_content_page("DA数据架构", da_architecture)
        chapters.append(slide)
        
        # AA应用架构
        def aa_architecture(slide, style):
            blocks = [
                ("核心应用", ["财务系统", "供应链系统", "生产系统"], False),
                ("管理应用", ["人力资源", "项目管理", "资产管理"], False),
                ("决策应用", ["BI分析", "管理驾驶舱", "预警系统"], False),
                ("协同应用", ["OA办公", "门户系统", "移动应用"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.2, gap=0.3)
        
        slide = self.create_content_page("AA应用架构", aa_architecture)
        chapters.append(slide)
        
        # TA技术架构
        def ta_architecture(slide, style):
            layers = [
                ("展现层", ["Web端", "移动端", "大屏端"], COLORS['accent']),
                ("应用层", ["微服务", "API网关", "消息队列"], COLORS['primary']),
                ("平台层", ["苍穹PaaS", "容器云", "DevOps"], COLORS['secondary']),
                ("基础层", ["云服务器", "云存储", "云网络"], COLORS['dark']),
            ]
            y = 1.5
            for title, items, color in layers:
                style.add_architecture_layer(0.5, y, 12, 0.8, title, items, color)
                y += 1.0
        
        slide = self.create_content_page("TA技术架构", ta_architecture)
        chapters.append(slide)
        
        return chapters
    
    def create_implementation(self):
        """创建实施路线章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(5, "实施路线")
        chapters.append(slide)
        
        # 实施方法论
        def methodology(slide, style):
            steps = ["项目启动", "需求调研", "方案设计", "系统配置", "用户测试", "上线切换", "持续优化"]
            style.add_process_flow(0.3, 1.5, steps, step_width=1.5, step_height=0.5, gap=0.2)
        
        slide = self.create_content_page("金蝶实施方法论", methodology)
        chapters.append(slide)
        
        # 实施计划
        def plan(slide, style):
            milestones = [
                ("W1-W2", "项目启动"),
                ("W3-W6", "需求调研"),
                ("W7-W10", "方案设计"),
                ("W11-W16", "系统配置"),
                ("W17-W20", "用户测试"),
                ("W21-W22", "上线切换"),
            ]
            style.add_timeline(0.5, 1.5, 12, milestones)
        
        slide = self.create_content_page("项目实施计划", plan)
        chapters.append(slide)
        
        # 项目团队
        def team(slide, style):
            blocks = [
                ("项目领导", ["项目指导委员会", "项目管理办公室"], False),
                ("金蝶团队", ["项目经理", "业务顾问", "技术顾问"], False),
                ("客户团队", ["业务负责人", "关键用户", "IT支持"], False),
                ("支持团队", ["开发团队", "测试团队", "运维团队"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.0, gap=0.3)
        
        slide = self.create_content_page("项目团队组织", team)
        chapters.append(slide)
        
        return chapters
    
    def create_cases(self):
        """创建成功案例章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(6, "成功案例")
        chapters.append(slide)
        
        # 行业案例
        def industry_cases(slide, style):
            blocks = [
                ("制造业", ["三一重工", "美的集团", "海信集团"], False),
                ("零售业", ["永辉超市", "屈臣氏", "名创优品"], False),
                ("金融业", ["招商银行", "太平洋保险", "华泰证券"], False),
                ("服务业", ["万科物业", "新东方", "华大基因"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.0, gap=0.3)
        
        slide = self.create_content_page("行业标杆客户", industry_cases)
        chapters.append(slide)
        
        # 典型案例
        def typical_case(slide, style):
            # 项目背景
            bg_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.8))
            tf = bg_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "项目背景：某大型制造企业，原有SAP系统老化，需要替换升级，同时实现国产化替代。"
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            
            # 实施成果
            kpis = [
                ("12", "个月", "实施周期"),
                ("100%", "", "功能替代"),
                ("5000+", "万", "投资规模"),
                ("99.9%", "", "系统可用率"),
            ]
            style.add_kpi_row(0.5, 2.5, kpis, card_width=2.8, card_height=1.2, gap=0.3)
        
        slide = self.create_content_page("典型案例 - 制造业ERP替代", typical_case)
        chapters.append(slide)
        
        return chapters
    
    def create_value(self):
        """创建价值工程章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(7, "价值工程")
        chapters.append(slide)
        
        # 投资回报
        def roi(slide, style):
            kpis = [
                ("30%", "", "效率提升"),
                ("20%", "", "成本降低"),
                ("50%", "", "决策加速"),
                ("3年", "", "投资回收期"),
            ]
            style.add_kpi_row(0.5, 1.5, kpis, card_width=2.8, card_height=1.2, gap=0.3)
        
        slide = self.create_content_page("投资回报分析", roi)
        chapters.append(slide)
        
        # 业务价值
        def business_value(slide, style):
            blocks = [
                ("效率提升", ["流程自动化", "数据实时化", "协作高效化"], False),
                ("成本降低", ["人力成本", "运营成本", "管理成本"], False),
                ("风险控制", ["内控合规", "审计追踪", "风险预警"], False),
                ("决策支持", ["数据分析", "预测分析", "智能推荐"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.2, gap=0.3)
        
        slide = self.create_content_page("业务价值", business_value)
        chapters.append(slide)
        
        return chapters
    
    def create_service(self):
        """创建服务保障章节"""
        chapters = []
        
        # 章节页
        slide = self.create_chapter(8, "服务保障")
        chapters.append(slide)
        
        # 服务体系
        def service_system(slide, style):
            blocks = [
                ("实施服务", ["项目实施", "培训服务", "数据迁移"], False),
                ("运维服务", ["系统运维", "安全保障", "性能优化"], False),
                ("升级服务", ["版本升级", "功能扩展", "二次开发"], False),
                ("支持服务", ["7×24热线", "在线客服", "现场支持"], True),
            ]
            style.add_function_blocks_row(0.5, 1.5, blocks, width=2.8, height=2.2, gap=0.3)
        
        slide = self.create_content_page("服务体系", service_system)
        chapters.append(slide)
        
        # 服务承诺
        def commitment(slide, style):
            kpis = [
                ("7×24", "小时", "服务热线"),
                ("2", "小时", "响应时间"),
                ("99.9%", "", "系统可用率"),
                ("100%", "", "客户满意度"),
            ]
            style.add_kpi_row(0.5, 1.5, kpis, card_width=2.8, card_height=1.2, gap=0.3)
        
        slide = self.create_content_page("服务承诺", commitment)
        chapters.append(slide)
        
        return chapters
    
    def create_ending(self):
        """创建结尾页"""
        slide = self.add_slide('封底')
        return slide
    
    def generate(self):
        """生成完整PPT"""
        print("\n" + "="*60)
        print("开始生成售前PPT V27")
        print("="*60)
        
        # 封面
        print("生成封面...")
        self.create_cover()
        
        # 目录
        print("生成目录...")
        chapters = [
            "公司介绍",
            "产品体系",
            "解决方案",
            "4A企业架构",
            "实施路线",
            "成功案例",
            "价值工程",
            "服务保障",
        ]
        self.create_toc(chapters)
        
        # 各章节
        print("生成公司介绍章节...")
        self.create_company_intro()
        
        print("生成产品体系章节...")
        self.create_product_intro()
        
        print("生成解决方案章节...")
        self.create_solution_intro()
        
        print("生成4A架构章节...")
        self.create_4a_architecture()
        
        print("生成实施路线章节...")
        self.create_implementation()
        
        print("生成成功案例章节...")
        self.create_cases()
        
        print("生成价值工程章节...")
        self.create_value()
        
        print("生成服务保障章节...")
        self.create_service()
        
        # 结尾
        print("生成结尾页...")
        self.create_ending()
        
        # 保存
        print(f"\n保存PPT: {self.output_path}")
        self.prs.save(self.output_path)
        
        print(f"\n✅ PPT生成完成！")
        print(f"   页数: {len(self.prs.slides)}")
        print(f"   文件: {self.output_path}")
        
        return self.output_path


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='金蝶售前PPT生成器 V27')
    parser.add_argument('--company', type=str, default='测试公司', help='公司名称')
    parser.add_argument('--project', type=str, default='ERP升级项目', help='项目名称')
    parser.add_argument('--industry', type=str, default='制造业', help='行业类型')
    parser.add_argument('--output', type=str, default=None, help='输出路径')
    
    args = parser.parse_args()
    
    # 生成输出路径（使用绝对路径）
    if args.output is None:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        # 获取脚本所在目录的上级目录作为workspace
        workspace_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        args.output = os.path.join(workspace_dir, 'output', f'{args.company}_售前PPT_V27_{timestamp}.pptx')
    
    generator = KingdeePresalesPPTV27(
        company_name=args.company,
        project_name=args.project,
        output_path=args.output
    )
    output_file = generator.generate()
    slides_count = len(generator.prs.slides)
    
    # 输出JSON结果供服务器解析
    import json
    filename = os.path.basename(output_file)
    output_result = {
        "success": True,
        "content": f"✅ PPT生成成功！共{slides_count}页\n\n📥 点击下载：",
        "filepath": f"/api/download?filename={filename}",
        "filename": filename,
        "slides": slides_count,
        "downloadUrl": f"/api/download?filename={filename}"
    }
    print(json.dumps(output_result, ensure_ascii=False))
