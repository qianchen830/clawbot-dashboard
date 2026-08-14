#!/usr/bin/env python3
"""
检查母版模板是否有现有幻灯片
"""
from pptx import Presentation

def check_master_template(ppt_path):
    """检查母版模板"""
    print(f"检查母版模板: {ppt_path}")
    print("=" * 80)
    
    prs = Presentation(ppt_path)
    
    print(f"\n总页数: {len(prs.slides)}")
    
    if len(prs.slides) > 0:
        print("\n警告: 母版模板包含现有幻灯片!")
        print("前10页标题:")
        for idx in range(min(10, len(prs.slides))):
            slide = prs.slides[idx]
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        title = text[:60]
                        break
            print(f"  [{idx+1:3d}] {slide.slide_layout.name:15s} - {title}")
    else:
        print("\n母版模板为空，适合作为模板使用。")

if __name__ == "__main__":
    master_path = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/ppt模板/ppt母版.pptx"
    check_master_template(master_path)
