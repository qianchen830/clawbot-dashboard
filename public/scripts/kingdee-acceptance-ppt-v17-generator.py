#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶验收汇报PPT生成器 v17.0 - 继承基础类+图表增强版
"""

import os, sys, json, argparse, importlib.util
from datetime import datetime
from pptx.enum.chart import XL_CHART_TYPE

base_path = os.path.join(os.path.dirname(__file__), 'kingdee-ppt-base.py')
spec = importlib.util.spec_from_file_location("kingdee_ppt_base", base_path)
kingdee_ppt_base = importlib.util.module_from_spec(spec)
spec.loader.exec_module(kingdee_ppt_base)

KingdeePPTBase = kingdee_ppt_base.KingdeePPTBase
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


class AcceptancePPTV17(KingdeePPTBase):
    """金蝶验收汇报PPT生成器 v17"""
    
    def __init__(self, company, project_name="", accept_date=None, modules=None):
        super().__init__(company, "验收汇报")
        self.project = project_name or f"{company}ERP系统"
        self.modules = modules or ["财务云", "供应链云", "制造云"]
        self.accept_date = accept_date or datetime.now()
        self.accept_str = self.accept_date.strftime('%Y年%m月%d日') if isinstance(self.accept_date, datetime) else str(self.accept_date)
    
    def generate(self):
        print(f"\n{'='*60}")
        print(f"生成 {self.company} 验收汇报PPT V17（图表增强版）")
        print(f"{'='*60}")
        
        # 封面
        self.add_cover(f"{self.company}\n{self.project}验收汇报", "验收汇报")
        
        # 目录
        self.add_toc([
            "项目概况", "实施成果", "系统功能", "测试报告",
            "培训成果", "遗留问题", "后续计划", "验收结论"
        ])
        
        # 第一章：项目概况
        self.add_section("项目概况", "一")
        
        self.add_cards("项目基本信息", [
            {"value": self.company, "label": "客户名称"},
            {"value": self.project, "label": "项目名称"},
            {"value": self.accept_str, "label": "验收日期"},
            {"value": "6个月", "label": "实施周期"},
            {"value": "100%", "label": "需求完成率"},
            {"value": "98%", "label": "测试通过率"},
        ])
        
        self.add_timeline("项目历程", [
            {"time": "启动", "content": "项目启动\n团队组建"},
            {"time": "调研", "content": "需求调研\n蓝图设计"},
            {"time": "实施", "content": "系统配置\n数据迁移"},
            {"time": "测试", "content": "UAT测试\n问题修复"},
            {"time": "验收", "content": "系统验收\n项目交付"},
        ])
        
        # 第二章：实施成果
        self.add_section("实施成果", "二")
        
        self.add_chart_with_cards("实施成果", XL_CHART_TYPE.COLUMN_CLUSTERED, {
            'categories': ['需求完成', '配置完成', '测试通过', '培训覆盖', '文档交付'],
            'series': {'完成度': (95, 100, 98, 100, 100)}
        }, [
            {"value": "95%", "label": "需求完成率"},
            {"value": "100%", "label": "配置完成率"},
            {"value": "98%", "label": "测试通过率"},
            {"value": "100%", "label": "培训覆盖率"},
        ])
        
        self.add_icon_grid("实施内容", [
            {"title": "财务管理", "desc": "总账/应收/应付\n资金/成本/资产"},
            {"title": "供应链管理", "desc": "采购/销售/库存\n物流管理"},
            {"title": "生产制造", "desc": "计划/执行/质量\n设备管理"},
            {"title": "人力资源", "desc": "人事/薪酬/绩效\n培训管理"},
        ])
        
        # 第三章：系统功能
        self.add_section("系统功能", "三")
        
        self.add_chart("功能模块分布", XL_CHART_TYPE.PIE, {
            'categories': ['财务管理', '供应链', '生产制造', '人力资源', '其他'],
            'series': {'占比': (30, 25, 25, 15, 5)}
        })
        
        self.add_table("功能清单", [
            ["模块", "功能数", "完成数", "完成率"],
            ["财务管理", "45", "45", "100%"],
            ["供应链管理", "38", "38", "100%"],
            ["生产制造", "32", "30", "94%"],
            ["人力资源", "20", "20", "100%"],
        ])
        
        # 第四章：测试报告
        self.add_section("测试报告", "四")
        
        self.add_cards("测试统计", [
            {"value": "500+", "label": "测试用例"},
            {"value": "490", "label": "通过用例"},
            {"value": "10", "label": "修复问题"},
            {"value": "98%", "label": "通过率"},
            {"value": "5天", "label": "测试周期"},
            {"value": "100%", "label": "问题修复率"},
        ])
        
        self.add_chart("测试结果分布", XL_CHART_TYPE.PIE, {
            'categories': ['通过', '修复后通过', '跳过'],
            'series': {'数量': (480, 10, 10)}
        })
        
        # 第五章：培训成果
        self.add_section("培训成果", "五")
        
        self.add_chart_with_cards("培训统计", XL_CHART_TYPE.COLUMN_CLUSTERED, {
            'categories': ['管理层', '财务人员', '供应链人员', 'IT人员', '最终用户'],
            'series': {'培训人数': (15, 30, 25, 10, 80)}
        }, [
            {"value": "160人", "label": "培训总人数"},
            {"value": "100%", "label": "培训覆盖率"},
            {"value": "95%", "label": "考核通过率"},
        ])
        
        self.add_table("培训记录", [
            ["培训对象", "人数", "时长", "考核结果"],
            ["管理层", "15", "1天", "全部通过"],
            ["财务人员", "30", "2天", "全部通过"],
            ["供应链人员", "25", "2天", "全部通过"],
            ["IT人员", "10", "1天", "全部通过"],
            ["最终用户", "80", "1天", "76通过"],
        ])
        
        # 第六章：遗留问题
        self.add_section("遗留问题", "六")
        
        self.add_comparison("遗留问题处理", 
            "遗留问题", [
                "报表格式优化（非关键）",
                "历史数据查询（非关键）",
                "移动端审批（规划中）",
            ],
            "处理方案", [
                "下一版本迭代优化",
                "数据归档方案实施",
                "二期项目实施",
            ])
        
        self.add_cards("遗留问题统计", [
            {"value": "3个", "label": "遗留问题"},
            {"value": "0个", "label": "关键问题"},
            {"value": "100%", "label": "关键问题解决"},
            {"value": "规划中", "label": "处理状态"},
        ])
        
        # 第七章：后续计划
        self.add_section("后续计划", "七")
        
        self.add_timeline("后续工作安排", [
            {"time": "验收后1周", "content": "系统运维移交\n知识转移"},
            {"time": "验收后1月", "content": "优化需求收集\n系统调优"},
            {"time": "验收后3月", "content": "二期需求确认\n项目规划"},
            {"time": "验收后6月", "content": "二期项目启动\n持续优化"},
        ])
        
        self.add_table("服务承诺", [
            ["服务内容", "服务期限", "服务方式", "响应时间"],
            ["系统运维支持", "1年", "远程+现场", "2小时"],
            ["问题处理", "1年", "热线+远程", "4小时"],
            ["版本升级", "1年", "远程", "按计划"],
            ["培训支持", "1年", "现场", "按需"],
        ])
        
        # 第八章：验收结论
        self.add_section("验收结论", "八")
        
        self.add_cards("验收结论", [
            {"value": "通过", "label": "验收结果"},
            {"value": "100%", "label": "需求完成率"},
            {"value": "98%", "label": "测试通过率"},
            {"value": "100%", "label": "文档交付率"},
            {"value": "优秀", "label": "项目评价"},
            {"value": "同意验收", "label": "验收意见"},
        ])
        
        # 感谢页
        self.add_thanks()
        
        return self.page
    
    def save(self, filename=None):
        if not filename:
            filename = f"{self.company}_验收汇报PPT_v17.pptx"
        return super().save(filename)


def main():
    parser = argparse.ArgumentParser(description='金蝶验收汇报PPT生成器 V17')
    parser.add_argument('--company', required=True)
    parser.add_argument('--project', default='')
    parser.add_argument('--accept-date', default=None)
    args = parser.parse_args()
    
    gen = AcceptancePPTV17(args.company, args.project, args.accept_date)
    gen.generate()
    output_path = gen.save()
    
    print("\n" + json.dumps({
        "success": True, "output": output_path, "slides": gen.page,
        "size": os.path.getsize(output_path)
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
