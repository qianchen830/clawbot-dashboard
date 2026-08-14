#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - 图表模块 v1.0
提供数据可视化图表生成功能
支持饼图、柱状图、时间线图、组织架构图
"""

import os
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# 导入公共模块
from kingdee_ppt_common import COLORS, create_presentation


# 图表配色方案
CHART_COLORS = [
    RGBColor(0, 82, 147),    # 金蝶蓝
    RGBColor(0, 112, 192),   # 辅助蓝
    RGBColor(255, 153, 0),   # 强调橙
    RGBColor(0, 176, 80),    # 成功绿
    RGBColor(255, 192, 0),   # 警告黄
    RGBColor(112, 48, 160),  # 紫色
    RGBColor(192, 0, 0),     # 红色
    RGBColor(0, 128, 128),   # 青色
]


def add_pie_chart_slide(prs: Presentation, title: str, 
                        data: Dict[str, float],
                        chart_title: str = "") -> Any:
    """
    添加饼图幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        data: 数据字典，键为标签，值为数值
        chart_title: 图表标题
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = list(data.keys())
    chart_data.add_series(chart_title or '数据', list(data.values()))
    
    # 添加饼图
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(6), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    # 设置图表样式
    chart.has_legend = True
    chart.legend.include_in_layout = False
    
    # 添加数据标签
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = False
    
    # 添加说明文字
    tb_desc = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(5), Inches(5))
    tf = tb_desc.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "数据说明"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    total = sum(data.values())
    for label, value in data.items():
        p = tf.add_paragraph()
        percentage = (value / total * 100) if total > 0 else 0
        p.text = f"• {label}: {value} ({percentage:.1f}%)"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(8)
    
    return slide


def add_bar_chart_slide(prs: Presentation, title: str,
                        categories: List[str],
                        series_data: Dict[str, List[float]],
                        chart_title: str = "") -> Any:
    """
    添加柱状图幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        categories: 分类列表
        series_data: 系列数据字典，键为系列名，值为数值列表
        chart_title: 图表标题
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)
    
    # 添加柱状图
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # 设置图表样式
    chart.has_legend = True
    chart.legend.include_in_layout = False
    
    return slide


def add_line_chart_slide(prs: Presentation, title: str,
                         categories: List[str],
                         series_data: Dict[str, List[float]],
                         chart_title: str = "") -> Any:
    """
    添加折线图幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        categories: 分类列表（如月份）
        series_data: 系列数据字典，键为系列名，值为数值列表
        chart_title: 图表标题
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)
    
    # 添加折线图
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    # 设置图表样式
    chart.has_legend = True
    chart.legend.include_in_layout = False
    
    return slide


def add_architecture_diagram_slide(prs: Presentation, title: str,
                                   layers: List[Dict[str, Any]]) -> Any:
    """
    添加架构图幻灯片（4A架构：BA/DA/AA/TA）
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        layers: 层级列表，每个层级包含：
               - name: 层级名称
               - components: 组件列表
               - color: 颜色（可选）
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 架构层布局
    layer_height = 1.2
    layer_gap = 0.2
    start_y = 1.3
    layer_width = 12.333
    layer_x = 0.5
    
    for i, layer in enumerate(layers):
        y = start_y + i * (layer_height + layer_gap)
        
        # 层级颜色
        layer_color = layer.get('color', CHART_COLORS[i % len(CHART_COLORS)])
        
        # 层级背景
        layer_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(layer_x), Inches(y),
            Inches(layer_width), Inches(layer_height)
        )
        layer_bg.fill.solid()
        layer_bg.fill.fore_color.rgb = layer_color
        layer_bg.line.fill.background()
        
        # 层级名称
        tb_name = slide.shapes.add_textbox(
            Inches(layer_x + 0.2), Inches(y + 0.1),
            Inches(2), Inches(0.4)
        )
        p = tb_name.text_frame.paragraphs[0]
        p.text = layer.get('name', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        
        # 组件列表
        components = layer.get('components', [])
        component_width = (layer_width - 2.5) / max(len(components), 1)
        
        for j, comp in enumerate(components):
            comp_x = layer_x + 2.3 + j * component_width
            comp_y = y + 0.5
            
            # 组件框
            comp_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(comp_x), Inches(comp_y),
                Inches(component_width - 0.1), Inches(0.55)
            )
            comp_bg.fill.solid()
            comp_bg.fill.fore_color.rgb = COLORS['white']
            comp_bg.line.color.rgb = layer_color
            
            # 组件文字
            tb_comp = slide.shapes.add_textbox(
                Inches(comp_x), Inches(comp_y + 0.1),
                Inches(component_width - 0.1), Inches(0.4)
            )
            p = tb_comp.text_frame.paragraphs[0]
            p.text = comp
            p.font.size = Pt(11)
            p.font.color.rgb = layer_color
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_org_chart_slide(prs: Presentation, title: str,
                        org_structure: Dict[str, Any]) -> Any:
    """
    添加组织架构图幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        org_structure: 组织架构字典，包含：
                      - name: 组织名称
                      - children: 子组织列表
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 绘制组织架构
    def draw_org_node(node, x, y, width, level=0):
        """递归绘制组织节点"""
        box_height = 0.6
        box_width = min(width, 2.5)
        
        # 节点框
        node_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        node_bg.fill.solid()
        node_bg.fill.fore_color.rgb = CHART_COLORS[level % len(CHART_COLORS)]
        node_bg.line.fill.background()
        
        # 节点文字
        tb_node = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.15),
            Inches(box_width), Inches(0.4)
        )
        p = tb_node.text_frame.paragraphs[0]
        p.text = node.get('name', '')
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 子节点
        children = node.get('children', [])
        if children:
            child_y = y + box_height + 0.3
            child_width = width / max(len(children), 1)
            
            for i, child in enumerate(children):
                child_x = x + i * child_width
                draw_org_node(child, child_x, child_y, child_width, level + 1)
    
    # 从顶层开始绘制
    root = org_structure
    start_x = 0.5
    start_y = 1.5
    total_width = 12.333
    
    draw_org_node(root, start_x, start_y, total_width)
    
    return slide


def add_process_flow_slide(prs: Presentation, title: str,
                           steps: List[Dict[str, str]]) -> Any:
    """
    添加流程图幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        steps: 步骤列表，每个步骤包含：
              - name: 步骤名称
              - description: 步骤描述
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 流程步骤布局
    step_width = 2.0
    step_height = 1.0
    gap = 0.3
    start_x = 0.5
    start_y = 1.8
    arrow_width = 0.3
    
    total_width = len(steps) * step_width + (len(steps) - 1) * (gap + arrow_width)
    start_x = (13.333 - total_width) / 2
    
    for i, step in enumerate(steps):
        x = start_x + i * (step_width + gap + arrow_width)
        
        # 步骤框
        step_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(start_y),
            Inches(step_width), Inches(step_height)
        )
        step_bg.fill.solid()
        step_bg.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
        step_bg.line.fill.background()
        
        # 步骤编号
        tb_num = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 0.1),
            Inches(step_width), Inches(0.3)
        )
        p = tb_num.text_frame.paragraphs[0]
        p.text = f"步骤 {i + 1}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 步骤名称
        tb_name = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 0.35),
            Inches(step_width), Inches(0.4)
        )
        p = tb_name.text_frame.paragraphs[0]
        p.text = step.get('name', '')
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 步骤描述
        if 'description' in step:
            tb_desc = slide.shapes.add_textbox(
                Inches(x), Inches(start_y + step_height + 0.2),
                Inches(step_width), Inches(3)
            )
            tf = tb_desc.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step.get('description', '')
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['text_dark']
            p.alignment = PP_ALIGN.CENTER
        
        # 箭头（除最后一个步骤）
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + step_width + 0.05), Inches(start_y + 0.35),
                Inches(arrow_width), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['accent']
            arrow.line.fill.background()
    
    return slide


def add_kpi_dashboard_slide(prs: Presentation, title: str,
                            kpis: List[Dict[str, Any]]) -> Any:
    """
    添加KPI仪表盘幻灯片
    
    Args:
        prs: Presentation对象
        title: 幻灯片标题
        kpis: KPI列表，每个KPI包含：
             - name: KPI名称
             - value: 当前值
             - target: 目标值（可选）
             - unit: 单位（可选）
             - trend: 趋势（up/down/flat，可选）
    
    Returns:
        创建的幻灯片对象
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # KPI卡片布局 - 3列
    cols = 3
    rows = (len(kpis) + cols - 1) // cols
    card_width = 3.8
    card_height = 2.2
    gap_x = 0.4
    gap_y = 0.3
    start_x = 0.7
    start_y = 1.3
    
    for i, kpi in enumerate(kpis[:9]):  # 最多显示9个KPI
        col = i % cols
        row = i // cols
        
        x = start_x + col * (card_width + gap_x)
        y = start_y + row * (card_height + gap_y)
        
        # KPI卡片背景
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['light_bg']
        card_bg.line.color.rgb = RGBColor(220, 220, 220)
        
        # KPI名称
        tb_name = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.1),
            Inches(card_width - 0.3), Inches(0.3)
        )
        p = tb_name.text_frame.paragraphs[0]
        p.text = kpi.get('name', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['text_light']
        
        # KPI值
        value = kpi.get('value', '')
        unit = kpi.get('unit', '')
        tb_value = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.5),
            Inches(card_width - 0.3), Inches(0.8)
        )
        p = tb_value.text_frame.paragraphs[0]
        p.text = f"{value}{unit}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 目标值和趋势
        if 'target' in kpi:
            target = kpi.get('target', '')
            tb_target = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 1.4),
                Inches(card_width - 0.3), Inches(0.3)
            )
            p = tb_target.text_frame.paragraphs[0]
            p.text = f"目标: {target}{unit}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['text_light']
        
        # 趋势指示器
        trend = kpi.get('trend', 'flat')
        if trend != 'flat':
            trend_symbol = '↑' if trend == 'up' else '↓'
            trend_color = COLORS['success'] if trend == 'up' else COLORS['error']
            
            tb_trend = slide.shapes.add_textbox(
                Inches(x + card_width - 0.6), Inches(y + 0.1),
                Inches(0.5), Inches(0.4)
            )
            p = tb_trend.text_frame.paragraphs[0]
            p.text = trend_symbol
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = trend_color
    
    return slide


# 模块信息
__version__ = "1.0.0"
__author__ = "ClawBot"
__description__ = "金蝶PPT生成器图表模块"


if __name__ == "__main__":
    # 测试代码
    print(f"金蝶PPT图表模块 v{__version__}")
    print(f"支持的图表类型:")
    print(f"  - 饼图: add_pie_chart_slide()")
    print(f"  - 柱状图: add_bar_chart_slide()")
    print(f"  - 折线图: add_line_chart_slide()")
    print(f"  - 架构图: add_architecture_diagram_slide()")
    print(f"  - 组织架构图: add_org_chart_slide()")
    print(f"  - 流程图: add_process_flow_slide()")
    print(f"  - KPI仪表盘: add_kpi_dashboard_slide()")
