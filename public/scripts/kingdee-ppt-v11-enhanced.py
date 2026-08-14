#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 - v11.0 专业增强版
支持行业特定内容，包含详细解决方案和案例
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'white': RGBColor(255, 255, 255),
    'light': RGBColor(240, 240, 240),
    'success': RGBColor(76, 175, 80),
    'warning': RGBColor(255, 152, 0),
    'error': RGBColor(244, 67, 54),
}

def add_title_slide(prs, title, subtitle="", date=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    
    if date:
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = date
        p3.font.size = Pt(18)
        p3.font.color.rgb = COLORS['white']
        p3.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith('###'):
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(8)
        elif item.strip():
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
    
    return slide

def add_value_card_slide(prs, title, items):
    """添加价值卡片页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 价值卡片
    card_width = Inches(5)
    card_height = Inches(1.5)
    card_spacing = Inches(0.5)
    cards_per_row = 2
    
    for i, item in enumerate(items):
        row = i // cards_per_row
        col = i % cards_per_row
        
        x = Inches(0.5) + col * (card_width + card_spacing)
        y = Inches(1.5) + row * (card_height + Inches(0.3))
        
        # 卡片背景
        card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, card_height)
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['light']
        card_bg.line.fill.background()
        
        # 卡片内容
        card_text = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), card_width - Inches(0.4), card_height - Inches(0.6))
        card_tf = card_text.text_frame
        card_tf.word_wrap = True
        
        p = card_tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def get_industry_content(industry, modules):
    """根据行业获取特定内容"""
    industry_content = {
        '制造业': {
            'pain_points': [
                "###生产计划不精准",
                "- 产能计算不准确",
                "- 物料需求预测偏差大",
                "- 生产进度跟踪困难",
                "",
                "###质量控制不完善",
                "- 质量数据分散",
                "- 质量追溯困难",
                "- 质量标准执行不严格",
                "",
                "###成本核算不准确",
                "- 成本归集不完整",
                "- 成本分摊不合理",
                "- 成本分析不及时"
            ],
            'solutions': [
                "###智能制造解决方案",
                "- 数字化生产计划",
                "- 智能排产与调度",
                "- 实时生产监控",
                "- 设备管理维护",
                "",
                "###质量管控解决方案",
                "- 全流程质量追溯",
                "- 质量数据分析",
                "- 质量预警机制",
                "- 质量标准管理",
                "",
                "###成本管理解决方案",
                "- 精准成本核算",
                "- 成本分析与预测",
                "- 成本控制优化",
                "- 盈利能力分析"
            ],
            'value_metrics': [
                "生产效率提升40%",
                "产品质量合格率提升15%",
                "生产成本降低20%",
                "设备利用率提升30%"
            ]
        },
        '零售业': {
            'pain_points': [
                "###库存管理混乱",
                "- 库存数据不准确",
                "- 缺货/积压严重",
                "- 库存周转率低",
                "",
                "###采购效率低下",
                "- 采购流程繁琐",
                "- 供应商管理困难",
                "- 采购成本控制不力",
                "",
                "###销售分析不足",
                "- 销售数据分散",
                "- 客户画像不清晰",
                "- 营销效果难评估"
            ],
            'solutions': [
                "###智慧零售解决方案",
                "- 智能库存管理",
                "- 自动化采购",
                "- 供应商协同",
                "- 采购成本优化",
                "",
                "###客户营销解决方案",
                "- 全渠道客户管理",
                "- 智能营销推荐",
                "- 客户价值分析",
                "- 营销效果追踪",
                "",
                "###门店运营解决方案",
                "- 统一商品管理",
                "- 智能定价策略",
                "- 门店绩效管理",
                "- 会员体系构建"
            ],
            'value_metrics': [
                "库存周转率提升50%",
                "采购成本降低25%",
                "销售转化率提升30%",
                "客户满意度提升40%"
            ]
        },
        '服务业': {
            'pain_points': [
                "###客户管理分散",
                "- 客户信息不完整",
                "- 客户需求不明确",
                "- 客户服务不一致",
                "",
                "###项目管控困难",
                "- 项目进度不透明",
                "- 资源分配不合理",
                "- 项目成本超支",
                "",
                "###财务管理复杂",
                "- 收入确认复杂",
                "- 成本分摊困难",
                "- 财务分析滞后"
            ],
            'solutions': [
                "###客户关系管理",
                "- 全生命周期客户管理",
                "- 客户需求分析",
                "- 服务质量监控",
                "- 客户满意度管理",
                "",
                "###项目管理解决方案",
                "- 项目全生命周期管理",
                "- 资源优化配置",
                "- 项目成本控制",
                "- 项目绩效评估",
                "",
                "###财务管理解决方案",
                "- 收入成本匹配",
                "- 项目成本核算",
                "- 财务分析报告",
                "- 风险预警机制"
            ],
            'value_metrics': [
                "客户满意度提升35%",
                "项目交付准时率提升40%",
                "运营成本降低20%",
                "客户留存率提升25%"
            ]
        }
    }
    
    # 默认制造业内容
    return industry_content.get(industry, industry_content['制造业'])

def generate_presales_ppt_v11(customer_info):
    """生成售前PPT - v11.0 专业增强版"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance', 'supply'])
    
    # 获取行业特定内容
    industry_content = get_industry_content(industry, modules)
    
    # 封面
    add_title_slide(prs, company_name, f"{industry}企业数字化转型解决方案", datetime.now().strftime('%Y年%m月%d日'))
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、企业概况与需求分析",
        "二、行业痛点与挑战", 
        "三、总体解决方案",
        "四、核心模块设计",
        "五、4A架构详解",
        "六、价值工程分析",
        "七、实施路线图",
        "八、成功案例分享",
        "九、项目团队与保障"
    ])
    
    # 企业概况
    add_section_slide(prs, "一、企业概况与需求分析")
    add_content_slide(prs, "企业基本情况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{customer_info.get('companySize', '中型企业')}",
        f"员工人数：{customer_info.get('employees', '待定')}人",
        f"年营业额：{customer_info.get('revenue', '待定')}万元",
        "",
        "###信息化现状：",
        "- 系统分散，数据孤岛",
        "- 业务流程不标准化",
        "- 管理决策滞后"
    ])
    
    # 行业痛点
    add_section_slide(prs, "二、行业痛点与挑战")
    add_content_slide(prs, f"{industry}行业核心痛点", industry_content['pain_points'])
    
    # 总体解决方案
    add_section_slide(prs, "三、总体解决方案")
    add_content_slide(prs, "数字化转型总体框架", [
        "###技术架构：",
        "- 金蝶云·星空平台",
        "- 云原生微服务架构",
        "- 多租户SaaS模式",
        "- 高可用集群部署",
        "",
        "###业务架构：",
        "- 财务管理一体化",
        "- 供应链协同管理",
        "- 生产制造智能化",
        "- 人力资源管理现代化",
        "",
        "###数据架构：",
        "- 统一数据平台",
        "- 数据治理体系",
        "- 智能分析应用"
    ])
    
    # 核心模块设计
    add_section_slide(prs, "四、核心模块设计")
    add_content_slide(prs, "核心业务模块", [
        "###财务管理模块：",
        "- 财务核算总账",
        "- 应收应付管理",
        "- 成本管理",
        "- 资金管理",
        "",
        "###供应链管理模块：",
        "- 采购管理",
        "- 销售管理",
        "- 库存管理",
        "",
        "###生产制造模块：",
        "- 生产计划管理",
        "- 生产执行管理",
        "- 质量管理",
        "",
        "###人力资源管理模块：",
        "- 人事管理",
        "- 薪酬管理",
        "- 绩效管理"
    ])
    
    # 4A架构详解
    add_section_slide(prs, "五、4A架构详解")
    add_content_slide(prs, "业务架构（BA）", [
        "###核心价值流：",
        "- 订单到收款（O2C）",
        "- 采购到付款（P2P）", 
        "- 计划到生产（P2M）",
        "- 研发到上市（R2M）",
        "",
        "###业务能力：",
        "- 财务管理能力",
        "- 供应链管理能力",
        "- 生产制造能力",
        "- 人力资源管理能力"
    ])
    
    add_content_slide(prs, "数据架构（DA）", [
        "###数据实体：",
        "- 客户、供应商、物料",
        "- 订单、发票、凭证",
        "- 产品、工艺、设备",
        "",
        "###数据服务：",
        "- 数据查询服务",
        "- 数据分析服务", 
        "- 数据治理服务"
    ])
    
    # 价值工程分析
    add_section_slide(prs, "六、价值工程分析")
    add_value_card_slide(prs, "核心价值指标", [
        "战略价值：提升企业竞争力",
        "经济价值：降低运营成本",
        "管理价值：提升管理效率",
        "创新价值：促进业务创新"
    ])
    
    add_content_slide(prs, f"{industry}行业价值收益", [
        "###量化收益指标：",
        "- 财务结账时间缩短50%",
        "- 库存周转率提升40%",
        "- 业务效率提升30%",
        "- 客户满意度提升35%",
        "",
        "###行业特定价值：",
        *industry_content['value_metrics']
    ])
    
    # 实施路线图
    add_section_slide(prs, "七、实施路线图")
    add_content_slide(prs, "分阶段实施计划", [
        "###第一阶段：基础建设（1-2个月）",
        "- 项目启动与团队组建",
        "- 需求调研与蓝图设计",
        "- 系统基础配置",
        "",
        "###第二阶段：核心实施（2-3个月）",
        "- 财务管理模块上线",
        "- 供应链管理模块上线",
        "- 数据迁移与验证",
        "",
        "###第三阶段：全面推广（1-2个月）",
        "- 生产制造模块上线",
        "- 人力资源管理模块上线",
        "- 用户培训与推广",
        "",
        "###第四阶段：持续优化（持续）",
        "- 系统功能优化",
        "- 业务流程优化",
        "- 数据价值挖掘"
    ])
    
    # 成功案例分享
    add_section_slide(prs, "八、成功案例分享")
    add_content_slide(prs, f"{industry}行业成功案例", [
        "###案例1：大型制造企业",
        "企业规模：5000+员工，年营收20亿",
        "实施范围：财务、供应链、生产制造",
        "项目成果：",
        "- 财务结账时间从15天缩短到5天",
        "- 库存周转率提升45%",
        "- 生产效率提升35%",
        "",
        "###案例2：零售连锁企业",
        "企业规模：1000+门店，年营收50亿",
        "实施范围：财务、供应链、门店管理",
        "项目成果：",
        "- 采购成本降低30%",
        "- 销售转化率提升40%",
        "- 客户满意度提升45%"
    ])
    
    # 项目团队与保障
    add_section_slide(prs, "九、项目团队与保障")
    add_content_slide(prs, "专业实施团队", [
        "###核心团队配置：",
        "- 项目总监：10年以上ERP实施经验",
        "- 业务顾问：8年以上行业经验",
        "- 技术专家：5年以上技术架构经验",
        "- 实施顾问：6年以上实施经验",
        "",
        "###服务保障体系：",
        "- 7×24小时技术支持",
        "- 定期项目回访",
        "- 持续系统优化",
        "- 培训认证服务"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", "金蝶软件（中国）有限公司")
    
    return prs

def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 3:
        print("用法: python3 kingdee-ppt-v11-enhanced.py --type <type> --companyName <name> [其他参数]")
        sys.exit(1)
    
    # 解析参数
    params = {}
    for i in range(1, len(sys.argv), 2):
        if i + 1 < len(sys.argv):
            key = sys.argv[i].lstrip('-')
            params[key] = sys.argv[i + 1]
    
    # 生成PPT
    ppt_type = params.get('type', 'presales')
    company_name = params.get('companyName', '企业名称')
    
    try:
        if ppt_type == 'presales':
            prs = generate_presales_ppt_v11(params)
        else:
            print(f"不支持的PPT类型: {ppt_type}")
            sys.exit(1)
        
        # 保存文件
        filename = f"{company_name}_{ppt_type}_v11.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        print(json.dumps({
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "slides": len(prs.slides),
            "type": ppt_type
        }, ensure_ascii=False))
        
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": str(e)
        }, ensure_ascii=False))

if __name__ == "__main__":
    main()