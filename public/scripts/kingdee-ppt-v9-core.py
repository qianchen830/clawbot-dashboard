#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - v9.0 核心版本
基于真实金蝶项目文档模板，支持100+页专业内容
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'light': RGBColor(240, 240, 240),
    'white': RGBColor(255, 255, 255),
}

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Content
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith('###'):
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(8)
        elif item.strip():
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
    
    return slide

print("✅ v9.0核心版本创建成功！")
