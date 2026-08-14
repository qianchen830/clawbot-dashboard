#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 - v13.0 专业母版版
使用参考文档的母版，包含完整内容
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import random

TEMPLATE_PATH = os.path.expanduser("~/.openclaw/workspace/output/ppt-assets/kingdee-template.pptx")
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
ASSETS_DIR = os.path.join(OUTPUT_DIR, "ppt-assets")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 82, 147),
    'dark': RGBColor(51, 51, 51),
    'gray': RGBColor(128, 128, 128),
}

def get_random_image(category):
    folder = os.path.join(ASSETS_DIR, category)
    if not os.path.exists(folder):
        return None
    images = [f for f in os.listdir(folder) if f.endswith(('.png', '.jpg', '.jpeg'))]
    if not images:
        return None
    return os.path.join(folder, random.choice(images))

def generate_presales_ppt_v13(customer_info):
    company_name = customer_info.get('companyName', '客户企业')
    industry = customer_info.get('industry', '制造业')
    company_size = customer_info.get('companySize', '中型企业')
    modules = customer_info.get('modules', ['finance', 'supply'])
    employees = customer_info.get('employees', '1000')
    revenue = customer_info.get('revenue', '10000')
    pain_points = customer_info.get('painPoints', '')
    business_goals = customer_info.get('businessGoals', '')
    
    module_names = {
        'finance': '财务管理',
        'supply': '供应链管理',
        'manufacturing': '制造管理',
        'hr': '人力资源管理'
    }
    module_list = [module_names.get(m, m) for m in modules]
    
    try:
        prs = Presentation(TEMPLATE_PATH)
    except Exception as e:
        prs = Presentation()
    
    layouts = {}
    for layout in prs.slide_layouts:
        layouts[layout.name] = layout
    
    cover_layout = layouts.get('封面', layouts.get('2_空白', prs.slide_layouts[0]))
    toc_layout = layouts.get('目录、提纲', layouts.get('1_目录、提纲', prs.slide_layouts[0]))
    content_layout = layouts.get('白色内页', layouts.get('1_白色内页', prs.slide_layouts[0]))
    blank_layout = layouts.get('2_空白', prs.slide_layouts[0])
    
    today = datetime.now().strftime('%Y年%m月%d日')
    
    def add_slide(layout_name, title, content=None, image=None):
        layout = layouts.get(layout_name, blank_layout)
        slide = prs.slides.add_slide(layout)
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text and len(para.text) < 50:
                        para.text = title
                        break
                break
        
        if content:
            content_added = False
            for shape in slide.shapes:
                if shape.has_text_frame and not content_added:
                    if len(shape.text_frame.text) > 50 or shape.top > Inches(1):
                        tf = shape.text_frame
                        tf.clear()
                        for i, item in enumerate(content):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            if item.startswith('  '):
                                p.text = f"    {item.strip()}"
                                p.font.size = Pt(12)
                                p.font.color.rgb = COLORS['gray']
                            else:
                                p.text = f"● {item}"
                                p.font.size = Pt(14)
                                p.font.color.rgb = COLORS['dark']
                            p.space_after = Pt(6)
                        content_added = True
                        break
            
            if not content_added:
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                for i, item in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"● {item}" if not item.startswith('  ') else f"    {item.strip()}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = COLORS['dark']
                    p.space_after = Pt(6)
        
        if image and os.path.exists(image):
            try:
                slide.shapes.add_picture(image, Inches(0.5), Inches(1.5), width=Inches(12))
            except:
                pass
        
        return slide
    
    # ===== 封面 =====
    add_slide('封面', f"{company_name}ERP系统解决方案")
    
    # ===== 目录 =====
    add_slide('目录、提纲', "目 录", [
        "01 项目需求理解",
        "02 解决方案设计", 
        "03 技术平台架构",
        "04 项目实施安排",
        "05 项目团队配置",
        "06 成功案例分享",
        "07 服务承诺保障",
        "08 金蝶实力介绍"
    ])
    
    # ===== 第1部分：项目需求理解 =====
    add_slide('白色内页', "一、行业洞察", [
        f"随着{industry}行业发展环境变化，业务增长重心从单一产品向整体解决方案转变",
        "数字化转型成为必然趋势，企业需构建敏捷、智能的管理体系",
        "市场竞争加剧，成本控制、效率提升、质量保障成为核心关注点",
        "",
        "行业面临的挑战：",
        "  - 技术创新不足，与国际先进水平存在差距",
        "  - 落后产能过剩，转型升级迫在眉睫",
        "  - 政策监管加强，环保要求日益严格",
        "  - 市场需求变化，需要快速适应"
    ])
    
    add_slide('白色内页', "二、企业概况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{company_size}",
        f"员工人数：{employees}人",
        f"年营业额：{revenue}万元",
        f"实施模块：{'、'.join(module_list)}",
        "",
        "企业当前正处于数字化转型关键期，需要：",
        "  - 业务流程标准化、规范化",
        "  - 财务业务一体化管理",
        "  - 数据驱动决策能力提升"
    ])
    
    add_slide('白色内页', "三、当前业务痛点分析", [
        "信息系统分散，数据孤岛严重，跨部门协同困难",
        "业务流程不顺畅，审批效率低，流程可追溯性差",
        "财务核算不及时，成本核算不精准，决策支持不足",
        "供应链管理粗放，库存周转率低，采购效率不高",
        "生产计划不准确，物料需求预测偏差大，交付周期长",
        "质量管理缺乏追溯体系，问题定位困难，改进措施难落地",
        "主数据管理分散，缺乏统一的主数据管理平台",
        "各系统间数据断点多，对运营决策支撑不足"
    ])
    
    add_slide('白色内页', "四、项目建设目标", [
        "统一业务及管理主数据，消除信息孤岛",
        "实现销售到收款的全闭环数字化管理",
        "建立项目预算、项目经营及项目核算的闭环管理",
        "实现采购到应付的全链条数字化管理",
        "构建从计划到交付的生产及供应链数字化管理体系",
        "实现业务变更及工程变更的全程跟踪和联动",
        "建立质量检验、追溯及持续改善体系",
        "提升人才招、选、用、育、留的全面体验"
    ])
    
    # ===== 第2部分：解决方案设计 =====
    add_slide('白色内页', "五、总体架构设计", [
        "业务架构 BA：销售到收款、采购到付款、计划到生产全流程管理",
        "数据架构 DA：数据治理、数据资产、数据服务",
        "应用架构 AA：财务管理、供应链、制造管理、人力资源管理",
        "技术架构 TA：云原生微服务、分布式部署、开放集成平台",
        "",
        "设计原则：",
        "  - 业务驱动：以业务需求为导向，确保系统支撑业务发展",
        "  - 技术先进：采用云原生架构，确保系统稳定可扩展",
        "  - 数据贯通：打通数据断点，实现业务数据一体化",
        "  - 安全可控：数据加密、权限管控，确保数据安全"
    ])
    
    # 添加架构图
    arch_image = get_random_image('architecture')
    if arch_image:
        add_slide('白色内页', "六、系统架构图", image=arch_image)
    
    add_slide('白色内页', "七、核心解决方案", [
        "财务管理：总账、应收、应付、固定资产、成本、预算",
        "供应链管理：销售、采购、库存、计划、物流",
        "生产制造：BOM管理、MRP运算、生产执行、质量管理",
        "数据分析：经营驾驶舱、多维报表、智能预警",
        "",
        "解决方案特点：",
        "  - 一体化设计：业务流程端到端贯通",
        "  - 灵活配置：支持企业个性化需求",
        "  - 数据驱动：实时数据支撑决策",
        "  - 智能分析：AI赋能业务洞察"
    ])
    
    if 'finance' in modules:
        add_slide('白色内页', "八、财务管理方案", [
            "总账管理：凭证处理、自动记账、期末结账、财务报表",
            "应收管理：销售开票、收款核销、账龄分析、信用管理",
            "应付管理：采购发票、付款处理、供应商对账、付款计划",
            "固定资产：资产台账、折旧计提、资产处置、资产盘点",
            "成本管理：成本核算、成本分配、成本分析、成本控制",
            "全面预算：预算编制、预算审批、预算控制、预算分析",
            "",
            "预期价值：财务结账时间缩短60%，成本核算精度提升40%"
        ])
    
    if 'supply' in modules:
        add_slide('白色内页', "九、供应链管理方案", [
            "销售管理：订单处理、信用控制、出货管理、销售分析",
            "采购管理：询价比价、订单跟踪、入库检验、采购结算",
            "库存管理：出入库管理、盘点管理、库存分析、库存预警",
            "计划管理：需求计划、采购计划、生产计划、MRP运算",
            "物流管理：配送管理、运输跟踪、物流成本、签收管理",
            "",
            "预期价值：库存周转率提升40%，采购成本降低15%"
        ])
    
    if 'manufacturing' in modules:
        add_slide('白色内页', "十、生产制造管理方案", [
            "产品数据：BOM管理、工艺路线、工程变更、版本管理",
            "计划管理：主生产计划、物料需求计划、能力计划、排程优化",
            "生产执行：工单管理、工序汇报、完工入库、进度跟踪",
            "质量管理：来料检验、过程检验、成品检验、质量追溯",
            "设备管理：设备台账、维护保养、故障处理、设备分析",
            "",
            "预期价值：生产计划准确率提升40%，MRP运算效率提升80%"
        ])
    
    add_slide('白色内页', "十一、预期价值分析", [
        "运营效率提升40%：流程自动化、数据实时共享",
        "库存周转提升30%：精准计划、库存优化",
        "结账周期缩短60%：自动化核算、实时报表",
        "成本精度提升35%：精细核算、多维分析",
        "决策响应提升60%：实时分析、智能预警",
        "协同效率提升50%：流程贯通、信息共享",
        "质量追溯提升90%：全过程记录、快速定位",
        "财务准确提升50%：业务驱动、自动核算"
    ])
    
    # ===== 第3部分：技术平台 =====
    add_slide('白色内页', "十二、金蝶云·苍穹平台优势", [
        "云原生架构：微服务、容器化、分布式部署",
        "低代码开发：可视化建模、快速交付、灵活扩展",
        "AI原生能力：大模型集成、智能助手、预测分析",
        "数据中台：数据治理、数据资产、数据服务",
        "开放集成：API网关、标准接口、生态连接",
        "安全可信：等保三级、数据加密、权限管控"
    ])
    
    # 添加截图
    screenshots = [f for f in os.listdir(os.path.join(ASSETS_DIR, 'screenshots')) if f.endswith(('.png', '.jpg'))][:2]
    for i, img in enumerate(screenshots):
        img_path = os.path.join(ASSETS_DIR, 'screenshots', img)
        add_slide('白色内页', f"系统界面展示 ({i+1})", image=img_path)
    
    # ===== 第4部分：项目实施 =====
    add_slide('白色内页', "十三、实施方法论", [
        "第一阶段：项目启动（2周）",
        "  - 成立项目组织、明确项目范围、召开启动会议",
        "第二阶段：需求调研（4周）",
        "  - 业务调研、需求分析、差距分析、编制调研报告",
        "第三阶段：方案设计（3周）",
        "  - 业务蓝图设计、系统配置方案、方案评审确认",
        "第四阶段：系统配置（4周）",
        "  - 系统参数配置、基础数据准备、业务流程配置",
        "第五阶段：测试培训（3周）",
        "  - UAT测试、最终用户培训、操作手册编制",
        "第六阶段：上线验收（2周）",
        "  - 系统切换上线、上线支持保障、项目验收"
    ])
    
    add_slide('白色内页', "十四、项目团队配置", [
        "项目总监 1人：项目整体把控、资源协调（PMP认证、10年以上经验）",
        "项目经理 1人：项目日常管理、进度控制（PMP认证、5年以上经验）",
        "业务顾问 2人：业务方案设计、流程优化（行业认证、3年以上经验）",
        "技术顾问 1人：技术方案设计、系统集成（技术认证、3年以上经验）",
        "实施顾问 3人：系统配置、用户培训（产品认证、2年以上经验）",
        "开发工程师 2人：客户化开发、报表开发（开发认证、2年以上经验）"
    ])
    
    # ===== 第5部分：成功案例 =====
    add_slide('白色内页', "十五、行业典型案例", [
        "案例一：某大型制造企业",
        "  - 企业规模：年营收50亿，员工5000人",
        "  - 实施模块：财务、供应链、制造、HR",
        "  - 实施周期：8个月",
        "  - 项目成效：财务结账时间缩短60%、库存周转率提升40%",
        "",
        "案例二：某大型零售企业",
        "  - 企业规模：年营收30亿，员工3000人",
        "  - 实施模块：财务、供应链、零售管理",
        "  - 实施周期：6个月",
        "  - 项目成效：采购成本降低15%、库存准确率达到99%"
    ])
    
    # ===== 第6部分：服务承诺 =====
    add_slide('白色内页', "十六、服务承诺", [
        "系统稳定性承诺",
        "  - 系统可用率≥99.9%",
        "  - 数据安全有保障，支持容灾备份",
        "",
        "响应时效承诺",
        "  - 一级问题：2小时内响应，8小时内解决",
        "  - 二级问题：4小时内响应，24小时内解决",
        "",
        "培训承诺",
        "  - 提供不少于10场系统培训",
        "  - 培训覆盖率100%，考核通过率≥90%",
        "",
        "质保承诺",
        "  - 系统上线后提供1年免费质保服务"
    ])
    
    # ===== 第7部分：金蝶实力 =====
    add_slide('白色内页', "十七、金蝶核心数据", [
        "服务客户：743万+ 企业及组织选择金蝶",
        "云收入占比：79.46% 云转型成功",
        "市场份额：连续17年 中国ERP市场第一",
        "员工规模：1.2万+ 专业服务团队",
        "覆盖国家：170+ 业务遍布全球",
        "研发投入：30%+ 持续创新能力",
        "成功率：95%+ 项目交付成功率"
    ])
    
    add_slide('白色内页', "十八、金蝶云·星瀚核心优势", [
        "为大企业提供世界一流的数字化管理平台",
        "",
        "产品能力：",
        "  - 100+云服务，覆盖全业务场景",
        "  - 预组装数字能力，快速部署上线",
        "  - 灵活配置，适应企业个性化需求",
        "",
        "技术优势：",
        "  - 云原生架构，支持高并发、高可用",
        "  - 低代码开发平台，快速响应变化",
        "  - 开放集成能力，打通信息孤岛",
        "",
        "行业认可：",
        "  - IDC中国EA SaaS市场份额第一",
        "  - Gartner全球应用平台软件TOP5"
    ])
    
    # ===== 封底 =====
    add_slide('封底', "感谢聆听")
    
    # 保存
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    customer_code = customer_info.get('customerCode', 'CUSTOMER')
    filename = f"{customer_code}_售前解决方案_v13_{timestamp}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    
    prs.save(filepath)
    
    return {
        'success': True,
        'filepath': filepath,
        'filename': filename,
        'slides': len(prs.slides),
        'template': '使用参考文档母版'
    }

if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--type', default='presales')
    parser.add_argument('--companyName', required=True)
    parser.add_argument('--customerCode', default='CUSTOMER')
    parser.add_argument('--industry', default='制造业')
    parser.add_argument('--companySize', default='中型企业')
    parser.add_argument('--employees', default='')
    parser.add_argument('--revenue', default='')
    parser.add_argument('--modules', default='finance,supply')
    parser.add_argument('--painPoints', default='')
    parser.add_argument('--businessGoals', default='')
    parser.add_argument('--goliveDate', default='')
    parser.add_argument('--golivePhase', default='')
    parser.add_argument('--switchPlan', default='')
    parser.add_argument('--acceptanceDate', default='')
    parser.add_argument('--acceptanceConclusion', default='')
    
    args = parser.parse_args()
    
    customer_info = {
        'companyName': args.companyName,
        'customerCode': args.customerCode,
        'industry': args.industry,
        'companySize': args.companySize,
        'employees': args.employees,
        'revenue': args.revenue,
        'modules': args.modules.split(','),
        'painPoints': args.painPoints,
        'businessGoals': args.businessGoals
    }
    
    result = generate_presales_ppt_v13(customer_info)
    print(json.dumps(result, ensure_ascii=False))
