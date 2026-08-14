#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - 增强版
基于真实金蝶项目文档模板，支持丰富内容和专业样式
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 输出目录
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 金蝶标准配色
COLORS = {
    'primary': RGBColor(0, 102, 153),      # 金蝶蓝
    'secondary': RGBColor(0, 153, 204),     # 浅蓝
    'accent': RGBColor(255, 153, 0),        # 橙色
    'text': RGBColor(51, 51, 51),           # 深灰
    'light': RGBColor(240, 240, 240),       # 浅灰背景
    'white': RGBColor(255, 255, 255),
}

def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, layout_type="bullet"):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容
    if layout_type == "bullet":
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理层级
            if item.startswith('###'):
                p.text = item.replace('###', '').strip()
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLORS['primary']
                p.space_before = Pt(15)
            elif item.startswith('-'):
                p.text = "• " + item[1:].strip()
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['text']
                p.level = 1
                p.space_before = Pt(8)
            elif item.strip():
                p.text = item
                p.font.size = Pt(20)
                p.font.color.rgb = COLORS['text']
                p.space_before = Pt(10)
    
    elif layout_type == "two_column":
        # 左列
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        
        # 右列
        right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        
        mid = len(content_list) // 2
        for i, item in enumerate(content_list[:mid]):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(8)
        
        for i, item in enumerate(content_list[mid:]):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(8)
    
    return slide

def add_toc_slide(prs, items):
    """添加目录幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目 录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 目录项
    content_box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9.333), Inches(5))
    tf = content_box.text_frame
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}、{item}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text']
        p.space_before = Pt(20)
    
    return slide

def add_section_slide(prs, section_title):
    """添加章节分隔幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_presales_ppt_enhanced(customer_info):
    """生成售前PPT - 增强版（40+页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    company_size = customer_info.get('companySize', '中型企业')
    employees = customer_info.get('employees', '')
    revenue = customer_info.get('revenue', '')
    
    # ========== 第一部分：开场（5页）==========
    add_title_slide(prs, company_name, "数字化转型解决方案")
    
    add_content_slide(prs, "感谢致辞", [
        f"感谢{company_name}领导及各位专家评委",
        "",
        "金蝶项目团队",
        f"{datetime.now().strftime('%Y年%m月%d日')}"
    ])
    
    add_content_slide(prs, "金蝶述标团队", [
        "项目经理：" + customer_info.get('projectManager', '待定'),
        "技术总监：待定",
        "业务顾问：待定",
        "开发经理：待定",
        "实施顾问：待定"
    ])
    
    add_toc_slide(prs, [
        "企业概况与需求分析",
        "解决方案设计",
        "业务架构蓝图",
        "4A架构设计",
        "价值工程分析",
        "实施路线图",
        "成功案例分享",
        "项目团队与保障"
    ])
    
    # ========== 第二部分：企业概况（8页）==========
    add_section_slide(prs, "一、企业概况与需求分析")
    
    add_content_slide(prs, "企业基本情况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{company_size}",
        f"员工人数：{employees}人" if employees else "员工人数：待定",
        f"年营业额：{revenue}万元" if revenue else "年营业额：待定",
        "",
        "企业特点：",
        "- 业务模式多元化",
        "- 组织架构复杂",
        "- 管理精细化需求高",
        "- 数字化转型迫切"
    ])
    
    add_content_slide(prs, "企业组织架构", [
        "总部职能中心：",
        "- 财务中心、人力资源中心、运营管理中心",
        "- 采购中心、销售中心、研发中心",
        "",
        "下属单位：",
        "- 分公司、子公司、事业部",
        "- 生产基地、销售网点、服务中心"
    ])
    
    add_content_slide(prs, "当前业务现状", [
        "核心业务流程：",
        "- 销售管理：订单管理、客户管理、渠道管理",
        "- 采购管理：供应商管理、采购订单、入库管理",
        "- 生产管理：生产计划、车间管理、质量管理",
        "- 库存管理：出入库管理、盘点管理、库存分析",
        "- 财务管理：总账、应收、应付、成本、资金",
        "",
        "现有系统：",
        "- ERP系统、OA系统、CRM系统、财务系统"
    ])
    
    add_content_slide(prs, "业务痛点分析", [
        "痛点一：信息孤岛严重",
        "- 各系统数据不互通，重复录入效率低",
        "- 数据口径不一致，报表统计困难",
        "",
        "痛点二：业务流程不畅",
        "- 审批流程冗长，决策效率低下",
        "- 流程断点多，业务协同困难",
        "",
        "痛点三：管理精细化不足",
        "- 成本核算不精准，盈利分析困难",
        "- 预算控制不到位，资金管理粗放",
        "",
        "痛点四：决策支持不足",
        "- 数据分析滞后，无法实时监控",
        "- 缺乏数据可视化，决策依据不足"
    ])
    
    add_content_slide(prs, "数字化转型需求", [
        "战略需求：",
        "- 支撑企业战略发展，提升核心竞争力",
        "- 实现业务数字化转型，提高运营效率",
        "",
        "业务需求：",
        "- 打通业务流程，实现端到端协同",
        "- 统一数据标准，实现数据共享",
        "- 精细化管理，提升运营效率",
        "",
        "技术需求：",
        "- 云原生架构，支持弹性扩展",
        "- 开放平台，支持业务创新",
        "- 移动办公，随时随地访问"
    ])
    
    add_content_slide(prs, "项目目标", [
        "总体目标：",
        "- 建立统一数字化平台，实现业财一体化",
        "- 提升运营效率30%以上",
        "- 降低运营成本15%以上",
        "",
        "具体目标：",
        "- 业务流程标准化率：95%以上",
        "- 数据准确率：99%以上",
        "- 系统可用性：99.9%以上",
        "- 用户满意度：90%以上"
    ])
    
    add_content_slide(prs, "实施范围", [
        "财务云：总账、应收、应付、固定资产、现金管理、成本管理",
        "供应链云：采购管理、库存管理、销售管理、物流管理",
        "制造云：生产计划、车间管理、质量管理、设备管理",
        "人力云：人事管理、薪酬管理、绩效管理、培训管理",
        "",
        "实施模块：",
        "- 一期：财务云 + 供应链云",
        "- 二期：制造云 + 人力云"
    ])
    
    # ========== 第三部分：解决方案（10页）==========
    add_section_slide(prs, "二、解决方案设计")
    
    add_content_slide(prs, "解决方案总体架构", [
        "金蝶云·星空架构：",
        "- 云原生架构，支持多租户",
        "- 微服务架构，弹性扩展",
        "- 开放平台，支持二次开发",
        "",
        "核心能力：",
        "- 业务中台：财务、供应链、制造、人力",
        "- 数据中台：数据治理、数据分析、数据可视化",
        "- 技术中台：开发平台、集成平台、运维平台"
    ])
    
    add_content_slide(prs, "财务云解决方案", [
        "总账管理：",
        "- 多组织核算，支持多会计准则",
        "- 自动生成凭证，提高核算效率",
        "- 实时财务报表，支持决策分析",
        "",
        "应收应付：",
        "- 采购付款一体化，提高资金效率",
        "- 销售收款一体化，降低坏账风险",
        "- 往来对账自动化，减少人工核对",
        "",
        "成本管理：",
        "- 精细化成本核算，支持多维度分析",
        "- 实时成本监控，及时发现问题",
        "- 成本预测分析，支持经营决策"
    ])
    
    add_content_slide(prs, "供应链云解决方案", [
        "采购管理：",
        "- 供应商全生命周期管理",
        "- 采购申请、订单、入库、结算全流程管理",
        "- 采购价格管理，降低采购成本",
        "",
        "库存管理：",
        "- 多仓库、多货位管理",
        "- 批次管理、序列号管理",
        "- 库存预警、安全库存管理",
        "",
        "销售管理：",
        "- 客户全生命周期管理",
        "- 销售订单、发货、开票、收款全流程管理",
        "- 价格管理、信用管理、促销管理"
    ])
    
    add_content_slide(prs, "制造云解决方案", [
        "生产计划：",
        "- MRP运算，自动生成生产计划",
        "- 物料需求计划，降低库存成本",
        "- 能力需求计划，优化资源利用",
        "",
        "车间管理：",
        "- 生产订单管理，实时跟踪生产进度",
        "- 工序管理，精细化生产控制",
        "- 质量检验，确保产品质量",
        "",
        "质量管理：",
        "- 来料检验、过程检验、成品检验",
        "- 质量追溯，支持质量分析",
        "- 不良品管理，降低质量成本"
    ])
    
    add_content_slide(prs, "人力云解决方案", [
        "人事管理：",
        "- 员工全生命周期管理",
        "- 组织架构、职位管理",
        "- 合同管理、档案管理",
        "",
        "薪酬管理：",
        "- 薪资核算自动化",
        "- 社保公积金管理",
        "- 个税计算、工资发放",
        "",
        "绩效管理：",
        "- KPI指标管理",
        "- 绩效考核、绩效分析",
        "- 绩效改进计划"
    ])
    
    add_content_slide(prs, "系统集成方案", [
        "集成架构：",
        "- API网关，统一接入管理",
        "- 消息队列，异步解耦",
        "- 数据中台，数据交换中心",
        "",
        "集成系统：",
        "- OA系统：审批流程、通知推送",
        "- CRM系统：客户数据同步",
        "- 银行系统：资金支付、银行对账",
        "- 税务系统：发票管理、税务申报"
    ])
    
    add_content_slide(prs, "数据迁移方案", [
        "迁移原则：",
        "- 完整性：数据不丢失、不重复",
        "- 准确性：数据准确、格式正确",
        "- 一致性：新旧数据一致",
        "- 可追溯：数据来源可查",
        "",
        "迁移内容：",
        "- 主数据：组织、客户、供应商、物料、科目",
        "- 期初数据：科目余额、库存余额、往来余额",
        "- 历史数据：历史单据、历史报表"
    ])
    
    add_content_slide(prs, "实施策略", [
        "实施原则：",
        "- 总体规划、分步实施",
        "- 先试点、后推广",
        "- 先核心、后扩展",
        "",
        "实施分期：",
        "- 一期（6个月）：财务云 + 供应链云",
        "- 二期（4个月）：制造云 + 人力云",
        "",
        "实施保障：",
        "- 成立项目组织，明确职责分工",
        "- 制定详细计划，定期跟踪进度",
        "- 建立沟通机制，及时解决问题"
    ])
    
    add_content_slide(prs, "风险管理", [
        "风险识别：",
        "- 需求变更风险",
        "- 数据质量风险",
        "- 用户接受风险",
        "- 系统集成风险",
        "",
        "风险应对：",
        "- 需求变更：建立变更控制流程",
        "- 数据质量：数据清洗、数据验证",
        "- 用户接受：加强培训、持续沟通",
        "- 系统集成：接口测试、联调验证"
    ])
    
    # ========== 第四部分：业务架构（5页）==========
    add_section_slide(prs, "三、业务架构蓝图")
    
    add_content_slide(prs, "业务能力地图", [
        "战略管理能力：战略规划、经营分析、绩效管理",
        "财务管理能力：资金管理、成本管理、预算管理、核算管理",
        "供应链能力：采购管理、库存管理、销售管理、物流管理",
        "生产制造能力：生产计划、车间管理、质量管理、设备管理",
        "人力资源能力：人事管理、薪酬管理、绩效管理、培训管理"
    ])
    
    add_content_slide(prs, "业务流程架构", [
        "端到端流程：",
        "- 订单到现金（O2C）",
        "- 采购到付款（P2P）",
        "- 计划到生产（P2M）",
        "- 记录到报告（R2R）",
        "",
        "核心业务流程：",
        "- 销售流程：订单→发货→开票→收款",
        "- 采购流程：申请→订单→入库→付款",
        "- 生产流程：计划→领料→生产→入库"
    ])
    
    add_content_slide(prs, "数据架构设计", [
        "主数据管理：",
        "- 组织主数据：公司、部门、岗位、人员",
        "- 业务主数据：客户、供应商、物料、科目",
        "",
        "数据标准：",
        "- 编码规则：统一编码体系",
        "- 数据字典：标准数据定义",
        "- 数据质量：数据验证规则"
    ])
    
    add_content_slide(prs, "应用架构设计", [
        "核心应用：",
        "- 财务云：总账、应收、应付、固定资产、现金管理",
        "- 供应链云：采购、库存、销售、物流",
        "- 制造云：生产计划、车间管理、质量管理",
        "- 人力云：人事、薪酬、绩效、培训",
        "",
        "扩展应用：",
        "- 预算管理、合并报表、资金管理",
        "- 条码管理、移动应用、数据分析"
    ])
    
    add_content_slide(prs, "技术架构设计", [
        "技术平台：",
        "- 云原生架构：容器化部署、微服务架构",
        "- 开发平台：低代码开发、快速定制",
        "- 集成平台：API网关、消息队列",
        "",
        "安全保障：",
        "- 数据加密、访问控制、审计日志",
        "- 灾备方案、备份策略"
    ])
    
    # ========== 第五部分：价值工程（5页）==========
    add_section_slide(prs, "四、价值工程分析")
    
    add_content_slide(prs, "价值目标", [
        "战略价值：",
        "- 支撑企业战略发展",
        "- 提升核心竞争力",
        "- 实现数字化转型",
        "",
        "业务价值：",
        "- 提升运营效率30%以上",
        "- 降低运营成本15%以上",
        "- 提高数据准确率99%以上",
        "",
        "管理价值：",
        "- 管理精细化、决策科学化",
        "- 流程标准化、协同高效化"
    ])
    
    add_content_slide(prs, "价值指标体系", [
        "财务指标：",
        "- 成本降低：15%以上",
        "- 资金周转：提升20%以上",
        "- 应收账款：周转天数减少30%",
        "",
        "运营指标：",
        "- 订单处理效率：提升40%以上",
        "- 库存周转率：提升25%以上",
        "- 生产效率：提升20%以上",
        "",
        "管理指标：",
        "- 决策响应时间：缩短50%以上",
        "- 数据准确率：99%以上",
        "- 流程效率：提升30%以上"
    ])
    
    add_content_slide(prs, "投资回报分析", [
        "项目投资：",
        "- 软件许可费：待定",
        "- 实施服务费：待定",
        "- 硬件设备费：待定",
        "",
        "预期收益：",
        "- 直接收益：成本降低、效率提升",
        "- 间接收益：管理提升、风险降低",
        "",
        "投资回报：",
        "- 投资回收期：2-3年",
        "- 投资回报率：200%以上"
    ])
    
    # ========== 第六部分：实施路线（5页）==========
    add_section_slide(prs, "五、实施路线图")
    
    add_content_slide(prs, "项目实施计划", [
        "第一阶段：项目启动（第1-2周）",
        "- 成立项目组织",
        "- 制定项目计划",
        "- 召开启动会议",
        "",
        "第二阶段：需求调研（第3-6周）",
        "- 业务调研",
        "- 需求分析",
        "- 蓝图设计",
        "",
        "第三阶段：系统配置（第7-14周）",
        "- 系统配置",
        "- 客户化开发",
        "- 系统测试"
    ])
    
    add_content_slide(prs, "项目实施计划（续）", [
        "第四阶段：测试培训（第15-20周）",
        "- UAT测试",
        "- 用户培训",
        "- 上线准备",
        "",
        "第五阶段：上线切换（第21-24周）",
        "- 数据迁移",
        "- 系统切换",
        "- 上线支持",
        "",
        "第六阶段：验收支持（第25-28周）",
        "- 系统验收",
        "- 持续优化",
        "- 运维支持"
    ])
    
    add_content_slide(prs, "项目组织架构", [
        "项目领导小组：",
        "- 组长：企业高层领导",
        "- 副组长：项目负责人",
        "- 成员：各部门负责人",
        "",
        "项目实施团队：",
        "- 项目经理：金蝶方",
        "- 业务顾问：金蝶方",
        "- 技术顾问：金蝶方",
        "- 关键用户：企业方"
    ])
    
    add_content_slide(prs, "项目成功要素", [
        "高层支持：",
        "- 项目获得高层重视和支持",
        "- 协调资源，推动项目进展",
        "",
        "目标明确：",
        "- 项目目标清晰、可衡量",
        "- 实施范围明确、可控",
        "",
        "团队保障：",
        "- 项目团队经验丰富",
        "- 关键用户全程参与",
        "",
        "有效沟通：",
        "- 建立沟通机制",
        "- 及时解决问题"
    ])
    
    # ========== 第七部分：成功案例（5页）==========
    add_section_slide(prs, "六、成功案例分享")
    
    add_content_slide(prs, "典型客户案例", [
        "案例一：某大型制造企业",
        "- 实施模块：财务云、供应链云、制造云",
        "- 实施周期：8个月",
        "- 项目成果：效率提升40%，成本降低20%",
        "",
        "案例二：某集团企业",
        "- 实施模块：财务云、人力云、预算管理",
        "- 实施周期：6个月",
        "- 项目成果：核算效率提升50%，资金周转提升30%"
    ])
    
    add_content_slide(prs, "行业成功案例", [
        "制造业：",
        "- 山东重工、潍柴雷沃、中车集团",
        "- 实施成果：生产效率提升30%，库存周转提升25%",
        "",
        "零售业：",
        "- 永辉超市、红旗连锁、家家悦",
        "- 实施成果：库存准确率99%，订单处理效率提升50%",
        "",
        "服务业：",
        "- 招商局、华侨城、中旅集团",
        "- 实施成果：财务共享效率提升60%，运营成本降低15%"
    ])
    
    add_content_slide(prs, "金蝶市场地位", [
        "市场份额：",
        "- 金蝶云连续3年市场份额第一",
        "- 服务企业超过743万家",
        "",
        "客户认可：",
        "- 客户满意度90%以上",
        "- 续约率95%以上",
        "",
        "行业认可：",
        "- IDC中国ERP市场第一",
        "- Gartner魔力象限领导者"
    ])
    
    # ========== 第八部分：项目保障（5页）==========
    add_section_slide(prs, "七、项目团队与保障")
    
    add_content_slide(prs, "项目团队介绍", [
        "金蝶项目团队：",
        "- 项目经理：10年以上实施经验",
        "- 业务顾问：5年以上行业经验",
        "- 技术顾问：精通金蝶产品技术",
        "",
        "项目经验：",
        "- 实施1000+企业成功案例",
        "- 覆盖制造业、零售业、服务业等多个行业"
    ])
    
    add_content_slide(prs, "项目质量保障", [
        "质量管理体系：",
        "- ISO9001质量管理体系认证",
        "- 项目管理规范",
        "",
        "质量控制措施：",
        "- 需求确认、设计评审",
        "- 系统测试、用户验收",
        "- 上线检查、验收评审"
    ])
    
    add_content_slide(prs, "售后服务承诺", [
        "服务内容：",
        "- 7×24小时技术支持",
        "- 系统升级、功能优化",
        "- 问题处理、性能优化",
        "",
        "服务响应：",
        "- 一般问题：4小时内响应",
        "- 紧急问题：1小时内响应",
        "",
        "服务保障：",
        "- 专属客户经理",
        "- 定期客户回访"
    ])
    
    add_content_slide(prs, "项目预期成果", [
        "业务层面：",
        "- 业务流程标准化，协同效率提升",
        "- 数据准确可靠，决策支持有力",
        "",
        "管理层面：",
        "- 管理精细化，运营成本降低",
        "- 风险可控化，合规性提升",
        "",
        "战略层面：",
        "- 支撑企业战略发展",
        "- 提升核心竞争力",
        "- 实现数字化转型"
    ])
    
    # ========== 结尾（3页）==========
    add_title_slide(prs, f"预祝{company_name}", "数字化转型项目取得圆满成功！")
    
    add_content_slide(prs, "联系方式", [
        "金蝶软件（中国）有限公司",
        "",
        "项目经理：" + customer_info.get('projectManager', '待定'),
        "联系电话：待定",
        "电子邮箱：待定",
        "",
        "感谢聆听！"
    ])
    
    # 保存文件
    filename = f"{customer_info.get('customerCode', '客户')}_售前解决方案_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    
    return {
        'success': True,
        'filepath': filepath,
        'filename': filename,
        'slides': len(prs.slides)
    }

# 测试
if __name__ == '__main__':
    test_data = {
        'companyName': '测试公司',
        'customerCode': 'CS',
        'industry': '制造业',
        'companySize': '中型企业',
        'employees': '500',
        'revenue': '10000',
        'projectManager': '张三'
    }
    result = generate_presales_ppt_enhanced(test_data)
    print(f"生成成功：{result['filename']}")
    print(f"幻灯片数量：{result['slides']}页")
