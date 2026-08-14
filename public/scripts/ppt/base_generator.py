# -*- coding: utf-8 -*-
"""
金蝶PPT基础生成器
提供所有PPT生成器的公共基类和样式管理
"""

import os
import sys
from abc import ABC, abstractmethod
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 金蝶官方配色
KINGDEE_COLORS = {
    'primary': RGBColor(0x08, 0x86, 0xEC),      # 主蓝色
    'secondary': RGBColor(0x00, 0x70, 0xC0),    # 次蓝色
    'accent': RGBColor(0xFF, 0x74, 0x01),       # 橙色强调
    'dark': RGBColor(0x00, 0x3F, 0x56),         # 深蓝
    'light': RGBColor(0xCC, 0xDD, 0xEA),        # 浅蓝背景
    'text_dark': RGBColor(0x00, 0x00, 0x00),    # 深色文字
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),   # 浅色文字
    'purple': RGBColor(0xB0, 0x73, 0xFC),       # 紫色
    'green': RGBColor(0x00, 0xB4, 0x2D),        # 绿色
    'red': RGBColor(0xE5, 0x39, 0x35),          # 红色
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x80, 0x80, 0x80),
    'light_gray': RGBColor(0xF0, 0xF0, 0xF0),
}

# 字体配置
FONTS = {
    'title': {'name': '微软雅黑', 'size': 44, 'bold': True},
    'subtitle': {'name': '微软雅黑', 'size': 28, 'bold': True},
    'heading1': {'name': '微软雅黑', 'size': 32, 'bold': True},
    'heading2': {'name': '微软雅黑', 'size': 24, 'bold': True},
    'heading3': {'name': '微软雅黑', 'size': 20, 'bold': True},
    'body': {'name': '微软雅黑', 'size': 16, 'bold': False},
    'caption': {'name': '微软雅黑', 'size': 12, 'bold': False},
}


@dataclass
class PPTConfig:
    """PPT配置"""
    company_name: str
    project_name: str
    output_path: str
    template_path: Optional[str] = None
    author: str = ""
    date: str = ""


class StyleManager:
    """样式管理器"""
    
    def __init__(self, colors: Dict = None, fonts: Dict = None):
        self.colors = colors or KINGDEE_COLORS
        self.fonts = fonts or FONTS
    
    def get_color(self, name: str) -> RGBColor:
        """获取颜色"""
        return self.colors.get(name, self.colors['text_dark'])
    
    def get_font(self, name: str) -> Dict:
        """获取字体配置"""
        return self.fonts.get(name, self.fonts['body'])
    
    def apply_font(self, text_frame, style_name: str, color_name: str = 'text_dark'):
        """应用字体样式"""
        font_config = self.get_font(style_name)
        color = self.get_color(color_name)
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = font_config['name']
            paragraph.font.size = Pt(font_config['size'])
            paragraph.font.bold = font_config.get('bold', False)
            paragraph.font.color.rgb = color


class LayoutManager:
    """布局管理器"""
    
    def __init__(self, presentation: Presentation):
        self.prs = presentation
        self.layouts = {}
        self._init_layouts()
    
    def _init_layouts(self):
        """初始化布局"""
        for i, layout in enumerate(self.prs.slide_master.slide_layouts):
            self.layouts[layout.name] = i
    
    def get_layout(self, name: str, default: int = 0) -> int:
        """获取布局索引"""
        # 尝试精确匹配
        if name in self.layouts:
            return self.layouts[name]
        
        # 尝试模糊匹配
        for layout_name, idx in self.layouts.items():
            if name in layout_name or layout_name in name:
                return idx
        
        return default
    
    def add_slide(self, layout_name: str = None, layout_index: int = None):
        """添加幻灯片"""
        if layout_index is not None:
            layout = self.prs.slide_master.slide_layouts[layout_index]
        elif layout_name:
            layout = self.prs.slide_master.slide_layouts[self.get_layout(layout_name)]
        else:
            layout = self.prs.slide_master.slide_layouts[0]
        
        return self.prs.slides.add_slide(layout)


class ShapeBuilder:
    """形状构建器"""
    
    def __init__(self, slide, style: StyleManager):
        self.slide = slide
        self.style = style
    
    def add_textbox(self, text: str, left: float, top: float, 
                    width: float, height: float, style_name: str = 'body',
                    color_name: str = 'text_dark', align: str = 'left') -> Any:
        """添加文本框"""
        textbox = self.slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = self._get_alignment(align)
        
        font_config = self.style.get_font(style_name)
        p.font.name = font_config['name']
        p.font.size = Pt(font_config['size'])
        p.font.bold = font_config.get('bold', False)
        p.font.color.rgb = self.style.get_color(color_name)
        
        return textbox
    
    def add_card(self, title: str, content: str, left: float, top: float,
                 width: float, height: float, bg_color: str = 'light') -> Any:
        """添加卡片"""
        # 背景
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.style.get_color(bg_color)
        shape.line.color.rgb = self.style.get_color('primary')
        
        # 标题
        self.add_textbox(title, left + 0.1, top + 0.1, width - 0.2, 0.4,
                        'heading3', 'primary', 'left')
        
        # 内容
        self.add_textbox(content, left + 0.1, top + 0.5, width - 0.2, height - 0.6,
                        'body', 'text_dark', 'left')
        
        return shape
    
    def add_icon_card(self, icon: str, title: str, content: str,
                      left: float, top: float, width: float = 2.5, height: float = 1.5):
        """添加图标卡片"""
        # 背景
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.style.get_color('white')
        shape.line.color.rgb = self.style.get_color('light')
        
        # 图标
        self.add_textbox(icon, left + 0.1, top + 0.1, 0.5, 0.5, 'heading2')
        
        # 标题
        self.add_textbox(title, left + 0.6, top + 0.1, width - 0.7, 0.3, 'heading3')
        
        # 内容
        self.add_textbox(content, left + 0.1, top + 0.5, width - 0.2, height - 0.6,
                        'caption', 'gray')
        
        return shape
    
    def add_timeline(self, items: List[Dict], left: float, top: float,
                     width: float = 9, item_width: float = 2):
        """添加时间线"""
        x = left
        
        for i, item in enumerate(items):
            # 圆点
            circle = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + item_width/2 - 0.15), Inches(top), Inches(0.3), Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.style.get_color('primary')
            circle.line.fill.background()
            
            # 连接线
            if i < len(items) - 1:
                line = self.slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x + item_width/2 + 0.15), Inches(top + 0.12),
                    Inches(item_width - 0.3), Inches(0.06)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = self.style.get_color('light')
                line.line.fill.background()
            
            # 标题
            self.add_textbox(item.get('title', ''), x, top + 0.4, item_width, 0.3,
                           'body', 'text_dark', 'center')
            
            # 内容
            self.add_textbox(item.get('content', ''), x, top + 0.7, item_width, 0.5,
                           'caption', 'gray', 'center')
            
            x += item_width
    
    def _get_alignment(self, align: str):
        """获取对齐方式"""
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
        }
        return align_map.get(align, PP_ALIGN.LEFT)


class BasePPTGenerator(ABC):
    """PPT生成器基类"""
    
    # 生成器元信息
    name: str = "base"
    version: str = "1.0"
    description: str = "基础生成器"
    pages: int = 0
    
    def __init__(self, config: PPTConfig):
        self.config = config
        self.output_path = Path(config.output_path)
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # 初始化组件
        self.style = StyleManager()
        self.prs = self._create_presentation()
        self.layout = LayoutManager(self.prs)
        
        # 记录幻灯片
        self.slides = []
        
        # 初始化日期
        if not config.date:
            config.date = datetime.now().strftime("%Y.%m.%d")
    
    def _create_presentation(self) -> Presentation:
        """创建演示文稿"""
        if self.config.template_path and Path(self.config.template_path).exists():
            return Presentation(self.config.template_path)
        
        # 使用默认模板
        default_template = Path(__file__).parent.parent / 'templates' / 'kingdee-template-clean.pptx'
        if default_template.exists():
            return Presentation(str(default_template))
        
        return Presentation()
    
    def add_slide(self, layout_name: str = None, layout_index: int = None):
        """添加幻灯片"""
        slide = self.layout.add_slide(layout_name, layout_index)
        self.slides.append(slide)
        return slide
    
    def create_shape_builder(self, slide) -> ShapeBuilder:
        """创建形状构建器"""
        return ShapeBuilder(slide, self.style)
    
    @abstractmethod
    def generate(self) -> str:
        """生成PPT - 子类实现"""
        pass
    
    def save(self) -> str:
        """保存PPT"""
        self.prs.save(str(self.output_path))
        return str(self.output_path)
    
    def get_info(self) -> Dict:
        """获取生成器信息"""
        return {
            'name': self.name,
            'version': self.version,
            'description': self.description,
            'pages': self.pages,
            'config': {
                'company_name': self.config.company_name,
                'project_name': self.config.project_name,
                'output_path': str(self.config.output_path),
            }
        }


class GeneratorRegistry:
    """生成器注册表"""
    
    _generators: Dict[str, type] = {}
    
    @classmethod
    def register(cls, name: str):
        """注册装饰器"""
        def decorator(generator_class):
            cls._generators[name] = generator_class
            return generator_class
        return decorator
    
    @classmethod
    def get(cls, name: str) -> Optional[type]:
        """获取生成器类"""
        return cls._generators.get(name)
    
    @classmethod
    def list(cls) -> List[Dict]:
        """列出所有生成器"""
        return [
            {'name': name, 'class': gen.__name__, 'info': gen(PPTConfig('', '', '')).get_info()}
            for name, gen in cls._generators.items()
        ]
    
    @classmethod
    def create(cls, name: str, config: PPTConfig) -> Optional[BasePPTGenerator]:
        """创建生成器实例"""
        generator_class = cls.get(name)
        if generator_class:
            return generator_class(config)
        return None


# 使用示例
if __name__ == '__main__':
    # 配置
    config = PPTConfig(
        company_name="测试公司",
        project_name="测试项目",
        output_path="output/test.pptx"
    )
    
    # 列出所有生成器
    print("可用生成器:")
    for gen in GeneratorRegistry.list():
        print(f"  - {gen['name']}: {gen['info']['description']}")
    
    print("\n基类测试完成")
