# -*- coding: utf-8 -*-
"""
金蝶PPT基础生成器 V1.0
所有PPT统一使用金蝶母版和配色
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 金蝶官方配色
KINGDEE_COLORS = {
    'primary': RGBColor(0x08, 0x86, 0xEC),      # 主蓝色
    'secondary': RGBColor(0x00, 0x70, 0xC0),    # 次蓝色
    'accent': RGBColor(0xFF, 0x74, 0x01),       # 橙色强调
    'dark': RGBColor(0x00, 0x3F, 0x56),         # 深蓝
    'light': RGBColor(0xCC, 0xDD, 0xEA),        # 浅蓝背景
    'text_dark': RGBColor(0x00, 0x00, 0x00),
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),
}

class KingdeePPTBase:
    """金蝶PPT基础类"""
    
    def __init__(self, output_path, template_name='kingdee-template-clean.pptx'):
        self.output_path = output_path
        template_path = os.path.join(os.path.dirname(__file__), '..', 'templates', template_name)
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.layouts = {}
        for i, layout in enumerate(self.prs.slide_master.slide_layouts):
            self.layouts[layout.name] = i
        
        self.colors = KINGDEE_COLORS
        self.slides_count = 0
    
    def add_slide(self, layout_name='3_白色内页'):
        """添加幻灯片"""
        for name, idx in self.layouts.items():
            if layout_name in name or name in layout_name:
                slide = self.prs.slides.add_slide(self.prs.slide_master.slide_layouts[idx])
                self.slides_count += 1
                return slide
        slide = self.prs.slides.add_slide(self.prs.slide_master.slide_layouts[2])
        self.slides_count += 1
        return slide
    
    def set_title(self, slide, title):
        """设置标题"""
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = title
                    shape.text_frame.paragraphs[0].font.size = Pt(24)
                    shape.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
                    return
        self.add_text(slide, 0.3, 0.2, 12, 0.6, title, 24, True, self.colors['dark'])
    
    def add_rect(self, slide, left, top, width, height, fill_color, text="", font_size=11, font_color=None):
        """添加圆角矩形"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or self.colors['text_dark']
            p.font.name = "微软雅黑"
        return shape
    
    def add_text(self, slide, left, top, width, height, text, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        """添加文本"""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.colors['text_dark']
        p.font.name = "微软雅黑"
        p.alignment = align
        return box
    
    def add_kpi(self, slide, left, top, width, height, value, unit, label):
        """添加KPI卡片"""
        self.add_rect(slide, left, top, width, height, self.colors['light'])
        self.add_text(slide, left+0.1, top+0.12, width-0.2, 0.5, f"{value}{unit}", 22, True, self.colors['primary'], PP_ALIGN.CENTER)
        self.add_text(slide, left+0.1, top+height-0.32, width-0.2, 0.3, label, 10, False, self.colors['text_dark'], PP_ALIGN.CENTER)
    
    def add_block(self, slide, left, top, width, height, title, items, accent=False):
        """添加功能模块"""
        bg = self.colors['accent'] if accent else self.colors['primary']
        self.add_rect(slide, left, top, width, 0.32, bg, title, 10, self.colors['text_light'])
        for i, item in enumerate(items[:int((height-0.38)/0.26)]):
            self.add_rect(slide, left, top+0.35+i*0.25, width, 0.23, self.colors['light'], item, 8)
    
    def add_chapter(self, num, title):
        """添加章节页"""
        slide = self.add_slide('目录、提纲')
        self.add_text(slide, 0.5, 2.3, 2, 1.5, f"{num:02d}", 72, True, self.colors['primary'])
        self.add_text(slide, 2.5, 2.6, 9, 1, title, 36, True, self.colors['dark'])
        return slide
    
    def save(self):
        """保存PPT"""
        self.prs.save(self.output_path)
        return self.output_path
