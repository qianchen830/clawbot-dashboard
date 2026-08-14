#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶验收汇报PPT生成器 - v12.0 专业版（25页）
基于真实金蝶项目PPT母版，整合高级模板库，数据驱动
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 配色方案
C = {
    'pri': RGBColor(0, 82, 147),
    'pri2': RGBColor(0, 112, 192),
    'pri_dk': RGBColor(0, 62, 117),
    'pri_lt': RGBColor(0, 142, 212),
    'acc': RGBColor(255, 153, 0),
    'txt': RGBColor(51, 51, 51),
    'txt2': RGBColor(102, 102, 102),
    'w': RGBColor(255, 255, 255),
    'lbg': RGBColor(245, 247, 250),
    'cbg': RGBColor(248, 250, 252),
    'ok': RGBColor(0, 176, 80),
    'ok2': RGBColor(46, 139, 87),
    'warn': RGBColor(255, 192, 0),
    'err': RGBColor(244, 67, 54),
    'err2': RGBColor(204, 51, 51),
    'ablue': RGBColor(51, 102, 204),
    'bd': RGBColor(220, 220, 220),
    'lblue': RGBColor(230, 243, 255),
    'lgrn': RGBColor(230, 255, 240),
    'lred': RGBColor(255, 230, 230),
    'lorng': RGBColor(255, 243, 224),
}

CMAP = {'green': C['ok2'], 'blue': C['ablue'], 'red': C['err2'], 'orange': C['acc'], 'primary': C['pri']}

def _tb(slide, x, y, w, h):
    """快捷创建文本框"""
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def _bg_shape(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    """快捷创建背景形状"""
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _header(slide, title, text_color=C['w']):
    """快捷添加标题栏"""
    _bg_shape(slide, 0, 0, 13.333, 1.0, C['pri'])
    tb = _tb(slide, 0.5, 0.25, 12.333, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = text_color

def _section_header(slide, title, text_color=C['w']):
    """章节页标题栏"""
    _bg_shape(slide, 0, 0, 13.333, 1.2, C['pri'])
    tb = _tb(slide, 0.5, 0.3, 12.333, 0.7)
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = text_color

def _add_para(tf, text, size=16, bold=False, color=None, space_before=0, indent=False):
    """快捷添加段落"""
    p = tf.add_paragraph()
    if text.startswith('###'):
        p.text = text.replace('###', '').strip()
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color or C['pri']
        p.space_before = Pt(15)
    elif text.startswith('- '):
        p.text = "• " + text[2:]
        p.font.size = Pt(size)
        p.font.color.rgb = color or C['txt']
        p.space_before = Pt(space_before or 8)
    elif text.strip():
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or C['txt']
        p.space_before = Pt(space_before)
    return p

def _bullet_list(slide, items, x=0.8, y=1.3, w=11.733, h=5.8):
    """快捷添加要点列表"""
    tb = _tb(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            if item.startswith('###'):
                tf.paragraphs[0].text = item.replace('###', '').strip()
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = C['pri']
            elif item.startswith('- '):
                tf.paragraphs[0].text = "• " + item[2:]
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = C['txt']
            else:
                tf.paragraphs[0].text = item
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = C['txt']
        else:
            _add_para(tf, item)
    return tb

# ============================================
# 基础幻灯片（保持原有接口兼容）
# ============================================

def add_cover_slide(prs, title, subtitle="", company="", date=""):
    """添加封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_shape(slide, 0, 0, 13.333, 7.5, C['pri'])
    # 装饰线条
    _bg_shape(slide, 0.8, 2.8, 0.1, 1.5, C['acc'])
    # 主标题
    tb = _tb(slide, 1, 2.8, 11.333, 1.2)
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = C['w']
    p.alignment = PP_ALIGN.LEFT
    # 副标题
    if subtitle:
        tb2 = _tb(slide, 1, 4.2, 11.333, 0.6)
        p = tb2.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = C['w']
        p.alignment = PP_ALIGN.LEFT
    # 底部信息
    parts = []
    if company:
        parts.append(f"企业：{company}")
    parts.append(date or datetime.now().strftime('%Y年%m月'))
    if parts:
        tb3 = _tb(slide, 1, 6.3, 11.333, 0.5)
        p = tb3.text_frame.paragraphs[0]
        p.text = "  |  ".join(parts)
        p.font.size = Pt(16)
        p.font.color.rgb = C['w']
    return slide

def add_contents_slide(prs, title="目录", sections=[]):
    """添加目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _section_header(slide, title)
    if not sections:
        sections = ["一、项目概述", "二、验收标准", "三、验收结果", "四、验收结论", "五、后续计划"]
    for i, sec in enumerate(sections, 1):
        y = 1.8 + (i - 1) * 0.8
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y), Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C['pri']
        circ.line.fill.background()
        tn = _tb(slide, 1, y, 0.5, 0.5)
        tn.text_frame.paragraphs[0].text = f"0{i}" if i < 10 else str(i)
        tn.text_frame.paragraphs[0].font.size = Pt(16)
        tn.text_frame.paragraphs[0].font.bold = True
        tn.text_frame.paragraphs[0].font.color.rgb = C['w']
        tn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tn.text_frame.paragraphs[0].space_before = Pt(6)
        ts = _tb(slide, 1.7, y, 10, 0.5)
        ts.text_frame.paragraphs[0].text = sec
        ts.text_frame.paragraphs[0].font.size = Pt(20)
        ts.text_frame.paragraphs[0].font.color.rgb = C['txt']
    return slide

def add_section_slide(prs, number, title, subtitle=""):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_shape(slide, 0, 0, 13.333, 7.5, C['pri2'])
    tn = _tb(slide, 1, 2.5, 2, 1.2)
    tn.text_frame.paragraphs[0].text = f"0{number}" if number < 10 else str(number)
    tn.text_frame.paragraphs[0].font.size = Pt(72)
    tn.text_frame.paragraphs[0].font.bold = True
    tn.text_frame.paragraphs[0].font.color.rgb = C['w']
    tt = _tb(slide, 3.5, 2.8, 9, 0.8)
    tt.text_frame.paragraphs[0].text = title
    tt.text_frame.paragraphs[0].font.size = Pt(44)
    tt.text_frame.paragraphs[0].font.bold = True
    tt.text_frame.paragraphs[0].font.color.rgb = C['w']
    if subtitle:
        ts = _tb(slide, 3.5, 3.8, 9, 0.6)
        ts.text_frame.paragraphs[0].text = subtitle
        ts.text_frame.paragraphs[0].font.size = Pt(20)
        ts.text_frame.paragraphs[0].font.color.rgb = C['w']
    return slide

def add_content_slide(prs, title, content_type="bullet", items=[], left_items=[], right_items=[]):
    """添加内容页（兼容旧接口）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    if content_type == "bullet":
        _bullet_list(slide, items)
    elif content_type == "two-column":
        for side, x in [(left_items, 0.5), (right_items, 6.8)]:
            tb = _tb(slide, x, 1.3, 6, 5.8)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = side.get('title', '')
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = C['pri']
            tf.paragraphs[0].space_before = Pt(10)
            for item in side.get('items', []):
                _add_para(tf, item, size=14, space_before=6)
    return slide

def add_data_dashboard_slide(prs, title, data_cards=[]):
    """添加数据看板页（兼容旧接口）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    if not data_cards:
        data_cards = [
            {"label": "功能完成度", "value": "100%", "status": "success"},
            {"label": "性能达标", "value": "98%", "status": "success"},
            {"label": "用户满意度", "value": "96%", "status": "success"},
            {"label": "数据准确率", "value": "99%", "status": "success"},
        ]
    cw, ch, sx, sy, gap = 2.8, 1.8, 0.7, 1.5, 0.3
    for i, card in enumerate(data_cards[:8]):
        col, row = i % 4, i // 4
        x = sx + col * (cw + gap)
        y = sy + row * (ch + gap)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C['lbg']
        bg.line.color.rgb = C['bd']
        sc = C['ok'] if card.get('status') == 'success' else C['warn']
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 2.2), Inches(y + 0.2), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = sc
        dot.line.fill.background()
        lb = _tb(slide, x + 0.1, y + 0.15, cw - 0.2, 0.3)
        lb.text_frame.paragraphs[0].text = card.get('label', '')
        lb.text_frame.paragraphs[0].font.size = Pt(12)
        lb.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        vl = _tb(slide, x + 0.1, y + 0.5, cw - 0.2, 0.6)
        vl.text_frame.paragraphs[0].text = card.get('value', '')
        vl.text_frame.paragraphs[0].font.size = Pt(32)
        vl.text_frame.paragraphs[0].font.bold = True
        vl.text_frame.paragraphs[0].font.color.rgb = C['pri']
    return slide

def add_checklist_slide(prs, title, checklist_items=[]):
    """添加检查清单页（兼容旧接口）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    _bullet_list(slide, checklist_items)
    return slide

def add_comparison_slide(prs, title, before_title="验收前", after_title="验收后", before_items=[], after_items=[]):
    """添加对比页（兼容旧接口）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    # 左侧
    _bg_shape(slide, 0.3, 1.3, 6.2, 5.8, C['lred'])
    tl = _tb(slide, 0.5, 1.5, 5.8, 0.5)
    tl.text_frame.paragraphs[0].text = f"❌ {before_title}"
    tl.text_frame.paragraphs[0].font.size = Pt(22)
    tl.text_frame.paragraphs[0].font.bold = True
    tl.text_frame.paragraphs[0].font.color.rgb = C['err2']
    cl = _tb(slide, 0.5, 2.2, 5.8, 4.5)
    cl.text_frame.word_wrap = True
    for i, item in enumerate(before_items):
        p = cl.text_frame.paragraphs[0] if i == 0 else cl.text_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = C['txt']
        p.space_before = Pt(8)
    # 右侧
    _bg_shape(slide, 6.8, 1.3, 6.2, 5.8, C['lgrn'])
    tr = _tb(slide, 7, 1.5, 5.8, 0.5)
    tr.text_frame.paragraphs[0].text = f"✅ {after_title}"
    tr.text_frame.paragraphs[0].font.size = Pt(22)
    tr.text_frame.paragraphs[0].font.bold = True
    tr.text_frame.paragraphs[0].font.color.rgb = C['ok2']
    cr = _tb(slide, 7, 2.2, 5.8, 4.5)
    cr.text_frame.word_wrap = True
    for i, item in enumerate(after_items):
        p = cr.text_frame.paragraphs[0] if i == 0 else cr.text_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = C['txt']
        p.space_before = Pt(8)
    return slide

def add_timeline_slide(prs, title, phases=[]):
    """添加时间轴页（兼容旧接口）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    if not phases:
        phases = [
            {"name": "启动准备", "time": "第1周", "work": "组建团队\n制定计划"},
            {"name": "功能测试", "time": "第2-3周", "work": "功能测试\n集成测试"},
            {"name": "用户验收", "time": "第4周", "work": "用户培训\n用户测试"},
            {"name": "正式验收", "time": "第5周", "work": "验收会议\n验收签字"}
        ]
    sx, yp, gap = 1.5, 2.5, 2.5
    _bg_shape(slide, sx, yp + 0.25, 10, 0.05, C['pri'])
    for i, ph in enumerate(phases):
        x = sx + i * gap
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(yp), Inches(0.6), Inches(0.6))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C['pri']
        circ.line.fill.background()
        tn = _tb(slide, x, yp + 0.8, 2, 0.4)
        tn.text_frame.paragraphs[0].text = ph.get('name', '')
        tn.text_frame.paragraphs[0].font.size = Pt(16)
        tn.text_frame.paragraphs[0].font.bold = True
        tn.text_frame.paragraphs[0].font.color.rgb = C['pri']
        tn.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tt = _tb(slide, x, yp + 1.2, 2, 0.3)
        tt.text_frame.paragraphs[0].text = ph.get('time', '')
        tt.text_frame.paragraphs[0].font.size = Pt(12)
        tt.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        tt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tw = _tb(slide, x, yp + 1.6, 2, 1.2)
        tw.text_frame.word_wrap = True
        tw.text_frame.paragraphs[0].text = ph.get('work', '')
        tw.text_frame.paragraphs[0].font.size = Pt(11)
        tw.text_frame.paragraphs[0].font.color.rgb = C['txt']
        tw.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_ending_slide(prs, title="验收结论", subtitle="项目验收通过", company="金蝶软件（中国）有限公司"):
    """添加结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_shape(slide, 0, 0, 13.333, 7.5, C['pri'])
    tb = _tb(slide, 1, 2.5, 11.333, 1.2)
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(56)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = C['w']
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if subtitle:
        tb2 = _tb(slide, 1, 4, 11.333, 0.6)
        tb2.text_frame.paragraphs[0].text = subtitle
        tb2.text_frame.paragraphs[0].font.size = Pt(24)
        tb2.text_frame.paragraphs[0].font.color.rgb = C['w']
        tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb3 = _tb(slide, 1, 6, 11.333, 0.5)
    tb3.text_frame.paragraphs[0].text = company
    tb3.text_frame.paragraphs[0].font.size = Pt(18)
    tb3.text_frame.paragraphs[0].font.color.rgb = C['w']
    tb3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

# ============================================
# 高级模板（从高级模板库整合）
# ============================================

def add_big_number_slide(prs, title, big_numbers):
    """大数字展示页 - 最多4个"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    n = min(len(big_numbers), 4)
    nw, hg = 2.8, 0.5
    tw = n * nw + (n - 1) * hg
    sx = (13.333 - tw) / 2
    for i, num in enumerate(big_numbers[:n]):
        x = sx + i * (nw + hg)
        y = 2.0
        clr = CMAP.get(num.get('color', 'blue'), C['pri'])
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(y), Inches(2), Inches(2))
        circ.fill.solid()
        circ.fill.fore_color.rgb = clr
        circ.line.fill.background()
        vt = _tb(slide, x + 0.4, y + 0.5, 2, 1.0)
        vt.text_frame.paragraphs[0].text = num.get('value', '0%')
        vt.text_frame.paragraphs[0].font.size = Pt(48)
        vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = C['w']
        vt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lt = _tb(slide, x, y + 2.3, nw, 0.5)
        lt.text_frame.paragraphs[0].text = num.get('label', '')
        lt.text_frame.paragraphs[0].font.size = Pt(18)
        lt.text_frame.paragraphs[0].font.bold = True
        lt.text_frame.paragraphs[0].font.color.rgb = C['txt']
        lt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        dt = _tb(slide, x, y + 2.8, nw, 0.5)
        dt.text_frame.paragraphs[0].text = num.get('desc', '')
        dt.text_frame.paragraphs[0].font.size = Pt(14)
        dt.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        dt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_kpi_card_slide(prs, title, kpi_data):
    """KPI指标卡片页 - 最多6个"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    cn = min(len(kpi_data), 6)
    cols = 3 if cn > 3 else cn
    cw, ch, hg, vg = 3.8, 2.5, 0.5, 0.4
    tw = cols * cw + (cols - 1) * hg
    sx = (13.333 - tw) / 2
    sy = 1.5
    for i, kpi in enumerate(kpi_data[:cn]):
        row, col = i // cols, i % cols
        x = sx + col * (cw + hg)
        y = sy + row * (ch + vg)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.fill.solid()
        card.fill.fore_color.rgb = C['cbg']
        card.line.color.rgb = C['bd']
        bar_clr = CMAP.get(kpi.get('color', 'blue'), C['pri'])
        _bg_shape(slide, x, y, 0.08, ch, bar_clr)
        lt = _tb(slide, x + 0.3, y + 0.3, cw - 0.5, 0.5)
        lt.text_frame.paragraphs[0].text = kpi.get('label', '')
        lt.text_frame.paragraphs[0].font.size = Pt(14)
        lt.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        vt = _tb(slide, x + 0.3, y + 0.9, cw - 0.5, 0.8)
        vt.text_frame.paragraphs[0].text = kpi.get('value', '')
        vt.text_frame.paragraphs[0].font.size = Pt(36)
        vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = C['txt']
        if kpi.get('change'):
            ct = _tb(slide, x + 0.3, y + 1.8, cw - 0.5, 0.4)
            ct.text_frame.paragraphs[0].text = kpi['change']
            ct.text_frame.paragraphs[0].font.size = Pt(14)
            ct.text_frame.paragraphs[0].font.bold = True
            pos = any(k in kpi['change'] for k in ['+', '↑', '缩短', '降低'])
            ct.text_frame.paragraphs[0].font.color.rgb = C['ok2'] if pos else C['err2']
    return slide

def add_table_slide(prs, title, headers, rows, highlight_cols=None):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    nc = len(headers)
    nr = len(rows) + 1
    tw = 12.0
    th = min(5.5, 0.6 + 0.5 * nr)
    tbl = slide.shapes.add_table(nr, nc, Inches(0.667), Inches(1.5), Inches(tw), Inches(th)).table
    cw_each = Inches(tw // nc)
    for j in range(nc):
        tbl.columns[j].width = cw_each
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C['pri']
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = C['w']
            p.alignment = PP_ALIGN.CENTER
    for i, rd in enumerate(rows):
        for j, cv in enumerate(rd):
            cell = tbl.cell(i + 1, j)
            cell.text = str(cv)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C['cbg'] if i % 2 == 0 else C['w']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = C['pri'] if (highlight_cols and j in highlight_cols) else C['txt']
                p.font.bold = bool(highlight_cols and j in highlight_cols)
                p.alignment = PP_ALIGN.CENTER
    return slide

def add_three_column_slide(prs, title, col1, col2, col3):
    """三栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    cw, hg = 3.8, 0.5
    tw = 3 * cw + 2 * hg
    sx = (13.333 - tw) / 2
    for idx, col in enumerate([col1, col2, col3]):
        x = sx + idx * (cw + hg)
        y = 1.5
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C['cbg']
        card.line.color.rgb = C['bd']
        icon = col.get('icon', '')
        tt = _tb(slide, x + 0.2, y + 0.2, cw - 0.4, 0.6)
        tt.text_frame.paragraphs[0].text = f"{icon} {col.get('title', '')}" if icon else col.get('title', '')
        tt.text_frame.paragraphs[0].font.size = Pt(20)
        tt.text_frame.paragraphs[0].font.bold = True
        tt.text_frame.paragraphs[0].font.color.rgb = C['pri']
        sep = _bg_shape(slide, x + 0.2, y + 0.8, cw - 0.4, 0.02, C['pri_lt'])
        it = _tb(slide, x + 0.2, y + 1.0, cw - 0.4, 4.2)
        tf = it.text_frame
        tf.word_wrap = True
        for j, item in enumerate(col.get('items', [])):
            if j == 0:
                if item.startswith('###'):
                    tf.paragraphs[0].text = item.replace('###', '').strip()
                    tf.paragraphs[0].font.size = Pt(15)
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.color.rgb = C['pri']
                else:
                    tf.paragraphs[0].text = ("• " + item[2:]) if item.startswith('- ') else item
                    tf.paragraphs[0].font.size = Pt(13)
                    tf.paragraphs[0].font.color.rgb = C['txt']
            else:
                _add_para(tf, item, size=13, space_before=4)
    return slide

def add_icon_card_slide(prs, title, cards):
    """图标卡片页 - 最多6个"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    cn = min(len(cards), 6)
    cw, ch, hg, vg = 3.8, 2.5, 0.5, 0.4
    tw = 3 * cw + 2 * hg
    sx = (13.333 - tw) / 2
    sy = 1.5
    for i, card in enumerate(cards[:cn]):
        row, col = i // 3, i % 3
        x = sx + col * (cw + hg)
        y = sy + row * (ch + vg)
        cs = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        cs.fill.solid()
        cs.fill.fore_color.rgb = C['w']
        cs.line.color.rgb = C['bd']
        _bg_shape(slide, x, y, cw, 0.05, C['pri'])
        ic = _tb(slide, x, y + 0.2, cw, 0.6)
        ic.text_frame.paragraphs[0].text = card.get('icon', '📌')
        ic.text_frame.paragraphs[0].font.size = Pt(32)
        ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tt = _tb(slide, x + 0.2, y + 0.9, cw - 0.4, 0.5)
        tt.text_frame.paragraphs[0].text = card.get('title', '')
        tt.text_frame.paragraphs[0].font.size = Pt(16)
        tt.text_frame.paragraphs[0].font.bold = True
        tt.text_frame.paragraphs[0].font.color.rgb = C['txt']
        tt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        dt = _tb(slide, x + 0.2, y + 1.4, cw - 0.4, 1.0)
        dt.text_frame.paragraphs[0].text = card.get('desc', '')
        dt.text_frame.paragraphs[0].font.size = Pt(12)
        dt.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        dt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_process_flow_slide(prs, title, steps):
    """流程图页 - 最多6步"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    sn = min(len(steps), 6)
    sw, aw = 1.8, 0.3
    tw = sn * sw + (sn - 1) * aw
    sx = (13.333 - tw) / 2
    sy = 2.5
    flow_colors = [C['pri'], C['pri_lt'], C['acc'], C['ok2'], C['ablue'], C['pri2']]
    for i, step in enumerate(steps[:sn]):
        x = sx + i * (sw + aw)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(sy), Inches(1.0), Inches(1.0))
        circ.fill.solid()
        circ.fill.fore_color.rgb = flow_colors[i % len(flow_colors)]
        circ.line.fill.background()
        nt = _tb(slide, x + 0.4, sy + 0.2, 1.0, 0.6)
        nt.text_frame.paragraphs[0].text = str(i + 1)
        nt.text_frame.paragraphs[0].font.size = Pt(32)
        nt.text_frame.paragraphs[0].font.bold = True
        nt.text_frame.paragraphs[0].font.color.rgb = C['w']
        nt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        nm = _tb(slide, x, sy + 1.3, sw, 0.6)
        nm.text_frame.paragraphs[0].text = step.get('name', '')
        nm.text_frame.paragraphs[0].font.size = Pt(16)
        nm.text_frame.paragraphs[0].font.bold = True
        nm.text_frame.paragraphs[0].font.color.rgb = C['txt']
        nm.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        dc = _tb(slide, x, sy + 1.9, sw, 1.5)
        dc.text_frame.word_wrap = True
        dc.text_frame.paragraphs[0].text = step.get('desc', '')
        dc.text_frame.paragraphs[0].font.size = Pt(12)
        dc.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        dc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < sn - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + sw + 0.05), Inches(sy + 0.3), Inches(0.2), Inches(0.4))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C['bd']
            arr.line.fill.background()
    return slide

def add_quote_slide(prs, title, quote, author):
    """引用页（客户评价）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_shape(slide, 0, 0, 13.333, 7.5, C['lblue'])
    if title:
        _section_header(slide, title)
    lq = _tb(slide, 0.5, 1.5, 1, 1)
    lq.text_frame.paragraphs[0].text = '"'
    lq.text_frame.paragraphs[0].font.size = Pt(120)
    lq.text_frame.paragraphs[0].font.color.rgb = C['pri_lt']
    qt = _tb(slide, 1.5, 2.5, 10.333, 2.5)
    qt.text_frame.word_wrap = True
    qt.text_frame.paragraphs[0].text = quote
    qt.text_frame.paragraphs[0].font.size = Pt(28)
    qt.text_frame.paragraphs[0].font.italic = True
    qt.text_frame.paragraphs[0].font.color.rgb = C['txt']
    qt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    at = _tb(slide, 1.5, 5.5, 10.333, 0.5)
    at.text_frame.paragraphs[0].text = f"— {author}"
    at.text_frame.paragraphs[0].font.size = Pt(18)
    at.text_frame.paragraphs[0].font.color.rgb = C['pri']
    at.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide

def add_stats_grid_slide(prs, title, stats):
    """统计网格页 - 最多12个"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    sn = min(len(stats), 12)
    cols = 4
    cw, ch, hg, vg = 2.8, 2.5, 0.3, 0.3
    tw = cols * cw + (cols - 1) * hg
    sx = (13.333 - tw) / 2
    sy = 1.5
    for i, st in enumerate(stats[:sn]):
        row, col = i // cols, i % cols
        x = sx + col * (cw + hg)
        y = sy + row * (ch + vg)
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
        cell.fill.solid()
        cell.fill.fore_color.rgb = C['w']
        cell.line.color.rgb = C['bd']
        ic = _tb(slide, x, y + 0.2, cw, 0.5)
        ic.text_frame.paragraphs[0].text = st.get('icon', '📊')
        ic.text_frame.paragraphs[0].font.size = Pt(28)
        ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        vt = _tb(slide, x, y + 0.8, cw, 0.8)
        vt.text_frame.paragraphs[0].text = st.get('value', '')
        vt.text_frame.paragraphs[0].font.size = Pt(32)
        vt.text_frame.paragraphs[0].font.bold = True
        vt.text_frame.paragraphs[0].font.color.rgb = C['pri']
        vt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lt = _tb(slide, x, y + 1.6, cw, 0.6)
        lt.text_frame.paragraphs[0].text = st.get('label', '')
        lt.text_frame.paragraphs[0].font.size = Pt(14)
        lt.text_frame.paragraphs[0].font.color.rgb = C['txt2']
        lt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

# ============================================
# 主生成函数 - 25页专业版
# ============================================

def generate_acceptance_ppt_v12(customer_info):
    """生成验收汇报PPT - v12.0 专业版（25页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    cn = customer_info.get('companyName', '企业名称')
    ad = customer_info.get('acceptanceDate', datetime.now().strftime('%Y年%m月'))
    ac = customer_info.get('acceptanceConclusion', '验收通过')
    proj = customer_info.get('projectName', '金蝶云·星空')
    
    # ==================== P1: 封面 ====================
    add_cover_slide(prs, f"{cn}\n{proj}项目验收汇报", ac, date=ad)
    
    # ==================== P2: 目录 ====================
    add_contents_slide(prs, "目录", [
        "一、项目概述", "二、项目成果", "三、功能验收",
        "四、项目价值", "五、经验总结", "六、后续支持计划"
    ])
    
    # ==================== 第一章：项目概述 (P3-P7) ====================
    add_section_slide(prs, 1, "项目概述", "项目背景 | 项目目标 | 项目范围 | 项目团队 | 实施历程")
    
    # P4: 项目背景与目标
    add_content_slide(prs, "项目背景与目标", "two-column",
        left_items={
            "title": "项目背景",
            "items": [
                "- 企业数字化转型战略驱动",
                "- 业务快速发展，管理效率亟待提升",
                "- 现有系统无法满足多业务协同需求",
                "- 数据孤岛严重，决策支持不足",
                "- 合规与风控要求不断提高"
            ]
        },
        right_items={
            "title": "项目目标",
            "items": [
                "- 实现业务流程标准化，提升运营效率",
                "- 建立统一数据平台，消除信息孤岛",
                "- 提升系统集成效率，实现业务协同",
                "- 强化内控体系，降低经营风险",
                "- 支撑管理决策，推动持续发展"
            ]
        })
    
    # P5: 项目范围
    add_icon_card_slide(prs, "项目实施范围", [
        {"icon": "💰", "title": "财务云", "desc": "总账、应收应付、固定资产、\n现金管理、成本管理、预算管理"},
        {"icon": "🔗", "title": "供应链云", "desc": "采购管理、库存管理、销售管理、\n物流管理、供应商协同"},
        {"icon": "👥", "title": "人力资源云", "desc": "人事管理、薪酬管理、\n绩效管理、培训管理"},
        {"icon": "🏭", "title": "制造云", "desc": "生产计划、车间管理、\n质量管理、设备管理"},
        {"icon": "📊", "title": "数据分析平台", "desc": "BI报表、经营分析、\n预警监控、移动看板"},
        {"icon": "🔧", "title": "技术平台", "desc": "云原生架构、API网关、\n数据集成、流程集成"},
    ])
    
    # P6: 项目团队
    add_three_column_slide(prs, "项目团队与组织保障", 
        {"title": "客户团队", "icon": "🏢", "items": [
            "###项目领导小组",
            "- 项目发起人：公司高层",
            "- 项目经理：业务负责人",
            "- 关键用户：各部门骨干",
            "",
            "###职责分工",
            "- 需求确认与变更审批",
            "- 用户测试与验收确认",
            "- 培训推广与上线支持"
        ]},
        {"title": "实施团队", "icon": "💻", "items": [
            "###金蝶实施团队",
            "- 项目经理：1人",
            "- 功能顾问：3人",
            "- 技术顾问：2人",
            "- 测试工程师：1人",
            "",
            "###实施方法",
            "- 采用金蝶KingdeeWay方法论",
            "- 敏捷迭代，分阶段交付"
        ]},
        {"title": "质量保障", "icon": "🛡️", "items": [
            "###质量管理体系",
            "- 周例会机制：每周进度同步",
            "- 阶段评审：里程碑检查",
            "- 风险管理：风险识别与应对",
            "",
            "###沟通机制",
            "- 每周项目周报",
            "- 月度高层汇报",
            "- 即时通讯群组"
        ]}
    )
    
    # P7: 实施历程
    add_timeline_slide(prs, "项目实施历程", [
        {"name": "项目启动", "time": "第1-2月", "work": "需求调研\n方案设计\n项目计划"},
        {"name": "蓝图设计", "time": "第3-4月", "work": "业务梳理\n蓝图确认\n方案评审"},
        {"name": "系统实现", "time": "第5-7月", "work": "系统配置\n开发测试\n数据迁移"},
        {"name": "上线切换", "time": "第8月", "work": "用户培训\n数据导入\n上线切换"},
        {"name": "运维支持", "time": "第9月起", "work": "系统运维\n持续优化\n验收准备"}
    ])
    
    # ==================== 第二章：项目成果 (P8-P12) ====================
    add_section_slide(prs, 2, "项目成果", "核心指标 | KPI详情 | 数据对比 | 业务提升")
    
    # P9: 核心成果概览（大数字模板）
    add_big_number_slide(prs, "核心成果概览", [
        {"value": "100%", "label": "功能完成度", "desc": "所有规划功能全部上线", "color": "green"},
        {"value": "50%", "label": "效率提升", "desc": "财务结账时间大幅缩短", "color": "blue"},
        {"value": "30%", "label": "成本降低", "desc": "运营管理成本显著下降", "color": "orange"},
        {"value": "96%", "label": "用户满意度", "desc": "用户培训满意度调研", "color": "green"},
    ])
    
    # P10: 核心KPI指标（KPI卡片模板）
    add_kpi_card_slide(prs, "核心KPI指标", [
        {"label": "财务结账周期", "value": "3天", "change": "缩短70%（原10天）", "color": "green"},
        {"label": "库存周转率", "value": "8.5次", "change": "↑35%", "color": "blue"},
        {"label": "订单处理时效", "value": "2小时", "change": "缩短60%", "color": "green"},
        {"label": "采购成本节约", "value": "12%", "change": "年节约200万+", "color": "orange"},
        {"label": "审批效率", "value": "即时", "change": "缩短90%（原3天）", "color": "green"},
        {"label": "数据准确率", "value": "99.5%", "change": "↑3.5%", "color": "blue"},
    ])
    
    # P11: 系统性能指标
    add_data_dashboard_slide(prs, "系统性能指标", [
        {"label": "平均响应时间", "value": "1.8秒", "status": "success"},
        {"label": "并发用户数", "value": "200人", "status": "success"},
        {"label": "数据处理能力", "value": "20000条/分", "status": "success"},
        {"label": "系统可用性", "value": "99.95%", "status": "success"},
        {"label": "月均在线时长", "value": "720小时", "status": "success"},
        {"label": "故障恢复时间", "value": "≤15分钟", "status": "success"},
        {"label": "数据备份", "value": "每日自动", "status": "success"},
        {"label": "安全等级", "value": "等保三级", "status": "success"},
    ])
    
    # P12: 项目成果对比
    add_comparison_slide(prs, "项目成果对比",
        before_title="项目实施前",
        after_title="项目实施后",
        before_items=[
            "财务结账周期10天，月结常需加班",
            "库存数据分散，周转率偏低",
            "手工开单，订单处理耗时长",
            "采购流程冗长，成本管控弱",
            "审批环节多，流程效率低",
            "数据不准确，决策支撑不足"
        ],
        after_items=[
            "财务结账3天完成，自动化率90%",
            "库存实时可视，周转率提升35%",
            "系统自动生成，处理时效2小时",
            "电子采购平台，成本降低12%",
            "移动审批即时完成，效率提升90%",
            "数据实时准确，BI决策看板赋能"
        ])
    
    # ==================== 第三章：功能验收 (P13-P17) ====================
    add_section_slide(prs, 3, "功能验收", "验收标准 | 功能验收 | 性能验收 | 安全验收 | 用户验收")
    
    # P14: 验收标准总览
    add_table_slide(prs, "验收标准总览",
        ["验收维度", "验收标准", "验收方法", "达成情况"],
        [
            ["功能完整性", "所有规划功能100%实现", "功能测试 + 用户确认", "✅ 通过"],
            ["系统性能", "响应≤3秒，并发≥100", "性能测试工具", "✅ 通过"],
            ["数据质量", "准确率≥99%", "数据校验 + 抽检", "✅ 通过"],
            ["系统安全", "等保三级，权限完善", "安全审计 + 渗透测试", "✅ 通过"],
            ["用户满意度", "培训满意度≥90%", "问卷调查", "✅ 通过"],
            ["文档交付", "全套文档齐全", "文档评审", "✅ 通过"],
        ],
        highlight_cols=[3]
    )
    
    # P15: 功能验收详情
    add_table_slide(prs, "功能验收详情",
        ["模块", "功能项", "测试用例", "通过率"],
        [
            ["总账", "凭证/账簿/报表", "45项", "100%"],
            ["应收应付", "发票/收付款/核销", "38项", "100%"],
            ["固定资产", "资产卡片/折旧/处置", "28项", "100%"],
            ["采购管理", "申请/订单/入库/结算", "42项", "100%"],
            ["销售管理", "报价/订单/发货/收款", "40项", "100%"],
            ["库存管理", "出入库/调拨/盘点", "35项", "100%"],
            ["薪酬管理", "社保/公积金/个税/工资", "30项", "100%"],
            ["报表分析", "财务/业务/自定义报表", "25项", "100%"],
        ],
        highlight_cols=[3]
    )
    
    # P16: 性能验收结果
    add_kpi_card_slide(prs, "性能验收结果", [
        {"label": "页面加载时间", "value": "0.8秒", "change": "标准：≤3秒 ✅", "color": "green"},
        {"label": "报表生成时间", "value": "2.5秒", "change": "标准：≤5秒 ✅", "color": "green"},
        {"label": "批量导入速度", "value": "15000条/分", "change": "标准：≥5000 ✅", "color": "green"},
        {"label": "并发用户支持", "value": "200人", "change": "标准：≥100 ✅", "color": "green"},
        {"label": "系统可用性", "value": "99.95%", "change": "标准：≥99.9% ✅", "color": "green"},
        {"label": "数据恢复时间", "value": "12分钟", "change": "标准：≤30分钟 ✅", "color": "green"},
    ])
    
    # P17: 安全验收 & 用户验收
    add_three_column_slide(prs, "安全与用户验收",
        {"title": "安全验收", "icon": "🔒", "items": [
            "- 等保三级认证通过",
            "- 数据传输SSL加密",
            "- 多层次权限管理体系",
            "- 操作日志完整审计",
            "- 数据备份与容灾机制",
            "- 防SQL注入/XSS攻击"
        ]},
        {"title": "用户验收", "icon": "👥", "items": [
            "- 培训覆盖率100%",
            "- 关键用户考核通过率98%",
            "- 操作手册交付齐全",
            "- 用户满意度96%",
            "- UAT测试全部通过",
            "- 上线运行稳定无重大问题"
        ]},
        {"title": "文档验收", "icon": "📄", "items": [
            "- 需求规格说明书",
            "- 蓝图设计文档",
            "- 系统配置文档",
            "- 测试报告",
            "- 用户操作手册",
            "- 运维管理手册"
        ]}
    )
    
    # ==================== 第四章：项目价值 (P18-P21) ====================
    add_section_slide(prs, 4, "项目价值", "业务价值 | 管理价值 | 技术价值 | 客户评价")
    
    # P19: 业务价值
    add_three_column_slide(prs, "业务价值",
        {"title": "财务价值", "icon": "💰", "items": [
            "###直接经济效益",
            "- 财务结账效率提升70%",
            "- 采购成本降低12%",
            "- 库存占用减少20%",
            "- 人力成本节约15%",
            "",
            "###年度收益估算",
            "- 直接经济效益：300万+",
            "- 管理效率提升折算：100万+"
        ]},
        {"title": "运营价值", "icon": "📈", "items": [
            "###效率提升",
            "- 订单处理效率提升60%",
            "- 审批流程效率提升90%",
            "- 报表生成效率提升80%",
            "- 信息查询效率提升95%",
            "",
            "###能力提升",
            "- 数据驱动决策能力",
            "- 风险预警与管控能力"
        ]},
        {"title": "战略价值", "icon": "🎯", "items": [
            "###数字化转型",
            "- 奠定数字化管理基础",
            "- 支撑业务快速扩张",
            "- 提升企业核心竞争力",
            "",
            "###可持续发展",
            "- 标准化流程体系",
            "- 知识沉淀与传承",
            "- 持续优化的技术平台"
        ]}
    )
    
    # P20: 管理价值
    add_icon_card_slide(prs, "管理价值", [
        {"icon": "📊", "title": "实时可视", "desc": "经营数据实时呈现\n管理驾驶舱全面覆盖"},
        {"icon": "⚡", "title": "流程优化", "desc": "端到端流程打通\n审批效率大幅提升"},
        {"icon": "🎯", "title": "精准决策", "desc": "BI分析辅助决策\n数据驱动科学管理"},
        {"icon": "🛡️", "title": "风险管控", "desc": "内控体系完善\n合规风控实时预警"},
        {"icon": "🔄", "title": "协同高效", "desc": "跨部门业务协同\n打破信息壁垒"},
        {"icon": "📐", "title": "标准规范", "desc": "统一编码标准\n统一业务流程规范"},
    ])
    
    # P21: 技术价值
    add_process_flow_slide(prs, "技术架构价值", [
        {"name": "云原生", "desc": "弹性扩展\n按需使用\n降低IT投入"},
        {"name": "微服务", "desc": "独立部署\n灵活迭代\n快速响应"},
        {"name": "数据中台", "desc": "统一标准\n实时同步\n分析赋能"},
        {"name": "集成平台", "desc": "API网关\n多系统打通\n生态连接"},
        {"name": "智能分析", "desc": "AI预警\n智能报表\n趋势预测"},
        {"name": "移动办公", "desc": "移动审批\n实时查看\n随时随地"},
    ])
    
    # P22: 客户评价
    add_quote_slide(prs, "客户评价",
        f"金蝶云·星空系统的成功上线，标志着{cn}数字化转型迈出了关键一步。"
        "系统不仅满足了当前业务管理需求，更为我们未来的发展奠定了坚实的技术基础。"
        "项目团队专业、高效，实施过程规范有序，我们对项目成果非常满意。",
        f"{cn} 项目负责人"
    )
    
    # ==================== 第五章：经验总结 (P23) ====================
    add_section_slide(prs, 5, "经验总结", "成功经验 | 改进方向")
    
    add_three_column_slide(prs, "项目经验总结",
        {"title": "成功经验", "icon": "✅", "items": [
            "###高层支持",
            "- 企业高层全程参与",
            "- 资源保障到位",
            "- 决策快速高效",
            "",
            "###方法论保障",
            "- 标准化实施方法论",
            "- 阶段性评审机制",
            "- 风险提前识别应对"
        ]},
        {"title": "关键成功因素", "icon": "🔑", "items": [
            "###团队协作",
            "- 客户与实施团队紧密配合",
            "- 关键用户深度参与",
            "- 跨部门协调顺畅",
            "",
            "###变更管理",
            "- 合理控制需求变更",
            "- 充分的沟通与培训",
            "- 稳步推进上线切换"
        ]},
        {"title": "改进方向", "icon": "📝", "items": [
            "###持续优化",
            "- 部分高级功能需进一步优化",
            "- 移动端应用体验待完善",
            "- 自定义报表功能增强",
            "",
            "###未来规划",
            "- 智能化应用探索",
            "- 产业链协同扩展",
            "- 数据资产深度利用"
        ]}
    )
    
    # ==================== 第六章：后续支持计划 (P24-P25) ====================
    add_section_slide(prs, 6, "后续支持计划", "优化计划 | 培训计划 | 运维保障")
    
    # P25: 后续支持计划
    add_table_slide(prs, "后续支持计划",
        ["阶段", "时间", "工作内容", "交付成果"],
        [
            ["系统优化", "上线后1-2月", "功能优化、性能调优、体验改进", "优化报告"],
            ["深度培训", "上线后2-3月", "高级功能培训、报表开发培训", "培训记录"],
            ["常规运维", "上线后1-6月", "日常运维、故障处理、版本升级", "运维周报"],
            ["持续改进", "长期", "需求收集、功能迭代、性能优化", "迭代计划"],
            ["知识转移", "上线后3月", "管理员认证培训、知识库建设", "转移确认书"],
            ["年度回访", "每年", "系统体检、需求回顾、升级建议", "回访报告"],
        ],
        highlight_cols=[]
    )
    
    # P26: 运维保障体系
    add_three_column_slide(prs, "运维保障体系",
        {"title": "服务保障", "icon": "🛠️", "items": [
            "###运维服务",
            "- 7×24小时在线监控",
            "- 4小时响应机制",
            "- 专属运维工程师",
            "- 每月运维报告",
            "",
            "###应急保障",
            "- 紧急故障15分钟响应",
            "- 备用方案快速切换",
            "- 数据恢复保障机制"
        ]},
        {"title": "升级保障", "icon": "📦", "items": [
            "###版本管理",
            "- 季度版本更新通知",
            "- 年度大版本升级",
            "- 升级前测试验证",
            "- 升级后回归测试",
            "",
            "###补丁管理",
            "- 安全补丁及时推送",
            "- 功能补丁评估安装"
        ]},
        {"title": "持续服务", "icon": "🤝", "items": [
            "###客户成功",
            "- 专属客户成功经理",
            "- 季度业务回访",
            "- 年度系统健康检查",
            "- 最佳实践分享",
            "",
            "###社区支持",
            "- 金蝶云社区",
            "- 知识库自助查询",
            "- 同行经验交流"
        ]}
    )
    
    # ==================== P27: 结束页 ====================
    add_ending_slide(prs, "感谢聆听", f"{cn} {proj}项目验收汇报", cn)
    
    # 保存
    filename = f"验收汇报_{cn}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    
    return {
        "success": True,
        "message": f"验收汇报PPT生成成功",
        "file": filename,
        "path": filepath,
        "pages": len(prs.slides)
    }


def main():
    """主函数 - 命令行入口"""
    import argparse
    
    parser = argparse.ArgumentParser(description='金蝶验收汇报PPT生成器 v12.0 专业版（25页）')
    parser.add_argument('--companyName', required=True, help='企业名称')
    parser.add_argument('--acceptanceDate', default='', help='验收日期')
    parser.add_argument('--acceptanceConclusion', default='验收通过', help='验收结论')
    parser.add_argument('--projectName', default='金蝶云·星空', help='项目名称')
    
    args = parser.parse_args()
    
    customer_info = {
        'companyName': args.companyName,
        'acceptanceDate': args.acceptanceDate,
        'acceptanceConclusion': args.acceptanceConclusion,
        'projectName': args.projectName,
    }
    
    result = generate_acceptance_ppt_v12(customer_info)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
