#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - 公共模块 v1.0
提取三个PPT生成器的公共代码，统一管理
包含：配色方案、公共函数、公共常量
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 输出目录
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 专业配色方案（金蝶官方色系）
COLORS = {
    'primary': RGBColor(0, 82, 147),      # 金蝶蓝
    'secondary': RGBColor(0, 112, 192),    # 辅助蓝
    'accent': RGBColor(255, 153, 0),       # 强调橙
    'text_dark': RGBColor(51, 51, 51),     # 深色文字
    'text_light': RGBColor(102, 102, 102), # 浅色文字
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(245, 247, 250),   # 浅色背景
    'success': RGBColor(0, 176, 80),       # 成功绿
    'warning': RGBColor(255, 192, 0),      # 警告黄
    'error': RGBColor(244, 67, 54),        # 错误红
}

# 备用配色方案（旧版本兼容）
COLORS_LEGACY = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'white': RGBColor(255, 255, 255),
    'light': RGBColor(240, 240, 240),
}


def create_presentation():
    """创建标准PPT对象"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_cover_slide(prs, title, subtitle="", speaker="", date="", colors=None):
    """
    添加封面页 - 专业母版版
    
    Args:
        prs: PPT对象
        title: 主标题
        subtitle: 副标题
        speaker: 演讲者
        date: 日期
        colors: 配色方案（默认使用COLORS）
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['primary']
    bg.line.fill.background()
    
    # 装饰线条
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(0.1), Inches(1.5))
    line1.fill.solid()
    line1.fill.fore_color.rgb = colors['accent']
    line1.line.fill.background()
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    # 底部信息
    bottom_info = []
    if speaker:
        bottom_info.append(f"演讲者：{speaker}")
    if date:
        bottom_info.append(date)
    else:
        bottom_info.append(datetime.now().strftime('%Y年%m月'))
    
    if bottom_info:
        tb_bottom = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11.333), Inches(0.5))
        p = tb_bottom.text_frame.paragraphs[0]
        p.text = "  |  ".join(bottom_info)
        p.font.size = Pt(16)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.LEFT
    
    return slide


def add_section_slide(prs, number, title, subtitle="", colors=None):
    """
    添加章节页
    
    Args:
        prs: PPT对象
        number: 章节编号
        title: 章节标题
        subtitle: 章节副标题
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['secondary']
    bg.line.fill.background()
    
    # 编号
    tb_num = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.2))
    p = tb_num.text_frame.paragraphs[0]
    p.text = f"0{number}" if number < 10 else str(number)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(3.5), Inches(2.8), Inches(9), Inches(0.8))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(3.5), Inches(3.8), Inches(9), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = colors['white']
    
    return slide


def add_contents_slide(prs, title="目录", sections=None, colors=None):
    """
    添加目录页
    
    Args:
        prs: PPT对象
        title: 标题
        sections: 章节列表
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    if sections is None:
        sections = ["一、项目背景与目标", "二、项目建设过程", "三、项目价值实现", "四、项目成果展示", "五、合作展望"]
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 目录内容
    y_pos = 1.8
    for i, section in enumerate(sections, 1):
        # 编号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors['primary']
        circle.line.fill.background()
        
        # 编号文字
        tb_num = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5))
        p = tb_num.text_frame.paragraphs[0]
        p.text = f"0{i}" if i < 10 else str(i)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
        tb_num.text_frame.paragraphs[0].space_before = Pt(6)
        
        # 章节文字
        tb_section = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos), Inches(10), Inches(0.5))
        p = tb_section.text_frame.paragraphs[0]
        p.text = section
        p.font.size = Pt(20)
        p.font.color.rgb = colors['text_dark']
        
        y_pos += 0.8
    
    return slide


def add_content_slide(prs, title, items=None, colors=None):
    """
    添加内容页
    
    Args:
        prs: PPT对象
        title: 标题
        items: 内容项列表（支持###小标题、-项目符号、普通文字）
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    if items is None:
        items = []
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 内容
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.8))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        
        if item.startswith('###'):
            # 小标题
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = colors['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            # 项目符号
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text_dark']
            p.space_before = Pt(8)
        elif item.strip():
            # 普通文字
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text_dark']
            p.space_before = Pt(10)
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, colors=None):
    """
    添加双栏内容页
    
    Args:
        prs: PPT对象
        title: 标题
        left_title: 左栏标题
        left_items: 左栏内容列表
        right_title: 右栏标题
        right_items: 右栏内容列表
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 左栏
    ltb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.8))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    
    # 左栏标题
    p = ltf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.space_before = Pt(10)
    
    # 左栏内容
    for item in left_items:
        p = ltf.add_paragraph()
        if item.startswith('-'):
            p.text = "• " + item[1:].strip()
        else:
            p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text_dark']
        p.space_before = Pt(6)
    
    # 右栏
    rtb = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.8))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    # 右栏标题
    p = rtf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.space_before = Pt(10)
    
    # 右栏内容
    for item in right_items:
        p = rtf.add_paragraph()
        if item.startswith('-'):
            p.text = "• " + item[1:].strip()
        else:
            p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text_dark']
        p.space_before = Pt(6)
    
    return slide


def add_data_dashboard_slide(prs, title, data_cards=None, colors=None):
    """
    添加数据看板页
    
    Args:
        prs: PPT对象
        title: 标题
        data_cards: 数据卡片列表，每个卡片包含：
            - label: 标签
            - value: 数值
            - change: 变化（可选）
            - color: 颜色类型（success/warning）
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    if data_cards is None:
        data_cards = [
            {"label": "财务结账时间", "value": "3天", "change": "↓70%", "color": "success"},
            {"label": "库存周转率", "value": "6次/年", "change": "↑50%", "color": "success"},
            {"label": "采购周期", "value": "7天", "change": "↓53%", "color": "success"},
            {"label": "用户满意度", "value": "92%", "change": "↑22%", "color": "success"},
        ]
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 卡片布局 - 4列
    card_width = 2.8
    card_height = 1.8
    start_x = 0.7
    start_y = 1.5
    gap = 0.3
    
    for i, card in enumerate(data_cards[:8]):  # 最多8个卡片
        col = i % 4
        row = i // 4
        
        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)
        
        # 卡片背景
        card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(x), Inches(y), 
                                         Inches(card_width), Inches(card_height))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = colors['light_bg']
        card_bg.line.color.rgb = RGBColor(220, 220, 220)
        
        # 标签
        tb_label = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15), 
                                            Inches(card_width - 0.2), Inches(0.3))
        p = tb_label.text_frame.paragraphs[0]
        p.text = card.get('label', '')
        p.font.size = Pt(12)
        p.font.color.rgb = colors['text_light']
        
        # 数值
        tb_value = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.5), 
                                            Inches(card_width - 0.2), Inches(0.6))
        p = tb_value.text_frame.paragraphs[0]
        p.text = card.get('value', '')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        
        # 变化
        if 'change' in card:
            tb_change = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.2), 
                                                 Inches(card_width - 0.2), Inches(0.3))
            p = tb_change.text_frame.paragraphs[0]
            p.text = card.get('change', '')
            p.font.size = Pt(14)
            change_color = colors['success'] if card.get('color') == 'success' else colors['warning']
            p.font.color.rgb = change_color
    
    return slide


def add_comparison_slide(prs, title, before_title="上线前", after_title="上线后", 
                         before_items=None, after_items=None, colors=None):
    """
    添加对比页（上线前 vs 上线后）
    
    Args:
        prs: PPT对象
        title: 标题
        before_title: 左侧标题
        after_title: 右侧标题
        before_items: 左侧内容列表
        after_items: 右侧内容列表
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    if before_items is None:
        before_items = []
    if after_items is None:
        after_items = []
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 左侧背景
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0.3), Inches(1.3), 
                                     Inches(6.2), Inches(5.8))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(255, 240, 240)
    left_bg.line.fill.background()
    
    # 左侧标题
    tb_left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    p = tb_left_title.text_frame.paragraphs[0]
    p.text = f"❌ {before_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    # 左侧内容
    ltb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    
    for i, item in enumerate(before_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text_dark']
        p.space_before = Pt(8)
    
    # 右侧背景
    right_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                      Inches(6.8), Inches(1.3), 
                                      Inches(6.2), Inches(5.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)
    right_bg.line.fill.background()
    
    # 右侧标题
    tb_right_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.5))
    p = tb_right_title.text_frame.paragraphs[0]
    p.text = f"✅ {after_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # 右侧内容
    rtb = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.8), Inches(4.5))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    for i, item in enumerate(after_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = colors['text_dark']
        p.space_before = Pt(8)
    
    return slide


def add_timeline_slide(prs, title, phases=None, colors=None):
    """
    添加时间轴页
    
    Args:
        prs: PPT对象
        title: 标题
        phases: 阶段列表，每个阶段包含：
            - name: 阶段名称
            - time: 时间
            - work: 工作内容
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    if phases is None:
        phases = [
            {"name": "项目启动", "time": "第1月", "work": "需求调研\n蓝图设计"},
            {"name": "系统实施", "time": "第2-3月", "work": "系统配置\n数据迁移"},
            {"name": "测试培训", "time": "第4-5月", "work": "系统测试\n用户培训"},
            {"name": "上线支持", "time": "第6月", "work": "系统上线\n运维支持"},
        ]
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = colors['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # 时间轴
    start_x = 1.5
    y_pos = 2.5
    gap = 2.5
    
    # 连接线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                  Inches(start_x), Inches(y_pos + 0.25), 
                                  Inches(10), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['primary']
    line.line.fill.background()
    
    for i, phase in enumerate(phases):
        x = start_x + i * gap
        
        # 圆点
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                        Inches(x + 0.7), Inches(y_pos), 
                                        Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors['primary']
        circle.line.fill.background()
        
        # 阶段名称
        tb_name = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 0.8), Inches(2), Inches(0.4))
        p = tb_name.text_frame.paragraphs[0]
        p.text = phase.get('name', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 时间
        tb_time = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 1.2), Inches(2), Inches(0.3))
        p = tb_time.text_frame.paragraphs[0]
        p.text = phase.get('time', '')
        p.font.size = Pt(12)
        p.font.color.rgb = colors['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 工作
        tb_work = slide.shapes.add_textbox(Inches(x), Inches(y_pos + 1.6), Inches(2), Inches(0.8))
        p = tb_work.text_frame.paragraphs[0]
        p.text = phase.get('work', '')
        p.font.size = Pt(11)
        p.font.color.rgb = colors['text_dark']
        p.alignment = PP_ALIGN.CENTER
        tb_work.text_frame.word_wrap = True
    
    return slide


def add_ending_slide(prs, title="谢谢", subtitle="期待与您合作", company="金蝶软件（中国）有限公司", colors=None):
    """
    添加结束页
    
    Args:
        prs: PPT对象
        title: 主标题
        subtitle: 副标题
        company: 公司信息
        colors: 配色方案
    
    Returns:
        slide: 创建的幻灯片
    """
    if colors is None:
        colors = COLORS
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['primary']
    bg.line.fill.background()
    
    # 主标题
    tb_title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.2))
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        tb_subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(0.6))
        p = tb_subtitle.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    # 公司信息
    tb_company = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.5))
    p = tb_company.text_frame.paragraphs[0]
    p.text = company
    p.font.size = Pt(18)
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def save_ppt(prs, filename, output_dir=None):
    """
    保存PPT文件
    
    Args:
        prs: PPT对象
        filename: 文件名
        output_dir: 输出目录（默认使用OUTPUT_DIR）
    
    Returns:
        dict: 包含文件信息的字典
    """
    if output_dir is None:
        output_dir = OUTPUT_DIR
    
    os.makedirs(output_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if not filename.endswith('.pptx'):
        filename = f"{filename}_{timestamp}.pptx"
    
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    
    return {
        "success": True,
        "filename": os.path.basename(filepath),
        "filepath": filepath,
        "pages": len(prs.slides)
    }


# 行业案例数据（扩展版 v2.0）
INDUSTRY_CASES = {
    "制造业": {
        "pain_points": ["生产计划难以协调，排产效率低", "库存积压严重，资金占用大", "成本核算困难，精细化程度低", "质量管理滞后，追溯困难"],
        "solutions": ["生产计划优化：APS高级排程系统", "精益库存管理：JIT准时制生产", "精细化成本管理：标准成本法", "质量追溯体系：全流程质量追溯"],
        "benefits": ["生产效率提升30%", "库存周转率提升40%", "成本降低15%", "质量合格率提升至98%"],
        "typical_customer": {
            "name": "某大型制造企业",
            "scale": "员工5000人，年营收50亿",
            "modules": "财务云、供应链云、制造云",
            "duration": "12个月",
            "results": ["财务结账从10天缩短到3天", "库存周转率提升40%", "生产计划准确性提升50%"]
        }
    },
    "零售业": {
        "pain_points": ["库存管理困难，缺货与积压并存", "销售数据分散，难以统一分析", "会员管理混乱，客户粘性低", "供应链效率低，响应速度慢"],
        "solutions": ["智能库存管理：库存预警、自动补货", "全渠道销售管理：线上线下统一", "会员营销系统：会员画像、精准营销", "供应链协同：供应商协同、物流跟踪"],
        "benefits": ["库存准确率提升至99%", "销售额增长20%", "会员复购率提升35%", "供应链效率提升50%"],
        "typical_customer": {
            "name": "某连锁零售企业",
            "scale": "员工1000人，门店200家",
            "modules": "财务云、供应链云、人力云",
            "duration": "8个月",
            "results": ["订单处理效率提升70%", "库存准确率提升至99%", "销售额增长20%"]
        }
    },
    "服务业": {
        "pain_points": ["项目成本失控，利润难以保障", "资源调度困难，利用率低", "客户服务不及时，满意度下降", "财务核算复杂，报表周期长"],
        "solutions": ["项目成本管理：项目预算、成本核算", "资源调度优化：资源池管理、智能调度", "客户服务系统：工单管理、服务跟踪", "财务自动化：自动记账、智能核算"],
        "benefits": ["项目利润率提升20%", "资源利用率提升25%", "客户满意度提升至95%", "财务效率提升60%"],
        "typical_customer": {
            "name": "某专业服务企业",
            "scale": "员工500人，年营收5亿",
            "modules": "财务云、项目管理、人力云",
            "duration": "6个月",
            "results": ["项目利润率提升20%", "资源利用率提升25%", "客户满意度提升至95%"]
        }
    },
    "金融": {
        "pain_points": ["风险管理不足，合规压力大", "业务创新慢，市场响应迟缓", "客户体验差，数字化程度低", "数据孤岛严重，难以整合分析"],
        "solutions": ["风险管控系统：风险评估、预警监控", "敏捷开发平台：快速迭代、灵活配置", "客户体验优化：全渠道服务、智能客服", "数据中台建设：数据整合、数据治理"],
        "benefits": ["风险识别率提升40%", "合规成本降低30%", "产品上线周期缩短50%", "客户满意度提升25%"],
        "typical_customer": {
            "name": "某城市商业银行",
            "scale": "员工3000人，资产规模500亿",
            "modules": "财务云、风险管理、数据中台",
            "duration": "18个月",
            "results": ["风险识别率提升40%", "合规成本降低30%", "产品上线周期缩短50%"]
        }
    },
    "医疗": {
        "pain_points": ["医疗资源紧张，调度困难", "患者体验差，等待时间长", "成本控制难，精细化管理不足", "数据安全风险，合规要求高"],
        "solutions": ["资源调度系统：床位管理、手术室排程", "患者服务优化：预约挂号、智能导诊", "成本精细管理：科室成本、项目成本", "数据安全体系：数据加密、访问控制"],
        "benefits": ["资源利用率提升20%", "患者满意度提升30%", "成本降低10%", "数据安全合规100%"],
        "typical_customer": {
            "name": "某三甲医院",
            "scale": "员工2000人，床位1500张",
            "modules": "财务云、供应链云、人力云",
            "duration": "12个月",
            "results": ["资源利用率提升20%", "患者满意度提升30%", "数据安全合规100%"]
        }
    },
    "教育": {
        "pain_points": ["教学资源分散，共享困难", "学生管理复杂，信息不透明", "财务流程繁琐，效率低下", "信息化程度低，数字化转型慢"],
        "solutions": ["教学资源平台：资源共享、在线备课", "学生管理系统：学籍管理、成绩管理", "财务自动化：预算管理、费用报销", "数字化转型：智慧校园、数据中台"],
        "benefits": ["教学效率提升25%", "管理效率提升40%", "财务效率提升50%", "信息化水平大幅提升"],
        "typical_customer": {
            "name": "某知名高校",
            "scale": "员工1500人，学生20000人",
            "modules": "财务云、人力云、学生管理",
            "duration": "10个月",
            "results": ["教学效率提升25%", "管理效率提升40%", "财务效率提升50%"]
        }
    },
    "物流": {
        "pain_points": ["仓储管理粗放，效率低下", "配送调度困难，成本高", "货物追踪不便，客户体验差", "成本核算复杂，精细化不足"],
        "solutions": ["智能仓储管理：WMS系统、库位优化", "配送调度优化：路径优化、车辆调度", "全程货物追踪：GPS定位、温度监控", "精细化成本核算：订单成本、线路成本"],
        "benefits": ["仓储效率提升40%", "配送成本降低20%", "货物准时率提升至95%", "客户满意度提升25%"],
        "typical_customer": {
            "name": "某大型物流企业",
            "scale": "员工3000人，网点500个",
            "modules": "财务云、供应链云、仓储管理",
            "duration": "12个月",
            "results": ["仓储效率提升40%", "配送成本降低20%", "货物准时率提升至95%"]
        }
    },
    "房地产": {
        "pain_points": ["项目管理复杂，进度难控", "成本超支严重，利润率下降", "销售管理粗放，回款困难", "客户关系维护不足，复购率低"],
        "solutions": ["项目管理系统：进度管理、质量管理", "成本精细管控：目标成本、动态成本", "销售管理升级：客户管理、销售漏斗", "客户关系维护：客户画像、精准营销"],
        "benefits": ["项目进度可控率提升至90%", "成本超支率降低60%", "销售回款率提升至95%", "客户复购率提升20%"],
        "typical_customer": {
            "name": "某知名房地产企业",
            "scale": "员工2000人，年营收100亿",
            "modules": "财务云、供应链云、项目管理",
            "duration": "18个月",
            "results": ["项目进度可控率提升至90%", "成本超支率降低60%", "销售回款率提升至95%"]
        }
    },
    "能源": {
        "pain_points": ["设备管理粗放，故障率高", "能耗管理困难，成本压力大", "安全风险高，合规要求严", "环保压力大，碳排放管控难"],
        "solutions": ["设备全生命周期管理：设备台账、故障预警", "能耗精细化管理：能耗监测、节能优化", "安全生产管理：隐患排查、应急管理", "环保合规管理：排放监测、碳足迹追踪"],
        "benefits": ["设备故障率降低40%", "能耗成本降低15%", "安全事故率降低50%", "环保合规率100%"],
        "typical_customer": {
            "name": "某大型能源企业",
            "scale": "员工5000人，年营收200亿",
            "modules": "财务云、设备管理、安全生产",
            "duration": "24个月",
            "results": ["设备故障率降低40%", "能耗成本降低15%", "环保合规率100%"]
        }
    },
    "建筑": {
        "pain_points": ["项目管理分散，进度难控", "成本核算粗放，利润不透明", "供应链管理弱，材料浪费严重", "安全管理难，事故风险高"],
        "solutions": ["项目集中管理：进度协同、资源调度", "成本精细核算：目标成本、动态成本", "供应链优化：集中采购、库存管控", "安全管理升级：隐患排查、安全培训"],
        "benefits": ["项目进度可控率提升85%", "成本超支率降低50%", "材料浪费率降低30%", "安全事故率降低60%"],
        "typical_customer": {
            "name": "某大型建筑企业",
            "scale": "员工3000人，年营收80亿",
            "modules": "财务云、供应链云、项目管理",
            "duration": "15个月",
            "results": ["项目进度可控率提升85%", "成本超支率降低50%", "安全事故率降低60%"]
        }
    }
}


def get_industry_case(industry):
    """
    获取行业案例数据
    
    Args:
        industry: 行业名称
    
    Returns:
        dict: 行业案例数据
    """
    return INDUSTRY_CASES.get(industry, INDUSTRY_CASES["制造业"])


# 模块信息
__version__ = "1.0.0"
__author__ = "ClawBot"
__description__ = "金蝶PPT生成器公共模块"


if __name__ == "__main__":
    # 测试代码
    print(f"金蝶PPT公共模块 v{__version__}")
    print(f"输出目录: {OUTPUT_DIR}")
    print(f"配色方案: {list(COLORS.keys())}")
    print(f"行业案例: {list(INDUSTRY_CASES.keys())}")
