#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 - 专业完整版 v12.0
直接使用参考文档的素材和内容
"""

import os
import json
import random
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
ASSETS_DIR = os.path.join(OUTPUT_DIR, "ppt-assets")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 82, 147),
    'secondary': RGBColor(0, 122, 194),
    'accent': RGBColor(255, 153, 0),
    'dark': RGBColor(51, 51, 51),
    'light': RGBColor(240, 245, 250),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(128, 128, 128),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def get_random_image(category):
    """获取随机图片"""
    folder = os.path.join(ASSETS_DIR, category)
    if not os.path.exists(folder):
        return None
    images = [f for f in os.listdir(folder) if f.endswith(('.png', '.jpg', '.jpeg'))]
    if not images:
        return None
    return os.path.join(folder, random.choice(images))

def add_slide_with_bg(prs, bg_image=None):
    """添加带背景的幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    if bg_image and os.path.exists(bg_image):
        try:
            pic = slide.shapes.add_picture(bg_image, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        except:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLORS['primary']
            bg.line.fill.background()
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['primary']
        bg.line.fill.background()
    
    return slide

def add_title_slide(prs, title, subtitle="", company=""):
    """封面页 - 使用背景图"""
    bg_image = get_random_image('backgrounds')
    slide = add_slide_with_bg(prs, bg_image)
    
    # 半透明遮罩
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS['primary']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = COLORS['white']
    
    # 底部
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8))
    tf3 = footer.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "金蝶软件（中国）有限公司"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLORS['white']
    p4 = tf3.add_paragraph()
    p4.text = datetime.now().strftime('%Y年%m月')
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(200, 200, 200)
    
    return slide

def add_image_slide(prs, title, image_path, subtitle=""):
    """图片页 - 直接使用架构图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 添加图片
    if os.path.exists(image_path):
        try:
            # 计算图片尺寸
            pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(12))
        except Exception as e:
            print(f"添加图片失败: {e}")
    
    return slide

def add_content_slide(prs, title, items, subtitle=""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 内容
    start_y = Inches(1.0)
    content_box = slide.shapes.add_textbox(Inches(0.5), start_y, Inches(12), Inches(6))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        if item.startswith('  '):
            p.text = f"    {item.strip()}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS['gray']
        else:
            p.text = f"● {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(5)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    cols = len(headers)
    col_width = min(Inches(12 / cols), Inches(4))
    table_width = col_width * cols
    table_height = Inches(0.4 * (len(rows) + 1))
    
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.0), table_width, table_height).table
    
    for i in range(cols):
        table.columns[i].width = col_width
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light']
    
    return slide

def add_section_slide(prs, title, number=""):
    """章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SLIDE_HEIGHT)
    left.fill.solid()
    left.fill.fore_color.rgb = COLORS['primary']
    left.line.fill.background()
    
    if number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(5), Inches(2.8), Inches(7.5), Inches(1.5))
    tf2 = title_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['dark']
    
    return slide

def add_thank_slide(prs):
    """感谢页"""
    bg_image = get_random_image('backgrounds')
    slide = add_slide_with_bg(prs, bg_image)
    
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS['primary']
    
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_presales_ppt_v12(customer_info):
    """生成使用参考文档素材的专业PPT"""
    company_name = customer_info.get('companyName', '客户企业')
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance', 'supply'])
    
    module_names = {'finance': '财务管理', 'supply': '供应链管理', 'manufacturing': '制造管理', 'hr': '人力资源管理'}
    module_list = [module_names.get(m, m) for m in modules]
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 封面
    add_title_slide(prs, f"{company_name}ERP系统解决方案", "数字化转型项目述标方案")
    
    # 目录
    add_content_slide(prs, "目 录", [
        "项目需求理解",
        "解决方案设计",
        "技术平台架构",
        "项目实施安排",
        "项目团队配置",
        "成功案例分享",
        "服务承诺保障",
        "金蝶实力介绍"
    ])
    
    # 第1部分：项目需求理解
    add_section_slide(prs, "项目需求理解", "01")
    
    add_content_slide(prs, "一、企业概况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"实施模块：{'、'.join(module_list)}",
        "",
        "企业当前正处于数字化转型关键期，需要通过ERP系统建设实现：",
        "  - 业务流程标准化、规范化",
        "  - 财务业务一体化管理",
        "  - 数据驱动决策能力提升"
    ])
    
    add_content_slide(prs, "二、当前业务痛点分析", [
        "信息系统分散，数据孤岛严重，跨部门协同困难",
        "业务流程不顺畅，审批效率低，流程可追溯性差",
        "财务核算不及时，成本核算不精准，决策支持不足",
        "供应链管理粗放，库存周转率低，采购效率不高",
        "生产计划不准确，物料需求预测偏差大，交付周期长",
        "质量管理缺乏追溯体系，问题定位困难，改进措施难落地",
        "主数据管理分散，缺乏统一的主数据管理平台",
        "各系统间数据断点多，对运营决策支撑不足"
    ])
    
    add_content_slide(prs, "三、项目建设目标", [
        "统一业务及管理主数据，消除信息孤岛",
        "实现销售到收款的全闭环数字化管理",
        "建立项目预算、项目经营及项目核算的闭环管理",
        "实现采购到应付的全链条数字化管理",
        "构建从计划到交付的生产及供应链数字化管理体系",
        "实现业务变更及工程变更的全程跟踪和联动",
        "建立质量检验、追溯及持续改善体系",
        "提升人才招、选、用、育、留的全面体验"
    ])
    
    # 第2部分：解决方案设计
    add_section_slide(prs, "解决方案设计", "02")
    
    # 添加架构图（如果有）
    arch_images = [f for f in os.listdir(os.path.join(ASSETS_DIR, 'architecture')) if f.endswith('.png')][:3]
    for i, img in enumerate(arch_images):
        img_path = os.path.join(ASSETS_DIR, 'architecture', img)
        add_image_slide(prs, f"系统架构图 ({i+1})", img_path)
    
    add_content_slide(prs, "总体架构设计", [
        "业务架构 BA：销售到收款、采购到付款、计划到生产全流程管理",
        "数据架构 DA：数据治理、数据资产、数据服务",
        "应用架构 AA：财务管理、供应链、制造管理、人力资源管理",
        "技术架构 TA：云原生微服务、分布式部署、开放集成平台"
    ])
    
    add_content_slide(prs, "核心解决方案", [
        "财务管理：总账、应收、应付、固定资产、成本、预算",
        "供应链管理：销售、采购、库存、计划、物流",
        "生产制造：BOM管理、MRP运算、生产执行、质量管理",
        "数据分析：经营驾驶舱、多维报表、智能预警"
    ])
    
    # 详细方案
    if 'finance' in modules:
        add_content_slide(prs, "财务管理方案", [
            "总账管理：凭证处理、自动记账、期末结账、财务报表",
            "应收管理：销售开票、收款核销、账龄分析、信用管理",
            "应付管理：采购发票、付款处理、供应商对账、付款计划",
            "固定资产：资产台账、折旧计提、资产处置、资产盘点",
            "成本管理：成本核算、成本分配、成本分析、成本控制",
            "",
            "预期价值：",
            "  - 财务结账时间缩短60%",
            "  - 成本核算精度提升40%",
            "  - 财务报表准确率提升50%"
        ])
    
    if 'supply' in modules:
        add_content_slide(prs, "供应链管理方案", [
            "销售管理：订单处理、信用控制、出货管理、销售分析",
            "采购管理：询价比价、订单跟踪、入库检验、采购结算",
            "库存管理：出入库管理、盘点管理、库存分析、库存预警",
            "计划管理：需求计划、采购计划、生产计划、MRP运算",
            "",
            "预期价值：",
            "  - 库存周转率提升40%",
            "  - 采购成本降低15%",
            "  - 订单交付准时率提升35%"
        ])
    
    # 第3部分：技术平台
    add_section_slide(prs, "技术平台架构", "03")
    
    add_content_slide(prs, "金蝶云·苍穹平台优势", [
        "云原生架构：微服务、容器化、分布式部署",
        "低代码开发：可视化建模、快速交付、灵活扩展",
        "AI原生能力：大模型集成、智能助手、预测分析",
        "数据中台：数据治理、数据资产、数据服务",
        "开放集成：API网关、标准接口、生态连接",
        "安全可信：等保三级、数据加密、权限管控"
    ])
    
    # 添加截图
    screenshots = [f for f in os.listdir(os.path.join(ASSETS_DIR, 'screenshots')) if f.endswith(('.png', '.jpg'))][:2]
    for i, img in enumerate(screenshots):
        img_path = os.path.join(ASSETS_DIR, 'screenshots', img)
        add_image_slide(prs, f"系统界面展示 ({i+1})", img_path)
    
    # 第4部分：项目实施
    add_section_slide(prs, "项目实施安排", "04")
    
    add_content_slide(prs, "实施方法论", [
        "第一阶段：项目启动（2周）",
        "  - 成立项目组织、明确项目范围、召开启动会议",
        "第二阶段：需求调研（4周）",
        "  - 业务调研、需求分析、差距分析、编制调研报告",
        "第三阶段：方案设计（3周）",
        "  - 业务蓝图设计、系统配置方案、方案评审确认",
        "第四阶段：系统配置（4周）",
        "  - 系统参数配置、基础数据准备、业务流程配置",
        "第五阶段：测试培训（3周）",
        "  - UAT测试、最终用户培训、操作手册编制",
        "第六阶段：上线验收（2周）",
        "  - 系统切换上线、上线支持保障、项目验收"
    ])
    
    add_table_slide(prs, "实施计划安排",
        ["阶段", "主要工作", "周期", "交付成果"],
        [
            ["项目启动", "组建团队、制定计划", "2周", "项目章程、实施计划"],
            ["需求调研", "业务调研、需求分析", "4周", "调研报告、需求清单"],
            ["方案设计", "蓝图设计、方案评审", "3周", "业务蓝图、配置方案"],
            ["系统配置", "参数配置、数据准备", "4周", "系统配置、基础数据"],
            ["测试培训", "UAT测试、用户培训", "3周", "测试报告、培训记录"],
            ["上线验收", "系统切换、项目验收", "2周", "上线报告、验收报告"]
        ]
    )
    
    # 第5部分：项目团队
    add_section_slide(prs, "项目团队配置", "05")
    
    add_table_slide(prs, "项目团队配置",
        ["角色", "职责", "人数", "资质要求"],
        [
            ["项目总监", "项目整体把控、资源协调", "1人", "PMP认证、10年以上经验"],
            ["项目经理", "项目日常管理、进度控制", "1人", "PMP认证、5年以上经验"],
            ["业务顾问", "业务方案设计、流程优化", "2人", "行业认证、3年以上经验"],
            ["技术顾问", "技术方案设计、系统集成", "1人", "技术认证、3年以上经验"],
            ["实施顾问", "系统配置、用户培训", "3人", "产品认证、2年以上经验"],
            ["开发工程师", "客户化开发、报表开发", "2人", "开发认证、2年以上经验"]
        ]
    )
    
    # 第6部分：成功案例
    add_section_slide(prs, "成功案例分享", "06")
    
    add_content_slide(prs, "行业典型案例", [
        "案例一：某大型制造企业",
        "  - 企业规模：年营收50亿，员工5000人",
        "  - 实施模块：财务、供应链、制造、HR",
        "  - 实施周期：8个月",
        "  - 项目成效：财务结账时间缩短60%、库存周转率提升40%",
        "",
        "案例二：某大型零售企业",
        "  - 企业规模：年营收30亿，员工3000人",
        "  - 实施模块：财务、供应链、零售管理",
        "  - 实施周期：6个月",
        "  - 项目成效：采购成本降低15%、库存准确率达到99%"
    ])
    
    # 第7部分：服务承诺
    add_section_slide(prs, "服务承诺保障", "07")
    
    add_content_slide(prs, "服务承诺", [
        "系统稳定性承诺",
        "  - 系统可用率≥99.9%",
        "  - 数据安全有保障，支持容灾备份",
        "",
        "响应时效承诺",
        "  - 一级问题：2小时内响应，8小时内解决",
        "  - 二级问题：4小时内响应，24小时内解决",
        "",
        "培训承诺",
        "  - 提供不少于10场系统培训",
        "  - 培训覆盖率100%，考核通过率≥90%",
        "",
        "质保承诺",
        "  - 系统上线后提供1年免费质保服务"
    ])
    
    # 第8部分：金蝶实力
    add_section_slide(prs, "金蝶实力介绍", "08")
    
    add_content_slide(prs, "金蝶核心数据", [
        "服务客户：743万+ 企业及组织选择金蝶",
        "云收入占比：79.46% 云转型成功",
        "市场份额：连续17年 中国ERP市场第一",
        "员工规模：1.2万+ 专业服务团队",
        "覆盖国家：170+ 业务遍布全球",
        "研发投入：30%+ 持续创新能力",
        "成功率：95%+ 项目交付成功率"
    ])
    
    # 感谢页
    add_thank_slide(prs)
    
    # 保存
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    customer_code = customer_info.get('customerCode', 'CUSTOMER')
    filename = f"{customer_code}_售前解决方案_v12_{timestamp}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    
    prs.save(filepath)
    
    return {
        'success': True,
        'filepath': filepath,
        'filename': filename,
        'slides': len(prs.slides)
    }

if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--type', default='presales')
    parser.add_argument('--companyName', required=True)
    parser.add_argument('--customerCode', default='CUSTOMER')
    parser.add_argument('--industry', default='制造业')
    parser.add_argument('--companySize', default='中型企业')
    parser.add_argument('--employees', default='')
    parser.add_argument('--revenue', default='')
    parser.add_argument('--modules', default='finance,supply')
    parser.add_argument('--painPoints', default='')
    parser.add_argument('--businessGoals', default='')
    parser.add_argument('--goliveDate', default='')
    parser.add_argument('--golivePhase', default='')
    parser.add_argument('--switchPlan', default='')
    parser.add_argument('--acceptanceDate', default='')
    parser.add_argument('--acceptanceConclusion', default='')
    
    args = parser.parse_args()
    
    customer_info = {
        'companyName': args.companyName,
        'customerCode': args.customerCode,
        'industry': args.industry,
        'companySize': args.companySize,
        'employees': args.employees,
        'revenue': args.revenue,
        'modules': args.modules.split(','),
        'painPoints': args.painPoints,
        'businessGoals': args.businessGoals
    }
    
    result = generate_presales_ppt_v12(customer_info)
    print(json.dumps(result, ensure_ascii=False))
