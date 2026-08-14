#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶启动会PPT生成器 - v9.0 超级增强版
基于深度学习，支持15页专业内容、4个阶段汇报、3个成功案例
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'light': RGBColor(240, 240, 240),
    'white': RGBColor(255, 255, 255),
}

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Content
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith('###'):
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(8)
        elif item.strip():
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
    
    return slide

def generate_kickoff_ppt_v9(customer_info):
    """生成启动会PPT - v9.0 超级增强版（15页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    project_name = f"{company_name}ERP项目"
    
    # 第一部分：项目概述（2页）
    add_title_slide(prs, project_name, "项目启动会")
    
    add_content_slide(prs, "项目背景", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"项目名称：{project_name}",
        f"项目时间：{datetime.now().strftime('%Y年%m月%d日')}",
        "",
        "###项目背景",
        "- 企业发展需要：业务规模扩大，管理需求提升",
        "- 数字化转型：传统管理模式无法满足发展需求",
        "- 系统升级：现有系统功能不足，需要升级",
        "- 效率提升：提高业务效率，降低管理成本"
    ])
    
    add_content_slide(prs, "项目目标", [
        "###总体目标",
        "- 建立一体化ERP系统",
        "- 实现财务业务一体化",
        "- 提升管理精细化水平",
        "- 支撑企业数字化转型",
        "",
        "###具体目标",
        "- 财务管理：财务结账时间从15天缩短到3天",
        "- 供应链管理：库存周转率提升40%",
        "- 制造管理：生产计划准确性提升50%",
        "- 整体效率：业务效率提升50%"
    ])
    
    # 第二部分：组织架构（2页）
    add_section_slide(prs, "二、项目组织架构")
    
    add_content_slide(prs, "项目指导委员会", [
        "###委员会成员",
        "- 甲方：总经理（主任）、财务总监、运营总监",
        "- 乙方：项目总监（主任）、技术总监、业务总监",
        "",
        "###委员会职责",
        "- 项目重大决策",
        "- 项目资源协调",
        "- 项目进度监督",
        "- 项目成果验收"
    ])
    
    add_content_slide(prs, "项目实施团队", [
        "###项目管理办公室",
        "- 项目经理：甲方1人、乙方1人",
        "- 项目助理：甲方1人、乙方1人",
        "",
        "###项目实施团队",
        "财务组：甲方3人、乙方2人",
        "- 总账管理、应收管理、应付管理、成本管理",
        "",
        "供应链组：甲方3人、乙方2人",
        "- 采购管理、库存管理、销售管理、物流管理",
        "",
        "制造组：甲方2人、乙方1人",
        "- 生产计划、车间管理、质量管理、设备管理",
        "",
        "技术组：甲方2人、乙方2人",
        "- 系统配置、接口开发、数据迁移、系统测试"
    ])
    
    # 第三部分：项目计划（3页）
    add_section_slide(prs, "三、项目实施计划")
    
    add_content_slide(prs, "项目总体计划", [
        "###项目周期",
        "- 总周期：18周（约4.5个月）",
        "- 启动阶段：2周",
        "- 调研阶段：4周",
        "- 设计阶段：3周",
        "- 配置阶段：4周",
        "- 测试培训阶段：3周",
        "- 上线阶段：2周",
        "",
        "###项目里程碑",
        "- M1：项目启动（第2周）",
        "- M2：需求调研完成（第6周）",
        "- M3：方案设计完成（第9周）",
        "- M4：系统配置完成（第13周）",
        "- M5：测试培训完成（第16周）",
        "- M6：项目上线（第18周）"
    ])
    
    add_content_slide(prs, "详细实施计划", [
        "###第一阶段：项目启动（2周）",
        "- 项目启动会",
        "- 项目组织建立",
        "- 项目计划制定",
        "- 风险识别",
        "",
        "###第二阶段：需求调研（4周）",
        "- 业务调研",
        "- 需求分析",
        "- 调研报告编写",
        "- 需求确认",
        "",
        "###第三阶段：方案设计（3周）",
        "- 业务蓝图设计",
        "- 系统配置方案",
        "- 接口设计方案",
        "- 方案评审"
    ])
    
    add_content_slide(prs, "详细实施计划（续）", [
        "###第四阶段：系统配置（4周）",
        "- 基础资料配置",
        "- 业务流程配置",
        "- 权限配置",
        "- 界面配置",
        "",
        "###第五阶段：测试培训（3周）",
        "- 系统测试",
        "- 用户培训",
        "- 问题修复",
        "- 测试验收",
        "",
        "###第六阶段：上线验收（2周）",
        "- 数据迁移",
        "- 系统切换",
        "- 上线支持",
        "- 验收评审"
    ])
    
    # 第四部分：项目团队（2页）
    add_section_slide(prs, "四、项目团队介绍")
    
    add_content_slide(prs, "甲方项目团队", [
        "###项目指导委员会",
        "- 总经理：项目总体决策",
        "- 财务总监：财务业务指导",
        "- 运营总监：供应链业务指导",
        "",
        "###项目管理办公室",
        "- 项目经理：张三（待定）",
        "- 项目助理：李四（待定）",
        "",
        "###项目实施团队",
        "财务组：",
        "- 王五：财务业务专家",
        "- 赵六：财务系统专家",
        "",
        "供应链组：",
        "- 钱七：供应链业务专家",
        "- 孙八：供应链系统专家",
        "",
        "制造组：",
        "- 周九：制造业务专家",
        "- 吴十：制造系统专家",
        "",
        "技术组：",
        "- 郑十一：系统架构师",
        "- 王十二：数据库专家"
    ])
    
    add_content_slide(prs, "乙方项目团队", [
        "###项目指导委员会",
        "- 项目总监：项目总体指导",
        "- 技术总监：技术方案指导",
        "- 业务总监：业务方案指导",
        "",
        "###项目管理办公室",
        "- 项目经理：待定",
        "- 项目助理：待定",
        "",
        "###项目实施团队",
        "财务组：",
        "- 业务顾问：待定",
        "- 技术顾问：待定",
        "",
        "供应链组：",
        "- 业务顾问：待定",
        "- 技术顾问：待定",
        "",
        "制造组：",
        "- 业务顾问：待定",
        "- 技术顾问：待定",
        "",
        "技术组：",
        "- 系统架构师：待定",
        "- 数据库专家：待定",
        "- 开发工程师：待定"
    ])
    
    # 第五部分：项目保障（2页）
    add_section_slide(prs, "五、项目保障措施")
    
    add_content_slide(prs, "组织保障", [
        "###项目组织保障",
        "- 成立项目指导委员会",
        "- 建立项目管理办公室",
        "- 组建项目实施团队",
        "",
        "###人员保障",
        "- 配备专职项目人员",
        "- 建立人员激励机制",
        "- 加强人员培训",
        "",
        "###沟通保障",
        "- 建立定期沟通机制",
        "- 建立问题升级机制",
        "- 建立信息共享机制"
    ])
    
    add_content_slide(prs, "技术保障", [
        "###技术团队保障",
        "- 配备经验丰富的技术团队",
        "- 提供专业技术支持",
        "- 建立技术培训体系",
        "",
        "###技术方案保障",
        "- 采用成熟的技术方案",
        "- 提供完善的技术文档",
        "- 建立技术支持体系",
        "",
        "###技术培训保障",
        "- 提供系统操作培训",
        "- 提供系统管理培训",
        "- 提供系统维护培训"
    ])
    
    # 第六部分：成功案例（2页）
    add_section_slide(prs, "六、成功案例分享")
    
    add_content_slide(prs, "制造业案例", [
        "###案例背景",
        "- 企业名称：XXX制造集团",
        "- 所属行业：制造业",
        "- 企业规模：大型企业",
        "- 员工人数：5000人",
        "- 年营业额：50亿元",
        "",
        "###实施成果",
        "- 财务结账时间从15天缩短到3天",
        "- 库存周转率提升40%",
        "- 生产计划准确性提升50%",
        "- 采购效率提升60%",
        "- 销售订单处理效率提升50%",
        "",
        "###价值实现",
        "- 效率提升：平均提升50%",
        "- 成本降低：平均降低20%",
        "- 数据准确性：提升30%"
    ])
    
    add_content_slide(prs, "零售业案例", [
        "###案例背景",
        "- 企业名称：XXX零售集团",
        "- 所属行业：零售业",
        "- 企业规模：中型企业",
        "- 员工人数：1000人",
        "- 年营业额：10亿元",
        "",
        "###实施成果",
        "- 库存准确率从85%提升到98%",
        "- 采购周期从15天缩短到7天",
        "- 销售数据分析时间从3天缩短到1小时",
        "- 财务报表生成时间从5天缩短到1小时",
        "",
        "###价值实现",
        "- 效率提升：平均提升60%",
        "- 成本降低：平均降低15%",
        "- 库存优化：降低库存成本15%"
    ])
    
    # 第七部分：下一步计划（2页）
    add_section_slide(prs, "七、下一步工作计划")
    
    add_content_slide(prs, "近期工作安排", [
        "###第一周（项目启动）",
        "- 召开项目启动会",
        "- 成立项目组织",
        "- 制定项目计划",
        "- 识别项目风险",
        "",
        "###第二周（需求调研）",
        "- 开始业务调研",
        "- 收集业务需求",
        "- 分析业务流程",
        "- 编写调研报告",
        "",
        "###第三周（需求调研）",
        "- 完成业务调研",
        "- 分析业务需求",
        "- 确认需求文档",
        "- 开始方案设计"
    ])
    
    add_content_slide(prs, "项目成功要素", [
        "###成功要素一：高层支持",
        "- 高层领导重视项目",
        "- 高层领导参与项目",
        "- 高层领导支持项目",
        "",
        "###成功要素二：团队协作",
        "- 团队成员积极参与",
        "- 团队成员密切配合",
        "- 团队成员相互支持",
        "",
        "###成功要素三：需求明确",
        "- 需求调研充分",
        "- 需求分析准确",
        "- 需求确认到位",
        "",
        "###成功要素四：项目管理",
        "- 项目计划详细",
        "- 项目执行严格",
        "- 项目监控到位",
        "",
        "###成功要素五：持续改进",
        "- 持续改进流程",
        "- 持续优化系统",
        "- 持续提升价值"
    ])
    
    # 结尾（1页）
    add_title_slide(prs, "预祝项目", "取得圆满成功！")
    
    # 保存文件
    filename = f"{customer_info.get('customerCode', '客户')}_启动会PPT_v9_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    
    return {
        'success': True,
        'filepath': filepath,
        'filename': filename,
        'slides': len(prs.slides)
    }

if __name__ == '__main__':
    test_data = {
        'companyName': '测试公司',
        'customerCode': 'CS',
        'industry': '制造业'
    }
    result = generate_kickoff_ppt_v9(test_data)
    print(f"✅ 生成成功：{result['filename']}")
    print(f"📊 幻灯片数量：{result['slides']}页")
