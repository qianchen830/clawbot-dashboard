#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶交付文档生成器
支持：PPT（售前、启动会、上线汇报、验收汇报）和 Word（调研、蓝图、客户化开发、开发集成、UAT测试、上线切换）
"""

import sys
import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# 输出目录
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

def create_ppt(title, slides_data, customer_info):
    """创建PPT文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        if 'title' in slide_info:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_info['title']
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.alignment = PP_ALIGN.CENTER
        
        # 添加内容
        if 'content' in slide_info:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for item in slide_info['content']:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.space_after = Pt(10)
                if item.startswith('###'):
                    p.font.bold = True
                    p.font.size = Pt(22)
    
    # 保存文件
    filename = f"{customer_info.get('customerCode', '客户')}_{title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    return filepath

def create_word(title, sections, customer_info):
    """创建Word文档"""
    doc = Document()
    
    # 标题
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 文档信息
    doc.add_paragraph(f"客户名称：{customer_info.get('companyName', '')}")
    doc.add_paragraph(f"文档版本：{customer_info.get('docVersion', 'V1.0')}")
    doc.add_paragraph(f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}")
    doc.add_paragraph("")
    
    # 各章节
    for section in sections:
        doc.add_heading(section['title'], level=1)
        for content in section.get('content', []):
            if isinstance(content, str):
                doc.add_paragraph(content)
            elif isinstance(content, dict):
                if content.get('type') == 'list':
                    for item in content['items']:
                        doc.add_paragraph(item, style='List Bullet')
                elif content.get('type') == 'table':
                    table = doc.add_table(rows=len(content['rows']), cols=len(content['cols']))
                    for i, row_data in enumerate(content['rows']):
                        for j, cell_data in enumerate(row_data):
                            table.rows[i].cells[j].text = cell_data
    
    # 保存文件
    filename = f"{customer_info.get('customerCode', '客户')}_{title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    doc.save(filepath)
    return filepath

def get_int(data, key, default=0):
    """安全获取整数值"""
    try:
        val = data.get(key, default)
        return int(val) if val else default
    except:
        return default

def generate_presales_ppt(customer_info, sections):
    """生成售前PPT - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    company_size = customer_info.get('companySize', '中型企业')
    modules = customer_info.get('modules', ['finance'])
    impl_months = get_int(customer_info, 'implMonths', 6)
    
    # 行业特点
    industry_features = {
        '制造业': ['生产计划管理', '成本精细核算', '供应链协同', '质量管理'],
        '零售业': ['全渠道销售', '库存优化', '会员管理', '供应链整合'],
        '服务业': ['项目管理', '人力资本', '财务共享', '客户服务'],
        '金融': ['风险管理', '合规管理', '资金管理', '财务核算'],
        '医疗': ['预算管理', '成本控制', '物资管理', '人力资源'],
        '教育': ['预算管理', '科研管理', '人力资源', '资产管理']
    }
    
    features = industry_features.get(industry, industry_features['制造业'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    
    slides = [
        {'title': '封面', 'content': [
            f"**{customer_info.get('companyName', '企业名称')}**",
            f"数字化转型解决方案",
            "",
            f"文档版本：{customer_info.get('docVersion', 'V1.0')}",
            f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}",
            "",
            "**金蝶软件（中国）有限公司**"
        ]},
        {'title': '目录', 'content': [
            '一、企业概况与需求分析',
            '二、解决方案设计',
            '三、业务架构蓝图',
            '四、4A架构设计',
            '五、价值工程分析',
            '六、实施路线图',
            '七、成功案例分享',
            '八、项目团队与保障'
        ]},
        {'title': '一、企业概况与需求分析', 'content': [
            f"### 企业基本信息",
            f"企业名称：{customer_info.get('companyName', '')}",
            f"所属行业：{industry}",
            f"企业规模：{company_size}",
            f"员工人数：{customer_info.get('employees', '')}人",
            f"年营业额：{customer_info.get('revenue', '')}万元",
            "",
            f"### 行业特点",
            *[f"- {f}" for f in features],
            "",
            f"### 当前痛点",
            customer_info.get('painPoints', '- 信息化程度不足\n- 业务流程不顺畅\n- 数据孤岛严重\n- 决策支持能力弱')
        ]},
        {'title': '二、解决方案设计', 'content': [
            f"### 实施范围",
            f"核心模块：{', '.join(selected_modules)}",
            f"用户规模：{customer_info.get('userCount', '')}人",
            f"实施周期：{customer_info.get('implMonths', '')}个月",
            "",
            f"### 解决方案亮点",
            *[f"- {f}数字化" for f in features[:3]],
            "- 业财一体化管理",
            "- 数据驱动决策",
            "- 移动办公支持"
        ]},
        {'title': '三、业务架构蓝图', 'content': [
            "### 战略层",
            "- 数字化转型战略",
            "- 业务创新战略",
            "",
            "### 业务层",
            "- 核心业务流程优化",
            "- 业务协同平台建设",
            "",
            "### 支撑层",
            "- 数据中台建设",
            "- 技术平台建设"
        ]},
        {'title': '四、4A架构设计', 'content': [
            "### 业务架构(BA)",
            "- 业务流程梳理",
            "- 业务能力规划",
            "",
            "### 数据架构(DA)",
            "- 主数据管理",
            "- 数据治理体系",
            "",
            "### 应用架构(AA)",
            f"- 核心应用：{', '.join(selected_modules)}",
            "- 应用集成方案",
            "",
            "### 技术架构(TA)",
            "- 云原生架构",
            "- 微服务架构"
        ]},
        {'title': '五、价值工程分析', 'content': [
            "### 价值目标",
            "- 提升运营效率 30%+",
            "- 降低运营成本 20%+",
            "- 缩短业务周期 40%+",
            "",
            "### 价值实现路径",
            "- 业务流程标准化",
            "- 数据资产化",
            "- 决策智能化"
        ]},
        {'title': '六、实施路线图', 'content': [
            f"### 第一阶段：项目启动（第1-2周）",
            "- 项目组织建立",
            "- 实施计划确认",
            "",
            f"### 第二阶段：需求调研（第3-6周）",
            "- 业务调研",
            "- 需求分析",
            "",
            f"### 第三阶段：方案设计（第7-12周）",
            "- 解决方案设计",
            "- 系统配置",
            "",
            f"### 第四阶段：上线验收（第{impl_months*4-4}周-{impl_months*4}周）",
            "- 系统测试",
            "- 用户培训",
            "- 上线支持"
        ]},
        {'title': '七、成功案例分享', 'content': [
            f"### 同行业案例",
            f"- 某大型{industry}企业数字化转型项目",
            f"- 项目周期：{impl_months}个月",
            "- 实施范围：财务云+供应链云",
            "",
            "### 项目成果",
            "- 业务效率提升35%",
            "- 运营成本降低25%",
            "- 数据准确率99%+"
        ]},
        {'title': '八、项目团队与保障', 'content': [
            "### 项目组织",
            "- 项目指导委员会",
            "- 项目管理办公室",
            "- 业务实施团队",
            "- 技术开发团队",
            "",
            "### 项目保障",
            "- 专业实施团队",
            "- 标准实施方法论",
            "- 完善的培训体系",
            "- 7×24小时技术支持"
        ]},
        {'title': '感谢聆听', 'content': [
            "",
            f"**{customer_info.get('companyName', '企业名称')}**",
            "数字化转型之旅，金蝶与您同行",
            "",
            "**金蝶软件（中国）有限公司**",
            "咨询热线：4008-830-830",
            "官网：www.kingdee.com"
        ]}
    ]
    
    return create_ppt('售前解决方案', slides, customer_info)

def generate_startup_ppt(customer_info):
    """生成启动会PPT - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    impl_months = get_int(customer_info, 'implMonths', 6)
    
    slides = [
        {'title': '封面', 'content': [
            f"**{customer_info.get('companyName', '企业名称')}**",
            "ERP项目启动会",
            "",
            f"项目经理：{customer_info.get('projectManager', '待定')}",
            f"启动日期：{datetime.now().strftime('%Y年%m月%d日')}",
            "",
            "**金蝶软件（中国）有限公司**"
        ]},
        {'title': '会议议程', 'content': [
            "一、项目背景与目标",
            "二、项目组织架构",
            "三、实施范围与计划",
            "四、项目沟通机制",
            "五、风险管理策略",
            "六、下一步工作安排"
        ]},
        {'title': '一、项目背景与目标', 'content': [
            f"### 项目背景",
            f"- {customer_info.get('companyName', '')}信息化建设需求",
            f"- {industry}行业数字化转型趋势",
            "- 企业管理升级内在要求",
            "",
            f"### 项目目标",
            "- 构建一体化数字管理平台",
            f"- 实现业务流程标准化",
            "- 提升决策支持能力",
            f"- 实施范围：{', '.join(selected_modules)}",
            f"- 项目周期：{impl_months}个月"
        ]},
        {'title': '二、项目组织架构', 'content': [
            "### 项目指导委员会",
            "- 项目总监（客户方）",
            "- 项目总监（金蝶方）",
            "",
            "### 项目管理办公室(PMO)",
            "- 项目经理（客户方）",
            "- 项目经理（金蝶方）",
            "",
            "### 业务实施团队",
            "- 财务顾问",
            "- 供应链顾问",
            "- 制造顾问",
            "- 技术顾问"
        ]},
        {'title': '三、实施范围与计划', 'content': [
            f"### 实施范围",
            f"- 核心模块：{', '.join(selected_modules)}",
            f"- 用户规模：{customer_info.get('userCount', '')}人",
            f"- 项目预算：{customer_info.get('budget', '')}万元",
            "",
            "### 实施计划",
            f"- 项目启动：第1-2周",
            f"- 需求调研：第3-6周",
            f"- 方案设计：第7-12周",
            f"- 系统配置：第13-20周",
            f"- 测试培训：第21-24周",
            f"- 上线验收：第{impl_months*4-4}-{impl_months*4}周"
        ]},
        {'title': '四、项目沟通机制', 'content': [
            "### 定期会议",
            "- 周例会：每周五下午",
            "- 月度汇报：每月最后一周",
            "- 里程碑评审：关键节点",
            "",
            "### 沟通渠道",
            "- 项目微信群",
            "- 项目邮箱",
            "- 项目管理平台"
        ]},
        {'title': '五、风险管理策略', 'content': [
            "### 主要风险",
            "- 需求变更风险",
            "- 数据迁移风险",
            "- 人员流动风险",
            "- 进度延期风险",
            "",
            "### 应对措施",
            "- 建立变更控制流程",
            "- 制定详细数据迁移方案",
            "- 建立知识转移机制",
            "- 设置进度预警机制"
        ]},
        {'title': '六、下一步工作安排', 'content': [
            "### 本周工作",
            "- 确认项目组织架构",
            "- 召开项目启动会",
            "- 制定详细实施计划",
            "",
            "### 下周工作",
            "- 开始业务调研",
            "- 收集基础数据",
            "- 确认系统环境"
        ]},
        {'title': '谢谢！', 'content': [
            "",
            "让我们携手共创数字化未来！",
            "",
            "**金蝶软件（中国）有限公司**",
            "咨询热线：4008-830-830"
        ]}
    ]
    
    return create_ppt('项目启动会', slides, customer_info)

def generate_golive_ppt(customer_info):
    """生成上线汇报PPT - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    impl_months = get_int(customer_info, 'implMonths', 6)
    golive_date = customer_info.get('goliveDate', datetime.now().strftime('%Y年%m月%d日'))
    
    slides = [
        {'title': '封面', 'content': [
            f"**{customer_info.get('companyName', '企业名称')}**",
            "ERP项目上线汇报",
            "",
            f"上线日期：{golive_date}",
            f"汇报日期：{datetime.now().strftime('%Y年%m月%d日')}",
            "",
            "**金蝶软件（中国）有限公司**"
        ]},
        {'title': '目录', 'content': [
            '一、项目概况',
            '二、上线准备',
            '三、上线计划',
            '四、应急预案',
            '五、后续支持',
            '六、总结与展望'
        ]},
        {'title': '一、项目概况', 'content': [
            f"### 项目基本信息",
            f"企业名称：{customer_info.get('companyName', '')}",
            f"所属行业：{industry}",
            f"企业规模：{customer_info.get('companySize', '中型企业')}",
            f"实施范围：{', '.join(selected_modules)}",
            f"项目周期：{impl_months}个月",
            f"上线日期：{golive_date}",
            "",
            f"### 项目成果",
            "- 系统建设完成",
            "- 数据迁移完成",
            "- 用户培训完成",
            "- 业务流程优化"
        ]},
        {'title': '二、上线准备', 'content': [
            f"### 2.1 系统准备",
            "- 系统环境搭建完成",
            "- 系统配置完成",
            "- 用户权限配置完成",
            "- 系统测试完成",
            "",
            f"### 2.2 数据准备",
            "- 基础数据导入完成",
            "- 期初数据核对完成",
            "- 数据备份完成",
            "- 数据验证完成",
            "",
            f"### 2.3 人员准备",
            "- 用户培训完成（{customer_info.get('userCount', 100)}人）",
            "- 运维团队就位",
            "- 业务支持团队就位"
        ]},
        {'title': '三、上线计划', 'content': [
            f"### 3.1 上线时间安排",
            f"- 数据迁移：上线前1天",
            f"- 系统切换：上线日00:00-02:00",
            f"- 用户培训：上线后1周",
            f"- 运维支持：上线后持续",
            "",
            f"### 3.2 上线步骤",
            "1. 数据迁移准备",
            "2. 系统切换执行",
            "3. 用户权限验证",
            "4. 业务流程验证",
            "5. 问题处理与支持",
            "",
            f"### 3.3 关键节点",
            f"- T-7：数据核对",
            f"- T-3：系统测试",
            f"- T-1：最终备份",
            f"- T日：系统切换",
            f"- T+1：运维支持"
        ]},
        {'title': '四、应急预案', 'content': [
            f"### 4.1 风险识别",
            "- 系统故障风险",
            "- 数据错误风险",
            "- 用户操作风险",
            "- 业务中断风险",
            "",
            f"### 4.2 应急措施",
            "- 系统回退方案",
            "- 数据恢复方案",
            "- 技术支持方案",
            "- 业务连续性方案",
            "",
            f"### 4.3 应急联系人",
            "- 项目经理：{customer_info.get('projectManager', '待定')}",
            "- 技术支持：4008-830-830",
            "- 业务支持：业务顾问团队"
        ]},
        {'title': '五、后续支持', 'content': [
            f"### 5.1 运维保障",
            "- 7×24小时技术支持",
            "- 问题响应机制（30分钟响应）",
            "- 系统优化计划",
            "",
            f"### 5.2 培训支持",
            "- 上线后培训（1周）",
            "- 岗位操作培训",
            "- 系统功能培训",
            "",
            f"### 5.3 优化建议",
            "- 业务流程持续优化",
            "- 系统功能持续完善",
            "- 数据质量持续提升"
        ]},
        {'title': '六、总结与展望', 'content': [
            f"### 6.1 项目总结",
            f"{customer_info.get('companyName', '')}ERP项目成功上线，",
            "实现了业务流程标准化，",
            "提升了管理效率，",
            "为数字化转型奠定了基础。",
            "",
            f"### 6.2 未来展望",
            "- 深化应用：全面推广使用",
            "- 持续优化：系统功能完善",
            "- 价值挖掘：数据价值提升",
            "- 数字化转型：业务创新"
        ]},
        {'title': '谢谢！', 'content': [
            "",
            "ERP项目上线成功！",
            "数字化转型新征程！",
            "",
            "**金蝶软件（中国）有限公司**",
            "咨询热线：4008-830-830"
        ]}
    ]
    
    return create_ppt('上线汇报', slides, customer_info)

def generate_acceptance_ppt(customer_info):
    """生成验收汇报PPT - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    impl_months = get_int(customer_info, 'implMonths', 6)
    acceptance_date = customer_info.get('acceptanceDate', datetime.now().strftime('%Y年%m月%d日'))
    acceptance_result = customer_info.get('acceptanceResult', '验收通过')
    
    slides = [
        {'title': '封面', 'content': [
            f"**{customer_info.get('companyName', '企业名称')}**",
            "ERP项目验收汇报",
            "",
            f"验收日期：{acceptance_date}",
            f"验收结论：{acceptance_result}",
            "",
            "**金蝶软件（中国）有限公司**"
        ]},
        {'title': '目录', 'content': [
            '一、项目概况',
            '二、项目成果',
            '三、验收依据',
            '四、验收结论',
            '五、遗留问题',
            '六、后续计划'
        ]},
        {'title': '一、项目概况', 'content': [
            f"### 项目基本信息",
            f"企业名称：{customer_info.get('companyName', '')}",
            f"所属行业：{industry}",
            f"企业规模：{customer_info.get('companySize', '中型企业')}",
            f"实施范围：{', '.join(selected_modules)}",
            f"项目周期：{impl_months}个月",
            f"验收日期：{acceptance_date}",
            f"验收结论：{acceptance_result}",
            "",
            f"### 项目组织",
            "- 项目指导委员会",
            "- 项目管理办公室",
            "- 业务实施团队",
            "- 技术开发团队"
        ]},
        {'title': '二、项目成果', 'content': [
            f"### 2.1 系统建设成果",
            "- 系统部署完成",
            "- 系统配置完成",
            "- 用户培训完成（{customer_info.get('userCount', 100)}人）",
            "- 数据迁移完成",
            "- 业务流程优化完成",
            "",
            f"### 2.2 业务成果",
            "- 业务流程标准化",
            "- 数据准确性提升（95%+）",
            "- 工作效率提升（30%+）",
            "- 管理水平提升",
            "",
            f"### 2.3 用户反馈",
            "- 用户满意度：{customer_info.get('acceptanceResult', '验收通过')}",
            "- 系统易用性：良好",
            "- 功能完整性：完整"
        ]},
        {'title': '三、验收依据', 'content': [
            f"### 3.1 验收文档",
            "- 系统功能验收测试报告",
            "- 用户培训验收报告",
            "- 数据迁移验收报告",
            "- 业务流程验收报告",
            "",
            f"### 3.2 验收标准",
            "- 功能完整性：100%",
            "- 系统稳定性：99.9%",
            "- 用户满意度：90%+",
            "- 数据准确性：99%+"
        ]},
        {'title': '四、验收结论', 'content': [
            f"### 4.1 验收结果",
            f"ERP项目验收结论：{acceptance_result}",
            "",
            f"### 4.2 验收意见",
            "- 系统功能完整，满足业务需求",
            "- 系统运行稳定，性能良好",
            "- 用户培训到位，能够熟练使用",
            "- 数据迁移准确，数据质量良好",
            "",
            f"### 4.3 验收建议",
            "- 持续优化系统功能",
            "- 加强用户使用培训",
            "- 定期系统维护"
        ]},
        {'title': '五、遗留问题', 'content': [
            f"### 5.1 问题清单",
            customer_info.get('acceptanceIssues', '- 无遗留问题'),
            "",
            f"### 5.2 解决计划",
            customer_info.get('acceptancePlan', '- 持续技术支持\n- 定期回访\n- 系统优化建议')
        ]},
        {'title': '六、后续计划', 'content': [
            f"### 6.1 运维保障",
            "- 7×24小时技术支持",
            "- 问题响应机制（30分钟响应）",
            "- 系统优化计划",
            "",
            f"### 6.2 培训支持",
            "- 上线后培训（1周）",
            "- 岗位操作培训",
            "- 系统功能培训",
            "",
            f"### 6.3 优化建议",
            "- 业务流程持续优化",
            "- 系统功能持续完善",
            "- 数据质量持续提升"
        ]},
        {'title': '谢谢！', 'content': [
            "",
            "ERP项目验收成功！",
            "数字化转型新征程！",
            "",
            "**金蝶软件（中国）有限公司**",
            "咨询热线：4008-830-830"
        ]}
    ]
    
    return create_ppt('验收汇报', slides, customer_info)

def generate_survey_word(customer_info):
    """生成调研纪要Word - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务管理', 'supply': '供应链管理', 'manufacturing': '制造管理', 'hr': '人力资源管理'}
    selected_modules = [module_names.get(m, m) for m in modules]
    
    sections = [
        {'title': '一、调研基本信息', 'content': [
            f"调研对象：{customer_info.get('companyName', '')}",
            f"所属行业：{industry}",
            f"调研日期：{datetime.now().strftime('%Y年%m月%d日')}",
            f"调研模块：{', '.join(selected_modules)}",
            f"调研人员：项目经理、业务顾问、技术顾问",
            f"受访人员：财务部、采购部、销售部、仓库、IT部",
            f"调研地点：{customer_info.get('companyName', '')}会议室"
        ]},
        {'title': '二、调研目的', 'content': [
            f"本次调研旨在全面了解{customer_info.get('companyName', '')}的业务现状和需求，为后续方案设计提供依据。具体目标包括：",
            "",
            "1. 了解企业现有业务流程和管理模式",
            "2. 识别现有系统存在的问题和痛点",
            "3. 明确企业数字化转型需求",
            "4. 确定系统实施范围和边界",
            "5. 收集基础数据和关键业务场景"
        ]},
        {'title': '三、调研内容', 'content': [
            f"### 3.1 财务管理现状",
            "- 现有财务系统使用情况",
            "- 会计科目体系",
            "- 成本核算方法",
            "- 预算管理流程",
            "- 财务报表体系",
            "",
            f"### 3.2 供应链管理现状",
            "- 采购业务流程",
            "- 销售业务流程",
            "- 库存管理流程",
            "- 供应商管理",
            "- 客户管理",
            "",
            f"### 3.3 数据管理现状",
            "- 主数据管理现状",
            "- 基础数据质量",
            "- 数据接口情况"
        ]},
        {'title': '四、存在问题', 'content': [
            f"### 4.1 业务流程问题",
            "- 部门间协作效率低",
            "- 流程审批周期长",
            "- 信息传递不及时",
            "",
            f"### 4.2 数据管理问题",
            "- 数据孤岛现象严重",
            "- 数据一致性差",
            "- 数据准确性待提升",
            "",
            f"### 4.3 系统应用问题",
            "- 系统功能不完善",
            "- 用户体验有待提升",
            "- 系统集成困难"
        ]},
        {'title': '五、需求分析', 'content': [
            f"### 5.1 业务需求",
            f"- 构建{industry}行业一体化管理平台",
            "- 实现业务流程标准化",
            "- 提升数据准确性及时效性",
            "",
            f"### 5.2 功能需求",
            *[f"- {m}功能模块建设" for m in selected_modules],
            "- 移动办公支持",
            "- 数据分析报表",
            "",
            f"### 5.3 技术需求",
            "- 云端部署",
            "- 数据安全保障",
            "- 系统高可用性"
        ]},
        {'title': '六、下一步计划', 'content': [
            "1. 完成调研纪要确认",
            "2. 编制需求分析报告",
            "3. 设计解决方案",
            "4. 方案评审确认",
            "5. 制定详细实施计划"
        ]},
        {'title': '七、附件', 'content': [
            "- 调研问卷",
            "- 业务流程图",
            "- 数据收集模板",
            "- 会议纪要"
        ]}
    ]
    
    return create_word('调研纪要', sections, customer_info)

def generate_blueprint_word(customer_info):
    """生成业务蓝图Word - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    
    sections = [
        {'title': '一、蓝图概述', 'content': [
            f"客户名称：{customer_info.get('companyName', '')}",
            f"所属行业：{industry}",
            f"实施范围：{', '.join(selected_modules)}",
            f"文档版本：{customer_info.get('docVersion', 'V1.0')}",
            f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}",
            "",
            f"本蓝图设计文档基于{customer_info.get('companyName', '')}业务调研结果，结合金蝶云星空产品功能，设计符合企业实际需求的业务解决方案。"
        ]},
        {'title': '二、业务架构设计', 'content': [
            f"### 2.1 组织架构",
            "- 公司层级：法人公司、核算组织、利润中心",
            "- 部门层级：成本中心、利润中心",
            "- 岗位层级：业务岗位、管理岗位",
            "",
            f"### 2.2 业务流程架构",
            "- 核心业务流程：采购到付款、销售到收款、生产到成本",
            "- 管理业务流程：预算管理、资金管理、成本管理",
            "- 支撑业务流程：人力资源管理、资产管理"
        ]},
        {'title': '三、系统架构设计', 'content': [
            f"### 3.1 应用架构",
            *[f"- {m}：核心业务模块" for m in selected_modules],
            "- 基础平台：用户管理、权限管理、流程管理",
            "- 数据分析：报表中心、BI分析",
            "",
            f"### 3.2 功能模块设计",
            "",
            f"#### 3.2.1 财务云",
            "- 总账管理：凭证、账簿、报表",
            "- 应收管理：收款、核销、账龄",
            "- 应付管理：付款、核销、账龄",
            "- 固定资产：卡片、折旧、处置",
            "- 成本管理：成本核算、成本分析",
            "",
            f"#### 3.2.2 供应链云",
            "- 采购管理：请购、订单、入库、结算",
            "- 销售管理：报价、订单、出库、结算",
            "- 库存管理：入库、出库、盘点、调拨",
            "- 质量管理：检验、判定、追溯"
        ]},
        {'title': '四、数据架构设计', 'content': [
            f"### 4.1 主数据管理",
            "- 客户主数据：客户档案、信用额度、结算规则",
            "- 供应商主数据：供应商档案、付款条件、结算规则",
            "- 物料主数据：物料档案、计价方法、库存策略",
            "- 会计科目：科目体系、辅助核算、外币核算",
            "",
            f"### 4.2 基础数据",
            "- 部门档案",
            "- 人员档案",
            "- 仓库档案",
            "- 结算方式",
            "- 收付款条件"
        ]},
        {'title': '五、业务流程设计', 'content': [
            f"### 5.1 采购业务流程",
            "采购申请 → 采购订单 → 收货入库 → 质量检验 → 采购结算 → 付款管理",
            "",
            f"### 5.2 销售业务流程",
            "销售报价 → 销售订单 → 发货出库 → 质量检验 → 销售结算 → 收款管理",
            "",
            f"### 5.3 生产业务流程",
            "销售预测 → 生产计划 → 生产订单 → 生产领料 → 生产入库 → 成本核算"
        ]},
        {'title': '六、接口架构设计', 'content': [
            f"### 6.1 系统集成方案",
            "- 银企直联：银行账户、资金流水、电子对账",
            "- 税务接口：发票开具、发票认证、纳税申报",
            "- OA集成：流程审批、消息通知、单据传递",
            "",
            f"### 6.2 数据迁移方案",
            "- 历史数据清洗",
            "- 基础数据导入",
            "- 期初数据迁移"
        ]},
        {'title': '七、实施建议', 'content': [
            "1. 分阶段实施，降低风险",
            "2. 强化培训，确保使用",
            "3. 完善制度，保障运行",
            "4. 持续优化，提升价值"
        ]}
    ]
    
    return create_word('业务蓝图设计', sections, customer_info)

def generate_devdesign_word(customer_info):
    """生成客户化开发Word"""
    sections = [
        {'title': '开发需求概述', 'content': [
            f"客户名称：{customer_info.get('companyName', '')}",
            f"文档版本：{customer_info.get('docVersion', 'V1.0')}",
            f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}"
        ]},
        {'title': '需求分析', 'content': [
            '功能需求：',
            '- 需求1：待补充',
            '- 需求2：待补充',
            '',
            '接口需求：',
            '- 接口1：待补充',
            '- 接口2：待补充'
        ]},
        {'title': '技术方案', 'content': [
            '开发环境：',
            '- 开发平台：金蝶云星空BOS平台',
            '- 开发语言：Java',
            '- 数据库：MySQL',
            '',
            '技术架构：',
            '- 前端：Vue.js',
            '- 后端：Spring Boot',
            '- 中间件：Redis、RabbitMQ'
        ]},
        {'title': '开发计划', 'content': [
            '- 需求确认：第1周',
            '- 设计评审：第2周',
            '- 开发实现：第3-6周',
            '- 测试验收：第7-8周'
        ]},
        {'title': '测试方案', 'content': [
            '单元测试：',
            '- 功能测试',
            '- 性能测试',
            '',
            '集成测试：',
            '- 接口测试',
            '- 回归测试'
        ]}
    ]
    
    return create_word('客户化开发设计说明书', sections, customer_info)

def generate_integration_word(customer_info):
    """生成开发集成Word"""
    sections = [
        {'title': '集成概述', 'content': [
            f"客户名称：{customer_info.get('companyName', '')}",
            f"文档版本：{customer_info.get('docVersion', 'V1.0')}",
            f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}"
        ]},
        {'title': '集成范围', 'content': [
            '内部系统集成：',
            '- OA系统',
            '- CRM系统',
            '- WMS系统',
            '',
            '外部系统对接：',
            '- 银企直联',
            '- 税务平台',
            '- 电子发票'
        ]},
        {'title': '技术方案', 'content': [
            '集成方式：',
            '- API接口',
            '- 数据库同步',
            '- 文件传输',
            '',
            '技术标准：',
            '- RESTful API',
            '- JSON数据格式',
            '- OAuth2.0认证'
        ]},
        {'title': '实施计划', 'content': [
            '- 接口设计：第1-2周',
            '- 开发实现：第3-6周',
            '- 联调测试：第7-8周',
            '- 上线部署：第9周'
        ]}
    ]
    
    return create_word('开发集成方案', sections, customer_info)

def generate_uat_word(customer_info):
    """生成UAT测试报告Word - 基于金蝶标准模板"""
    industry = customer_info.get('industry', '制造业')
    modules = customer_info.get('modules', ['finance'])
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 'hr': '人力云'}
    selected_modules = [module_names.get(m, m) for m in modules]
    impl_months = get_int(customer_info, 'implMonths', 6)
    
    sections = [
        {'title': '一、测试概述', 'content': [
            f"客户名称：{customer_info.get('companyName', '')}",
            f"测试模块：{', '.join(selected_modules)}",
            f"测试日期：{datetime.now().strftime('%Y年%m月%d日')}",
            f"测试周期：{impl_months // 2}周",
            f"测试人员：业务骨干、关键用户、IT人员",
            "",
            "UAT测试（用户验收测试）是系统上线前的关键环节，旨在验证系统是否满足业务需求，确保系统功能完整、运行稳定。"
        ]},
        {'title': '二、测试环境', 'content': [
            f"### 2.1 服务器环境",
            "- 操作系统：Linux CentOS 7.x",
            "- 数据库：MySQL 8.0",
            "- 应用服务器：Tomcat 9.x",
            "- 中间件：Redis 6.x",
            "",
            f"### 2.2 客户端环境",
            "- 浏览器：Chrome 90+ / Edge 90+",
            "- 分辨率：1920×1080以上",
            "- 网络：企业内网",
            "",
            f"### 2.3 测试数据",
            "- 基础数据：已导入完成",
            "- 业务数据：已准备测试场景",
            "- 期初数据：已核对无误"
        ]},
        {'title': '三、测试范围', 'content': [
            f"### 3.1 功能测试范围",
            *[f"- {m}模块功能测试" for m in selected_modules],
            "- 系统集成测试",
            "- 权限管理测试",
            "",
            f"### 3.2 性能测试范围",
            "- 并发用户测试（100用户并发）",
            "- 大数据量测试（10万条记录）",
            "- 响应时间测试（<3秒）"
        ]},
        {'title': '四、测试用例', 'content': [
            f"### 4.1 财务云测试用例",
            "| 序号 | 测试项 | 测试步骤 | 预期结果 | 实际结果 | 状态 |",
            "|------|--------|----------|----------|----------|------|",
            "| 1 | 凭证录入 | 新增凭证→录入分录→保存 | 凭证保存成功 | 凭证保存成功 | 通过 |",
            "| 2 | 凭证审核 | 选择凭证→点击审核 | 审核成功 | 审核成功 | 通过 |",
            "| 3 | 凭证记账 | 选择凭证→点击记账 | 记账成功 | 记账成功 | 通过 |",
            "",
            f"### 4.2 供应链云测试用例",
            "| 序号 | 测试项 | 测试步骤 | 预期结果 | 实际结果 | 状态 |",
            "|------|--------|----------|----------|----------|------|",
            "| 1 | 采购订单 | 新增订单→录入明细→保存 | 订单保存成功 | 订单保存成功 | 通过 |",
            "| 2 | 采购入库 | 选择订单→生成入库单 | 入库成功 | 入库成功 | 通过 |",
            "| 3 | 采购结算 | 选择入库单→生成结算单 | 结算成功 | 结算成功 | 通过 |"
        ]},
        {'title': '五、测试结果', 'content': [
            f"### 5.1 测试统计",
            f"- 测试用例总数：{customer_info.get('userCount', 100) * 2}个",
            f"- 通过用例数：{customer_info.get('userCount', 100) * 2 - 5}个",
            f"- 失败用例数：5个",
            f"- 通过率：{((customer_info.get('userCount', 100) * 2 - 5) / (customer_info.get('userCount', 100) * 2) * 100):.1f}%",
            "",
            f"### 5.2 问题分类",
            "- 功能问题：3个（已修复）",
            "- 性能问题：1个（已优化）",
            "- 界面问题：1个（已调整）",
            "",
            f"### 5.3 测试结论",
            "系统功能完整，运行稳定，满足业务需求，建议上线。"
        ]},
        {'title': '六、遗留问题', 'content': [
            "| 序号 | 问题描述 | 优先级 | 负责人 | 计划完成时间 | 状态 |",
            "|------|----------|--------|--------|--------------|------|",
            "| 1 | 报表导出格式优化 | 低 | 开发团队 | 上线后1周 | 待处理 |",
            "| 2 | 移动端部分功能适配 | 低 | 开发团队 | 上线后2周 | 待处理 |"
        ]},
        {'title': '七、测试结论', 'content': [
            f"经UAT测试验证，{customer_info.get('companyName', '')}ERP系统：",
            "",
            "1. 功能完整：核心业务流程全部实现",
            "2. 性能稳定：系统响应时间满足要求",
            "3. 数据准确：测试数据核对无误",
            "4. 用户满意：关键用户验收通过",
            "",
            "**建议：系统具备上线条件，可以进入上线阶段。**"
        ]},
        {'title': '八、签字确认', 'content': [
            "",
            f"客户方项目负责人：______________    日期：______________",
            "",
            f"金蝶方项目经理：______________    日期：______________",
            "",
            f"测试负责人：______________    日期：______________"
        ]}
    ]
    
    return create_word('UAT测试报告', sections, customer_info)

def generate_golive_word(customer_info):
    """生成上线切换Word"""
    sections = [
        {'title': '上线概述', 'content': [
            f"客户名称：{customer_info.get('companyName', '')}",
            f"上线日期：{datetime.now().strftime('%Y年%m月%d日')}",
            f"实施模块：{', '.join(customer_info.get('modules', []))}"
        ]},
        {'title': '上线准备', 'content': [
            '系统准备：',
            '- 系统环境已搭建',
            '- 系统配置已完成',
            '- 用户权限已配置',
            '',
            '数据准备：',
            '- 基础数据已导入',
            '- 期初数据已核对',
            '- 数据备份已完成',
            '',
            '人员准备：',
            '- 用户培训已完成',
            '- 运维团队已就位'
        ]},
        {'title': '上线计划', 'content': [
            '上线步骤：',
            '1. 数据迁移',
            '2. 系统切换',
            '3. 用户培训',
            '4. 运维支持',
            '',
            '时间安排：',
            '- T-7：数据核对',
            '- T-3：系统测试',
            '- T-1：最终备份',
            '- T日：系统切换',
            '- T+1：运维支持'
        ]},
        {'title': '应急预案', 'content': [
            '风险识别：',
            '- 系统故障',
            '- 数据错误',
            '- 用户操作失误',
            '',
            '应急措施：',
            '- 系统回退',
            '- 数据恢复',
            '- 技术支持'
        ]},
        {'title': '运维保障', 'content': [
            '支持团队：',
            '- 项目经理',
            '- 技术顾问',
            '- 业务顾问',
            '',
            '联系方式：',
            '- 7×24小时热线',
            '- 微信群支持'
        ]}
    ]
    
    return create_word('上线切换方案', sections, customer_info)

def main():
    if len(sys.argv) < 3:
        print(json.dumps({'error': '参数不足', 'usage': 'python kingdee-generator.py <type> <json_data>'}))
        sys.exit(1)
    
    doc_type = sys.argv[1]
    try:
        customer_info = json.loads(sys.argv[2])
    except json.JSONDecodeError:
        print(json.dumps({'error': 'JSON解析失败'}))
        sys.exit(1)
    
    result = {'error': '未知类型'}
    
    try:
        if doc_type == 'presales_ppt':
            filepath = generate_presales_ppt(customer_info, customer_info.get('pptSections', []))
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'startup_ppt':
            filepath = generate_startup_ppt(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'golive_ppt':
            filepath = generate_golive_ppt(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'acceptance_ppt':
            filepath = generate_acceptance_ppt(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'survey_word':
            filepath = generate_survey_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'blueprint_word':
            filepath = generate_blueprint_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'devdesign_word':
            filepath = generate_devdesign_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'integration_word':
            filepath = generate_integration_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'uat_word':
            filepath = generate_uat_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
        elif doc_type == 'golive_word':
            filepath = generate_golive_word(customer_info)
            result = {'success': True, 'filepath': filepath, 'filename': os.path.basename(filepath)}
    except Exception as e:
        result = {'error': str(e)}
    
    print(json.dumps(result, ensure_ascii=False))

if __name__ == '__main__':
    main()
