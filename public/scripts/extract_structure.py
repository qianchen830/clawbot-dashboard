#!/usr/bin/env python3
"""
分析参考PPT的详细内容结构，提取章节标题和页面内容要点
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import re

def extract_content_structure(ppt_path):
    """提取PPT的详细内容结构"""
    print(f"正在提取内容结构: {ppt_path}")
    
    prs = Presentation(ppt_path)
    
    structure = []
    current_section = None
    
    for idx, slide in enumerate(prs.slides, 1):
        # 提取所有文本
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        
        # 提取标题（通常是第一个文本框）
        title = texts[0] if texts else ""
        
        # 检测章节标题
        section_keywords = ['目 录', '目录', '一、', '二、', '三、', '四、', '五、', '六、', '七、', '八、', '九、', '十、', 
                          '第一章', '第二章', '第三章', '第四章', '第五章']
        is_section = any(kw in title for kw in section_keywords)
        
        # 检测主要章节
        main_sections = {
            '目 录': '目录页',
            '目录': '目录页',
        }
        
        # 分析页面内容
        page_info = {
            'index': idx,
            'title': title[:100] if title else "",
            'layout': slide.slide_layout.name if slide.slide_layout else "Unknown",
            'texts': texts[:10],  # 只保存前10个文本
            'is_section': is_section,
            'shape_count': len(slide.shapes),
            'image_count': sum(1 for s in slide.shapes if s.shape_type == 13),
            'table_count': sum(1 for s in slide.shapes if s.has_table),
            'chart_count': sum(1 for s in slide.shapes if s.has_chart)
        }
        
        structure.append(page_info)
    
    # 输出结构化内容
    print("\n" + "=" * 80)
    print("章节结构分析:")
    print("=" * 80)
    
    section_count = 0
    for info in structure:
        if info['layout'] in ['封面', '封底', '1_目录、提纲']:
            section_count += 1
            print(f"\n[{info['index']:3d}] {info['layout']}: {info['title'][:60]}")
    
    return structure

if __name__ == "__main__":
    ppt_path = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/售前文件/中煤科工ERP重庆研究院ERP升级项目述标文件V7.0.pptx"
    structure = extract_content_structure(ppt_path)
    
    # 输出统计
    print("\n" + "=" * 80)
    print("结构统计:")
    print("=" * 80)
    print(f"总页数: {len(structure)}")
    print(f"目录页数: {sum(1 for s in structure if '目录' in s['layout'])}")
    print(f"内容页数: {sum(1 for s in structure if '白色内页' in s['layout'])}")
    print(f"封面/封底: {sum(1 for s in structure if s['layout'] in ['封面', '封底'])}")
