#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器基础类 V1.0
- 统一的样式和颜色
- 图表、卡片、时间线、表格等视觉元素
- 所有PPT生成器共用
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
import os

# 金蝶品牌色
COLORS = {
    'primary': RGBColor(0x1F, 0x4E, 0x79),       # 主色深蓝
    'kingdee_blue': RGBColor(0x00, 0x6C, 0xB8),  # 金蝶蓝
    'accent': RGBColor(0xED, 0x7D, 0x31),        # 强调橙
    'success': RGBColor(0x70, 0xAD, 0x47),       # 成功绿
    'warning': RGBColor(0xFF, 0xC0, 0x00),       # 警告黄
    'danger': RGBColor(0xC0, 0x00, 0x00),        # 危险红
    'text_primary': RGBColor(0x33, 0x33, 0x33),
    'text_secondary': RGBColor(0x66, 0x66, 0x66),
    'text_light': RGBColor(0x99, 0x99, 0x99),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'bg_light': RGBColor(0xF2, 0xF6, 0xFA),
    'bg_card': RGBColor(0xFA, 0xFB, 0xFC),
    'border': RGBColor(0xD0, 0xD8, 0xE0),
}

TEMPLATE_PATH = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/ppt模板/ppt母版.pptx"

class KingdeePPTBase:
    """金蝶PPT生成器基础类"""
    
    def __init__(self, company, doc_type="项目"):
        self.company = company
        self.doc_type = doc_type
        self.prs = None
        self.page = 0
        self.layout_map = {}
        self.C = COLORS
        
        if os.path.exists(TEMPLATE_PATH):
            self.prs = Presentation(TEMPLATE_PATH)
            self._clear_slides()
            for i, layout in enumerate(self.prs.slide_layouts):
                self.layout_map[layout.name] = i
            print(f"加载母版成功，布局: {list(self.layout_map.keys())}")
        else:
            self.prs = Presentation()
    
    def _clear_slides(self):
        """清除母版示例幻灯片"""
        pres_part = self.prs.part
        sldIdLst = self.prs._element.find(qn('p:sldIdLst'))
        if sldIdLst is not None:
            for sld_id in list(sldIdLst):
                rId = sld_id.get(qn('r:id'))
                if rId:
                    try:
                        pres_part.drop_rel(rId)
                    except:
                        pass
                sldIdLst.remove(sld_id)
            self.prs.slides._sldIdLst = sldIdLst
    
    def _get_layout(self, name):
        """获取布局"""
        if name in self.layout_map:
            return self.prs.slide_layouts[self.layout_map[name]]
        for key in ['白色内页', '内页', '3_白色内页']:
            if key in self.layout_map:
                return self.prs.slide_layouts[self.layout_map[key]]
        return self.prs.slide_layouts[0]
    
    def _add_slide(self, layout_name="白色内页"):
        """添加幻灯片"""
        slide = self.prs.slides.add_slide(self._get_layout(layout_name))
        self.page += 1
        return slide
    
    # ========== 基础元素 ==========
    
    def _add_title_bar(self, slide, title, subtitle=None):
        """标题栏"""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.85))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.C['primary']
        bg.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = '微软雅黑'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        
        if subtitle:
            sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12), Inches(0.25))
            p = sb.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.name = '微软雅黑'
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
    
    def _add_footer(self, slide):
        """页脚"""
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(3), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self.C['kingdee_blue']
        line.line.fill.background()
        
        fb = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(4), Inches(0.3))
        p = fb.text_frame.paragraphs[0]
        p.text = f"金蝶软件（中国）有限公司  |  第 {self.page} 页"
        p.font.name = '微软雅黑'
        p.font.size = Pt(9)
        p.font.color.rgb = self.C['text_light']
        p.alignment = PP_ALIGN.RIGHT
    
    # ========== 页面类型 ==========
    
    def add_cover(self, title, subtitle=None):
        """封面页"""
        slide = self._add_slide("封面")
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.C['primary']
        bg.line.fill.background()
        
        # 装饰线
        deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(1.5), Inches(0.08))
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.C['accent']
        deco.line.fill.background()
        
        # 标题
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(title.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = '微软雅黑'
            p.font.size = Pt(42) if i == 0 else Pt(28)
            p.font.bold = True if i == 0 else False
            p.font.color.rgb = self.C['white']
        
        if subtitle:
            sb = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(6), Inches(0.5))
            p = sb.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.name = '微软雅黑'
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
        
        # 公司
        cb = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(6), Inches(0.4))
        p = cb.text_frame.paragraphs[0]
        p.text = "金蝶软件（中国）有限公司"
        p.font.name = '微软雅黑'
        p.font.size = Pt(14)
        p.font.color.rgb = self.C['white']
        
        # 日期
        from datetime import datetime
        db = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(6), Inches(0.3))
        p = db.text_frame.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月")
        p.font.name = '微软雅黑'
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
        
        print(f"[{self.page:3d}] 封面")
        return slide
    
    def add_toc(self, items):
        """目录页"""
        slide = self._add_slide("目录、提纲")
        self._add_title_bar(slide, "目 录")
        
        half = (len(items) + 1) // 2
        for i, item in enumerate(items, 1):
            col = 0 if i <= half else 1
            row = (i - 1) if i <= half else (i - half - 1)
            x = Inches(0.8 + col * 6.2)
            y = Inches(1.3 + row * 0.9)
            
            # 编号
            num_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = self.C['kingdee_blue']
            num_bg.line.fill.background()
            
            num_tb = slide.shapes.add_textbox(x, y + Inches(0.08), Inches(0.5), Inches(0.4))
            p = num_tb.text_frame.paragraphs[0]
            p.text = str(i)
            p.font.name = '微软雅黑'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.C['white']
            p.alignment = PP_ALIGN.CENTER
            
            # 文本
            item_tb = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.1), Inches(5.5), Inches(0.4))
            p = item_tb.text_frame.paragraphs[0]
            p.text = item
            p.font.name = '微软雅黑'
            p.font.size = Pt(16)
            p.font.color.rgb = self.C['text_primary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] 目录")
        return slide
    
    def add_section(self, title, num=None):
        """章节页"""
        slide = self._add_slide("目录、提纲")
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.C['primary']
        bg.line.fill.background()
        
        if num:
            num_tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(3), Inches(0.8))
            p = num_tb.text_frame.paragraphs[0]
            p.text = f"第{num}章"
            p.font.name = '微软雅黑'
            p.font.size = Pt(20)
            p.font.color.rgb = self.C['accent']
        
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = '微软雅黑'
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        
        deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(2), Inches(0.06))
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.C['kingdee_blue']
        deco.line.fill.background()
        
        print(f"[{self.page:3d}] 章节: {title}")
        return slide
    
    def add_thanks(self):
        """感谢页"""
        slide = self._add_slide("封底")
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.C['primary']
        bg.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = "感谢聆听"
        p.font.name = '微软雅黑'
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        p.alignment = PP_ALIGN.CENTER
        
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.8))
        p = sb.text_frame.paragraphs[0]
        p.text = "THANK YOU"
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5))
        p = cb.text_frame.paragraphs[0]
        p.text = "金蝶软件（中国）有限公司"
        p.font.name = '微软雅黑'
        p.font.size = Pt(16)
        p.font.color.rgb = self.C['white']
        p.alignment = PP_ALIGN.CENTER
        
        print(f"[{self.page:3d}] 感谢页")
        return slide
    
    # ========== 视觉元素 ==========
    
    def add_cards(self, title, cards, cols=3, subtitle=None):
        """卡片页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        cw, ch = 3.9, 1.5
        for i, card in enumerate(cards):
            col, row = i % cols, i // cols
            x = Inches(0.5 + col * (cw + 0.25))
            y = Inches(1.3 + row * (ch + 0.2))
            
            # 卡片背景
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(cw), Inches(ch))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.C['bg_card']
            bg.line.color.rgb = self.C['border']
            
            # 左侧装饰
            deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), Inches(ch))
            deco.fill.solid()
            deco.fill.fore_color.rgb = self.C['kingdee_blue']
            deco.line.fill.background()
            
            # 数值
            vb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.25), Inches(cw - 0.4), Inches(0.7))
            p = vb.text_frame.paragraphs[0]
            p.text = str(card.get('value', ''))
            p.font.name = '微软雅黑'
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.C['kingdee_blue']
            
            # 标签
            lb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.0), Inches(cw - 0.4), Inches(0.35))
            p = lb.text_frame.paragraphs[0]
            p.text = card.get('label', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(12)
            p.font.color.rgb = self.C['text_secondary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [卡片x{len(cards)}]")
        return slide
    
    def add_chart(self, title, chart_type, data, subtitle=None):
        """图表页 - 带真实数据"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        # 创建图表数据
        chart_data = CategoryChartData()
        chart_data.categories = data['categories']
        
        for sn, sv in data['series'].items():
            chart_data.add_series(sn, sv)
        
        # 添加图表
        chart_shape = slide.shapes.add_chart(chart_type, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6), chart_data)
        chart = chart_shape.chart
        
        # 设置图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [图表: {len(data['categories'])}类别, {len(data['series'])}系列]")
        return slide
    
    def add_chart_with_cards(self, title, chart_type, chart_data, cards):
        """图表+卡片混合"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title)
        
        # 图表（左侧）
        cd = CategoryChartData()
        cd.categories = chart_data['categories']
        for sn, sv in chart_data['series'].items():
            cd.add_series(sn, sv)
        
        chart_shape = slide.shapes.add_chart(chart_type, Inches(0.5), Inches(1.1), Inches(7.5), Inches(4.5), cd)
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # 卡片（右侧）
        for i, card in enumerate(cards[:4]):
            y = Inches(1.1 + i * 1.3)
            
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), y, Inches(4.5), Inches(1.1))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.C['bg_card']
            bg.line.color.rgb = self.C['border']
            
            deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), y, Inches(0.06), Inches(1.1))
            deco.fill.solid()
            deco.fill.fore_color.rgb = self.C['kingdee_blue']
            deco.line.fill.background()
            
            vb = slide.shapes.add_textbox(Inches(8.5), y + Inches(0.15), Inches(2), Inches(0.55))
            p = vb.text_frame.paragraphs[0]
            p.text = str(card.get('value', ''))
            p.font.name = '微软雅黑'
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = self.C['kingdee_blue']
            
            lb = slide.shapes.add_textbox(Inches(8.5), y + Inches(0.7), Inches(4), Inches(0.3))
            p = lb.text_frame.paragraphs[0]
            p.text = card.get('label', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(11)
            p.font.color.rgb = self.C['text_secondary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [图表+卡片]")
        return slide
    
    def add_table(self, title, data, subtitle=None):
        """表格页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        rows, cols = len(data), len(data[0])
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5)).table
        
        # 设置列宽
        col_widths = [2.5, 4.5, 2.5, 2.8] if cols == 4 else [int(Inches(12.3) / cols)] * cols
        for i, w in enumerate(col_widths[:cols]):
            table.columns[i].width = Inches(w)
        
        # 填充数据
        for ri, rd in enumerate(data):
            for ci, ct in enumerate(rd):
                cell = table.rows[ri].cells[ci]
                cell.text = str(ct)
                for p in cell.text_frame.paragraphs:
                    p.font.name = '微软雅黑'
                    p.font.size = Pt(11)
                    p.alignment = PP_ALIGN.CENTER
                    
                    if ri == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.C['primary']
                        p.font.color.rgb = self.C['white']
                        p.font.bold = True
                    else:
                        if ri % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.C['bg_light']
                        p.font.color.rgb = self.C['text_primary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [表格: {rows}x{cols}]")
        return slide
    
    def add_timeline(self, title, items, subtitle=None):
        """时间线页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        start_y = 1.8
        
        # 主线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            Inches(0.8), Inches(start_y + 0.15), Inches(11.7), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = self.C['border']
        line.line.fill.background()
        
        for i, item in enumerate(items):
            x = Inches(0.8 + i * 2.5)
            
            # 圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), Inches(start_y), Inches(0.35), Inches(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.C['kingdee_blue']
            dot.line.fill.background()
            
            # 内圆
            inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.93), Inches(start_y + 0.08), Inches(0.19), Inches(0.19))
            inner.fill.solid()
            inner.fill.fore_color.rgb = self.C['white']
            inner.line.fill.background()
            
            # 时间
            tb1 = slide.shapes.add_textbox(x, Inches(start_y + 0.45), Inches(2.2), Inches(0.35))
            p = tb1.text_frame.paragraphs[0]
            p.text = item.get('time', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.C['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 内容卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                x, Inches(start_y + 0.9), Inches(2.2), Inches(2.0))
            card.fill.solid()
            card.fill.fore_color.rgb = self.C['bg_card']
            card.line.color.rgb = self.C['border']
            
            tb2 = slide.shapes.add_textbox(x + Inches(0.1), Inches(start_y + 1.0), Inches(2.0), Inches(1.8))
            tf = tb2.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get('content', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(10)
            p.font.color.rgb = self.C['text_primary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [时间线x{len(items)}]")
        return slide
    
    def add_icon_grid(self, title, items, cols=4, subtitle=None):
        """图标网格页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        cw, ch = 2.9, 1.5
        for i, item in enumerate(items):
            col, row = i % cols, i // cols
            x = Inches(0.5 + col * (cw + 0.2))
            y = Inches(1.3 + row * (ch + 0.15))
            
            # 卡片
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(cw), Inches(ch))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.C['bg_card']
            bg.line.color.rgb = self.C['border']
            
            # 图标
            icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4))
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = self.C['kingdee_blue']
            icon_bg.line.fill.background()
            
            # 标题
            tb = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.18), Inches(cw - 0.8), Inches(0.35))
            p = tb.text_frame.paragraphs[0]
            p.text = item.get('title', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.C['primary']
            
            # 描述
            desc = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.65), Inches(cw - 0.3), Inches(0.75))
            tf = desc.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get('desc', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(9)
            p.font.color.rgb = self.C['text_secondary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [图标网格x{len(items)}]")
        return slide
    
    def add_comparison(self, title, left_title, left_items, right_title, right_items, subtitle=None):
        """对比页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title, subtitle)
        
        # 左侧标题
        left_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.9), Inches(0.5))
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = self.C['accent']
        left_header.line.fill.background()
        
        left_title_tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.28), Inches(5.7), Inches(0.4))
        p = left_title_tb.text_frame.paragraphs[0]
        p.text = left_title
        p.font.name = '微软雅黑'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        
        # 左侧内容
        for i, item in enumerate(left_items[:8]):
            y = Inches(1.85 + i * 0.55)
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1), Inches(0.1), Inches(0.1))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self.C['accent']
            bullet.line.fill.background()
            
            tb = slide.shapes.add_textbox(Inches(0.9), y, Inches(5.5), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = item
            p.font.name = '微软雅黑'
            p.font.size = Pt(11)
            p.font.color.rgb = self.C['text_primary']
        
        # 右侧标题
        right_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.5))
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = self.C['success']
        right_header.line.fill.background()
        
        right_title_tb = slide.shapes.add_textbox(Inches(7), Inches(1.28), Inches(5.7), Inches(0.4))
        p = right_title_tb.text_frame.paragraphs[0]
        p.text = right_title
        p.font.name = '微软雅黑'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        
        # 右侧内容
        for i, item in enumerate(right_items[:8]):
            y = Inches(1.85 + i * 0.55)
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), y + Inches(0.1), Inches(0.1), Inches(0.1))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self.C['success']
            bullet.line.fill.background()
            
            tb = slide.shapes.add_textbox(Inches(7.3), y, Inches(5.5), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = item
            p.font.name = '微软雅黑'
            p.font.size = Pt(11)
            p.font.color.rgb = self.C['text_primary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [对比页]")
        return slide
    
    def add_content_with_cards(self, title, content, cards):
        """内容+卡片混合页"""
        slide = self._add_slide("白色内页")
        self._add_title_bar(slide, title)
        
        # 内容区域（左侧）
        lines = content.strip().split('\n')
        y_offset = 1.2
        
        for line in lines[:12]:
            if not line.strip():
                y_offset += 0.15
                continue
            
            is_title = line.startswith('【') and line.endswith('】')
            is_bullet = line.startswith('•')
            
            if is_title:
                title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Inches(0.5), Inches(y_offset), Inches(7.5), Inches(0.4))
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = self.C['bg_light']
                title_bg.line.fill.background()
                
                icon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                    Inches(0.5), Inches(y_offset), Inches(0.06), Inches(0.4))
                icon.fill.solid()
                icon.fill.fore_color.rgb = self.C['kingdee_blue']
                icon.line.fill.background()
                
                tb = slide.shapes.add_textbox(Inches(0.65), Inches(y_offset + 0.06), Inches(7.3), Inches(0.3))
                p = tb.text_frame.paragraphs[0]
                p.text = line
                p.font.name = '微软雅黑'
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.C['primary']
                y_offset += 0.5
            elif is_bullet:
                bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                    Inches(0.6), Inches(y_offset + 0.08), Inches(0.1), Inches(0.1))
                bullet.fill.solid()
                bullet.fill.fore_color.rgb = self.C['kingdee_blue']
                bullet.line.fill.background()
                
                tb = slide.shapes.add_textbox(Inches(0.8), Inches(y_offset), Inches(7.2), Inches(0.35))
                p = tb.text_frame.paragraphs[0]
                p.text = line[1:].strip()
                p.font.name = '微软雅黑'
                p.font.size = Pt(11)
                p.font.color.rgb = self.C['text_primary']
                y_offset += 0.35
            else:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(7.5), Inches(0.35))
                p = tb.text_frame.paragraphs[0]
                p.text = line
                p.font.name = '微软雅黑'
                p.font.size = Pt(11)
                p.font.color.rgb = self.C['text_primary']
                y_offset += 0.35
        
        # 卡片（右侧）
        for i, card in enumerate(cards[:4]):
            y = Inches(1.3 + i * 1.3)
            
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), y, Inches(4.5), Inches(1.1))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.C['bg_card']
            bg.line.color.rgb = self.C['border']
            
            deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), y, Inches(0.06), Inches(1.1))
            deco.fill.solid()
            deco.fill.fore_color.rgb = self.C['kingdee_blue']
            deco.line.fill.background()
            
            vb = slide.shapes.add_textbox(Inches(8.5), y + Inches(0.15), Inches(2), Inches(0.55))
            p = vb.text_frame.paragraphs[0]
            p.text = str(card.get('value', ''))
            p.font.name = '微软雅黑'
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = self.C['kingdee_blue']
            
            lb = slide.shapes.add_textbox(Inches(8.5), y + Inches(0.7), Inches(4), Inches(0.3))
            p = lb.text_frame.paragraphs[0]
            p.text = card.get('label', '')
            p.font.name = '微软雅黑'
            p.font.size = Pt(11)
            p.font.color.rgb = self.C['text_secondary']
        
        self._add_footer(slide)
        print(f"[{self.page:3d}] {title} [内容+卡片]")
        return slide
    
    def save(self, filename):
        """保存PPT"""
        output_path = os.path.join(os.path.expanduser("~/.openclaw/workspace/output"), filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        
        file_size = os.path.getsize(output_path)
        print(f"\n保存成功: {output_path}")
        print(f"文件大小: {file_size / 1024:.1f} KB")
        print(f"总页数: {self.page} 页")
        
        return output_path
