#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 - 专业版 v10.0
参考真实述标文档，包含专业布局、表格、图表
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 专业配色方案
COLORS = {
    'primary': RGBColor(0, 82, 147),      # 金蝶蓝
    'secondary': RGBColor(0, 122, 194),   # 浅蓝
    'accent': RGBColor(255, 153, 0),      # 橙色强调
    'dark': RGBColor(51, 51, 51),         # 深灰文字
    'light': RGBColor(240, 245, 250),     # 浅蓝背景
    'white': RGBColor(255, 255, 255),
    'success': RGBColor(0, 176, 80),      # 绿色
    'warning': RGBColor(255, 192, 0),     # 黄色
}

# 幻灯片尺寸（16:9）
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """添加封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf3 = footer.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"金蝶软件（中国）有限公司 | {datetime.now().strftime('%Y年%m月')}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLORS['white']
    p3.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, number=""):
    """添加章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 左侧色块
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), SLIDE_HEIGHT)
    left.fill.solid()
    left.fill.fore_color.rgb = COLORS['primary']
    left.line.fill.background()
    
    # 章节编号
    if number:
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # 右侧标题
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.8), Inches(8), Inches(1.5))
    tf2 = title_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = COLORS['dark']
    
    return slide

def add_content_slide(prs, title, content_items, has_icons=False):
    """添加内容页（带要点）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        
        p.text = f"● {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 计算表格尺寸
    cols = len(headers)
    table_width = min(Inches(12), Inches(0.5 + cols * 2))
    table_height = Inches(0.4 + len(rows) * 0.5)
    
    # 创建表格
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.3), table_width, table_height).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # 填充数据
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER
            # 隔行变色
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light']
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.6))
    tf_lt = left_title_box.text_frame
    p_lt = tf_lt.paragraphs[0]
    p_lt.text = left_title
    p_lt.font.size = Pt(20)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLORS['secondary']
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.5), Inches(5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf_l.paragraphs[0]
        else:
            p = tf_l.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(5.5), Inches(0.6))
    tf_rt = right_title_box.text_frame
    p_rt = tf_rt.paragraphs[0]
    p_rt.text = right_title
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLORS['secondary']
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.9), Inches(5.5), Inches(5))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf_r.paragraphs[0]
        else:
            p = tf_r.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(8)
    
    return slide

def add_chart_slide(prs, title, chart_type, categories, values, series_name=""):
    """添加图表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 创建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name or '数据', values)
    
    # 选择图表类型
    if chart_type == 'bar':
        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    elif chart_type == 'line':
        chart_type_enum = XL_CHART_TYPE.LINE
    elif chart_type == 'pie':
        chart_type_enum = XL_CHART_TYPE.PIE
    else:
        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    
    # 添加图表
    chart = slide.shapes.add_chart(
        chart_type_enum,
        Inches(1), Inches(1.5),
        Inches(11), Inches(5),
        chart_data
    ).chart
    
    # 设置图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    return slide

def generate_presales_ppt_v10(customer_info):
    """生成专业版售前PPT"""
    company_name = customer_info.get('companyName', '客户企业')
    industry = customer_info.get('industry', '制造业')
    company_size = customer_info.get('companySize', '中型企业')
    modules = customer_info.get('modules', ['finance', 'supply'])
    employees = customer_info.get('employees', '1000')
    revenue = customer_info.get('revenue', '10000')
    pain_points = customer_info.get('painPoints', '')
    business_goals = customer_info.get('businessGoals', '')
    
    # 模块名称映射
    module_names = {
        'finance': '财务管理',
        'supply': '供应链管理',
        'manufacturing': '制造管理',
        'hr': '人力资源管理'
    }
    module_list = [module_names.get(m, m) for m in modules]
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    today = datetime.now().strftime('%Y年%m月%d日')
    
    # ========== 第1部分：封面 ==========
    add_title_slide(prs, f"{company_name}ERP系统解决方案", "数字化转型项目实施方案")
    
    # ========== 第2部分：目录 ==========
    toc_items = [
        "01 项目需求理解",
        "02 解决方案设计",
        "03 实施方案规划",
        "04 项目团队配置",
        "05 成功案例分享",
        "06 服务承诺保障"
    ]
    add_content_slide(prs, "目 录", toc_items)
    
    # ========== 第3部分：项目需求理解 ==========
    add_section_slide(prs, "项目需求理解", "01")
    
    # 企业概况
    overview_items = [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{company_size}",
        f"员工人数：{employees}人",
        f"年营业额：{revenue}万元",
        f"实施模块：{'、'.join(module_list)}"
    ]
    add_content_slide(prs, "一、企业概况", overview_items)
    
    # 业务痛点
    if pain_points:
        pain_items = [p.strip() for p in pain_points.split('\n') if p.strip()][:8]
        add_content_slide(prs, "二、当前业务痛点", pain_items)
    else:
        pain_items_default = [
            "信息系统分散，数据孤岛严重",
            "业务流程不顺畅，协同效率低",
            "财务核算不及时，决策支持不足",
            "供应链管理粗放，库存成本高",
            "生产计划不准确，交付周期长"
        ]
        add_content_slide(prs, "二、当前业务痛点", pain_items_default)
    
    # 建设目标
    if business_goals:
        goal_items = [g.strip() for g in business_goals.split('\n') if g.strip()][:8]
        add_content_slide(prs, "三、项目建设目标", goal_items)
    else:
        goal_items_default = [
            "构建统一数字化平台，消除信息孤岛",
            "优化业务流程，提升运营效率",
            "实现财务业务一体化，提升决策质量",
            "降低运营成本，提高企业竞争力",
            "建立数据治理体系，实现数据驱动"
        ]
        add_content_slide(prs, "三、项目建设目标", goal_items_default)
    
    # ========== 第4部分：解决方案设计 ==========
    add_section_slide(prs, "解决方案设计", "02")
    
    # 总体架构
    add_two_column_slide(prs, "一、总体架构设计",
        "业务架构", [
            "销售到收款全流程管理",
            "采购到付款全流程管理",
            "生产计划与执行管理",
            "库存与物流管理",
            "财务核算与管理会计",
            "成本管理与分析"
        ],
        "技术架构", [
            "云原生微服务架构",
            "分布式部署支持",
            "多租户多组织",
            "弹性扩展能力",
            "高可用高安全",
            "开放集成平台"
        ]
    )
    
    # 实施模块
    module_details = {
        'finance': [
            "总账管理：凭证处理、期末结账、报表生成",
            "应收管理：销售开票、收款核销、账龄分析",
            "应付管理：采购发票、付款处理、供应商对账",
            "固定资产：资产台账、折旧计提、资产处置",
            "成本管理：成本核算、成本分析、成本控制"
        ],
        'supply': [
            "销售管理：订单处理、价格管理、出货管理",
            "采购管理：询价比价、订单跟踪、入库检验",
            "库存管理：出入库管理、盘点管理、库存分析",
            "计划管理：需求计划、采购计划、生产计划"
        ],
        'manufacturing': [
            "生产计划：主生产计划、物料需求计划",
            "生产执行：工单管理、工序汇报、完工入库",
            "质量管理：来料检验、过程检验、成品检验",
            "设备管理：设备台账、维护保养、故障处理"
        ],
        'hr': [
            "组织管理：组织架构、岗位体系、编制管理",
            "人事管理：员工档案、入离职管理、合同管理",
            "薪酬管理：薪酬核算、个税计算、社保管理",
            "绩效管理：绩效考核、结果应用、能力发展"
        ]
    }
    
    for i, module in enumerate(modules[:3], 1):
        module_name = module_names.get(module, module)
        details = module_details.get(module, [])
        if details:
            add_content_slide(prs, f"{['二', '三', '四'][i-1]}、{module_name}方案", details)
    
    # 价值分析
    add_table_slide(prs, "三、预期价值分析",
        ["价值维度", "指标项", "预期提升", "说明"],
        [
            ["运营效率", "业务处理效率", "提升40%", "流程自动化、数据实时共享"],
            ["运营效率", "库存周转率", "提升30%", "精准计划、库存优化"],
            ["财务管理", "结账周期", "缩短50%", "自动化核算、实时报表"],
            ["财务管理", "成本核算精度", "提升35%", "精细核算、多维分析"],
            ["决策支持", "数据准确性", "提升40%", "数据治理、主数据管理"],
            ["决策支持", "决策响应速度", "提升60%", "实时分析、智能预警"]
        ]
    )
    
    # ========== 第5部分：实施方案 ==========
    add_section_slide(prs, "实施方案规划", "03")
    
    # 实施方法论
    add_content_slide(prs, "一、实施方法论", [
        "金蝶云实施方法论：六阶段实施法",
        "第一阶段：项目启动（2周）- 成立组织、明确范围",
        "第二阶段：需求调研（4周）- 业务调研、需求分析",
        "第三阶段：方案设计（3周）- 蓝图设计、方案确认",
        "第四阶段：系统配置（4周）- 参数配置、流程设置",
        "第五阶段：测试培训（3周）- UAT测试、用户培训",
        "第六阶段：上线验收（2周）- 系统切换、项目验收"
    ])
    
    # 实施计划表
    add_table_slide(prs, "二、实施计划安排",
        ["阶段", "主要工作", "周期", "交付成果"],
        [
            ["项目启动", "组建团队、制定计划", "2周", "项目章程、实施计划"],
            ["需求调研", "业务调研、差距分析", "4周", "调研报告、需求清单"],
            ["方案设计", "蓝图设计、方案评审", "3周", "业务蓝图、配置方案"],
            ["系统配置", "系统设置、数据准备", "4周", "系统配置、基础数据"],
            ["测试培训", "UAT测试、操作培训", "3周", "测试报告、培训记录"],
            ["上线验收", "系统切换、验收评审", "2周", "上线报告、验收报告"]
        ]
    )
    
    # ========== 第6部分：项目团队 ==========
    add_section_slide(prs, "项目团队配置", "04")
    
    # 团队配置表
    add_table_slide(prs, "项目团队配置",
        ["角色", "职责", "人数", "资质要求"],
        [
            ["项目总监", "项目整体把控、资源协调", "1人", "PMP认证、10年以上经验"],
            ["项目经理", "项目日常管理、进度控制", "1人", "PMP认证、5年以上经验"],
            ["业务顾问", "业务方案设计、流程优化", "2人", "行业认证、3年以上经验"],
            ["技术顾问", "技术方案设计、系统集成", "1人", "技术认证、3年以上经验"],
            ["实施顾问", "系统配置、用户培训", "3人", "产品认证、2年以上经验"],
            ["开发工程师", "客户化开发、报表开发", "2人", "开发认证、2年以上经验"]
        ]
    )
    
    # ========== 第7部分：成功案例 ==========
    add_section_slide(prs, "成功案例分享", "05")
    
    # 案例介绍
    add_two_column_slide(prs, "行业典型案例",
        "案例一：某大型制造企业", [
            "企业规模：年营收50亿，员工5000人",
            "实施模块：财务、供应链、制造、HR",
            "实施周期：8个月",
            "项目成效：",
            "  - 财务结账时间缩短60%",
            "  - 库存周转率提升40%",
            "  - 生产计划准确率提升35%"
        ],
        "案例二：某连锁零售企业", [
            "企业规模：年营收20亿，员工3000人",
            "实施模块：财务、供应链、零售管理",
            "实施周期：6个月",
            "项目成效：",
            "  - 采购成本降低15%",
            "  - 库存准确率达到99%",
            "  - 销售预测准确率提升25%"
        ]
    )
    
    # ========== 第8部分：服务承诺 ==========
    add_section_slide(prs, "服务承诺保障", "06")
    
    # 服务承诺
    add_content_slide(prs, "服务承诺", [
        "系统稳定性承诺：系统可用率≥99.9%",
        "响应时效承诺：一级问题2小时内响应，8小时内解决",
        "培训承诺：提供不少于10场系统培训，培训覆盖率100%",
        "文档承诺：提供完整的操作手册和管理文档",
        "质保承诺：系统上线后提供1年免费质保服务",
        "升级承诺：提供系统版本升级服务和技术支持"
    ])
    
    # 联系方式
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf2 = contact.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"金蝶软件（中国）有限公司\n{today}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    customer_code = customer_info.get('customerCode', 'CUSTOMER')
    filename = f"{customer_code}_售前解决方案_v10_{timestamp}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    
    prs.save(filepath)
    
    return {
        'success': True,
        'filepath': filepath,
        'filename': filename,
        'slides': len(prs.slides)
    }

if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='生成售前PPT')
    parser.add_argument('--type', default='presales', help='PPT类型')
    parser.add_argument('--companyName', required=True)
    parser.add_argument('--customerCode', default='CUSTOMER')
    parser.add_argument('--industry', default='制造业')
    parser.add_argument('--companySize', default='中型企业')
    parser.add_argument('--employees', default='')
    parser.add_argument('--revenue', default='')
    parser.add_argument('--modules', default='finance,supply')
    parser.add_argument('--painPoints', default='')
    parser.add_argument('--businessGoals', default='')
    parser.add_argument('--goliveDate', default='', help='兼容参数')
    parser.add_argument('--golivePhase', default='', help='兼容参数')
    parser.add_argument('--switchPlan', default='', help='兼容参数')
    parser.add_argument('--acceptanceDate', default='', help='兼容参数')
    parser.add_argument('--acceptanceConclusion', default='', help='兼容参数')
    
    args = parser.parse_args()
    
    customer_info = {
        'companyName': args.companyName,
        'customerCode': args.customerCode,
        'industry': args.industry,
        'companySize': args.companySize,
        'employees': args.employees,
        'revenue': args.revenue,
        'modules': args.modules.split(','),
        'painPoints': args.painPoints,
        'businessGoals': args.businessGoals
    }
    
    result = generate_presales_ppt_v10(customer_info)
    print(json.dumps(result, ensure_ascii=False))
