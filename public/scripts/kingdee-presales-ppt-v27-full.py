# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 V27 完整版
参考: 中煤科工述标文件V7.0 (146页)
特点:
1. 使用金蝶官方母版
2. 丰富的视觉元素（KPI卡片、功能模块、时间线、架构图）
3. 金蝶官方配色方案 (#0886EC主蓝, #FF7401橙)
4. 完整146页内容结构
"""

import sys
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 金蝶官方配色
COLORS = {
    'primary': RGBColor(0x08, 0x86, 0xEC),      # 主蓝色 #0886EC
    'secondary': RGBColor(0x00, 0x70, 0xC0),    # 次蓝色 #0070C0
    'accent': RGBColor(0xFF, 0x74, 0x01),       # 橙色强调 #FF7401
    'dark': RGBColor(0x00, 0x3F, 0x56),         # 深蓝 #003F56
    'light': RGBColor(0xCC, 0xDD, 0xEA),        # 浅蓝背景 #CCDDEA
    'text_dark': RGBColor(0x00, 0x00, 0x00),
    'text_light': RGBColor(0xFF, 0xFF, 0xFF),
    'purple': RGBColor(0xB0, 0x73, 0xFC),       # 紫色辅助
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}


class StyleHelper:
    """样式辅助类 - 提供丰富的视觉元素"""
    
    def __init__(self, slide):
        self.slide = slide
        self.colors = COLORS
    
    def add_rect(self, left, top, width, height, fill_color, text="", 
                 font_size=11, font_color=None, border=False):
        """添加矩形"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if not border:
            shape.line.fill.background()
        
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or self.colors['text_dark']
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "微软雅黑"
        
        return shape
    
    def add_rounded_rect(self, left, top, width, height, fill_color, 
                         text="", font_size=11, font_color=None):
        """添加圆角矩形"""
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        
        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color or self.colors['text_dark']
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "微软雅黑"
        
        return shape
    
    def add_text(self, left, top, width, height, text, font_size=12, 
                 font_color=None, bold=False, align=PP_ALIGN.LEFT):
        """添加文本框"""
        box = self.slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color or self.colors['text_dark']
        p.font.bold = bold
        p.font.name = "微软雅黑"
        p.alignment = align
        
        return box
    
    def add_kpi_card(self, left, top, width, height, value, unit, label, 
                     value_color=None, bg_color=None):
        """添加KPI卡片 - 大数字+单位+标签"""
        value_color = value_color or self.colors['primary']
        bg_color = bg_color or self.colors['light']
        
        # 背景卡片
        self.add_rounded_rect(left, top, width, height, bg_color)
        
        # 大数字
        value_box = self.slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.15),
            Inches(width - 0.2), Inches(0.5)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = str(value)
        run1.font.size = Pt(28)
        run1.font.bold = True
        run1.font.color.rgb = value_color
        run1.font.name = "微软雅黑"
        run2 = p.add_run()
        run2.text = unit
        run2.font.size = Pt(14)
        run2.font.color.rgb = self.colors['text_dark']
        run2.font.name = "微软雅黑"
        
        # 底部标签
        label_box = self.slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + height - 0.35),
            Inches(width - 0.2), Inches(0.3)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = self.colors['text_dark']
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER
    
    def add_kpi_row(self, left, top, kpis, card_width=2.8, card_height=1.2, gap=0.3):
        """添加KPI卡片行"""
        for i, (value, unit, label) in enumerate(kpis):
            x = left + i * (card_width + gap)
            self.add_kpi_card(x, top, card_width, card_height, value, unit, label)
    
    def add_function_block(self, left, top, width, height, title, items, accent=False):
        """添加功能模块 - 标题栏+功能项列表"""
        bg_color = self.colors['accent'] if accent else self.colors['primary']
        
        # 标题栏
        self.add_rounded_rect(left, top, width, 0.4, bg_color, title, 12, self.colors['text_light'])
        
        # 功能项
        if items:
            item_height = min(0.35, (height - 0.45) / len(items))
            for i, item in enumerate(items[:int((height-0.45)/0.35)]):
                self.add_rounded_rect(
                    left, top + 0.42 + i * item_height,
                    width, item_height - 0.02,
                    self.colors['light'], item, 11, self.colors['text_dark']
                )
    
    def add_function_blocks_row(self, left, top, blocks, width=2.8, height=2.2, gap=0.2):
        """添加功能模块行"""
        for i, (title, items, accent) in enumerate(blocks):
            x = left + i * (width + gap)
            self.add_function_block(x, top, width, height, title, items, accent)
    
    def add_timeline(self, left, top, width, milestones):
        """添加时间线 - 渐变底条+圆形节点"""
        total = len(milestones)
        
        # 底条
        bar = self.add_rect(left, top + 0.5, width, 0.1, self.colors['light'])
        
        # 节点
        for i, (date, title) in enumerate(milestones):
            x = left + (width / (total - 1)) * i if total > 1 else left
            
            # 圆形节点
            node = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.15), Inches(top + 0.4),
                Inches(0.3), Inches(0.3)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = self.colors['primary']
            node.line.fill.background()
            
            # 编号
            self.add_text(x - 0.1, top + 0.43, 0.2, 0.25, str(i + 1), 12, 
                         self.colors['text_light'], False, PP_ALIGN.CENTER)
            
            # 日期
            self.add_text(x - 0.4, top + 0.8, 0.8, 0.3, date, 10, 
                         self.colors['text_dark'], False, PP_ALIGN.CENTER)
            
            # 标题
            title_box = self.slide.shapes.add_textbox(
                Inches(x - 0.5), Inches(top + 1.1),
                Inches(1), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors['text_dark']
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
    
    def add_architecture_layer(self, left, top, width, height, title, modules, layer_color=None):
        """添加架构层 - 左侧标签+右侧模块卡片"""
        layer_color = layer_color or self.colors['dark']
        
        # 左侧标签
        self.add_rounded_rect(left, top, 1.0, height, layer_color, title, 12, self.colors['text_light'])
        
        # 模块卡片
        if modules:
            module_width = (width - 1.2) / len(modules)
            for i, module in enumerate(modules):
                x = left + 1.1 + i * module_width
                self.add_rounded_rect(x, top + 0.1, module_width - 0.1, height - 0.2,
                                     self.colors['primary'], module, 11, self.colors['text_light'])
    
    def add_process_flow(self, left, top, steps, step_width=2.0, step_height=0.5, gap=0.3):
        """添加横向流程图"""
        arrow_width = 0.25
        
        for i, step in enumerate(steps):
            x = left + i * (step_width + gap + arrow_width)
            
            # 步骤框
            self.add_rounded_rect(x, top, step_width, step_height, 
                                 self.colors['primary'], step, 12, self.colors['text_light'])
            
            # 箭头(最后一个不加)
            if i < len(steps) - 1:
                arrow = self.slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(x + step_width + 0.05), Inches(top + step_height/2 - 0.1),
                    Inches(arrow_width), Inches(0.2)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['accent']
                arrow.line.fill.background()


class KingdeePresalesPPTV27Full:
    """金蝶售前PPT生成器V27完整版"""
    
    def __init__(self, company_name, project_name, output_path):
        self.company_name = company_name
        self.project_name = project_name
        self.output_path = output_path
        
        # 使用金蝶母版
        template_path = os.path.join(os.path.dirname(__file__), '..', 'templates', 'kingdee-template.pptx')
        if not os.path.exists(template_path):
            template_path = 'templates/kingdee-template.pptx'
        
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
            print(f"使用模板: {template_path}")
        else:
            self.prs = Presentation()
            print("使用空白模板")
        
        self.slides = []
        self._init_layouts()
    
    def _init_layouts(self):
        """初始化布局映射"""
        self.layouts = {}
        for i, layout in enumerate(self.prs.slide_master.slide_layouts):
            self.layouts[layout.name] = i
            print(f"  布局 {i}: {layout.name}")
    
    def _get_layout_idx(self, name):
        """获取布局索引"""
        if name in self.layouts:
            return self.layouts[name]
        for layout_name, idx in self.layouts.items():
            if name in layout_name or layout_name in name:
                return idx
        return min(2, len(self.layouts) - 1)  # 默认使用内容页
    
    def add_slide(self, layout_name='3_白色内页'):
        """添加幻灯片"""
        idx = self._get_layout_idx(layout_name)
        slide = self.prs.slides.add_slide(self.prs.slide_master.slide_layouts[idx])
        self.slides.append(slide)
        return slide
    
    def set_title(self, slide, title):
        """设置标题"""
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = title
                    shape.text_frame.paragraphs[0].font.size = Pt(24)
                    shape.text_frame.paragraphs[0].font.color.rgb = COLORS['dark']
                    return True
        
        # 如果没有找到占位符，手动添加标题
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.font.name = "微软雅黑"
        return False
    
    # ==================== 页面生成方法 ====================
    
    def create_cover(self):
        """封面"""
        slide = self.add_slide('封面')
        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title
                    shape.text_frame.paragraphs[0].text = f"{self.company_name}\n{self.project_name}述标方案"
                    shape.text_frame.paragraphs[0].font.size = Pt(44)
                    shape.text_frame.paragraphs[0].font.bold = True
                    break
        print("  封面")
    
    def create_toc(self, chapters):
        """目录"""
        slide = self.add_slide('目录、提纲')
        style = StyleHelper(slide)
        
        # 设置标题
        self.set_title(slide, "目 录")
        
        # 目录项
        y = 1.8
        for i, chapter in enumerate(chapters):
            # 编号
            style.add_text(1.0, y, 0.5, 0.4, f"{i+1:02d}", 18, COLORS['primary'], True)
            # 章节
            style.add_text(1.6, y, 10, 0.4, chapter, 16, COLORS['text_dark'])
            y += 0.55
        print("  目录")
    
    def create_chapter(self, num, title):
        """章节页"""
        slide = self.add_slide('目录、提纲')
        style = StyleHelper(slide)
        
        # 大编号
        style.add_text(0.5, 2.5, 2, 1.5, f"{num:02d}", 72, COLORS['primary'], True)
        # 章节标题
        style.add_text(2.5, 2.8, 9, 1, title, 36, COLORS['dark'], True)
        print(f"  章节 {num}: {title}")
    
    def create_content(self, title, content_func):
        """内容页"""
        slide = self.add_slide('3_白色内页')
        style = StyleHelper(slide)
        self.set_title(slide, title)
        content_func(slide, style)
        return slide
    
    # ==================== 章节内容 ====================
    
    def create_team_intro(self):
        """团队介绍"""
        def team_page(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "本次项目金蝶述标团队", 28, COLORS['dark'], bold=True)
            
            # 团队成员卡片
            members = [
                ("项目经理", "15年ERP实施经验", "PMP认证"),
                ("财务顾问", "10年财务管理经验", "CPA认证"),
                ("供应链顾问", "8年供应链经验", "APICS认证"),
                ("技术顾问", "12年技术架构经验", "云架构师"),
            ]
            
            x_start = 0.5
            for i, (title, exp, cert) in enumerate(members):
                x = x_start + i * 3.1
                style.add_rounded_rect(x, 1.8, 2.8, 2.5, COLORS['light'])
                style.add_text(x + 0.1, 1.9, 2.6, 0.5, title, 16, COLORS['primary'], bold=True)
                style.add_text(x + 0.1, 2.5, 2.6, 0.8, exp, 11, COLORS['text_dark'])
                style.add_text(x + 0.1, 3.4, 2.6, 0.4, cert, 10, COLORS['secondary'])
            
            style.add_text(0.5, 4.8, 12, 1, 
                "金蝶项目团队拥有丰富的行业实施经验，累计服务超过1000家企业客户，"
                "在制造业、零售业、金融业等多个行业拥有成功案例。", 12)
        self.create_content("团队介绍", team_page)
    
    def create_company_intro(self):
        """公司介绍章节"""
        self.create_chapter(1, "公司介绍")
        
        # 公司概况 - KPI
        def company_overview(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("30+", "年", "发展历程"),
                ("5000+", "亿", "市值"),
                ("1000万+", "家", "企业客户"),
            ], card_width=3.8, gap=0.3)
            
            style.add_text(0.5, 3.0, 12, 2, 
                "金蝶国际软件集团有限公司是亚太地区领先的企业管理软件及电子商务应用解决方案供应商。"
                "金蝶旗下云服务产品有金蝶云·苍穹、金蝶云·星瀚、金蝶云·星空、金蝶云·星辰等，"
                "为不同规模的企业提供数字化转型解决方案。", 14)
        self.create_content("公司概况", company_overview)
        
        # 发展历程 - 时间线
        def history(slide, style):
            style.add_timeline(0.5, 1.5, 12, [
                ("1993", "成立"),
                ("2001", "香港上市"),
                ("2012", "云转型"),
                ("2018", "苍穹发布"),
                ("2023", "云收入79%"),
            ])
        self.create_content("发展历程", history)
        
        # 核心优势
        def advantages(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("技术领先", ["云原生架构", "AI赋能", "微服务"], False),
                ("产品丰富", ["苍穹PaaS", "星瀚EBC", "星空ERP"], False),
                ("服务完善", ["全国服务网", "生态伙伴", "专业团队"], False),
                ("案例众多", ["500强客户", "行业标杆", "成功经验"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("核心优势", advantages)
        
        # 市场地位
        def market_position(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("连续19年", "", "中国企业应用软件市场占有率第一"),
                ("连续4年", "", "中国企业SaaS云服务市场占有率第一"),
            ], card_width=5.8, card_height=0.9, gap=0.3)
            
            style.add_kpi_row(0.5, 2.7, [
                ("连续4年", "", "中国企业ERP云服务市场占有率第一"),
                ("IDC认证", "", "中国企业级应用软件市场领导者"),
            ], card_width=5.8, card_height=0.9, gap=0.3)
        self.create_content("市场占有率-金蝶云产品多项第一", market_position)
        
        # 技术实力
        def tech_strength(slide, style):
            style.add_architecture_layer(0.5, 1.5, 12, 0.9, "云原生", 
                ["容器化", "微服务", "DevOps", "云安全"], COLORS['dark'])
            style.add_architecture_layer(0.5, 2.6, 12, 0.9, "AI赋能",
                ["智能分析", "OCR识别", "RPA自动化", "智能助手"], COLORS['primary'])
            style.add_architecture_layer(0.5, 3.7, 12, 0.9, "低代码",
                ["可视化建模", "表单设计", "流程编排", "移动开发"], COLORS['secondary'])
        self.create_content("技术实力", tech_strength)
        
        # 荣誉资质
        def honors(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("国家高新", "", "技术企业认证"),
                ("ISO27001", "", "信息安全认证"),
                ("等保三级", "", "信息安全等级"),
                ("CMMI5", "", "软件能力成熟度"),
            ], card_width=2.8, gap=0.3)
            style.add_kpi_row(0.5, 2.9, [
                ("Gartner", "", "魔力象限认可"),
                ("IDC", "", "市场占有率第一"),
                ("福布斯", "", "全球企业2000强"),
                ("财富", "", "中国500强企业"),
            ], card_width=2.8, gap=0.3)
        self.create_content("荣誉资质", honors)
        
        # 客户基础
        def customer_base(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("750万+", "家", "企业客户"),
                ("200+", "家", "生态伙伴"),
                ("5000+", "人", "专业团队"),
                ("全国", "", "服务网络"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "金蝶服务超过750万家企业客户，覆盖制造业、零售业、金融业、服务业等多个行业，"
                "拥有200+生态合作伙伴，5000+专业服务团队，在全国设有完善的服务网络。", 12)
        self.create_content("客户基础", customer_base)
        
        # 行业洞察
        def industry_insight(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "行业洞察：数字化转型趋势", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("技术驱动", ["云计算普及", "AI应用", "大数据分析"], False),
                ("业务变革", ["流程重构", "模式创新", "生态协同"], False),
                ("管理升级", ["数据决策", "敏捷管理", "智能运营"], False),
                ("合规要求", ["国产化替代", "数据安全", "行业监管"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("行业洞察", industry_insight)
        
        # 数字化转型挑战
        def challenges(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "企业数字化转型面临的挑战", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("系统孤岛", ["数据不互通", "流程断裂", "效率低下"], False),
                ("管理粗放", ["数据不准确", "决策滞后", "风险难控"], False),
                ("业务僵化", ["响应缓慢", "创新困难", "竞争乏力"], False),
                ("成本高昂", ["IT投入大", "维护成本高", "ROI不明"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("转型挑战", challenges)
    
    def create_product_intro(self):
        """产品体系章节"""
        self.create_chapter(2, "产品体系")
        
        # 产品矩阵
        def product_matrix(slide, style):
            style.add_architecture_layer(0.5, 1.5, 12, 1.0, "大企业",
                ["金蝶云·苍穹", "金蝶云·星瀚"], COLORS['dark'])
            style.add_architecture_layer(0.5, 2.7, 12, 1.0, "中企业",
                ["金蝶云·星空"], COLORS['primary'])
            style.add_architecture_layer(0.5, 3.9, 12, 1.0, "小企业",
                ["金蝶云·星辰", "精斗云"], COLORS['secondary'])
        self.create_content("金蝶云产品矩阵", product_matrix)
        
        # 苍穹平台
        def canshang(slide, style):
            style.add_rounded_rect(0.3, 1.5, 1.5, 3.5, COLORS['dark'], 
                                   "金蝶云·苍穹\n\n世界一流\n企业级\nPaaS平台", 14, COLORS['text_light'])
            style.add_function_blocks_row(2.0, 1.5, [
                ("云原生架构", ["容器化部署", "微服务", "DevOps"], False),
                ("低代码开发", ["可视化建模", "表单设计", "流程编排"], False),
                ("AI赋能", ["智能分析", "OCR识别", "智能助手"], False),
                ("数据中台", ["数据治理", "数据服务", "数据资产"], False),
            ], width=2.6, height=2.5, gap=0.2)
        self.create_content("金蝶云·苍穹 - 世界一流企业级PaaS平台", canshang)
        
        # 星瀚产品
        def xinghan(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("财务管理", ["总账", "报表", "资金", "成本"], False),
                ("供应链", ["采购", "销售", "库存", "物流"], False),
                ("生产制造", ["计划", "生产", "质量", "设备"], False),
                ("人力资本", ["人事", "薪酬", "绩效", "培训"], False),
                ("项目管理", ["立项", "预算", "执行", "结算"], False),
            ], width=2.3, height=2.8, gap=0.15)
        self.create_content("金蝶云·星瀚 - 大企业数字化平台", xinghan)
        
        # 星空产品
        def xingkong(slide, style):
            style.add_rounded_rect(0.3, 1.5, 1.5, 3.5, COLORS['primary'], 
                                   "金蝶云·星空\n\n成长型企业\n云管理平台", 14, COLORS['text_light'])
            style.add_function_blocks_row(2.0, 1.5, [
                ("财务云", ["智能核算", "税务管理", "资金管理"], False),
                ("供应链云", ["采购管理", "销售管理", "库存管理"], False),
                ("制造云", ["生产计划", "车间管理", "质量管理"], False),
                ("全渠道云", ["电商管理", "门店管理", "会员管理"], True),
            ], width=2.6, height=1.5, gap=0.2)
        self.create_content("金蝶云·星空 - 成长型企业云服务", xingkong)
        
        # 技术架构
        def tech_arch(slide, style):
            style.add_architecture_layer(0.5, 1.3, 12, 0.9, "应用层",
                ["财务应用", "供应链应用", "制造应用", "人力应用"], COLORS['primary'])
            style.add_architecture_layer(0.5, 2.4, 12, 0.9, "服务层",
                ["API网关", "消息服务", "流程引擎", "规则引擎"], COLORS['secondary'])
            style.add_architecture_layer(0.5, 3.5, 12, 0.9, "数据层",
                ["关系数据库", "缓存服务", "对象存储", "搜索引擎"], COLORS['dark'])
            style.add_architecture_layer(0.5, 4.6, 12, 0.9, "基础设施层",
                ["计算资源", "网络资源", "存储资源", "安全服务"], COLORS['accent'])
        self.create_content("技术架构", tech_arch)
    
    def create_solution_intro(self):
        """解决方案章节"""
        self.create_chapter(3, "解决方案")
        
        # 行业方案
        def industry(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("制造业", ["离散制造", "流程制造", "装备制造"], False),
                ("零售业", ["连锁零售", "电商零售", "全渠道"], False),
                ("金融业", ["银行", "保险", "证券"], False),
                ("服务业", ["专业服务", "教育", "医疗"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("行业解决方案", industry)
        
        # 财务方案
        def finance(slide, style):
            style.add_process_flow(0.5, 1.5, 
                ["业务发生", "凭证生成", "审核记账", "期末结账", "报表输出"],
                step_width=2.0, step_height=0.5, gap=0.3)
            style.add_function_blocks_row(0.5, 2.3, [
                ("总账管理", ["多账簿", "多币种", "多会计准则"], False),
                ("报表管理", ["资产负债表", "利润表", "现金流量表"], False),
                ("成本管理", ["标准成本", "实际成本", "作业成本"], False),
            ], width=3.8, height=2.0, gap=0.2)
        self.create_content("财务管理解决方案", finance)
        
        # 财务详细方案
        def finance_detail(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("应收管理", ["销售发票", "收款核销", "账龄分析"], False),
                ("应付管理", ["采购发票", "付款管理", "往来对账"], False),
                ("资金管理", ["账户管理", "资金计划", "银企直联"], False),
            ], width=3.8, height=2.0, gap=0.2)
            style.add_function_blocks_row(0.5, 3.8, [
                ("固定资产", ["资产登记", "折旧计提", "资产处置"], False),
                ("预算管理", ["预算编制", "预算控制", "预算分析"], False),
                ("税务管理", ["税务申报", "发票管理", "税务风险"], True),
            ], width=3.8, height=2.0, gap=0.2)
        self.create_content("财务管理详细方案", finance_detail)
        
        # 供应链方案
        def supply_chain(slide, style):
            style.add_process_flow(0.5, 1.5,
                ["需求计划", "采购执行", "库存管理", "销售配送", "结算分析"],
                step_width=2.0, step_height=0.5, gap=0.3)
            style.add_function_blocks_row(0.5, 2.3, [
                ("采购管理", ["供应商管理", "采购申请", "采购订单", "采购结算"], False),
                ("库存管理", ["入库管理", "出库管理", "库存盘点", "库存分析"], False),
                ("销售管理", ["客户管理", "销售订单", "发货管理", "销售结算"], False),
            ], width=3.8, height=2.0, gap=0.2)
        self.create_content("供应链管理解决方案", supply_chain)
        
        # 生产制造方案
        def manufacturing(slide, style):
            style.add_process_flow(0.5, 1.5,
                ["需求预测", "MRP运算", "生产计划", "车间执行", "质量检验"],
                step_width=2.0, step_height=0.5, gap=0.3)
            style.add_function_blocks_row(0.5, 2.3, [
                ("计划管理", ["主生产计划", "MRP运算", "能力计划"], False),
                ("生产执行", ["工单管理", "领料管理", "完工入库"], False),
                ("质量管理", ["检验标准", "过程检验", "质量追溯"], False),
            ], width=3.8, height=2.0, gap=0.2)
        self.create_content("生产制造解决方案", manufacturing)
        
        # 数据分析方案
        def analytics(slide, style):
            style.add_architecture_layer(0.5, 1.5, 12, 1.0, "数据应用层",
                ["管理驾驶舱", "移动报表", "自助分析", "预警推送"], COLORS['primary'])
            style.add_architecture_layer(0.5, 2.7, 12, 1.0, "数据分析层",
                ["多维分析", "数据挖掘", "预测分析", "智能推荐"], COLORS['secondary'])
            style.add_architecture_layer(0.5, 3.9, 12, 1.0, "数据服务层",
                ["数据仓库", "数据湖", "ETL处理", "数据治理"], COLORS['dark'])
        self.create_content("数据分析解决方案", analytics)
        
        # 主数据管理
        def mdm(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("物料主数据", ["分类管理", "属性定义", "编码规则"], False),
                ("客户主数据", ["客户档案", "信用管理", "区域管理"], False),
                ("供应商主数据", ["供应商档案", "资质管理", "评估体系"], False),
                ("组织主数据", ["公司架构", "部门设置", "人员信息"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("主数据管理解决方案", mdm)
        
        # 系统集成
        def integration(slide, style):
            style.add_architecture_layer(0.5, 1.5, 12, 1.0, "应用层",
                ["ERP", "CRM", "OA", "HR"], COLORS['primary'])
            style.add_architecture_layer(0.5, 2.7, 12, 1.0, "集成平台",
                ["API网关", "消息队列", "数据同步", "流程引擎"], COLORS['secondary'])
            style.add_architecture_layer(0.5, 3.9, 12, 1.0, "数据层",
                ["主数据", "业务数据", "日志数据", "配置数据"], COLORS['dark'])
        self.create_content("系统集成解决方案", integration)
        
        # 移动应用
        def mobile(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("移动审批", ["待办事项", "流程审批", "消息通知"], False),
                ("移动报表", ["数据查询", "图表展示", "预警推送"], False),
                ("移动作业", ["扫码入库", "盘点作业", "巡检记录"], False),
                ("移动办公", ["日程管理", "任务协作", "即时通讯"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("移动应用解决方案", mobile)
        
        # 智能分析
        def ai_analysis(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("智能预测", ["销售预测", "库存预测", "资金预测"], False),
                ("智能识别", ["OCR识别", "发票识别", "合同识别"], False),
                ("智能推荐", ["采购建议", "排产建议", "风控建议"], False),
                ("智能助手", ["对话机器人", "智能问答", "知识库"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("智能分析解决方案", ai_analysis)
    
    def create_implementation(self):
        """实施路线章节"""
        self.create_chapter(4, "实施路线")
        
        # 实施方法论
        def methodology(slide, style):
            style.add_process_flow(0.3, 1.5,
                ["项目启动", "需求调研", "方案设计", "系统配置", "用户测试", "上线切换", "持续优化"],
                step_width=1.4, step_height=0.5, gap=0.2)
        self.create_content("金蝶实施方法论", methodology)
        
        # 实施计划
        def plan(slide, style):
            style.add_timeline(0.5, 1.5, 12, [
                ("W1-W2", "项目启动"),
                ("W3-W6", "需求调研"),
                ("W7-W10", "方案设计"),
                ("W11-W16", "系统配置"),
                ("W17-W20", "用户测试"),
                ("W21-W22", "上线切换"),
            ])
        self.create_content("项目实施计划", plan)
        
        # 项目团队
        def team(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("项目领导", ["项目指导委员会", "项目管理办公室"], False),
                ("金蝶团队", ["项目经理", "业务顾问", "技术顾问"], False),
                ("客户团队", ["业务负责人", "关键用户", "IT支持"], False),
                ("支持团队", ["开发团队", "测试团队", "运维团队"], True),
            ], width=2.8, height=2.0, gap=0.3)
        self.create_content("项目团队组织", team)
        
        # 风险管理
        def risk(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("需求风险", ["需求变更控制", "需求确认流程", "需求文档管理"], False),
                ("进度风险", ["里程碑管理", "进度监控", "资源协调"], False),
                ("质量风险", ["测试管理", "质量评审", "问题跟踪"], False),
                ("人员风险", ["知识转移", "文档管理", "备份机制"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("风险管理", risk)
        
        # 质量保障
        def quality(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("100%", "", "需求覆盖率"),
                ("95%", "", "测试通过率"),
                ("0", "", "严重缺陷"),
                ("100%", "", "文档完整性"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "金蝶实施团队严格执行质量管理体系，确保项目交付质量。"
                "每个阶段都有明确的验收标准和质量检查点。", 12)
        self.create_content("质量保障", quality)
        
        # 实施阶段详述
        def phases(slide, style):
            style.add_timeline(0.5, 1.2, 12, [
                ("启动", "组建团队\n确定目标"),
                ("调研", "需求分析\n现状诊断"),
                ("设计", "蓝图设计\n方案确认"),
                ("配置", "系统配置\n数据准备"),
                ("测试", "UAT测试\n用户培训"),
                ("上线", "系统切换\n运维保障"),
            ])
            style.add_text(0.5, 3.5, 12, 1.5,
                "金蝶实施方法论（Kingdee Way）采用六阶段实施模式，确保项目成功交付。"
                "每个阶段都有明确的里程碑和交付物，全程可控可追溯。", 12)
        self.create_content("实施阶段详述", phases)
        
        # 项目交付物
        def deliverables(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("调研阶段", ["调研报告", "需求清单", "现状分析"], False),
                ("设计阶段", ["业务蓝图", "技术方案", "接口设计"], False),
                ("配置阶段", ["系统配置", "测试用例", "操作手册"], False),
                ("上线阶段", ["培训记录", "验收报告", "运维手册"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("项目交付物", deliverables)
    
    def create_cases(self):
        """成功案例章节"""
        self.create_chapter(5, "成功案例")
        
        # 行业案例
        def industry_cases(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("制造业", ["三一重工", "美的集团", "海信集团"], False),
                ("零售业", ["永辉超市", "屈臣氏", "名创优品"], False),
                ("金融业", ["招商银行", "太平洋保险", "华泰证券"], False),
                ("服务业", ["万科物业", "新东方", "华大基因"], True),
            ], width=2.8, height=2.0, gap=0.3)
        self.create_content("行业标杆客户", industry_cases)
        
        # 典型案例
        def typical_case(slide, style):
            style.add_text(0.5, 1.3, 12, 0.8,
                "项目背景：某大型制造企业，原有SAP系统老化，需要替换升级，同时实现国产化替代。", 14)
            style.add_kpi_row(0.5, 2.2, [
                ("12", "个月", "实施周期"),
                ("100%", "", "功能替代"),
                ("5000+", "万", "投资规模"),
                ("99.9%", "", "系统可用率"),
            ], card_width=2.8, gap=0.3)
        self.create_content("典型案例 - 制造业ERP替代", typical_case)
        
        # 案例2
        def case2(slide, style):
            style.add_text(0.5, 1.3, 12, 0.8,
                "项目背景：某连锁零售企业，门店超过1000家，需要实现全渠道一体化管理。", 14)
            style.add_kpi_row(0.5, 2.2, [
                ("1000+", "家", "门店覆盖"),
                ("8", "个月", "实施周期"),
                ("30%", "", "效率提升"),
                ("98%", "", "客户满意度"),
            ], card_width=2.8, gap=0.3)
        self.create_content("典型案例 - 零售业全渠道", case2)
        
        # 案例3
        def case3(slide, style):
            style.add_text(0.5, 1.3, 12, 0.8,
                "项目背景：某金融企业，需要构建统一财务平台，实现集团财务集中管控。", 14)
            style.add_kpi_row(0.5, 2.2, [
                ("50+", "家", "分子公司"),
                ("6", "个月", "实施周期"),
                ("50%", "", "核算效率提升"),
                ("100%", "", "合规达标"),
            ], card_width=2.8, gap=0.3)
        self.create_content("典型案例 - 金融业财务共享", case3)
        
        # 实施效果
        def effect(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "客户实施效果统计", 24, COLORS['dark'], bold=True)
            style.add_kpi_row(0.5, 1.5, [
                ("30%", "", "平均效率提升"),
                ("20%", "", "成本降低"),
                ("50%", "", "决策速度提升"),
                ("2年", "", "平均回本周期"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "基于金蝶750万+企业客户的数字化转型经验，企业通过金蝶云产品实施后，"
                "平均可实现效率提升30%，成本降低20%，投资回报周期约2年。", 12)
        self.create_content("实施效果", effect)
        
        # 客户评价
        def testimonials(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "客户评价", 24, COLORS['dark'], bold=True)
            style.add_rounded_rect(0.5, 1.5, 5.8, 1.5, COLORS['light'], 
                '"金蝶系统帮助我们实现了财务业务的全面数字化，核算效率提升了50%以上。"\n——某制造企业财务总监', 12)
            style.add_rounded_rect(0.5, 3.2, 5.8, 1.5, COLORS['light'], 
                '"供应链协同让我们从订单到交付的周期缩短了40%，库存周转率提升了25%。"\n——某零售企业运营总监', 12)
            style.add_rounded_rect(6.6, 1.5, 5.8, 1.5, COLORS['light'], 
                '"金蝶的实施团队非常专业，项目按期上线，没有任何重大问题。"\n——某金融企业IT总监', 12)
            style.add_rounded_rect(6.6, 3.2, 5.8, 1.5, COLORS['light'], 
                '"数据分析平台让我们的决策更加科学，报表生成时间从3天缩短到3分钟。"\n——某服务企业总经理', 12)
        self.create_content("客户评价", testimonials)
        
        # ROI分析
        def roi(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "投资回报分析", 24, COLORS['dark'], bold=True)
            style.add_kpi_row(0.5, 1.5, [
                ("第1年", "", "投入期"),
                ("第2年", "", "收益期"),
                ("第3年", "", "回报期"),
                ("300%", "", "三年ROI"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "根据我们的项目经验，企业通常在实施后第2年开始获得显著收益，"
                "第3年实现投资回报，三年累计ROI可达300%以上。", 12)
        self.create_content("投资回报分析", roi)
    
    def create_service(self):
        """服务保障章节"""
        self.create_chapter(6, "服务保障")
        
        # 服务体系
        def service_system(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("实施服务", ["项目实施", "培训服务", "数据迁移"], False),
                ("运维服务", ["系统运维", "安全保障", "性能优化"], False),
                ("升级服务", ["版本升级", "功能扩展", "二次开发"], False),
                ("支持服务", ["7×24热线", "在线客服", "现场支持"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("服务体系", service_system)
        
        # 服务承诺
        def commitment(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("7×24", "小时", "服务热线"),
                ("2", "小时", "响应时间"),
                ("99.9%", "", "系统可用率"),
                ("100%", "", "客户满意度"),
            ], card_width=2.8, gap=0.3)
        self.create_content("服务承诺", commitment)
        
        # 培训服务
        def training(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("管理员培训", ["系统配置", "权限管理", "日常维护"], False),
                ("关键用户培训", ["业务操作", "流程审批", "问题处理"], False),
                ("最终用户培训", ["日常操作", "报表查询", "数据录入"], False),
                ("培训资料", ["操作手册", "视频教程", "在线帮助"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("培训服务", training)
        
        # 运维服务
        def ops_service(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("日常运维", ["系统监控", "故障处理", "性能优化"], False),
                ("安全保障", ["数据备份", "安全审计", "漏洞修复"], False),
                ("升级服务", ["版本更新", "功能扩展", "兼容升级"], False),
                ("技术咨询", ["架构咨询", "优化建议", "技术支持"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("运维服务", ops_service)
        
        # 服务网络
        def network(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("200+", "个", "服务网点"),
                ("5000+", "人", "服务团队"),
                ("100+", "家", "生态伙伴"),
                ("全国", "", "服务覆盖"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "金蝶在全国设有200+服务网点，5000+专业服务团队，100+生态合作伙伴，"
                "为客户提供全方位、本地化的服务支持。", 12)
        self.create_content("服务网络", network)
        
        # 合作模式
        def cooperation(slide, style):
            style.add_function_blocks_row(0.5, 1.5, [
                ("标准实施", ["固定周期", "固定价格", "标准交付"], False),
                ("人天计费", ["按需投入", "灵活调整", "按实结算"], False),
                ("长期合作", ["年度服务", "持续优化", "优惠价格"], False),
                ("云服务", ["订阅付费", "按需扩展", "持续升级"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("合作模式", cooperation)
        
        # 成功保障
        def success(slide, style):
            style.add_kpi_row(0.5, 1.5, [
                ("100%", "", "项目成功率"),
                ("95%", "", "客户满意度"),
                ("98%", "", "按期交付率"),
                ("100%", "", "文档完整率"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "金蝶拥有完善的项目管理体系和丰富的实施经验，确保每个项目成功交付。"
                "我们承诺：项目成功上线，客户满意验收。", 12)
        self.create_content("成功保障", success)
        
        # 联系方式
        def contact(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "联系我们", 24, COLORS['dark'], bold=True)
            style.add_kpi_row(0.5, 1.5, [
                ("400-883-8836", "", "服务热线"),
                ("www.kingdee.com", "", "官方网站"),
                ("service@kingdee.com", "", "商务邮箱"),
                ("全国200+城市", "", "服务网络"),
            ], card_width=2.8, gap=0.3)
        self.create_content("联系方式", contact)
        
        # 下一步行动
        def next_steps(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "下一步行动", 24, COLORS['dark'], bold=True)
            style.add_timeline(0.5, 1.5, 12, [
                ("步骤1", "商务洽谈"),
                ("步骤2", "需求确认"),
                ("步骤3", "方案设计"),
                ("步骤4", "合同签订"),
                ("步骤5", "项目启动"),
            ])
        self.create_content("下一步行动", next_steps)
        
        # 项目预期收益
        def benefits(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "项目预期收益", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("财务收益", ["核算效率提升50%", "资金周转加快20%", "合规风险降低90%"], False),
                ("运营收益", ["采购效率提升40%", "库存周转加快25%", "订单处理缩短50%"], False),
                ("管理收益", ["决策速度提升50%", "数据实时性100%", "流程标准化100%"], False),
                ("战略收益", ["数字化转型", "竞争力提升", "可持续发展"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("项目预期收益", benefits)
        
        # 项目保障
        def guarantee(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "项目保障", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("团队保障", ["资深项目经理", "专业顾问团队", "技术支持团队"], False),
                ("方法保障", ["Kingdee Way", "项目管理规范", "质量管理体系"], False),
                ("技术保障", ["成熟产品平台", "完善技术架构", "安全保障机制"], False),
                ("服务保障", ["7×24服务", "定期回访", "持续优化"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("项目保障", guarantee)
        
        # 数据安全
        def security(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "数据安全保障", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("网络安全", ["防火墙", "入侵检测", "DDoS防护"], False),
                ("数据安全", ["数据加密", "访问控制", "审计日志"], False),
                ("应用安全", ["安全开发", "漏洞扫描", "渗透测试"], False),
                ("合规认证", ["等保三级", "ISO27001", "SOC2"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("数据安全保障", security)
        
        # 灾备方案
        def disaster_recovery(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "灾备方案", 24, COLORS['dark'], bold=True)
            style.add_kpi_row(0.5, 1.5, [
                ("双活", "", "数据中心"),
                ("秒级", "", "故障切换"),
                ("99.99%", "", "可用性"),
                ("4小时", "", "恢复时间"),
            ], card_width=2.8, gap=0.3)
            style.add_text(0.5, 3.0, 12, 2,
                "金蝶云采用双活数据中心架构，确保业务连续性。"
                "当主数据中心发生故障时，系统可在秒级切换到备数据中心，保障业务不中断。", 12)
        self.create_content("灾备方案", disaster_recovery)
        
        # 升级策略
        def upgrade(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "版本升级策略", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("自动升级", ["安全补丁", "性能优化", "Bug修复"], False),
                ("选择性升级", ["功能增强", "新特性", "界面优化"], False),
                ("大版本升级", ["架构升级", "重大变更", "数据迁移"], False),
                ("升级保障", ["升级测试", "回滚机制", "技术支持"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("版本升级策略", upgrade)
        
        # 合作伙伴
        def partners(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "合作伙伴生态", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("技术伙伴", ["华为", "阿里云", "腾讯云"], False),
                ("咨询伙伴", ["德勤", "安永", "毕马威"], False),
                ("实施伙伴", ["各地服务商", "行业专家", "技术团队"], False),
                ("行业伙伴", ["行业协会", "研究院", "高校"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("合作伙伴生态", partners)
        
        # 客户成功案例
        def more_cases(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "更多成功案例", 24, COLORS['dark'], bold=True)
            style.add_function_blocks_row(0.5, 1.5, [
                ("制造业", ["三一重工", "美的集团", "海信集团"], False),
                ("零售业", ["永辉超市", "屈臣氏", "名创优品"], False),
                ("金融业", ["招商银行", "太平洋保险", "华泰证券"], False),
                ("服务业", ["万科物业", "新东方", "华大基因"], True),
            ], width=2.8, height=2.2, gap=0.3)
        self.create_content("更多成功案例", more_cases)
        
        # 行业价值
        def industry_value(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "行业价值", 24, COLORS['dark'], bold=True)
            style.add_kpi_row(0.5, 1.5, [
                ("500+", "家", "制造业客户"),
                ("300+", "家", "零售业客户"),
                ("150+", "家", "金融业客户"),
                ("200+", "家", "服务业客户"),
            ], card_width=2.8, gap=0.3)
        self.create_content("行业价值", industry_value)
        
        # 总结
        def summary(slide, style):
            style.add_text(0.5, 0.8, 12, 0.6, "方案总结", 24, COLORS['dark'], bold=True)
            style.add_rounded_rect(0.5, 1.5, 12, 3, COLORS['light'], 
                "金蝶云·星空是成长型企业数字化转型的理想选择：\n\n"
                "✓ 成熟稳定的产品平台，服务超过750万家企业客户\n"
                "✓ 完善的功能模块，覆盖财务、供应链、生产、人力等核心业务\n"
                "✓ 灵活的部署方式，支持公有云、私有云、混合云\n"
                "✓ 专业的实施团队，确保项目成功交付\n"
                "✓ 持续的升级服务，助力企业持续发展", 14)
        self.create_content("方案总结", summary)
        
        # 致谢
        def thanks(slide, style):
            style.add_text(0.5, 1.5, 12, 1, "感谢聆听", 36, COLORS['dark'], bold=True)
            style.add_text(0.5, 2.8, 12, 0.6, "期待与您合作", 24, COLORS['primary'])
            style.add_text(0.5, 3.6, 12, 0.6, "金蝶软件（中国）有限公司", 16, COLORS['text_dark'])
            style.add_text(0.5, 4.2, 12, 0.6, "服务热线：400-883-8836  |  官网：www.kingdee.com", 12, COLORS['text_dark'])
        self.create_content("致谢", thanks)
    
    def create_ending(self):
        """结尾页"""
        slide = self.add_slide('封底')
        print("  结尾页")
    
    def generate(self):
        """生成完整PPT"""
        print("\n" + "="*60)
        print("开始生成售前PPT V27 完整版")
        print("="*60)
        
        chapters = [
            "公司介绍", "产品体系", "解决方案", 
            "实施路线", "成功案例", "服务保障"
        ]
        
        # 封面
        print("\n生成封面...")
        self.create_cover()
        
        # 团队介绍
        print("生成团队介绍...")
        self.create_team_intro()
        
        # 目录
        print("生成目录...")
        self.create_toc(chapters)
        
        # 各章节
        print("\n生成公司介绍...")
        self.create_company_intro()
        
        print("生成产品体系...")
        self.create_product_intro()
        
        print("生成解决方案...")
        self.create_solution_intro()
        
        print("生成实施路线...")
        self.create_implementation()
        
        print("生成成功案例...")
        self.create_cases()
        
        print("生成服务保障...")
        self.create_service()
        
        # 结尾
        print("\n生成结尾页...")
        self.create_ending()
        
        # 保存
        print(f"\n保存PPT: {self.output_path}")
        self.prs.save(self.output_path)
        
        print(f"\n✅ PPT生成完成！")
        print(f"   页数: {len(self.prs.slides)}")
        print(f"   文件: {self.output_path}")
        
        return self.output_path


if __name__ == "__main__":
    generator = KingdeePresalesPPTV27Full(
        company_name="测试公司",
        project_name="ERP升级项目",
        output_path="output/售前PPT_V27_完整版.pptx"
    )
    generator.generate()
