#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 V2
基于中煤科工述标文件模板样式，采用4A架构思路
"""

import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from copy import deepcopy
import os

# 模板路径
TEMPLATE_PATH = "/mnt/d/Kingdee文档/自动化交付工具/售前文件/中煤科工ERP重庆研究院ERP升级项目述标文件V7.0.pptx"

# ==================== 行业案例库 ====================
INDUSTRY_CASES = {
    "制造业": [
        {
            "name": "潍柴雷沃",
            "desc": "山东重工集团旗下，农业装备行业领军企业",
            "scale": "200亿营收，1.1万员工",
            "modules": ["财务云", "供应链云", "制造云", "成本管理"],
            "results": ["MRP运算从6-8小时缩短到1小时", "668项功能完全满足81%", "18套系统打通，258个接口"],
            "pain_points": ["SAP系统老旧", "多组织协同难", "成本核算复杂"]
        },
        {
            "name": "中车集团",
            "desc": "全球领先的轨道交通装备制造商",
            "scale": "超2000亿营收，17万员工",
            "modules": ["集团财务", "全球司库", "集中采购", "供应链协同"],
            "results": ["财务共享中心上线", "资金集中管理", "采购成本降低15%"],
            "pain_points": ["集团管控弱", "资金分散", "采购不透明"]
        },
        {
            "name": "招商局集团",
            "desc": "中央企业，综合性企业集团",
            "scale": "万亿资产，世界500强",
            "modules": ["价值创造型财务", "司库管理", "合并报表"],
            "results": ["运营财务→业务财务→战略财务转型", "资金使用效率提升30%"],
            "pain_points": ["财务转型需求", "资金管理复杂"]
        }
    ],
    "军工": [
        {
            "name": "某军工研究所",
            "desc": "国防尖端科技研究院",
            "scale": "涉密项目",
            "modules": ["项目管理", "生产制造", "供应链", "财务", "资产"],
            "results": ["全栈信创通过分保测评（机密级）", "飞腾芯片+银河麒麟OS+达梦数据库"],
            "pain_points": ["信创要求", "机密级安全", "项目管理复杂"]
        },
        {
            "name": "某项目型军工企业",
            "desc": "国防装备制造企业",
            "scale": "大型军工",
            "modules": ["项目制造", "质量", "成本", "财务"],
            "results": ["项目全生命周期管理", "成本精细化核算"],
            "pain_points": ["项目型生产", "质量追溯难"]
        }
    ],
    "煤炭": [
        {
            "name": "中煤科工重庆研究院",
            "desc": "煤矿安全技术国家工程研究中心",
            "scale": "高新技术企业",
            "modules": ["财务云", "供应链云", "制造云", "项目管理", "人力资源管理"],
            "results": ["信创替代SAP", "多组织多业态管控", "业财一体化"],
            "pain_points": ["SAP系统老旧", "主数据分散", "生产计划不统一"]
        }
    ]
}

# ==================== 4A架构知识库 ====================
ARCHITECTURE_KNOWLEDGE = {
    "BA": {
        "name": "业务架构",
        "desc": "描述企业业务运作方式的架构",
        "components": ["价值流", "业务能力", "业务流程", "组织架构", "业务对象"],
        "value_stream": ["从订单到收款", "从采购到付款", "从计划到生产", "从需求到交付"],
        "capabilities": {
            "财务域": ["总账管理", "应收管理", "应付管理", "资产管理", "成本管理", "预算管理", "资金管理", "合并报表"],
            "供应链域": ["采购管理", "库存管理", "销售管理", "物流管理", "供应商管理", "客户管理"],
            "制造域": ["生产计划", "生产执行", "质量管理", "设备管理", "工艺管理"],
            "人力域": ["组织管理", "人事管理", "薪酬管理", "绩效管理", "招聘管理", "培训管理"]
        }
    },
    "DA": {
        "name": "数据架构",
        "desc": "描述企业数据资产管理的架构",
        "components": ["主数据", "业务数据", "分析数据", "数据治理", "数据服务"],
        "master_data": ["客户主数据", "供应商主数据", "物料主数据", "会计科目", "组织主数据", "人员主数据"],
        "value": ["统一数据标准", "消除数据孤岛", "提升数据质量", "支撑数据决策"]
    },
    "AA": {
        "name": "应用架构",
        "desc": "描述企业应用系统及相互关系的架构",
        "components": ["核心ERP", "业务系统", "集成平台", "移动应用", "数据分析"],
        "systems": {
            "核心系统": ["金蝶云·星瀚", "财务管理", "供应链管理", "制造管理"],
            "周边系统": ["OA办公", "MES生产", "WMS仓储", "PLM研发", "CRM客户"],
            "集成方案": ["API网关", "ESB总线", "数据中台", "消息队列"]
        }
    },
    "TA": {
        "name": "技术架构",
        "desc": "描述企业技术基础设施的架构",
        "components": ["云平台", "数据库", "中间件", "安全技术", "运维监控"],
        "features": ["云原生架构", "微服务架构", "分布式部署", "容器化部署", "高可用架构"],
        "security": ["等保三级", "分保测评", "信创适配", "数据加密"]
    }
}

# ==================== 价值工程指标 ====================
VALUE_INDICATORS = {
    "财务价值": [
        {"name": "财务结账效率", "before": "15天", "after": "3天", "improve": "提升80%"},
        {"name": "资金周转效率", "before": "45天", "after": "30天", "improve": "提升33%"},
        {"name": "报表准确率", "before": "95%", "after": "99%", "improve": "提升4%"},
        {"name": "应收账款周转", "before": "60天", "after": "45天", "improve": "缩短15天"}
    ],
    "供应链价值": [
        {"name": "库存周转率", "before": "6次/年", "after": "9次/年", "improve": "提升50%"},
        {"name": "采购周期", "before": "15天", "after": "7天", "improve": "缩短53%"},
        {"name": "订单交付率", "before": "85%", "after": "95%", "improve": "提升10%"},
        {"name": "库存准确率", "before": "85%", "after": "98%", "improve": "提升13%"}
    ],
    "生产价值": [
        {"name": "生产计划准确性", "before": "70%", "after": "90%", "improve": "提升20%"},
        {"name": "生产效率", "before": "基准", "after": "提升", "improve": "提升30%"},
        {"name": "质量合格率", "before": "95%", "after": "99%", "improve": "提升4%"},
        {"name": "MRP运算时间", "before": "6-8小时", "after": "1小时", "improve": "缩短87%"}
    ],
    "管理价值": [
        {"name": "审批效率", "before": "3天", "after": "0.5天", "improve": "提升83%"},
        {"name": "数据一致性", "before": "80%", "after": "99%", "improve": "提升19%"},
        {"name": "决策响应时间", "before": "7天", "after": "1天", "improve": "缩短86%"}
    ]
}


def create_presentation_from_template():
    """从模板创建演示文稿"""
    if os.path.exists(TEMPLATE_PATH):
        try:
            prs = Presentation(TEMPLATE_PATH)
            print(f"成功加载模板: {TEMPLATE_PATH}")
            print(f"模板幻灯片数: {len(prs.slides)}")
            return prs
        except Exception as e:
            print(f"加载模板失败: {e}")
    
    # 如果模板不存在，创建空白演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def copy_slide(prs, slide_index):
    """复制模板中的幻灯片"""
    if slide_index < len(prs.slides):
        source_slide = prs.slides[slide_index]
        # 创建新幻灯片
        blank_layout = prs.slide_layouts[6]  # 空白布局
        new_slide = prs.slides.add_slide(blank_layout)
        
        # 复制形状（简化版本，python-pptx不支持完整复制）
        for shape in source_slide.shapes:
            if hasattr(shape, 'text_frame'):
                # 复制文本框
                try:
                    txBox = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    if hasattr(shape, 'text_frame'):
                        for para in shape.text_frame.paragraphs:
                            p = txBox.text_frame.add_paragraph()
                            p.text = para.text
                            p.font.size = para.font.size
                            p.font.bold = para.font.bold
                except:
                    pass
        return new_slide
    return None


def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.8))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RgbColor(102, 102, 102)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title, number=""):
    """添加章节标题幻灯片"""
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(13.33), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0, 51, 102)
    shape.line.fill.background()
    
    # 编号
    if number:
        num_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(1), Inches(1))
        num_frame = num_shape.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = number
        num_para.font.size = Pt(60)
        num_para.font.bold = True
        num_para.font.color.rgb = RgbColor(255, 255, 255)
    
    # 标题
    title_shape = slide.shapes.add_textbox(Inches(2.5), Inches(3.2), Inches(10), Inches(1.2))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)
    
    return slide


def add_content_slide(prs, title, content_list, subtitle=""):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RgbColor(0, 51, 102)
    title_bg.line.fill.background()
    
    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)
    
    # 副标题
    if subtitle:
        sub_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
        sub_frame = sub_shape.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(16)
        sub_para.font.color.rgb = RgbColor(102, 102, 102)
        start_y = 2.0
    else:
        start_y = 1.5
    
    # 内容
    content_shape = slide.shapes.add_textbox(Inches(0.8), Inches(start_y), Inches(11.5), Inches(5))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        para.text = f"• {item}"
        para.font.size = Pt(18)
        para.font.color.rgb = RgbColor(51, 51, 51)
        para.space_after = Pt(12)
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏幻灯片"""
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RgbColor(0, 51, 102)
    title_bg.line.fill.background()
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)
    
    # 左栏
    left_title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    left_title_frame = left_title_shape.text_frame
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.text = left_title
    left_title_para.font.size = Pt(20)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = RgbColor(0, 51, 102)
    
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(14)
        para.space_after = Pt(8)
    
    # 右栏
    right_title_shape = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.8), Inches(0.5))
    right_title_frame = right_title_shape.text_frame
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.text = right_title
    right_title_para.font.size = Pt(20)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = RgbColor(0, 51, 102)
    
    right_content = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.8), Inches(4.5))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(14)
        para.space_after = Pt(8)
    
    return slide


def add_table_slide(prs, title, headers, rows):
    """添加表格幻灯片"""
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RgbColor(0, 51, 102)
    title_bg.line.fill.background()
    
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)
    
    # 表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5)).table
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(0, 51, 102)
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RgbColor(255, 255, 255)
        para.font.bold = True
        para.font.size = Pt(12)
    
    # 数据行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(240, 240, 240)
    
    return slide


def generate_presales_ppt_v2(data):
    """生成售前PPT V2 - 基于4A架构思路"""
    prs = create_presentation_from_template()
    
    # 如果使用模板，先清除所有幻灯片（保留母版）
    if os.path.exists(TEMPLATE_PATH):
        # 使用模板的布局创建新幻灯片
        slide_count = len(prs.slides)
        print(f"使用模板，原有幻灯片数: {slide_count}")
        # 我们将添加新的幻灯片到末尾
    
    # 转换数据类型
    months = int(data.get('implMonths', 6) or 6)
    employees = data.get('employees', '') or ''
    revenue = data.get('revenue', '') or ''
    budget = data.get('budget', '') or ''
    user_count = data.get('userCount', '') or ''
    company_name = data.get('companyName', '客户企业')
    industry = data.get('industry', '制造业')
    
    # 获取行业案例
    industry_cases = INDUSTRY_CASES.get(industry, INDUSTRY_CASES.get("制造业", []))
    
    # ==================== 第1部分：封面 ====================
    add_title_slide(prs, 
        f"{company_name}ERP项目解决方案",
        "金蝶软件（中国）有限公司")
    
    # ==================== 第2部分：目录 ====================
    add_content_slide(prs, "目录", [
        "01 项目理解与需求分析",
        "02 解决方案设计（4A架构）",
        "03 业务架构（BA）- 业务蓝图",
        "04 数据架构（DA）- 数据治理",
        "05 应用架构（AA）- 系统集成",
        "06 技术架构（TA）- 平台支撑",
        "07 价值工程与预期收益",
        "08 成功案例",
        "09 项目实施安排",
        "10 服务保障承诺"
    ])
    
    # ==================== 第3部分：项目理解 ====================
    add_section_slide(prs, "项目理解与需求分析", "01")
    
    # 企业概况
    add_content_slide(prs, "企业概况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{data.get('companySize', '待确认')}",
        f"员工人数：{employees}人" if employees else "员工人数：待确认",
        f"年营业额：{revenue}万元" if revenue else "年营业额：待确认",
        f"用户规模：{user_count}用户" if user_count else "用户规模：待确认",
        f"管控模式：{data.get('controlMode', '待确认')}"
    ], "客户基本信息")
    
    # 行业洞察
    industry_insights = {
        "制造业": [
            "制造业数字化转型加速，智能制造成为核心方向",
            "多组织、多工厂协同管理需求迫切",
            "成本精细化核算要求越来越高",
            "供应链协同能力成为竞争关键"
        ],
        "煤炭": [
            "煤炭行业机械化程度已达99%以上",
            "从单一产品向整体解决方案转变",
            "安全生产、环保监管要求提升",
            "技术创新和管理创新双轮驱动"
        ],
        "军工": [
            "信创替代成为必选项",
            "机密级信息系统要求",
            "项目型生产管理复杂",
            "质量追溯要求严格"
        ]
    }
    add_content_slide(prs, "行业洞察", 
        industry_insights.get(industry, industry_insights["制造业"]),
        "行业发展趋势分析")
    
    # 痛点分析
    pain_points = data.get('painPoints', '')
    if pain_points:
        pain_list = [p.strip() for p in pain_points.split('\n') if p.strip()][:6]
    else:
        pain_list = [
            "主数据管理分散，缺乏统一的主数据管理平台",
            "各生产单位计划不统一，生产组织协同难",
            "采购模式多样化，供应商交付协同难",
            "未建立标准成本体系，业务财务未一体化",
            "业务价值链跨越多系统，数据有断点"
        ]
    add_content_slide(prs, "核心痛点分析", pain_list, "企业面临的核心问题")
    
    # 建设目标
    business_goals = data.get('businessGoals', '')
    if business_goals:
        goal_list = [g.strip() for g in business_goals.split('\n') if g.strip()][:6]
    else:
        goal_list = [
            "建立统一的主数据管理体系",
            "实现业财一体化管理",
            "提升生产协同效率",
            "强化供应链管控能力",
            "构建数据决策分析平台"
        ]
    add_content_slide(prs, "建设目标", goal_list, "项目预期达成的目标")
    
    # ==================== 第4部分：解决方案设计（4A架构） ====================
    add_section_slide(prs, "解决方案设计", "02")
    
    add_content_slide(prs, "总体设计思路", [
        "基于4A企业架构方法论，构建企业数字化能力",
        "业务架构（BA）：梳理价值流、业务能力、业务流程",
        "数据架构（DA）：建立主数据体系、数据治理平台",
        "应用架构（AA）：设计系统集成方案、应用部署架构",
        "技术架构（TA）：搭建云原生平台、安全保障体系",
        "价值工程：量化预期收益，支撑投资决策"
    ], "基于4A架构的顶层设计")
    
    # 实施范围
    modules = data.get('modules', ['finance']) or ['finance']
    if isinstance(modules, str):
        modules = [m.strip() for m in modules.split(',') if m.strip()]
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 
                    'hr': '人力云', 'plm': 'PLM云', 'pm': '项目管理'}
    module_list = [module_names.get(m, m) for m in modules]
    
    add_two_column_slide(prs, "实施范围",
        "实施模块", module_list + [
            f"用户规模：{user_count}用户" if user_count else "用户规模：待确认",
            f"实施周期：{months}个月",
            f"项目预算：{budget}万元" if budget else "项目预算：待确认"
        ],
        "实施阶段", [
            f"第一阶段（M1-M{months//3}）：核心模块上线",
            f"第二阶段（M{months//3+1}-M{months*2//3}）：扩展模块推广",
            f"第三阶段（M{months*2//3+1}-M{months}）：优化完善",
            "验收阶段：系统验收与知识转移"
        ])
    
    # ==================== 第5部分：业务架构（BA） ====================
    add_section_slide(prs, "业务架构（BA）", "03")
    
    # 业务架构蓝图
    ba = ARCHITECTURE_KNOWLEDGE['BA']
    add_content_slide(prs, "业务架构设计", [
        f"价值流：{', '.join(ba['value_stream'][:4])}",
        "业务能力地图：L1-L3三级能力体系",
        "业务流程：端到端流程设计",
        "组织架构：适配业务的组织设计",
        "业务对象：统一业务语言"
    ], ba['desc'])
    
    # 业务能力地图
    for domain, capabilities in list(ba['capabilities'].items())[:3]:
        add_content_slide(prs, f"{domain}能力地图", capabilities, f"{domain}业务能力体系")
    
    # ==================== 第6部分：数据架构（DA） ====================
    add_section_slide(prs, "数据架构（DA）", "04")
    
    da = ARCHITECTURE_KNOWLEDGE['DA']
    add_content_slide(prs, "数据架构设计", [
        f"主数据管理：{', '.join(da['master_data'])}",
        "数据治理：数据标准、数据质量、数据安全",
        "数据服务：数据采集、数据清洗、数据分析",
        "数据分析：经营分析、决策支持、预测预警"
    ], da['desc'])
    
    add_two_column_slide(prs, "主数据管理体系",
        "主数据范围", da['master_data'],
        "数据价值", da['value'])
    
    # ==================== 第7部分：应用架构（AA） ====================
    add_section_slide(prs, "应用架构（AA）", "05")
    
    aa = ARCHITECTURE_KNOWLEDGE['AA']
    add_content_slide(prs, "应用架构设计", [
        "核心ERP：金蝶云·星瀚",
        "业务系统：OA、MES、WMS、PLM、CRM等",
        "集成平台：API网关、ESB总线、数据中台",
        "移动应用：移动审批、移动报表、移动作业",
        "数据分析：BI平台、数据仓库、管理驾驶舱"
    ], aa['desc'])
    
    add_two_column_slide(prs, "系统集成方案",
        "核心系统", list(aa['systems']['核心系统']) + list(aa['systems']['周边系统'])[:3],
        "集成方式", aa['systems']['集成方案'])
    
    # ==================== 第8部分：技术架构（TA） ====================
    add_section_slide(prs, "技术架构（TA）", "06")
    
    ta = ARCHITECTURE_KNOWLEDGE['TA']
    add_content_slide(prs, "技术架构设计", [
        "云原生架构：微服务、容器化、分布式",
        "高可用架构：双活部署、容灾备份",
        "安全技术：等保三级、数据加密、访问控制",
        "运维监控：全链路监控、智能运维",
        "信创适配：国产数据库、国产操作系统、国产芯片"
    ], ta['desc'])
    
    add_two_column_slide(prs, "技术特性与安全",
        "技术特性", ta['features'],
        "安全保障", ta['security'])
    
    # ==================== 第9部分：价值工程 ====================
    add_section_slide(prs, "价值工程与预期收益", "07")
    
    # 价值指标表格
    value_rows = []
    for category, indicators in list(VALUE_INDICATORS.items())[:2]:
        for ind in indicators[:2]:
            value_rows.append([category, ind['name'], ind['before'], ind['after'], ind['improve']])
    
    add_table_slide(prs, "预期价值指标", 
        ["价值领域", "指标名称", "实施前", "实施后", "提升幅度"],
        value_rows)
    
    # ==================== 第10部分：成功案例 ====================
    add_section_slide(prs, "成功案例", "08")
    
    for case in industry_cases[:2]:
        add_two_column_slide(prs, f"案例：{case['name']}",
            "项目背景", [
                case['desc'],
                f"规模：{case['scale']}",
                f"实施模块：{', '.join(case['modules'])}",
                "痛点：" + "、".join(case['pain_points'])
            ],
            "项目成果", case['results'])
    
    # ==================== 第11部分：实施安排 ====================
    add_section_slide(prs, "项目实施安排", "09")
    
    m1 = max(1, months // 4)
    m2 = max(m1 + 1, months // 2)
    m3 = max(m2 + 1, int(months * 0.8))
    
    add_content_slide(prs, "实施计划", [
        f"项目启动（第1周）：组建团队、明确计划、召开启动会",
        f"需求调研（第2-{m1}周）：业务调研、需求梳理、蓝图设计",
        f"系统配置（第{m1+1}-{m2}周）：系统配置、数据准备、单元测试",
        f"系统测试（第{m2+1}-{m3}周）：集成测试、UAT测试、用户培训",
        f"上线切换（第{m3+1}-{months}周）：数据迁移、上线切换、支持保障",
        f"项目验收：验收测试、项目总结、知识转移"
    ], f"实施周期：{months}个月")
    
    # ==================== 第12部分：服务保障 ====================
    add_section_slide(prs, "服务保障承诺", "10")
    
    add_content_slide(prs, "服务保障", [
        "项目团队：资深项目经理+专业顾问团队",
        "实施方法论：Kingdee Way标准化实施方法论",
        "质量保障：里程碑验收、阶段性评审、问题跟踪",
        "知识转移：培训计划、操作手册、视频教程",
        "售后服务：免费服务期、专属服务团队、绿色通道",
        "持续优化：定期回访、版本升级、功能优化"
    ], "全方位服务保障")
    
    # ==================== 结束页 ====================
    add_title_slide(prs, "感谢聆听", f"{company_name}ERP项目解决方案\n\n金蝶软件（中国）有限公司")
    
    return prs


def main():
    """主函数"""
    if len(sys.argv) > 1:
        try:
            data = json.loads(sys.argv[1])
        except:
            data = {
                "companyName": "测试企业",
                "customerCode": "TEST",
                "industry": "制造业"
            }
    else:
        data = {
            "companyName": "测试企业",
            "customerCode": "TEST",
            "industry": "制造业",
            "companySize": "中型企业",
            "employees": "1000",
            "revenue": "50000",
            "budget": "500",
            "userCount": "500",
            "implMonths": "8",
            "controlMode": "集团管控",
            "modules": ["finance", "supply", "manufacturing"],
            "painPoints": "1、主数据管理分散\n2、生产计划不统一\n3、成本核算不精细",
            "businessGoals": "1、建立统一主数据平台\n2、实现业财一体化\n3、提升运营效率"
        }
    
    # 生成PPT
    prs = generate_presales_ppt_v2(data)
    
    # 保存文件
    output_dir = os.path.expanduser("~/.openclaw/workspace/output")
    os.makedirs(output_dir, exist_ok=True)
    
    customer_code = data.get('customerCode', data.get('companyName', '客户')[:4])
    output_path = os.path.join(output_dir, f"{customer_code}_售前