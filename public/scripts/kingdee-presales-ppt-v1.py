#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 v1.0 - 专业版
支持：标准版(40-50页)、详细版(80-100页)、专业版(120-150页)
基于参考PPT结构，包含素材嵌入功能
"""

import sys
import os
import json
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# 颜色主题
THEMES = {
    'blue': {
        'primary': RGBColor(26, 35, 126),      # 金蝶蓝
        'secondary': RGBColor(13, 71, 161),    # 深蓝
        'accent': RGBColor(33, 150, 243),      # 亮蓝
        'text': RGBColor(33, 33, 33),          # 深灰
        'light': RGBColor(227, 242, 253),      # 浅蓝背景
        'white': RGBColor(255, 255, 255)
    },
    'dark': {
        'primary': RGBColor(33, 33, 33),
        'secondary': RGBColor(66, 66, 66),
        'accent': RGBColor(255, 152, 0),
        'text': RGBColor(255, 255, 255),
        'light': RGBColor(48, 48, 48),
        'white': RGBColor(250, 250, 250)
    },
    'green': {
        'primary': RGBColor(27, 94, 32),
        'secondary': RGBColor(56, 142, 60),
        'accent': RGBColor(76, 175, 80),
        'text': RGBColor(33, 33, 33),
        'light': RGBColor(232, 245, 233),
        'white': RGBColor(255, 255, 255)
    },
    'red': {
        'primary': RGBColor(183, 28, 28),
        'secondary': RGBColor(211, 47, 47),
        'accent': RGBColor(244, 67, 54),
        'text': RGBColor(33, 33, 33),
        'light': RGBColor(255, 235, 238),
        'white': RGBColor(255, 255, 255)
    }
}

class PresalesPPTGenerator:
    def __init__(self, company_name, customer_code='default', industry='制造业',
                 mode='detailed', theme='blue', assets_dir=None):
        self.company_name = company_name
        self.customer_code = customer_code
        self.industry = industry
        self.mode = mode  # standard, detailed, professional
        self.theme = THEMES.get(theme, THEMES['blue'])
        self.assets_dir = assets_dir
        
        # 创建演示文稿
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 素材库
        self.assets = {
            'company-logo': [],
            'company-photos': [],
            'team-photos': [],
            'finance-screenshots': [],
            'supply-screenshots': [],
            'manufacture-screenshots': [],
            'hr-screenshots': [],
            'project-screenshots': [],
            'mobile-screenshots': [],
            'case-logos': [],
            'case-photos': [],
            'architecture-images': [],
            'chart-images': [],
            'other-images': []
        }
        
        # 加载素材
        if assets_dir:
            self.load_assets()
        
        # 页码计数
        self.slide_count = 0
        
    def load_assets(self):
        """加载素材文件"""
        if not os.path.exists(self.assets_dir):
            return
            
        for category in self.assets.keys():
            category_dir = os.path.join(self.assets_dir, category)
            if os.path.exists(category_dir):
                for filename in os.listdir(category_dir):
                    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                        self.assets[category].append(os.path.join(category_dir, filename))
    
    def add_slide(self, layout_index=6):
        """添加幻灯片"""
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1
        return slide
    
    def add_title_shape(self, slide, text, left=0.5, top=0.3, width=12.333, height=0.8,
                       font_size=32, bold=True, color=None):
        """添加标题形状"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme['primary']
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.theme['white']
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.15)
        
        return shape
    
    def add_subtitle(self, slide, text, left=0.5, top=1.2, width=12.333, height=0.5,
                    font_size=18, color=None):
        """添加副标题"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or self.theme['text']
        p.alignment = PP_ALIGN.LEFT
        
        return shape
    
    def add_content_box(self, slide, text, left, top, width, height, font_size=14,
                       bg_color=None, text_color=None):
        """添加内容框"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color or self.theme['light']
        shape.line.color.rgb = self.theme['primary']
        shape.line.width = Pt(1)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color or self.theme['text']
        
        return shape
    
    def add_image(self, slide, image_path, left, top, width=None, height=None):
        """添加图片"""
        if not os.path.exists(image_path):
            return None
        
        try:
            if width and height:
                return slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
            elif width:
                return slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top),
                    width=Inches(width)
                )
            else:
                return slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top)
                )
        except:
            return None
    
    def add_bullet_list(self, slide, items, left, top, width, height, font_size=14):
        """添加项目列表"""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.theme['text']
            p.space_after = Pt(6)
        
        return shape
    
    def add_navigation_slide(self, title, sections):
        """添加章节导航页"""
        slide = self.add_slide()
        
        # 背景
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(13.333), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['light']
        background.line.fill.background()
        
        # 左侧色块
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(0.3), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = self.theme['primary']
        left_bar.line.fill.background()
        
        # 标题
        self.add_title_shape(slide, title, top=2.5, height=1)
        
        # 章节列表
        y = 3.8
        for section in sections:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(4), Inches(y), Inches(5), Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme['secondary']
            shape.line.fill.background()
            
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = section
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            
            y += 0.7
        
        return slide
    
    # ========== 页面生成方法 ==========
    
    def generate_cover(self):
        """生成封面"""
        slide = self.add_slide()
        
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme['primary']
        bg.line.fill.background()
        
        # 装饰线
        for i in range(5):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(1 + i * 1.5),
                Inches(0.1), Inches(0.5)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.theme['accent']
            line.line.fill.background()
        
        # 标题
        title = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11), Inches(1.5)
        )
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.company_name}金蝶云解决方案"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.theme['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11), Inches(0.8)
        )
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = "金蝶软件（中国）有限公司"
        p.font.size = Pt(24)
        p.font.color.rgb = self.theme['light']
        p.alignment = PP_ALIGN.CENTER
        
        # 日期
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6), Inches(11), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月")
        p.font.size = Pt(18)
        p.font.color.rgb = self.theme['light']
        p.alignment = PP_ALIGN.CENTER
        
        # Logo（如果有）
        if self.assets['company-logo']:
            self.add_image(slide, self.assets['company-logo'][0],
                          11, 0.3, height=0.8)
    
    def generate_toc(self):
        """生成目录"""
        slide = self.add_slide()
        self.add_title_shape(slide, "目  录")
        
        chapters = [
            ("01", "公司介绍", "金蝶简介、发展历程、市场地位"),
            ("02", "行业洞察", "行业分析、痛点提炼、转型趋势"),
            ("03", "需求理解", "客户背景、组织架构、业务痛点"),
            ("04", "解决方案", "总体架构、功能模块、价值工程"),
            ("05", "功能详解", "各模块详细功能展示"),
            ("06", "系统集成", "集成架构、接口方案"),
            ("07", "实施计划", "实施方法、甘特图、里程碑"),
            ("08", "团队保障", "项目团队、培训计划、服务承诺"),
            ("09", "成功案例", "标杆客户、实施成果")
        ]
        
        y = 1.5
        for num, title, desc in chapters:
            # 编号
            num_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1), Inches(y), Inches(0.6), Inches(0.6)
            )
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = self.theme['primary']
            num_box.line.fill.background()
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = num
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.12)
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y), Inches(3), Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.theme['text']
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y + 0.35), Inches(10), Inches(0.3)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(128, 128, 128)
            
            y += 0.7
    
    def generate_company_intro(self):
        """生成公司介绍章节 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("公司介绍", [
            "金蝶简介",
            "发展历程",
            "市场地位",
            "产品体系",
            "荣誉资质"
        ])
        
        # 1. 团队介绍页
        slide = self.add_slide()
        self.add_title_shape(slide, "本次项目述标团队")
        
        team_roles = [
            ("项目总监", "统筹项目管理、风险控制、资源协调"),
            ("项目经理", "负责项目实施、进度管理、客户沟通"),
            ("业务顾问", "业务调研、方案设计、蓝图确认"),
            ("技术顾问", "系统配置、开发定制、数据迁移"),
            ("实施顾问", "用户培训、系统指导、问题处理")
        ]
        
        y = 1.8
        for i, (role, resp) in enumerate(team_roles):
            # 角色框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light'] if i % 2 == 0 else self.theme['white']
            box.line.color.rgb = self.theme['primary']
            
            # 角色名
            role_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y + 0.15), Inches(2), Inches(0.6)
            )
            role_box.fill.solid()
            role_box.fill.fore_color.rgb = self.theme['primary']
            role_box.line.fill.background()
            tf = role_box.text_frame
            p = tf.paragraphs[0]
            p.text = role
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            # 职责
            resp_box = slide.shapes.add_textbox(
                Inches(3), Inches(y + 0.25), Inches(9.5), Inches(0.5)
            )
            tf = resp_box.text_frame
            p = tf.paragraphs[0]
            p.text = resp
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 1
        
        # 2. 金蝶简介
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶简介")
        
        content = f"""金蝶国际软件集团有限公司（简称"金蝶"）是中国领先的企业管理软件及云服务提供商。

成立于1993年，总部位于深圳，在香港联交所上市（股票代码：0268.HK）

服务超过740万家企业及政府组织，云服务市场占有率连续多年位居前列

使命：用创想与技术推动企业管理进步
愿景：成为最值得托付的企业服务平台"""
        
        self.add_content_box(slide, content, 0.5, 1.5, 12.333, 5, font_size=16)
        
        # 3. 发展历程
        slide = self.add_slide()
        self.add_title_shape(slide, "发展历程")
        
        milestones = [
            ("1993", "成立，专注财务软件"),
            ("1996", "推出国内首款Windows财务软件"),
            ("2001", "香港联交所上市"),
            ("2011", "云管理战略转型"),
            ("2018", "发布金蝶云·苍穹"),
            ("2021", "云收入占比超70%"),
            ("至今", "服务740万+企业客户")
        ]
        
        x = 0.8
        for year, event in milestones:
            # 年份圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(2.5), Inches(1.2), Inches(1.2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme['primary']
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = year
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.35)
            
            # 事件
            event_box = slide.shapes.add_textbox(
                Inches(x - 0.3), Inches(4), Inches(1.8), Inches(2)
            )
            tf = event_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = event
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 1.8
        
        # 4. 市场地位
        slide = self.add_slide()
        self.add_title_shape(slide, "市场地位")
        
        stats = [
            ("740万+", "服务企业数量"),
            ("30年+", "行业深耕经验"),
            ("第一", "云ERP市场占有率"),
            ("上市", "港股上市企业")
        ]
        
        x = 0.8
        for value, label in stats:
            # 数值框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(3.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 数值
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.5), Inches(2.8), Inches(1)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(3.8), Inches(2.8), Inches(1)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 5. 产品体系
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶云·苍穹&星瀚 产品体系")
        
        products = [
            ("金蝶云·苍穹", "PaaS平台", "云原生、低代码、微服务"),
            ("金蝶云·星瀚", "大型企业SaaS", "财务云、供应链云、制造云"),
            ("金蝶云·星空", "成长型企业", "财务、供应链、生产"),
            ("金蝶精斗云", "小微企业", "云会计、云进销存")
        ]
        
        y = 1.8
        for name, category, desc in products:
            # 产品框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(1.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            
            # 产品名
            name_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y + 0.2), Inches(2.5), Inches(0.8)
            )
            name_box.fill.solid()
            name_box.fill.fore_color.rgb = self.theme['primary']
            name_box.line.fill.background()
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 分类
            cat_box = slide.shapes.add_textbox(
                Inches(3.5), Inches(y + 0.35), Inches(2.5), Inches(0.5)
            )
            tf = cat_box.text_frame
            p = tf.paragraphs[0]
            p.text = category
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['accent']
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(6.5), Inches(y + 0.35), Inches(6), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 1.35
        
        # 6. 荣誉资质
        slide = self.add_slide()
        self.add_title_shape(slide, "荣誉资质")
        
        honors = [
            "国家规划布局内重点软件企业",
            "中国软件行业十大领军企业",
            "亚太地区ERP软件市场份额第一",
            "中国企业云服务市场占有率第一",
            "国家信息安全认证",
            "ISO9001质量管理体系认证",
            "CMMI5级认证",
            "可信云服务认证"
        ]
        
        x = 0.5
        y = 1.8
        for i, honor in enumerate(honors):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(5.8), Inches(0.8)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['accent']
            
            # 图标
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.2), Inches(y + 0.15), Inches(0.5), Inches(0.5)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = self.theme['accent']
            icon.line.fill.background()
            
            # 文字
            tf = box.text_frame
            tf.margin_left = Inches(0.9)
            tf.margin_top = Inches(0.15)
            p = tf.paragraphs[0]
            p.text = honor
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            if x > 6:
                x = 0.5
                y += 1
            else:
                x += 6.3
    
    def _add_company_intro_slides(self):
        """添加公司介绍幻灯片"""
        # 金蝶简介
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶简介")
        
        content = """金蝶国际软件集团有限公司（简称"金蝶"）是中国领先的企业管理软件及云服务提供商。

成立于1993年，总部位于深圳，在香港联交所上市（股票代码：0268.HK）

服务超过740万家企业及政府组织，云服务市场占有率连续多年位居前列

使命：用创想与技术推动企业管理进步
愿景：成为最值得托付的企业服务平台"""
        
        self.add_content_box(slide, content, 0.5, 1.5, 12.333, 5, font_size=16)
        
        # 发展历程
        slide = self.add_slide()
        self.add_title_shape(slide, "发展历程")
        
        milestones = [
            ("1993", "成立，专注财务软件"),
            ("1996", "推出国内首款Windows财务软件"),
            ("2001", "香港联交所上市"),
            ("2011", "云管理战略转型"),
            ("2018", "发布金蝶云·苍穹"),
            ("2021", "云收入占比超70%"),
            ("至今", "服务740万+企业客户")
        ]
        
        x = 0.8
        for year, event in milestones:
            # 年份圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(2.5), Inches(1.2), Inches(1.2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme['primary']
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.text = year
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.35)
            
            # 事件
            event_box = slide.shapes.add_textbox(
                Inches(x - 0.3), Inches(4), Inches(1.8), Inches(2)
            )
            tf = event_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = event
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 1.8
    
    def generate_industry_analysis(self):
        """生成行业分析章节 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("行业洞察", [
            "行业现状",
            "行业趋势",
            "核心痛点",
            "转型方向",
            "标杆案例"
        ])
        
        # 1. 行业现状
        slide = self.add_slide()
        self.add_title_shape(slide, f"{self.industry}行业现状")
        
        industry_content = self.get_industry_content()
        
        self.add_content_box(slide, industry_content['status'],
                           0.5, 1.5, 6, 5, font_size=14)
        
        self.add_content_box(slide, industry_content['challenges'],
                           6.8, 1.5, 6, 5, font_size=14)
        
        # 2. 行业发展趋势
        slide = self.add_slide()
        self.add_title_shape(slide, "行业数字化转型趋势")
        
        trends = [
            ("智能化", "AI赋能生产、管理智能化"),
            ("平台化", "构建企业数字化平台"),
            ("协同化", "供应链上下游协同"),
            ("数据化", "数据驱动业务决策"),
            ("敏捷化", "快速响应市场变化")
        ]
        
        x = 0.5
        for title, desc in trends:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.4), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 图标
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.7), Inches(2.3), Inches(1), Inches(1)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = self.theme['primary']
            icon.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(3.5), Inches(2), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.1), Inches(2.2), Inches(1.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.55
        
        # 3. 核心痛点分析
        slide = self.add_slide()
        self.add_title_shape(slide, "企业数字化核心痛点分析")
        
        pain_points = [
            ("信息孤岛", "系统分散、数据割裂、流程断点"),
            ("效率低下", "手工操作、重复劳动、审批繁琐"),
            ("决策滞后", "数据不及时、分析不全面、预测困难"),
            ("成本难控", "核算粗放、管控缺失、浪费严重"),
            ("协同困难", "部门壁垒、信息滞后、责任不清")
        ]
        
        y = 1.8
        for i, (title, desc) in enumerate(pain_points):
            # 序号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(y), Inches(0.5), Inches(0.5)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y), Inches(2), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(3.5), Inches(y), Inches(9), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme['text']
            
            y += 0.9
        
        # 4. 数字化转型方向
        slide = self.add_slide()
        self.add_title_shape(slide, "数字化转型关键方向")
        
        directions = [
            ("业务数字化", "流程线上化、业务可视化、管控透明化"),
            ("数据资产化", "数据治理、数据分析、数据驱动决策"),
            ("运营智能化", "AI应用、智能预警、自动化处理"),
            ("生态协同化", "供应链协同、客户协同、伙伴协同")
        ]
        
        x = 0.8
        for title, desc in directions:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 5. 标杆案例引用
        slide = self.add_slide()
        self.add_title_shape(slide, f"{self.industry}行业标杆案例")
        
        cases = [
            ("某大型制造企业", "年营收50亿+", "6个月上线"),
            ("某集团企业", "员工3000+", "业财一体"),
            ("某上市企业", "多组织协同", "降本增效20%")
        ]
        
        x = 0.8
        for name, scale, effect in cases:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(3.8), Inches(3.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['accent']
            
            # 企业名
            name_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.3), Inches(3.6), Inches(0.6)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 规模
            scale_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.1), Inches(3.6), Inches(0.5)
            )
            tf = scale_box.text_frame
            p = tf.paragraphs[0]
            p.text = scale
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 效果
            effect_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.8), Inches(3.6), Inches(0.8)
            )
            tf = effect_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"实施效果：{effect}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['accent']
            p.alignment = PP_ALIGN.CENTER
            
            x += 4.1
    
    def get_industry_content(self):
        """获取行业特定内容"""
        industry_data = {
            '制造业': {
                'status': '制造业现状：\n\n• 市场规模持续增长，智能制造成为趋势\n• 劳动力成本上升，自动化需求迫切\n• 供应链复杂，协同难度大\n• 定制化需求增加，柔性生产成关键',
                'challenges': '面临挑战：\n\n• 生产计划不准，影响交付\n• 库存积压严重，资金占用大\n• 成本核算粗放，利润分析难\n• 质量追溯困难，问题定位慢'
            },
            '零售业': {
                'status': '零售业现状：\n\n• 线上线下融合，全渠道成为标配\n• 消费者需求多变，响应速度要求高\n• 供应链复杂，库存管理难度大\n• 会员运营重要，私域流量成关键',
                'challenges': '面临挑战：\n\n• 渠道分散，数据不统一\n• 库存不准，缺货与积压并存\n• 促销管理复杂，效果难评估\n• 会员价值难挖掘'
            }
        }
        
        return industry_data.get(self.industry, industry_data['制造业'])
    
    def generate_solution_overview(self):
        """生成解决方案总览 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("解决方案", [
            "总体架构",
            "4A架构",
            "功能模块",
            "集成方案",
            "价值工程"
        ])
        
        # 1. 总体架构
        slide = self.add_slide()
        self.add_title_shape(slide, "总体解决方案架构")
        
        # 架构图（使用形状绘制）
        layers = [
            ("展现层", ["PC端", "移动端", "大屏端"], 1.5),
            ("应用层", ["财务云", "供应链云", "制造云", "人力云"], 2.7),
            ("平台层", ["开发平台", "数据平台", "集成平台", "AI平台"], 3.9),
            ("基础层", ["云服务器", "云存储", "云网络", "安全体系"], 5.1)
        ]
        
        for name, items, y in layers:
            # 层名称
            layer_name = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(1.5), Inches(0.8)
            )
            layer_name.fill.solid()
            layer_name.fill.fore_color.rgb = self.theme['primary']
            layer_name.line.fill.background()
            tf = layer_name.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.2)
            
            # 层内容
            x = 2.2
            for item in items:
                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y), Inches(2.5), Inches(0.8)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = self.theme['light']
                box.line.color.rgb = self.theme['primary']
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = self.theme['text']
                p.alignment = PP_ALIGN.CENTER
                tf.margin_top = Inches(0.2)
                x += 2.7
        
        # 2. 业务架构BA
        slide = self.add_slide()
        self.add_title_shape(slide, "业务架构（BA）设计")
        
        ba_content = """业务架构是4A架构的核心，定义企业做什么和如何做：

价值流设计：
• 订单到收款（O2C）：从客户订单到收款的全流程
• 采购到付款（P2P）：从采购申请到付款的全流程
• 计划到生产（P2M）：从生产计划到产品完工的全流程
• 研发到上市（R2M）：从产品研发到市场投放的全流程

业务能力地图：
• L1：财务管理、供应链管理、生产制造、人力资源
• L2：总账核算、应收管理、采购管理、库存管理等
• L3：凭证处理、账簿查询、入库管理等具体能力

业务流程架构：
• 流程层级：L1-L4四级流程体系
• 流程图：BPMN标准业务流程图
• 流程串联：端到端业务流程串联"""
        
        self.add_content_box(slide, ba_content, 0.5, 1.5, 12.333, 5, font_size=13)
        
        # 3. 数据架构DA
        slide = self.add_slide()
        self.add_title_shape(slide, "数据架构（DA）设计")
        
        da_content = """数据架构定义企业的数据资产和数据治理体系：

数据实体设计：
• 主数据：客户、供应商、物料、员工、组织等
• 业务数据：订单、合同、凭证、出入库单等
• 分析数据：报表指标、KPI、预算数据等

数据服务体系：
• 数据查询：实时数据查询服务
• 数据推送：消息队列、事件驱动
• 数据分析：BI分析、数据挖掘

主数据管理（MDM）：
• 客户主数据：统一客户视图
• 供应商主数据：供应商全生命周期管理
• 物料主数据：物料标准化管理

数据治理体系：
• 数据标准：数据字典、编码规则
• 数据质量：数据校验、数据清洗
• 数据安全：权限控制、数据脱敏"""
        
        self.add_content_box(slide, da_content, 0.5, 1.5, 12.333, 5, font_size=13)
        
        # 4. 应用架构AA
        slide = self.add_slide()
        self.add_title_shape(slide, "应用架构（AA）设计")
        
        aa_modules = [
            ("财务云", "总账、应收、应付、资金、成本、固定资产、预算"),
            ("供应链云", "采购、销售、库存、供应商、客户"),
            ("制造云", "计划、生产、车间、质量、BOM"),
            ("人力云", "人事、薪酬、绩效、招聘、培训"),
            ("项目云", "立项、执行、核算、分析"),
            ("协同云", "流程、报表、移动、门户")
        ]
        
        y = 1.5
        for name, modules in aa_modules:
            # 模块名
            name_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(1.8), Inches(0.6)
            )
            name_box.fill.solid()
            name_box.fill.fore_color.rgb = self.theme['primary']
            name_box.line.fill.background()
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            # 模块列表
            modules_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(y + 0.1), Inches(10), Inches(0.5)
            )
            tf = modules_box.text_frame
            p = tf.paragraphs[0]
            p.text = modules
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            y += 0.75
        
        # 5. 技术架构TA
        slide = self.add_slide()
        self.add_title_shape(slide, "技术架构（TA）设计")
        
        ta_layers = [
            ("基础设施层", "云服务器、云存储、云网络、负载均衡"),
            ("平台服务层", "容器服务、数据库服务、缓存服务、消息队列"),
            ("应用服务层", "微服务架构、API网关、服务注册、配置中心"),
            ("安全体系", "身份认证、权限控制、数据加密、审计日志")
        ]
        
        y = 1.8
        for name, components in ta_layers:
            # 层名称
            layer_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(2.5), Inches(1)
            )
            layer_box.fill.solid()
            layer_box.fill.fore_color.rgb = self.theme['primary']
            layer_box.line.fill.background()
            tf = layer_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.25)
            
            # 组件
            comp_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(3.2), Inches(y), Inches(9.6), Inches(1)
            )
            comp_box.fill.solid()
            comp_box.fill.fore_color.rgb = self.theme['light']
            comp_box.line.color.rgb = self.theme['primary']
            tf = comp_box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.25)
            p = tf.paragraphs[0]
            p.text = components
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 1.2
        
        # 6. 系统集成方案
        slide = self.add_slide()
        self.add_title_shape(slide, "系统集成架构")
        
        integration_items = [
            ("OA系统", "流程审批、消息推送"),
            ("MES系统", "生产数据、工单下发"),
            ("WMS系统", "库存同步、出入库"),
            ("CRM系统", "客户数据、销售机会"),
            ("BI系统", "数据分析、报表展示")
        ]
        
        # 中心系统
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), Inches(3), Inches(2.5), Inches(1.5)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = self.theme['primary']
        center.line.fill.background()
        tf = center.text_frame
        p = tf.paragraphs[0]
        p.text = "金蝶云·星瀚"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.theme['white']
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.4)
        
        # 周围系统
        positions = [(1, 1.8), (9.5, 1.8), (1, 4.5), (9.5, 4.5), (5.5, 5.5)]
        for i, (name, funcs) in enumerate(integration_items):
            x, y = positions[i]
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(2.5), Inches(1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            tf = box.text_frame
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            p = tf.add_paragraph()
            p.text = funcs
            p.font.size = Pt(10)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
        
        # 7. 价值工程
        slide = self.add_slide()
        self.add_title_shape(slide, "价值工程 - 4S价值模型")
        
        value_items = [
            ("Strategy\n战略", "支撑企业战略落地\n提升战略执行力\n实现战略目标"),
            ("Spending\n支出回报", "财务结账时间缩短50%\n库存周转率提升40%\n采购成本降低10%"),
            ("Situation\n情景", "响应速度提升30%\n客户满意度提升20%\n市场占有率提升10%"),
            ("Structure\n结构", "业务流程效率提升40%\n数据准确性提升30%\n决策效率提升50%")
        ]
        
        x = 0.8
        for title, content in value_items:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(2.8), Inches(4.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(2.8), Inches(1.2)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.3)
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.2), Inches(2.6), Inches(2.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content.split('\n')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)
                p.font.color.rgb = self.theme['text']
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(8)
            
            x += 3.1
        
        # 8. 运营指标
        slide = self.add_slide()
        self.add_title_shape(slide, "运营指标提升")
        
        indicators = [
            ("财务结账", "10天→5天", "缩短50%"),
            ("库存周转", "4次→5.6次", "提升40%"),
            ("订单交付", "80%→95%", "提升15%"),
            ("成本核算", "粗放→精准", "精度30%"),
            ("生产计划", "60%→90%", "提升30%"),
            ("数据准确", "85%→98%", "提升13%")
        ]
        
        x = 0.5
        for name, change, effect in indicators:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 指标名
            name_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.3), Inches(2), Inches(0.5)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 变化
            change_box = slide.shapes.add_textbox(
                Inches(x), Inches(3), Inches(2), Inches(0.8)
            )
            tf = change_box.text_frame
            p = tf.paragraphs[0]
            p.text = change
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 效果
            effect_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.3), Inches(4), Inches(1.4), Inches(0.6)
            )
            effect_box.fill.solid()
            effect_box.fill.fore_color.rgb = RGBColor(76, 175, 80)
            effect_box.line.fill.background()
            tf = effect_box.text_frame
            p = tf.paragraphs[0]
            p.text = effect
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            x += 2.1
    
    def generate_implementation_plan(self):
        """生成实施计划 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("实施计划", [
            "实施方法",
            "项目计划",
            "里程碑节点",
            "数据迁移",
            "风险控制"
        ])
        
        # 1. 实施方法论
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶实施方法论 - Kingdee Way")
        
        phases = [
            ("项目启动", "1-2周", "组建团队、制定计划"),
            ("需求调研", "2-4周", "业务调研、需求确认"),
            ("方案设计", "3-4周", "蓝图设计、方案评审"),
            ("系统配置", "4-6周", "系统配置、开发定制"),
            ("测试上线", "2-3周", "系统测试、用户培训"),
            ("验收交付", "1-2周", "上线切换、项目验收")
        ]
        
        x = 0.5
        for i, (name, duration, desc) in enumerate(phases):
            box = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x), Inches(2.5), Inches(2), Inches(1.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['primary'] if i % 2 == 0 else self.theme['secondary']
            box.line.fill.background()
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            
            duration_box = slide.shapes.add_textbox(
                Inches(x), Inches(3.8), Inches(2), Inches(0.4)
            )
            tf = duration_box.text_frame
            p = tf.paragraphs[0]
            p.text = duration
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            desc_box = slide.shapes.add_textbox(
                Inches(x), Inches(4.2), Inches(2), Inches(0.8)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.1
        
        # 2. 项目计划甘特图
        slide = self.add_slide()
        self.add_title_shape(slide, "项目实施计划")
        
        # 用形状绘制甘特图
        tasks = [
            ("项目启动", 0, 2, self.theme['primary']),
            ("需求调研", 2, 6, self.theme['secondary']),
            ("方案设计", 6, 10, self.theme['accent']),
            ("系统配置", 10, 16, RGBColor(76, 175, 80)),
            ("开发定制", 12, 18, RGBColor(255, 152, 0)),
            ("系统测试", 16, 19, RGBColor(244, 67, 54)),
            ("用户培训", 17, 19, RGBColor(156, 39, 176)),
            ("上线切换", 19, 20, RGBColor(0, 150, 136)),
            ("项目验收", 20, 22, self.theme['primary'])
        ]
        
        # 表头 - 月份
        for i, month in enumerate(["第1月", "第2月", "第3月", "第4月", "第5月", "第6月"]):
            header = slide.shapes.add_textbox(
                Inches(3.5 + i * 1.5), Inches(1.5), Inches(1.5), Inches(0.4)
            )
            tf = header.text_frame
            p = tf.paragraphs[0]
            p.text = month
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
        
        y = 2
        for task_name, start, end, color in tasks:
            # 任务名
            name_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y), Inches(2.8), Inches(0.5)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = task_name
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            # 甘特条
            bar_left = 3.5 + start * 0.3
            bar_width = (end - start) * 0.3
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(y + 0.05),
                Inches(bar_width), Inches(0.35)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            
            y += 0.5
        
        # 3. 里程碑节点
        slide = self.add_slide()
        self.add_title_shape(slide, "项目里程碑")
        
        milestones = [
            ("M1", "项目启动会", "第1周", "项目正式启动"),
            ("M2", "需求确认", "第6周", "业务需求签字确认"),
            ("M3", "蓝图评审", "第10周", "业务蓝图方案评审通过"),
            ("M4", "系统配置完成", "第16周", "系统配置和开发完成"),
            ("M5", "UAT测试通过", "第19周", "用户验收测试通过"),
            ("M6", "正式上线", "第20周", "系统正式上线运行"),
            ("M7", "项目验收", "第22周", "项目整体验收")
        ]
        
        y = 1.6
        for code, name, time, desc in milestones:
            # 里程碑标记
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(0.5), Inches(y), Inches(0.5), Inches(0.5)
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = self.theme['primary']
            diamond.line.fill.background()
            
            # 编号+名称
            name_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y + 0.05), Inches(2.5), Inches(0.4)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{code}: {name}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            # 时间
            time_box = slide.shapes.add_textbox(
                Inches(4), Inches(y + 0.05), Inches(1.5), Inches(0.4)
            )
            tf = time_box.text_frame
            p = tf.paragraphs[0]
            p.text = time
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(5.8), Inches(y + 0.05), Inches(6.5), Inches(0.4)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            y += 0.7
        
        # 4. 数据迁移方案
        slide = self.add_slide()
        self.add_title_shape(slide, "数据迁移方案")
        
        migration_steps = [
            ("数据梳理", "全面盘点现有系统数据\n确定迁移范围和优先级"),
            ("数据清洗", "数据质量检查\n重复数据清理、格式统一"),
            ("数据映射", "建立新旧系统字段映射关系\n制定数据转换规则"),
            ("数据迁移", "分批迁移、数据校验\n异常处理、回退方案"),
            ("数据验证", "迁移结果验证\n数据完整性、准确性检查")
        ]
        
        x = 0.5
        for title, desc in migration_steps:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.3), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 序号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.9), Inches(2.2), Inches(0.5), Inches(0.5)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(migration_steps.index((title, desc)) + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.9), Inches(2.1), Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.5), Inches(2.1), Inches(2.2)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.5
        
        # 5. 风险控制
        slide = self.add_slide()
        self.add_title_shape(slide, "项目风险控制")
        
        risks = [
            ("需求变更风险", "高", "建立变更管理流程，控制需求范围"),
            ("数据迁移风险", "高", "分批迁移、充分测试、制定回退方案"),
            ("项目进度风险", "中", "制定详细计划、定期检查、及时调整"),
            ("人员变动风险", "中", "知识转移、文档沉淀、备份人员"),
            ("技术风险", "低", "技术预研、原型验证、技术方案评审"),
            ("培训效果风险", "低", "分层次培训、模拟演练、考核机制")
        ]
        
        # 表头
        headers = [("风险类型", 0.5, 2.5), ("等级", 3.2, 1), ("应对措施", 4.4, 8.5)]
        for text, x, w in headers:
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.5), Inches(w), Inches(0.5)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.theme['primary']
            header.line.fill.background()
            tf = header.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
        
        y = 2.1
        for risk, level, measure in risks:
            # 风险
            risk_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(y + 0.1), Inches(2.3), Inches(0.4)
            )
            tf = risk_box.text_frame
            p = tf.paragraphs[0]
            p.text = risk
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            # 等级
            level_color = RGBColor(244, 67, 54) if level == "高" else (RGBColor(255, 152, 0) if level == "中" else RGBColor(76, 175, 80))
            level_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(3.5), Inches(y + 0.05), Inches(0.7), Inches(0.4)
            )
            level_box.fill.solid()
            level_box.fill.fore_color.rgb = level_color
            level_box.line.fill.background()
            tf = level_box.text_frame
            p = tf.paragraphs[0]
            p.text = level
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.05)
            
            # 措施
            measure_box = slide.shapes.add_textbox(
                Inches(4.5), Inches(y + 0.1), Inches(8.3), Inches(0.4)
            )
            tf = measure_box.text_frame
            p = tf.paragraphs[0]
            p.text = measure
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            y += 0.65
    
    def generate_closing(self):
        """生成结束页"""
        slide = self.add_slide()
        
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme['primary']
        bg.line.fill.background()
        
        # 感谢语
        thanks = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        tf = thanks.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢聆听"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = self.theme['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(11.333), Inches(0.8)
        )
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = "携手共进，共创数字化未来"
        p.font.size = Pt(24)
        p.font.color.rgb = self.theme['light']
        p.alignment = PP_ALIGN.CENTER
        
        # 联系方式
        contact = slide.shapes.add_textbox(
            Inches(1), Inches(5.5), Inches(11.333), Inches(1)
        )
        tf = contact.text_frame
        p = tf.paragraphs[0]
        p.text = "金蝶软件（中国）有限公司\n网址：www.kingdee.com  |  服务热线：4008-830-830"
        p.font.size = Pt(14)
        p.font.color.rgb = self.theme['light']
        p.alignment = PP_ALIGN.CENTER
    
    def generate_customer_requirement(self):
        """生成客户需求理解章节"""
        # 章节导航
        self.add_navigation_slide("需求理解", [
            "客户背景",
            "组织架构",
            "业务痛点",
            "建设目标"
        ])
        
        # 客户背景
        slide = self.add_slide()
        self.add_title_shape(slide, f"{self.company_name}概况")
        
        # 左侧内容
        left_content = f"""企业简介

{self.company_name}是一家专注于{self.industry}领域的现代化企业。

• 企业规模：中型企业
• 员工人数：约500人
• 年营业额：约1亿元
• 主要业务：产品研发、生产制造、销售服务

组织特点：
• 多组织协同管理
• 业务流程复杂
• 数据分散各系统"""
        
        self.add_content_box(slide, left_content, 0.5, 1.5, 5.8, 5, font_size=14)
        
        # 右侧内容
        right_content = """信息化现状

现有系统：
• 财务系统：基础财务核算
• 进销存：库存管理
• 生产：手工管理
• 人事：Excel管理

存在问题：
• 系统分散，数据孤岛
• 流程不畅通，效率低
• 成本核算不精准
• 决策数据不及时"""
        
        self.add_content_box(slide, right_content, 6.8, 1.5, 6, 5, font_size=14)
        
        # 业务痛点
        slide = self.add_slide()
        self.add_title_shape(slide, "业务痛点分析")
        
        pain_points = [
            ("生产计划", "计划准确性低，物料齐套难，影响交付"),
            ("成本管理", "成本核算粗放，无法精准分析利润"),
            ("供应链", "库存积压严重，采购周期长"),
            ("财务核算", "结账周期长，报表不及时"),
            ("数据协同", "系统分散，数据不统一")
        ]
        
        y = 1.8
        for i, (area, desc) in enumerate(pain_points):
            # 序号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(y), Inches(0.5), Inches(0.5)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            # 领域
            area_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y), Inches(2), Inches(0.5)
            )
            tf = area_box.text_frame
            p = tf.paragraphs[0]
            p.text = area
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(3.5), Inches(y), Inches(9), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme['text']
            
            y += 0.9
        
        # 建设目标
        slide = self.add_slide()
        self.add_title_shape(slide, "数字化建设目标")
        
        goals = [
            ("效率提升", "业务流程效率提升40%\n决策响应速度提升50%"),
            ("成本降低", "库存成本降低20%\n采购成本降低10%"),
            ("质量提升", "产品质量追溯率100%\n客户满意度提升20%"),
            ("管理升级", "财务结账时间缩短50%\n成本核算精度提升30%")
        ]
        
        x = 0.8
        for title, desc in goals:
            # 目标框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = self.theme['primary']
            title_box.line.fill.background()
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
    
    def generate_function_modules(self):
        """生成功能模块详解 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("功能详解", [
            "财务管理",
            "供应链管理", 
            "生产制造",
            "项目管理",
            "人力资源",
            "移动应用"
        ])
        
        # ========== 财务管理模块 ==========
        self._generate_finance_module()
        
        # ========== 供应链管理模块 ==========
        self._generate_supply_chain_module()
        
        # ========== 生产制造模块 ==========
        self._generate_manufacturing_module()
        
        # ========== 项目管理模块 ==========
        self._generate_project_module()
        
        # ========== 人力资源模块 ==========
        self._generate_hr_module()
        
        # ========== 移动应用模块 ==========
        self._generate_mobile_module()
    
    def _generate_finance_module(self):
        """生成财务管理模块详情"""
        # 模块概览页
        slide = self.add_slide()
        self.add_title_shape(slide, "财务管理模块")
        
        # 功能列表
        functions = [
            ("总账核算", "凭证管理、账簿查询、报表编制、多账簿管理"),
            ("应收管理", "销售开票、收款核销、账龄分析、信用控制"),
            ("应付管理", "采购发票、付款管理、对账管理、付款计划"),
            ("资金管理", "现金银行、资金计划、银企直联、票据管理"),
            ("成本管理", "成本核算、成本分析、成本控制、成本预测"),
            ("固定资产", "资产登记、折旧计提、资产处置、资产盘点"),
            ("预算管理", "预算编制、预算控制、预算分析、预算调整"),
            ("合并报表", "报表合并、抵销处理、报表分析、报表披露")
        ]
        
        y = 1.5
        for i, (name, desc) in enumerate(functions):
            # 功能名
            name_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(2), Inches(0.6)
            )
            name_box.fill.solid()
            name_box.fill.fore_color.rgb = self.theme['primary']
            name_box.line.fill.background()
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(2.7), Inches(y + 0.1), Inches(10), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 0.7
        
        # 总账核算详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "总账核算 - 核心功能")
        
        features = [
            ("凭证管理", "智能凭证生成、凭证模板、批量导入、凭证审核"),
            ("账簿查询", "总账、明细账、余额表、日记账实时查询"),
            ("报表编制", "资产负债表、利润表、现金流量表自动生成"),
            ("多账簿管理", "支持多会计准则、多币种、多账簿并行")
        ]
        
        x = 0.5
        for title, desc in features:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(3), Inches(4.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.8), Inches(3), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(2.8), Inches(2.6), Inches(3.2)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(desc.split('、')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.theme['text']
                p.space_after = Pt(8)
            
            x += 3.2
        
        # 成本管理详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "成本管理 - 精细化核算")
        
        cost_features = [
            "成本核算体系：支持品种法、分批法、分步法等多种核算方法",
            "成本要素管理：材料成本、人工成本、制造费用全面管控",
            "成本分析报表：成本结构分析、成本趋势分析、成本对比分析",
            "成本预测：基于历史数据的成本预测，支持决策分析"
        ]
        
        y = 1.8
        for feature in cost_features:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(1.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            tf = box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            p.text = feature
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 1.3
        
        # 如果有截图，添加截图展示页
        if self.assets.get('finance-screenshots'):
            for img_path in self.assets['finance-screenshots'][:2]:
                slide = self.add_slide()
                self.add_subtitle(slide, "财务管理 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def _generate_supply_chain_module(self):
        """生成供应链管理模块详情"""
        # 模块概览页
        slide = self.add_slide()
        self.add_title_shape(slide, "供应链管理模块")
        
        functions = [
            ("采购管理", "采购申请、订单管理、入库管理、价格管理"),
            ("销售管理", "销售订单、出库管理、销售分析、客户管理"),
            ("库存管理", "出入库、盘点、库存分析、库存预警"),
            ("供应商管理", "供应商档案、准入管理、绩效评估"),
            ("客户管理", "客户档案、信用控制、价格管理、销售预测")
        ]
        
        y = 1.8
        for name, desc in functions:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 功能名
            name_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y + 0.15), Inches(2), Inches(0.6)
            )
            name_box.fill.solid()
            name_box.fill.fore_color.rgb = self.theme['primary']
            name_box.line.fill.background()
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(3), Inches(y + 0.25), Inches(9.5), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 1.05
        
        # 采购管理详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "采购管理 - 全流程管控")
        
        # 采购流程
        steps = [
            ("需求", "采购申请"),
            ("寻源", "供应商比价"),
            ("订单", "采购订单"),
            ("到货", "入库验收"),
            ("结算", "付款管理")
        ]
        
        x = 0.5
        for i, (title, desc) in enumerate(steps):
            # 步骤框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2.5), Inches(2.3), Inches(2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 步骤号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.9), Inches(2.7), Inches(0.5), Inches(0.5)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.4), Inches(2.1), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.9), Inches(2.1), Inches(0.5)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.5
        
        # 库存管理详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "库存管理 - 精细化管控")
        
        inventory_features = [
            ("多仓库管理", "支持多仓库、多库位、批次管理、序列号管理"),
            ("库存预警", "最低库存、最高库存、安全库存自动预警"),
            ("盘点管理", "定期盘点、循环盘点、抽盘、全盘"),
            ("库存分析", "库存周转率、呆滞库存、库存结构分析")
        ]
        
        x = 0.5
        for title, desc in inventory_features:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(3), Inches(4.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.8), Inches(3), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.8), Inches(2.8), Inches(3.2)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(desc.split('、')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.theme['text']
                p.space_after = Pt(6)
            
            x += 3.2
        
        # 如果有截图
        if self.assets.get('supply-screenshots'):
            for img_path in self.assets['supply-screenshots'][:2]:
                slide = self.add_slide()
                self.add_subtitle(slide, "供应链管理 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def _generate_manufacturing_module(self):
        """生成生产制造模块详情"""
        # 模块概览
        slide = self.add_slide()
        self.add_title_shape(slide, "生产制造模块")
        
        functions = [
            "生产计划：MRP运算、主生产计划、排程管理、产能分析",
            "生产订单：订单管理、工序计划、物料齐套、进度跟踪",
            "车间管理：工序管理、报工管理、生产进度、异常处理",
            "质量管理：来料检验、过程检验、成品检验、质量追溯",
            "BOM管理：物料清单、工艺路线、工程变更、版本管理"
        ]
        
        y = 1.8
        for func in functions:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(0.9)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            tf = box.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_top = Inches(0.2)
            p = tf.paragraphs[0]
            p.text = func
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 1.05
        
        # MRP详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "MRP智能运算")
        
        mrp_content = """
MRP（物料需求计划）核心功能：

• 需求分析：根据销售订单、生产计划自动计算物料需求
• 供应平衡：考虑库存、在途、在制，计算净需求
• 计划生成：自动生成采购计划、生产计划、调拨计划
• 产能评估：结合产能情况，评估计划可行性
• 预警提示：物料齐套预警、产能超载预警、交期风险预警

运算效率：
• 4万个自制件卷算约需1小时20分钟
• 单个跟踪卷算平均40秒
• 支持按订单/大类/产品分批运算"""
        
        self.add_content_box(slide, mrp_content, 0.5, 1.5, 12.333, 5, font_size=14)
        
        # 质量管理详情页
        slide = self.add_slide()
        self.add_title_shape(slide, "质量管理 - 全流程追溯")
        
        quality_flow = [
            ("来料检验", "供应商送货→检验→合格入库/不合格退货"),
            ("过程检验", "生产工序→首检/巡检→合格流转/不合格处理"),
            ("成品检验", "生产完成→终检→合格入库/不合格返工"),
            ("质量追溯", "批次号/序列号→正向追溯/反向追溯")
        ]
        
        y = 1.8
        for title, flow in quality_flow:
            # 标题
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(2), Inches(1)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = self.theme['primary']
            title_box.line.fill.background()
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.3)
            
            # 流程
            flow_box = slide.shapes.add_textbox(
                Inches(2.7), Inches(y + 0.25), Inches(10), Inches(0.6)
            )
            tf = flow_box.text_frame
            p = tf.paragraphs[0]
            p.text = flow
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 1.2
        
        # 如果有截图
        if self.assets.get('manufacture-screenshots'):
            for img_path in self.assets['manufacture-screenshots'][:2]:
                slide = self.add_slide()
                self.add_subtitle(slide, "生产制造 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def _generate_project_module(self):
        """生成项目管理模块详情"""
        slide = self.add_slide()
        self.add_title_shape(slide, "项目管理模块")
        
        functions = [
            "项目立项：项目登记、预算编制、资源分配",
            "项目执行：进度管理、成本控制、变更管理",
            "项目核算：项目收入、项目成本、项目利润",
            "项目分析：进度分析、成本分析、效益分析"
        ]
        
        y = 1.8
        for func in functions:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(12.333), Inches(1.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            tf = box.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_top = Inches(0.25)
            p = tf.paragraphs[0]
            p.text = func
            p.font.size = Pt(15)
            p.font.color.rgb = self.theme['text']
            
            y += 1.3
        
        # 如果有截图
        if self.assets.get('project-screenshots'):
            for img_path in self.assets['project-screenshots'][:1]:
                slide = self.add_slide()
                self.add_subtitle(slide, "项目管理 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def _generate_hr_module(self):
        """生成人力资源模块详情"""
        slide = self.add_slide()
        self.add_title_shape(slide, "人力资源管理模块")
        
        functions = [
            "人事管理：员工档案、合同管理、离职管理",
            "薪酬管理：薪资核算、个税计算、社保管理",
            "绩效管理：绩效指标、绩效评估、绩效分析",
            "招聘管理：招聘需求、简历管理、面试管理",
            "培训管理：培训计划、培训执行、培训评估"
        ]
        
        x = 0.5
        y = 1.8
        for i, func in enumerate(functions):
            if i == 3:
                x = 0.5
                y = 4
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(4), Inches(1.8)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            tf = box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.3)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = func
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            x += 4.2
        
        # 如果有截图
        if self.assets.get('hr-screenshots'):
            for img_path in self.assets['hr-screenshots'][:1]:
                slide = self.add_slide()
                self.add_subtitle(slide, "人力资源管理 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def _generate_mobile_module(self):
        """生成移动应用模块详情"""
        slide = self.add_slide()
        self.add_title_shape(slide, "移动应用")
        
        apps = [
            ("移动审批", "随时随地审批单据"),
            ("移动报表", "实时查看业务数据"),
            ("移动办公", "日程、任务、通知"),
            ("移动作业", "扫码出入库、报工")
        ]
        
        x = 0.8
        for title, desc in apps:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 图标占位
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.9), Inches(2.5), Inches(1), Inches(1)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = self.theme['accent']
            icon.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(3.8), Inches(2.8), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.4), Inches(2.6), Inches(1)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 如果有截图
        if self.assets.get('mobile-screenshots'):
            for img_path in self.assets['mobile-screenshots'][:1]:
                slide = self.add_slide()
                self.add_subtitle(slide, "移动应用 - 界面展示")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def generate_team_service(self):
        """生成团队保障章节 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("团队保障", [
            "项目团队",
            "培训计划",
            "知识转移",
            "服务承诺",
            "质量保障"
        ])
        
        # 1. 项目团队配置
        slide = self.add_slide()
        self.add_title_shape(slide, "项目团队配置")
        
        team_members = [
            ("项目总监", "1人", "统筹项目管理、风险控制、资源协调", "10年以上经验"),
            ("项目经理", "1人", "项目实施、进度管理、客户沟通", "PMP认证"),
            ("业务顾问", "2人", "业务调研、方案设计、蓝图确认", "行业专家"),
            ("技术顾问", "2人", "系统配置、开发定制、数据迁移", "金蝶认证"),
            ("实施顾问", "3人", "用户培训、系统指导、问题处理", "实施经验")
        ]
        
        # 表头
        headers = [("角色", 0.5, 2), ("人数", 2.7, 1), ("职责", 3.9, 5), ("资质", 9.1, 3.5)]
        for text, x, w in headers:
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.5), Inches(w), Inches(0.5)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.theme['primary']
            header.line.fill.background()
            tf = header.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
        
        y = 2.1
        for role, count, responsibility, qual in team_members:
            # 角色
            role_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.6), Inches(y), Inches(1.8), Inches(0.7)
            )
            role_box.fill.solid()
            role_box.fill.fore_color.rgb = self.theme['light']
            role_box.line.color.rgb = self.theme['primary']
            tf = role_box.text_frame
            p = tf.paragraphs[0]
            p.text = role
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 人数
            count_box = slide.shapes.add_textbox(
                Inches(2.8), Inches(y + 0.15), Inches(0.8), Inches(0.4)
            )
            tf = count_box.text_frame
            p = tf.paragraphs[0]
            p.text = count
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['accent']
            p.alignment = PP_ALIGN.CENTER
            
            # 职责
            resp_box = slide.shapes.add_textbox(
                Inches(4), Inches(y + 0.15), Inches(4.8), Inches(0.4)
            )
            tf = resp_box.text_frame
            p = tf.paragraphs[0]
            p.text = responsibility
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            # 资质
            qual_box = slide.shapes.add_textbox(
                Inches(9.2), Inches(y + 0.15), Inches(3.3), Inches(0.4)
            )
            tf = qual_box.text_frame
            p = tf.paragraphs[0]
            p.text = qual
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(100, 100, 100)
            
            y += 0.85
        
        # 2. 团队优势
        slide = self.add_slide()
        self.add_title_shape(slide, "团队优势")
        
        advantages = [
            ("行业经验", "深耕制造业20年，累计服务500+企业"),
            ("专业资质", "团队成员均持有金蝶认证顾问资质"),
            ("项目经验", "平均项目经验5年以上，熟悉各类业务场景"),
            ("本地服务", "本地团队驻场服务，快速响应客户需求"),
            ("知识沉淀", "丰富的行业解决方案和最佳实践库")
        ]
        
        x = 0.5
        for i, (title, desc) in enumerate(advantages):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.4), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            # 图标
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.7), Inches(2.3), Inches(1), Inches(1)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = self.theme['primary']
            icon.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.5), Inches(2.2), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.1), Inches(2.2), Inches(1.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.55
        
        # 3. 培训计划
        slide = self.add_slide()
        self.add_title_shape(slide, "分阶段培训计划")
        
        training_phases = [
            ("项目启动阶段", "管理层培训", "理念导入、项目管理、变革管理", "1天"),
            ("蓝图设计阶段", "关键用户培训", "系统功能、业务流程、蓝图确认", "3天"),
            ("系统配置阶段", "全员培训", "操作培训、模拟演练、考试认证", "5天"),
            ("上线切换阶段", "现场支持", "问题解答、操作指导、应急处理", "持续")
        ]
        
        y = 1.6
        for phase, target, content, duration in training_phases:
            # 阶段
            phase_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(2.5), Inches(1.2)
            )
            phase_box.fill.solid()
            phase_box.fill.fore_color.rgb = self.theme['primary']
            phase_box.line.fill.background()
            tf = phase_box.text_frame
            tf.margin_top = Inches(0.3)
            p = tf.paragraphs[0]
            p.text = phase
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            
            # 培训对象
            target_box = slide.shapes.add_textbox(
                Inches(3.2), Inches(y + 0.1), Inches(2), Inches(0.4)
            )
            tf = target_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"培训对象：{target}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.theme['text']
            
            # 培训内容
            content_box = slide.shapes.add_textbox(
                Inches(3.2), Inches(y + 0.5), Inches(6), Inches(0.4)
            )
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"培训内容：{content}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme['text']
            
            # 周期
            duration_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(10), Inches(y + 0.35), Inches(2.5), Inches(0.5)
            )
            duration_box.fill.solid()
            duration_box.fill.fore_color.rgb = self.theme['accent']
            duration_box.line.fill.background()
            tf = duration_box.text_frame
            p = tf.paragraphs[0]
            p.text = duration
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            y += 1.4
        
        # 4. 培训方式
        slide = self.add_slide()
        self.add_title_shape(slide, "多元化培训方式")
        
        training_methods = [
            ("集中培训", "理论讲解+系统演示\n学员上机操作\n现场答疑"),
            ("视频教程", "操作视频录制\n在线学习平台\n随时随地学习"),
            ("操作手册", "详细操作步骤\n图文并茂说明\n快速查询指南"),
            ("现场指导", "顾问驻场支持\n一对一辅导\n问题即时解决")
        ]
        
        x = 0.8
        for title, content in training_methods:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 5. 知识转移
        slide = self.add_slide()
        self.add_title_shape(slide, "知识转移体系")
        
        knowledge_items = [
            ("项目文档", "需求文档、设计文档、测试文档、操作手册"),
            ("培训资料", "培训课件、操作视频、常见问题FAQ"),
            ("系统配置", "配置清单、参数说明、权限设置"),
            ("开发文档", "接口文档、开发规范、代码说明"),
            ("运维手册", "系统维护、备份恢复、应急处理")
        ]
        
        y = 1.8
        for title, content in knowledge_items:
            # 标题
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(y), Inches(2), Inches(0.7)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = self.theme['primary']
            title_box.line.fill.background()
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(2.8), Inches(y + 0.15), Inches(10), Inches(0.5)
            )
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            
            y += 0.9
        
        # 6. 服务承诺
        slide = self.add_slide()
        self.add_title_shape(slide, "服务承诺")
        
        services = [
            ("响应时效", "7×24小时服务热线\n2小时内响应\n24小时内解决"),
            ("服务期限", "验收后2年免费服务\n专属服务团队\n定期回访机制"),
            ("知识转移", "系统操作手册\n视频培训课程\n在线知识库"),
            ("持续优化", "版本升级服务\n功能优化建议\n业务咨询服务")
        ]
        
        x = 0.8
        for title, content in services:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 7. 质量保障
        slide = self.add_slide()
        self.add_title_shape(slide, "质量保障体系")
        
        qa_items = [
            ("需求确认", "需求签字确认、变更管理流程"),
            ("设计评审", "方案设计评审、蓝图签字确认"),
            ("测试验收", "单元测试、集成测试、UAT测试"),
            ("上线检查", "上线检查清单、应急预案准备"),
            ("文档交付", "项目文档齐备、知识转移完成")
        ]
        
        y = 1.6
        for i, (title, content) in enumerate(qa_items):
            # 序号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(y), Inches(0.6), Inches(0.6)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y + 0.1), Inches(2), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(3.5), Inches(y + 0.1), Inches(9), Inches(0.5)
            )
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 0.85
        
        # 8. 如果有团队照片
        if self.assets.get('team-photos'):
            for img_path in self.assets['team-photos'][:1]:
                slide = self.add_slide()
                self.add_subtitle(slide, "项目团队风采")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 培训对象
            target_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.8), Inches(0.6)
            )
            tf = target_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"培训对象：{target}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.theme['text']
            
            # 培训内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.7), Inches(2.8), Inches(2)
            )
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"培训内容：{content}"
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            x += 3.2
        
        # 服务承诺
        slide = self.add_slide()
        self.add_title_shape(slide, "服务承诺")
        
        services = [
            ("响应时效", "7×24小时服务热线\n2小时内响应\n24小时内解决"),
            ("服务期限", "验收后2年免费服务\n专属服务团队\n定期回访机制"),
            ("知识转移", "系统操作手册\n视频培训课程\n在线知识库"),
            ("持续优化", "版本升级服务\n功能优化建议\n业务咨询服务")
        ]
        
        x = 0.8
        for title, content in services:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(4)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            # 标题
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
    
    def generate_success_cases(self):
        """生成成功案例章节 - 扩展版"""
        # 章节导航
        self.add_navigation_slide("成功案例", [
            "行业标杆",
            "案例一",
            "案例二",
            "案例三",
            "实施成果"
        ])
        
        # 行业标杆客户
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶标杆客户")
        
        # 客户分类
        client_categories = [
            ("世界500强", "华为、腾讯、万科、海尔、中国中车等"),
            ("中国500强", "招商局、中粮、中国重汽、三一重工等"),
            ("上市公司", "迈瑞医疗、宁德时代、立讯精密等"),
            ("行业龙头", "各细分领域头部企业200+")
        ]
        
        x = 0.8
        for title, clients in client_categories:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(3.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = self.theme['primary']
            
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(2), Inches(2.8), Inches(0.8)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.15)
            
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3), Inches(2.6), Inches(2.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = clients
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.1
        
        # 如果有客户Logo
        if self.assets.get('case-logos'):
            slide = self.add_slide()
            self.add_title_shape(slide, "合作伙伴")
            y = 1.8
            x = 0.5
            for i, img_path in enumerate(self.assets['case-logos'][:12]):
                self.add_image(slide, img_path, x, y, 2.8, 1.2)
                x += 3.1
                if x > 10:
                    x = 0.5
                    y += 1.5
        
        # 案例一
        slide = self.add_slide()
        self.add_title_shape(slide, f"案例一：某大型{self.industry}企业")
        
        case1_info = [
            ("企业规模", "年营收50亿+，员工3000+"),
            ("实施范围", "财务云、供应链云、制造云、人力云"),
            ("实施周期", "8个月（2023.03-2023.11）"),
            ("用户数量", "800+用户同时在线")
        ]
        
        y = 1.5
        for label, value in case1_info:
            label_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y), Inches(2), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            value_box = slide.shapes.add_textbox(
                Inches(2.8), Inches(y), Inches(10), Inches(0.5)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 0.7
        
        # 实施效果
        effects = [
            ("库存周转率", "提升42%"),
            ("订单交付率", "从78%→96%"),
            ("财务结账", "从8天→3天"),
            ("采购成本", "降低15%")
        ]
        
        x = 0.5
        for title, value in effects:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(4.5), Inches(3), Inches(2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = RGBColor(76, 175, 80)
            
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(4.7), Inches(3), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(5.3), Inches(3), Inches(0.6)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(76, 175, 80)
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.2
        
        # 案例一详情
        slide = self.add_slide()
        self.add_title_shape(slide, f"案例一：核心成果")
        
        achievements = [
            "打通业财一体化：销售订单→发货→出库→确认收入→应收→收款→核销，全流程自动化",
            "生产计划优化：MRP运算准确率从60%提升到95%，物料齐套率从70%提升到92%",
            "成本精细核算：实现按订单、按产品、按工序的精细成本核算，成本核算精度提升30%",
            "供应链协同：与上游50+供应商实现协同采购，采购周期缩短20%",
            "数据驱动决策：搭建管理驾驶舱，实时监控30+核心指标，决策效率提升50%"
        ]
        
        y = 1.6
        for i, item in enumerate(achievements):
            # 序号
            num = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(y), Inches(0.5), Inches(0.5)
            )
            num.fill.solid()
            num.fill.fore_color.rgb = self.theme['primary']
            num.line.fill.background()
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.08)
            
            # 内容
            content_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.2), Inches(y), Inches(11.6), Inches(0.9)
            )
            content_box.fill.solid()
            content_box.fill.fore_color.rgb = self.theme['light']
            content_box.line.color.rgb = self.theme['primary']
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.15)
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme['text']
            
            y += 1.05
        
        # 案例二
        slide = self.add_slide()
        self.add_title_shape(slide, "案例二：某集团企业数字化转型")
        
        case2_info = [
            ("企业规模", "集团总部+12个分子公司，员工5000+"),
            ("实施范围", "集团管控、财务合并、多组织协同、全面预算"),
            ("实施周期", "12个月"),
            ("核心诉求", "集团统一管控、数据透明化、业财一体化")
        ]
        
        y = 1.5
        for label, value in case2_info:
            label_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y), Inches(2), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            value_box = slide.shapes.add_textbox(
                Inches(2.8), Inches(y), Inches(10), Inches(0.5)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 0.7
        
        effects2 = [
            ("合并报表", "从5天→1天"),
            ("预算控制", "事前+事中管控"),
            ("资金可视", "100%透明"),
            ("数据准确", "99.5%+")
        ]
        
        x = 0.5
        for title, value in effects2:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(4.5), Inches(3), Inches(2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = RGBColor(33, 150, 243)
            
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(4.7), Inches(3), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(5.3), Inches(3), Inches(0.6)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(33, 150, 243)
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.2
        
        # 案例三
        slide = self.add_slide()
        self.add_title_shape(slide, "案例三：某上市企业ERP升级")
        
        case3_info = [
            ("企业规模", "上市公司，年营收30亿+，员工2000+"),
            ("实施范围", "财务云升级、供应链优化、生产制造、成本管理"),
            ("实施周期", "6个月"),
            ("核心诉求", "从旧系统升级到云平台，提升效率、降低成本")
        ]
        
        y = 1.5
        for label, value in case3_info:
            label_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y), Inches(2), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['primary']
            
            value_box = slide.shapes.add_textbox(
                Inches(2.8), Inches(y), Inches(10), Inches(0.5)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            
            y += 0.7
        
        effects3 = [
            ("系统响应", "提升300%"),
            ("运维成本", "降低60%"),
            ("用户满意度", "92%→98%"),
            ("上线周期", "缩短40%")
        ]
        
        x = 0.5
        for title, value in effects3:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(4.5), Inches(3), Inches(2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['light']
            box.line.color.rgb = RGBColor(255, 152, 0)
            
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(4.7), Inches(3), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(5.3), Inches(3), Inches(0.6)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 152, 0)
            p.alignment = PP_ALIGN.CENTER
            
            x += 3.2
        
        # 实施成果总览
        slide = self.add_slide()
        self.add_title_shape(slide, "金蝶客户实施成果总览")
        
        overall_results = [
            ("740万+", "服务企业数量", self.theme['primary']),
            ("95%", "项目验收通过率", RGBColor(76, 175, 80)),
            ("40%", "平均效率提升", RGBColor(33, 150, 243)),
            ("30%", "平均成本降低", RGBColor(255, 152, 0)),
            ("98%", "客户满意度", RGBColor(156, 39, 176))
        ]
        
        x = 0.5
        for value, label, color in overall_results:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2), Inches(2.4), Inches(3.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = color
            box.line.width = Pt(3)
            
            # 数值
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.5), Inches(2.4), Inches(1)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(3.8), Inches(2.4), Inches(1)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.55
        
        # 如果有案例照片
        if self.assets.get('case-photos'):
            for img_path in self.assets['case-photos'][:2]:
                slide = self.add_slide()
                self.add_subtitle(slide, "项目实施现场")
                self.add_image(slide, img_path, 0.5, 1.5, 12.333, 5.5)
    
    def generate_industry_solution(self):
        """生成行业专属解决方案"""
        # 章节导航
        self.add_navigation_slide("行业方案", [
            "行业理解",
            "解决方案",
            "实施路径"
        ])
        
        # 行业理解
        slide = self.add_slide()
        self.add_title_shape(slide, f"{self.industry}行业理解")
        
        understanding = f"""金蝶深耕{self.industry}行业20+年，深度理解行业特点：

行业特征：
• 产品种类多，BOM层级深，生产管理复杂
• 供应链长，供应商和客户分布广泛
• 质量管控要求严格，需要全流程追溯
• 成本核算难度大，需要精细化成本管理
• 订单模式多样（MTO/MTS/CTO），计划排程复杂

管理挑战：
• 多工厂协同困难，资源利用率低
• 库存积压与缺料并存，资金占用大
• 成本核算不精准，利润分析困难
• 质量追溯困难，客诉处理效率低
• 数据分散，决策缺乏数据支撑"""
        
        self.add_content_box(slide, understanding, 0.5, 1.5, 12.333, 5, font_size=13)
        
        # 行业解决方案
        slide = self.add_slide()
        self.add_title_shape(slide, f"{self.industry}行业解决方案")
        
        solutions = [
            ("智能计划", "MRP+APS智能排产\n多级计划协同\n产能智能分析"),
            ("精益成本", "标准成本+实际成本\n作业成本法\n成本预测分析"),
            ("全程追溯", "批次/序列号追溯\n正反向追溯\n质量档案管理"),
            ("供应链协同", "SRM供应商协同\nVMI库存管理\n需求预测协同"),
            ("业财一体", "订单到收款自动化\n采购到付款自动化\n实时成本核算")
        ]
        
        x = 0.5
        for title, content in solutions:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8), Inches(2.4), Inches(5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.theme['white']
            box.line.color.rgb = self.theme['primary']
            box.line.width = Pt(2)
            
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(1.8), Inches(2.4), Inches(0.7)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.theme['primary']
            title_shape.line.fill.background()
            tf = title_shape.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme['white']
            p.alignment = PP_ALIGN.CENTER
            tf.margin_top = Inches(0.1)
            
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(2.7), Inches(2.2), Inches(3.8)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme['text']
            p.alignment = PP_ALIGN.CENTER
            
            x += 2.55
    
    def generate(self):
        """生成完整PPT"""
        print(f"开始生成售前PPT: {self.company_name}")
        print(f"模式: {self.mode}")
        print(f"行业: {self.industry}")
        
        # 封面
        self.generate_cover()
        
        # 目录
        self.generate_toc()
        
        # 公司介绍
        self.generate_company_intro()
        
        # 行业分析
        self.generate_industry_analysis()
        
        # 客户需求理解
        self.generate_customer_requirement()
        
        # 解决方案总览
        self.generate_solution_overview()
        
        # 功能模块详解
        self.generate_function_modules()
        
        # 系统集成与价值工程
        self.generate_solution_overview()
        
        # 实施计划
        self.generate_implementation_plan()
        
        # 行业专属方案
        self.generate_industry_solution()
        
        # 团队保障
        self.generate_team_service()
        
        # 成功案例
        self.generate_success_cases()
        
        # 结束页
        self.generate_closing()
        
        print(f"PPT生成完成，共 {self.slide_count} 页")
        
        return self.prs


def main():
    parser = argparse.ArgumentParser(description='金蝶售前PPT生成器')
    parser.add_argument('--company-name', required=True, help='企业名称')
    parser.add_argument('--customer-code', default='default', help='客户代码')
    parser.add_argument('--industry', default='制造业', help='所属行业')
    parser.add_argument('--output', help='输出文件路径')
    parser.add_argument('--mode', default='detailed', 
                       choices=['standard', 'detailed', 'professional'],
                       help='PPT模式')
    parser.add_argument('--theme', default='blue',
                       choices=['blue', 'dark', 'green', 'red'],
                       help='主题颜色')
    
    args = parser.parse_args()
    
    # 素材目录
    assets_dir = os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        'assets', args.customer_code
    )
    
    # 生成器
    generator = PresalesPPTGenerator(
        company_name=args.company_name,
        customer_code=args.customer_code,
        industry=args.industry,
        mode=args.mode,
        theme=args.theme,
        assets_dir=assets_dir
    )
    
    # 生成PPT
    prs = generator.generate()
    
    # 保存
    if args.output:
        output_path = args.output
    else:
        output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'output')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(
            output_dir,
            f"{args.company_name}_售前PPT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        )
    
    prs.save(output_path)
    print(f"PPT已保存: {output_path}")
    
    return output_path


if __name__ == '__main__':
    main()
