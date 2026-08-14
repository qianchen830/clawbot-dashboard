#!/usr/bin/env python3
"""
分析金蝶售前PPT参考文件的结构
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict

def analyze_ppt(ppt_path):
    """分析PPT文件的完整结构"""
    print(f"正在分析: {ppt_path}")
    print("=" * 80)
    
    prs = Presentation(ppt_path)
    
    total_slides = len(prs.slides)
    print(f"\n总页数: {total_slides}")
    print(f"幻灯片尺寸: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    
    # 分析每一页
    slide_info = []
    section_pages = defaultdict(list)
    
    for idx, slide in enumerate(prs.slides, 1):
        # 获取页面标题
        title = None
        title_text = ""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) > 0:
                    # 第一个有文本的形状可能是标题
                    if not title_text:
                        title_text = text[:100]  # 只取前100个字符
                    break
        
        # 获取幻灯片布局名称
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        
        # 统计形状
        shape_count = len(slide.shapes)
        text_shapes = sum(1 for s in slide.shapes if s.has_text_frame)
        image_shapes = sum(1 for s in slide.shapes if s.shape_type == 13)  # Picture
        table_shapes = sum(1 for s in slide.shapes if s.has_table)
        chart_shapes = sum(1 for s in slide.shapes if s.has_chart)
        
        info = {
            'index': idx,
            'title': title_text,
            'layout': layout_name,
            'shapes': shape_count,
            'text_shapes': text_shapes,
            'image_shapes': image_shapes,
            'table_shapes': table_shapes,
            'chart_shapes': chart_shapes
        }
        slide_info.append(info)
        
        # 检测章节标题（通常是短文本、居中、字体较大）
        if title_text and len(title_text) < 30:
            # 可能是章节标题
            if any(keyword in title_text for keyword in ['章', '部分', '篇', '第', '一、', '二、', '三、', '四、', '五、', '六、', '七、', '八、']):
                section_pages[title_text].append(idx)
    
    # 输出详细分析
    print("\n" + "=" * 80)
    print("页面详细分析:")
    print("=" * 80)
    
    for info in slide_info:
        print(f"\n第{info['index']:3d}页: {info['title'][:50]}")
        print(f"  布局: {info['layout']}")
        print(f"  形状数: {info['shapes']} (文本:{info['text_shapes']}, 图片:{info['image_shapes']}, 表格:{info['table_shapes']}, 图表:{info['chart_shapes']})")
    
    # 章节统计
    if section_pages:
        print("\n" + "=" * 80)
        print("检测到的章节:")
        print("=" * 80)
        for section, pages in section_pages.items():
            print(f"  {section}: 页码 {pages}")
    
    # 统计布局使用情况
    layout_stats = defaultdict(int)
    for info in slide_info:
        layout_stats[info['layout']] += 1
    
    print("\n" + "=" * 80)
    print("布局使用统计:")
    print("=" * 80)
    for layout, count in sorted(layout_stats.items(), key=lambda x: x[1], reverse=True):
        print(f"  {layout}: {count}次")
    
    # 统计内容类型
    total_images = sum(info['image_shapes'] for info in slide_info)
    total_tables = sum(info['table_shapes'] for info in slide_info)
    total_charts = sum(info['chart_shapes'] for info in slide_info)
    
    print("\n" + "=" * 80)
    print("内容统计:")
    print("=" * 80)
    print(f"  总图片数: {total_images}")
    print(f"  总表格数: {total_tables}")
    print(f"  总图表数: {total_charts}")
    
    return slide_info, section_pages, layout_stats

if __name__ == "__main__":
    ppt_path = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/售前文件/中煤科工ERP重庆研究院ERP升级项目述标文件V7.0.pptx"
    analyze_ppt(ppt_path)
