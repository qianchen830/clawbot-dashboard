#!/usr/bin/env python3
"""
分析PPT母版模板的布局
"""
from pptx import Presentation
from pptx.util import Inches, Pt

def analyze_master(ppt_path):
    """分析母版模板"""
    print(f"正在分析母版模板: {ppt_path}")
    print("=" * 80)
    
    prs = Presentation(ppt_path)
    
    print(f"\n幻灯片尺寸: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    
    # 分析每个母版
    print(f"\n母版名称: {prs.slide_master.name if hasattr(prs.slide_master, 'name') else 'N/A'}")
    
    # 分析每个布局
    print("\n" + "=" * 80)
    print("布局列表:")
    print("=" * 80)
    
    for layout_idx, layout in enumerate(prs.slide_layouts):
        print(f"\n布局 {layout_idx + 1}: {layout.name}")
        print(f"  母版: {layout.slide_master.name if layout.slide_master else 'None'}")
        
        # 分析布局中的占位符
        if hasattr(layout, 'placeholders'):
            print(f"  占位符数量: {len(layout.placeholders)}")
            for ph_idx, ph in enumerate(layout.placeholders):
                try:
                    print(f"    占位符 {ph_idx + 1}: {ph.name} (类型: {ph.type})")
                except:
                    pass
    
    return prs.slide_layouts

if __name__ == "__main__":
    master_path = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/ppt模板/ppt母版.pptx"
    layouts = analyze_master(master_path)
    
    print("\n" + "=" * 80)
    print("布局名称索引:")
    print("=" * 80)
    layout_map = {layout.name: idx for idx, layout in enumerate(layouts)}
    for name, idx in sorted(layout_map.items()):
        print(f"  '{name}': 布局索引 {idx}")
