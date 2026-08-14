#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器基础类 V2.1
正确使用母版占位符，保留母版样式，添加高级样式方法
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
import os
from datetime import datetime

# 金蝶品牌色
COLORS = {
    'primary': RGBColor(0x1F, 0x4E, 0x79),
    'kingdee_blue': RGBColor(0x00, 0x6C, 0xB8),
    'accent': RGBColor(0xED, 0x7D, 0x31),
    'success': RGBColor(0x70, 0xAD, 0x47),
    'warning': RGBColor(0xFF, 0xC0, 0x00),
    'danger': RGBColor(0xC0, 0x00, 0x00),
    'text_primary': RGBColor(0x33, 0x33, 0x33),
    'text_secondary': RGBColor(0x66, 0x66, 0x66),
    'text_light': RGBColor(0x99, 0x99, 0x99),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'bg_light': RGBColor(0xF2, 0xF6, 0xFA),
    'bg_card': RGBColor(0xFA, 0xFB, 0xFC),
    'border': RGBColor(0xD0, 0xD8, 0xE0),
}

TEMPLATE_PATH = "/mnt/d/Kingdee文档/自动化交付工具/参考文档（模板）/ppt模板/ppt母版.pptx"


class KingdeePPTBaseV2:
    """金蝶PPT生成器基础类 V2.1"""
    
    def __init__(self, company, doc_type="项目"):
        self.company = company
        self.doc_type = doc_type
        self.C = COLORS
        self.page = 0
        
        if os.path.exists(TEMPLATE_PATH):
            self.prs = Presentation(TEMPLATE_PATH)
            self._clear_slides()
            self.layouts = {}
            for i, layout in enumerate(self.prs.slide_layouts):
                self.layouts[layout.name] = i
            print(f"加载母版成功，布局: {list(self.layouts.keys())}")
        else:
            self.prs = Presentation()
            self.layouts = {}
    
    def _clear_slides(self):
        rIdList = []
        for slide in self.prs.slides:
            rId = self.prs.part.relate_to(slide.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            rIdList.append(rId)
        for rId in rIdList:
            self.prs.part.drop_rel(rId)
        self.prs.slides._sldIdLst.clear()
    
    def _get_layout(self, name):
        if name in self.layouts:
            return self.prs.slide_layouts[self.layouts[name]]
        for key, idx in self.layouts.items():
            if name in key or key in name:
                return self.prs.slide_layouts[idx]
        return self.prs.slide_layouts[0]
    
    def _add_slide(self, layout_name):
        layout = self._get_layout(layout_name)
        slide = self.prs.slides.add_slide(layout)
        self.page += 1
        return slide
    
    def _fill_placeholder(self, slide, idx, text, font_size=None, bold=False, color=None):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == idx:
                shape.text = text
                if font_size:
                    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
                shape.text_frame.paragraphs[0].font.bold = bold
                if color:
                    shape.text_frame.paragraphs[0].font.color.rgb = color
                return True
        return False
    
    # ========== 基础页面方法 ==========
    
    def add_cover(self, title, subtitle=None):
        slide = self._add_slide("封面")
        self._fill_placeholder(slide, 0, title, font_size=36, bold=True, color=COLORS['white'])
        if subtitle:
            self._fill_placeholder(slide, 12, subtitle, font_size=18, color=COLORS['kingdee_blue'])
        self._fill_placeholder(slide, 13, "金蝶软件（中国）有限公司", font_size=14)
        self._fill_placeholder(slide, 14, datetime.now().strftime("%Y年%m月"), font_size=12)
        print(f"[{self.page:3d}] 封面")
        return slide
    
    def add_toc(self, items):
        slide = self._add_slide("目录、提纲")
        self._fill_placeholder(slide, 0, "目 录", font_size=28, bold=True, color=COLORS['primary'])
        for i, item in enumerate(items):
            if i < 16:
                self._fill_placeholder(slide, 13 + i, f"{i+1}. {item}", font_size=14)
        print(f"[{self.page:3d}] 目录")
        return slide
    
    def add_section(self, title, num=None):
        slide = self._add_slide("3_白色内页")
        self._fill_placeholder(slide, 0, f"第{num}章 {title}" if num else title, font_size=32, bold=True, color=COLORS['primary'])
        print(f"[{self.page:3d}] 章节: {title}")
        return slide
    
    def add_content(self, title, content):
        slide = self._add_slide("3_白色内页")
        self._fill_placeholder(slide, 0, title, font_size=24, bold=True, color=COLORS['primary'])
        if content:
            self._fill_placeholder(slide, 12, content, font_size=12)
        print(f"[{self.page:3d}] {title}")
        return slide
    
    def add_thanks(self):
        slide = self._add_slide("封底")
        print(f"[{self.page:3d}] 感谢页")
        return slide
    
    # ========== 高级样式方法 ==========
    
    def add_styled_title(self, slide, title, subtitle=None):
        title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.C['primary']
        title_bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = '微软雅黑'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.C['white']
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(0.3))
            p = sub_box.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.name = '微软雅黑'
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
    
    def add_card(self, slide, x, y, w, h, title, value):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = self.C['bg_card']
        card.line.color.rgb = self.C['border']
        
        value_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(h*0.5))
        p = value_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = '微软雅黑'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.C['primary']
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+h*0.5), Inches(w-0.2), Inches(h*0.4))
        p = label_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = '微软雅黑'
        p.font.size = Pt(11)
        p.font.color.rgb = self.C['text_secondary']
        p.alignment = PP_ALIGN.CENTER
    
    def add_cards_row(self, slide, cards, y=1.5, cols=4, card_w=2.8, card_h=1.2, gap=0.3):
        total_w = cols * card_w + (cols - 1) * gap
        start_x = (13.333 - total_w) / 2
        for i, card in enumerate(cards):
            col = i % cols
            x = start_x + col * (card_w + gap)
            self.add_card(slide, x, y, card_w, card_h, card.get('label', ''), card.get('value', ''))
        return len(cards)
    
    # ========== 装饰元素 ==========
    
    def add_decorative_line(self, slide, x, y, w, color=None, h=0.02):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        line.fill.solid()
        line.fill.fore_color.rgb = color or self.C['kingdee_blue']
        line.line.fill.background()
        return line
    
    def add_decorative_corner(self, slide, corner='top-right', size=0.5):
        if corner == 'top-right':
            shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(13.333-size), Inches(0), Inches(size), Inches(size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.C['kingdee_blue']
            shape.line.fill.background()
        elif corner == 'bottom-left':
            shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(7.5-size), Inches(size), Inches(size))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.C['kingdee_blue']
            shape.line.fill.background()
        return shape
    
    def add_page_number(self, slide, num=None):
        if num is None:
            num = self.page
        box = slide.shapes.add_textbox(Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(10)
        p.font.color.rgb = self.C['text_light']
        p.alignment = PP_ALIGN.RIGHT
        return box
    
    # ========== 时间线样式 ==========
    
    def add_timeline(self, slide, items, y=1.8, item_w=2.5, item_h=1.0):
        n = len(items)
        if n == 0:
            return
        total_w = n * item_w + (n - 1) * 0.3
        start_x = (13.333 - total_w) / 2
        
        # 主线
        line_y = y + item_h / 2
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(line_y), Inches(total_w), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.C['kingdee_blue']
        line.line.fill.background()
        
        for i, item in enumerate(items):
            x = start_x + i * (item_w + 0.3)
            
            # 圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + item_w/2 - 0.1), Inches(line_y - 0.05), Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.C['primary']
            dot.line.color.rgb = self.C['white']
            dot.line.width = Pt(2)
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(item_w), Inches(0.4))
            p = title_box.text_frame.paragraphs[0]
            p.text = item.get('phase', f'阶段{i+1}')
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.C['primary']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            desc_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(item_w), Inches(item_h - 0.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get('desc', '')
            p.font.size = Pt(10)
            p.font.color.rgb = self.C['text_secondary']
            p.alignment = PP_ALIGN.CENTER
    
    # ========== 图标卡片 ==========
    
    def add_icon_card(self, slide, x, y, w, h, icon, title, desc=None, color=None):
        color = color or self.C['primary']
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = self.C['white']
        card.line.color.rgb = self.C['border']
        card.line.width = Pt(1)
        
        # 图标背景圆
        icon_size = min(w, h) * 0.35
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + (w-icon_size)/2), Inches(y + 0.15), Inches(icon_size), Inches(icon_size))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = color
        icon_bg.line.fill.background()
        
        # 图标文字
        icon_box = slide.shapes.add_textbox(Inches(x + (w-icon_size)/2), Inches(y + 0.15 + icon_size*0.25), Inches(icon_size), Inches(icon_size*0.5))
        p = icon_box.text_frame.paragraphs[0]
        p.text = icon
        p.font.size = Pt(int(icon_size * 20))
        p.font.color.rgb = self.C['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15 + icon_size + 0.1), Inches(w), Inches(0.4))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.C['text_primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        if desc:
            desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15 + icon_size + 0.45), Inches(w - 0.2), Inches(h - 0.15 - icon_size - 0.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(9)
            p.font.color.rgb = self.C['text_secondary']
            p.alignment = PP_ALIGN.CENTER
    
    def add_icon_cards_row(self, slide, cards, y=1.5, cols=4, card_w=2.8, card_h=1.8, gap=0.3, colors=None):
        total_w = cols * card_w + (cols - 1) * gap
        start_x = (13.333 - total_w) / 2
        default_colors = [self.C['primary'], self.C['kingdee_blue'], self.C['accent'], self.C['success']]
        
        for i, card in enumerate(cards):
            col = i % cols
            x = start_x + col * (card_w + gap)
            color = (colors[i] if colors and i < len(colors) else default_colors[i % len(default_colors)])
            self.add_icon_card(slide, x, y, card_w, card_h, 
                             card.get('icon', '●'), card.get('title', ''), card.get('desc', ''), color)
        return len(cards)
    
    # ========== 数据展示 ==========
    
    def add_kpi_row(self, slide, kpis, y=1.2):
        cards = []
        for kpi in kpis:
            cards.append({
                'value': kpi.get('value', '-'),
                'label': kpi.get('label', '')
            })
        return self.add_cards_row(slide, cards, y=y, cols=len(kpis), card_w=2.5, card_h=1.0)
    
    def save(self, filename=None):
        if not filename:
            filename = f"{self.company}_{self.doc_type}_{datetime.now().strftime('%Y%m%d')}.pptx"
        output_dir = os.path.expanduser("~/.openclaw/workspace/output")
        os.makedirs(output_dir, exist_ok=True)
        filepath = os.path.join(output_dir, filename)
        self.prs.save(filepath)
        size = os.path.getsize(filepath) / 1024
        print(f"\n保存成功: {filepath}")
        print(f"文件大小: {size:.1f} KB")
        print(f"总页数: {self.page} 页")
        return {"success": True, "filename": filename, "output": filepath, "slides": self.page, "size": size * 1024}


if __name__ == "__main__":
    ppt = KingdeePPTBaseV2("测试公司", "售前方案")
    ppt.add_cover("测试公司\n数字化管理系统项目", "售前方案")
    ppt.add_toc(["公司介绍", "产品体系", "解决方案", "实施计划"])
    ppt.add_section("公司介绍", "一")
    ppt.add_content("公司简介", "测试公司是一家领先的...")
    ppt.add_thanks()
    result = ppt.save()
    print(result)
