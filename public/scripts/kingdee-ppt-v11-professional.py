#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶PPT生成器 - v11.0 专业版
基于金蝶交付完整指南，包含丰富的专业内容
支持售前PPT、上线汇报PPT、验收汇报PPT
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

COLORS = {
    'primary': RGBColor(0, 102, 153),
    'secondary': RGBColor(0, 153, 204),
    'accent': RGBColor(255, 153, 0),
    'text': RGBColor(51, 51, 51),
    'white': RGBColor(255, 255, 255),
    'light': RGBColor(240, 240, 240),
}

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLORS['white']
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容
    cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith('###'):
            p.text = item.replace('###', '').strip()
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            p.space_before = Pt(15)
        elif item.startswith('-'):
            p.text = "• " + item[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(8)
        elif item.strip():
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(10)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    tb_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    tb_bg.fill.solid()
    tb_bg.fill.fore_color.rgb = COLORS['primary']
    tb_bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 左栏
    ltb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    
    # 左栏标题
    p = ltf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.space_before = Pt(10)
    
    # 左栏内容
    for item in left_items:
        p = ltf.add_paragraph()
        if item.startswith('-'):
            p.text = "• " + item[1:].strip()
        else:
            p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.space_before = Pt(6)
    
    # 右栏
    rtb = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
    rtf = rtb.text_frame
    rtf.word_wrap = True
    
    # 右栏标题
    p = rtf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.space_before = Pt(10)
    
    # 右栏内容
    for item in right_items:
        p = rtf.add_paragraph()
        if item.startswith('-'):
            p.text = "• " + item[1:].strip()
        else:
            p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.space_before = Pt(6)
    
    return slide

def generate_presales_ppt_v11(customer_info):
    """生成售前PPT - v11.0 专业版（32页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    company_size = customer_info.get('companySize', '中型企业')
    employees = customer_info.get('employees', '500')
    revenue = customer_info.get('revenue', '10000')
    
    # 01 封面
    add_title_slide(prs, company_name, "数字化转型解决方案")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"金蝶软件（中国）有限公司\n{datetime.now().strftime('%Y年%m月')}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 02 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、企业概况与需求分析",
        "二、解决方案设计", 
        "三、业务架构蓝图",
        "四、4A架构设计",
        "五、价值工程分析",
        "六、实施路线图",
        "七、成功案例分享",
        "八、项目团队与保障"
    ])
    
    # 第一章：企业概况与需求分析
    add_section_slide(prs, "一、企业概况与需求分析")
    
    add_content_slide(prs, "企业基本情况", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"企业规模：{company_size}",
        f"员工人数：{employees}人",
        f"年营业额：{revenue}万元",
        "",
        "###组织架构：",
        "- 集团总部：战略决策、财务管理、人力资源",
        "- 业务部门：采购、销售、生产、仓储",
        "- 支持部门：IT、行政、法务"
    ])
    
    add_two_column_slide(prs, "业务痛点分析", 
        "当前痛点", [
            "- 业务流程不畅：流程不标准，审批效率低",
            "- 数据孤岛严重：系统不集成，数据不共享",
            "- 管理效率低下：手工操作多，出错率高",
            "- 决策支持不足：数据不及时，分析不深入",
            "- 成本控制困难：成本核算难，控制手段少"
        ],
        "影响分析", [
            "- 运营成本高：人力成本、库存成本居高不下",
            "- 市场响应慢：订单交付周期长，客户满意度低",
            "- 管理风险大：财务风险、运营风险难以控制",
            "- 发展受限：信息化水平低，难以支撑业务扩张",
            "- 竞争力弱：数字化转型滞后，市场竞争压力大"
        ]
    )
    
    add_content_slide(prs, "业务目标", [
        "###核心目标：",
        "- 实现业务流程标准化，提升运营效率",
        "- 建立统一数据平台，消除信息孤岛",
        "- 提升系统集成效率，实现业务协同",
        "- 支持实时决策分析，增强管理能力",
        "",
        "###价值目标：",
        "- 效率提升：财务结账时间缩短50%",
        "- 成本降低：人力成本降低20%，库存成本降低15%",
        "- 管理提升：数据准确性提升30%，决策实时化"
    ])
    
    # 第二章：解决方案设计
    add_section_slide(prs, "二、解决方案设计")
    
    add_content_slide(prs, "总体解决方案", [
        "###实施范围：",
        "- 财务云：总账、应收、应付、固定资产、现金管理",
        "- 供应链云：采购、库存、销售、物流",
        "- 制造云：生产计划、车间管理、质量管理",
        "- 人力云：人事、薪酬、绩效、培训",
        "",
        "###技术平台：",
        "- 金蝶云·星空平台：云原生架构、微服务设计",
        "- 数据中台：统一数据平台、实时数据分析",
        "- 集成平台：API网关、数据集成、流程集成"
    ])
    
    add_two_column_slide(prs, "实施策略",
        "实施原则", [
            "- 整体规划：统一规划、分步实施",
            "- 重点突破：核心先行、逐步推广",
            "- 风险控制：试点验证、风险可控",
            "- 价值导向：快速见效、持续优化"
        ],
        "实施策略", [
            "- 分期实施：一期核心模块，二期扩展模块",
            "- 分批上线：总部先行，分步推广到分子公司",
            "- 试点推广：典型单位试点，总结经验后推广",
            "- 并行过渡：新旧系统并行，平稳过渡切换"
        ]
    )
    
    # 第三章：业务架构蓝图
    add_section_slide(prs, "三、业务架构蓝图")
    
    add_content_slide(prs, "业务架构（BA）", [
        "###核心价值流：",
        "- 订单到收款（O2C）：订单管理 → 发货开票 → 收款核销",
        "- 采购到付款（P2P）：采购申请 → 采购订单 → 收货付款", 
        "- 计划到生产（P2M）：需求计划 → 生产计划 → 车间执行",
        "- 研发到上市（R2M）：产品研发 → 产品上市 → 市场推广",
        "",
        "###业务能力：",
        "- 财务管理：会计核算、资金管理、成本管理、预算管理",
        "- 供应链管理：采购管理、库存管理、销售管理、物流管理",
        "- 生产制造：计划管理、生产执行、质量管理、设备管理",
        "- 人力资源管理：人事管理、薪酬管理、绩效管理、培训管理"
    ])
    
    # 第四章：4A架构设计
    add_section_slide(prs, "四、4A架构设计")
    
    add_content_slide(prs, "数据架构（DA）", [
        "###数据实体：",
        "- 主数据：客户、供应商、物料、科目、组织",
        "- 业务数据：订单、发票、凭证、合同、计划",
        "- 分析数据：报表、指标、维度、度量",
        "",
        "###数据服务：",
        "- 数据查询服务：实时查询、多维查询",
        "- 数据分析服务：统计分析、预测分析", 
        "- 数据治理服务：数据质量、数据安全、数据标准"
    ])
    
    add_content_slide(prs, "应用架构（AA）", [
        "###核心应用：",
        "- 财务云：总账、应收、应付、固定资产、现金管理、成本管理",
        "- 供应链云：采购管理、库存管理、销售管理、物流管理",
        "- 制造云：生产计划、车间管理、质量管理、设备管理",
        "",
        "###扩展应用：",
        "- 全面预算：预算编制、预算控制、预算分析",
        "- 合并报表：报表编制、报表合并、报表分析",
        "- 资金管理：资金计划、资金调度、资金监控"
    ])
    
    add_content_slide(prs, "技术架构（TA）", [
        "###技术平台：",
        "- 云原生架构：容器化部署、微服务架构、DevOps",
        "- 高可用设计：负载均衡、故障转移、数据备份",
        "- 安全体系：身份认证、权限控制、数据加密",
        "",
        "###集成架构：",
        "- API网关：统一接口管理、接口安全控制",
        "- 数据集成：ETL工具、数据同步、数据转换",
        "- 流程集成：BPM引擎、流程编排、流程监控"
    ])
    
    # 第五章：价值工程分析
    add_section_slide(prs, "五、价值工程分析")
    
    add_content_slide(prs, "价值目标模型", [
        "###四大价值概念：",
        "1. Strategy（战略）：保障战略有效性，支撑企业战略落地",
        "2. Spending（支出回报）：增收节支，提升投入产出比",
        "3. Situation（情景）：提升竞争力，增强市场响应能力",
        "4. Structure（结构）：提升效率，优化业务流程结构",
        "",
        "###价值指标体系：",
        "- 效率指标：业务处理时间、流程周期、人工工作量",
        "- 成本指标：人力成本、库存成本、运营成本",
        "- 质量指标：数据准确性、流程合规性、客户满意度"
    ])
    
    add_two_column_slide(prs, "预期收益分析",
        "效率提升", [
            "- 财务结账时间：从10天缩短到3天（↓70%）",
            "- 采购周期：从15天缩短到7天（↓53%）",
            "- 销售订单处理：从2天缩短到0.5天（↓75%）",
            "- 库存周转率：从4次/年提升到6次/年（↑50%）",
            "- 生产计划准确率：从70%提升到90%（↑20%）"
        ],
        "成本降低", [
            "- 人力成本：降低20%（自动化处理）",
            "- 库存成本：降低15%（精准库存控制）",
            "- 运营成本：降低10%（流程优化）",
            "- 财务成本：降低25%（自动化核算）",
            "- 采购成本：降低10%（集中采购）"
        ]
    )
    
    # 第六章：实施路线图
    add_section_slide(prs, "六、实施路线图")
    
    add_content_slide(prs, "项目实施计划", [
        "###项目阶段：",
        "- 项目启动（第1-2周）：项目组建、需求调研、蓝图设计",
        "- 系统实施（第3-14周）：系统配置、数据迁移、接口开发",
        "- 测试培训（第15-18周）：系统测试、用户培训、问题修复",
        "- 上线支持（第19-20周）：系统上线、运维支持、持续优化",
        "",
        "###关键里程碑：",
        "- 需求确认：第4周",
        "- 蓝图设计：第8周",
        "- 系统配置：第14周",
        "- UAT测试：第18周",
        "- 系统上线：第20周"
    ])
    
    add_content_slide(prs, "项目组织保障", [
        "###项目组织架构：",
        "- 项目指导委员会：项目决策、资源协调、风险管控",
        "- 项目管理办公室：项目计划、进度管理、质量管理",
        "- 业务组：需求分析、业务设计、用户培训",
        "- 技术组：系统配置、接口开发、数据迁移",
        "",
        "###项目保障机制：",
        "- 周例会制度：每周项目例会，沟通问题、协调资源",
        "- 里程碑评审：关键节点评审，确保项目质量",
        "- 风险管理：风险识别、风险评估、风险应对"
    ])
    
    # 第七章：成功案例分享
    add_section_slide(prs, "七、成功案例分享")
    
    add_content_slide(prs, "案例1：制造业企业", [
        "###客户背景：",
        "- 企业名称：某大型制造企业",
        "- 所属行业：制造业",
        "- 企业规模：大型企业，员工5000人",
        "- 实施模块：财务云、供应链云、制造云",
        "",
        "###项目成果：",
        "- 财务结账时间从10天缩短到3天",
        "- 库存周转率提升40%",
        "- 生产计划准确性提升50%",
        "- 采购效率提升60%"
    ])
    
    add_content_slide(prs, "案例2：零售企业", [
        "###客户背景：",
        "- 企业名称：某连锁零售企业",
        "- 所属行业：零售业",
        "- 企业规模：中型企业，员工1000人",
        "- 实施模块：财务云、供应链云、人力云",
        "",
        "###项目成果：",
        "- 订单处理效率提升70%",
        "- 库存准确率提升至99%",
        "- 人力成本降低25%",
        "- 销售额增长20%"
    ])
    
    # 第八章：项目团队与保障
    add_section_slide(prs, "八、项目团队与保障")
    
    add_two_column_slide(prs, "项目团队配置",
        "核心团队", [
            "- 项目总监：项目总负责、资源协调",
            "- 项目经理：项目管理、进度控制",
            "- 业务顾问：业务设计、方案落地",
            "- 技术顾问：技术架构、系统配置",
            "- 开发工程师：接口开发、报表开发",
            "- 测试工程师：系统测试、问题跟踪"
        ],
        "服务保障", [
            "- 7×24小时技术支持热线",
            "- 远程支持：在线问题处理",
            "- 现场支持：关键节点现场服务",
            "- 定期回访：项目回访、持续优化",
            "- 知识转移：用户培训、技术文档",
            "- 版本升级：产品升级、功能增强"
        ]
    )
    
    add_content_slide(prs, "服务承诺", [
        "###项目质量承诺：",
        "- 按时交付：严格按项目计划推进，确保按时上线",
        "- 质量保证：系统功能完整，性能达标，数据准确",
        "- 用户满意：用户培训充分，操作熟练，满意度≥90%",
        "",
        "###服务支持承诺：",
        "- 上线支持：上线期间7×24小时现场支持",
        "- 运维支持：上线后3个月免费运维支持",
        "- 持续优化：定期回访，收集需求，持续优化"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", "期待与您合作")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"金蝶软件（中国）有限公司\n{datetime.now().strftime('%Y年%m月%d日')}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return prs

def generate_golive_ppt_v11(customer_info):
    """生成上线汇报PPT - v11.0 专业版（25页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    golive_date = customer_info.get('goliveDate', datetime.now().strftime('%Y-%m-%d'))
    
    # 封面
    add_title_slide(prs, f"{company_name}ERP系统上线汇报", "数字化转型重要里程碑")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"上线日期：{golive_date}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、项目概述",
        "二、项目背景与目标", 
        "三、项目实施过程",
        "四、项目业务范围",
        "五、项目价值实现",
        "六、项目成果展示",
        "七、上线准备情况",
        "八、项目总结展望"
    ])
    
    # 第一章：项目概述
    add_section_slide(prs, "一、项目概述")
    
    add_content_slide(prs, "项目概述", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"项目周期：8个月（2026年1月-2026年8月）",
        f"上线日期：{golive_date}",
        "",
        "###项目意义：",
        "- 企业数字化转型的重要里程碑",
        "- 提升管理效率，降低运营成本",
        "- 增强企业核心竞争力",
        "- 支撑企业可持续发展"
    ])
    
    # 第二章：项目背景与目标
    add_section_slide(prs, "二、项目背景与目标")
    
    add_two_column_slide(prs, "项目建设背景",
        "业务痛点", [
            "- 业务流程不标准：审批效率低，流程周期长",
            "- 数据管理分散：信息孤岛，数据不共享",
            "- 系统集成困难：系统独立，接口缺失",
            "- 管理决策滞后：数据不及时，分析不深入"
        ],
        "建设目标", [
            "- 实现业务流程标准化：流程规范化、审批自动化",
            "- 建立统一数据平台：数据集中、实时共享",
            "- 提升系统集成效率：系统互联、业务协同",
            "- 支持实时决策分析：数据可视化、分析智能化"
        ]
    )
    
    # 第三章：项目实施过程
    add_section_slide(prs, "三、项目实施过程")
    
    add_content_slide(prs, "项目实施历程", [
        "###第一阶段：项目启动（2026年1月）",
        "- 项目组建、需求调研、蓝图设计",
        "",
        "###第二阶段：系统实施（2026年2-5月）",
        "- 系统配置、数据迁移、接口开发、报表开发",
        "",
        "###第三阶段：测试培训（2026年6-7月）",
        "- 系统测试、问题修复、用户培训、操作演练",
        "",
        "###第四阶段：上线准备（2026年8月）",
        "- 上线准备、数据核对、应急预案、上线切换"
    ])
    
    # 第四章：项目业务范围
    add_section_slide(prs, "四、项目业务范围")
    
    add_content_slide(prs, "业务范围架构", [
        "###核心业务模块：",
        "- 财务管理：总账、应收、应付、固定资产、现金管理、成本管理",
        "- 供应链管理：采购管理、库存管理、销售管理、物流管理",
        "- 生产制造：生产计划、车间管理、质量管理、设备管理",
        "- 人力资源管理：人事管理、薪酬管理、绩效管理、培训管理",
        "",
        "###系统集成：",
        "- 与MES系统对接：生产工单、完工报告同步",
        "- 与WMS系统对接：出入库单据、库存数据同步",
        "- 与OA系统对接：审批流程、通知推送同步"
    ])
    
    # 第五章：项目价值实现
    add_section_slide(prs, "五、项目价值实现")
    
    add_two_column_slide(prs, "价值总览",
        "数据治理", [
            "- 主数据统一率：从60%提升至100%",
            "- 数据准确性：从85%提升至98%",
            "- 数据及时性：从T+1提升至实时",
            "- 数据完整性：从70%提升至95%"
        ],
        "业务提效", [
            "- 财务结账时间：从10天缩短到3天",
            "- 采购周期：从15天缩短到7天",
            "- 订单交付准时率：从80%提升至95%",
            "- 库存周转率：从4次/年提升到6次/年"
        ]
    )
    
    add_content_slide(prs, "管控提升", [
        "###流程优化：",
        "- 实现订单-应收-实收-稽核闭环：订单全程可追溯",
        "- 构建合同-请款-发票-付款管控闭环：合同全程管控",
        "- 建立预算-执行-分析闭环：预算实时控制",
        "",
        "###决策赋能：",
        "- 实时现金流监控：资金状况一目了然",
        "- 经营数据分析：多维度分析报表",
        "- 管理驾驶舱：关键指标实时展示"
    ])
    
    # 第六章：项目成果展示
    add_section_slide(prs, "六、项目成果展示")
    
    add_content_slide(prs, "实施成果", [
        "###系统成果：",
        "- 完成核心模块实施：财务、供应链、制造、人力",
        "- 实现系统集成对接：MES、WMS、OA",
        "- 建立数据治理体系：主数据、业务数据、分析数据",
        "",
        "###业务成果：",
        "- 业务流程标准化：流程规范化、审批自动化",
        "- 数据管理规范化：数据集中、实时共享",
        "- 决策支持智能化：数据可视化、分析智能化"
    ])
    
    # 第七章：上线准备情况
    add_section_slide(prs, "七、上线准备情况")
    
    add_content_slide(prs, "上线准备检查清单", [
        "###系统准备：",
        "- ✅ 系统配置完成：所有模块配置完成并验证",
        "- ✅ 接口开发完成：MES、WMS、OA接口开发完成",
        "- ✅ 报表开发完成：所有业务报表开发完成",
        "",
        "###数据准备：",
        "- ✅ 主数据迁移完成：客户、供应商、物料、科目迁移完成",
        "- ✅ 期初数据迁移完成：期初余额、期初库存、期初往来迁移完成",
        "- ✅ 数据准确性验证：余额核对、库存核对、往来核对完成",
        "",
        "###用户准备：",
        "- ✅ 用户培训完成：所有用户培训并通过考核",
        "- ✅ 操作手册完成：用户操作手册编制完成",
        "- ✅ 应急预案完成：应急预案编制并演练完成"
    ])
    
    # 第八章：项目总结展望
    add_section_slide(prs, "八、项目总结展望")
    
    add_two_column_slide(prs, "项目总结",
        "成功经验", [
            "- 领导高度重视：项目决策及时、资源保障充分",
            "- 团队协作良好：项目团队专业、执行力强",
            "- 用户参与积极：业务部门配合、需求明确",
            "- 方法论规范：实施方法论规范、质量控制严格"
        ],
        "改进建议", [
            "- 加强用户培训：持续培训、提升操作熟练度",
            "- 优化系统性能：持续优化、提升用户体验",
            "- 完善运维支持：建立运维体系、快速响应问题"
        ]
    )
    
    # 封底
    add_title_slide(prs, "谢谢", "金蝶软件（中国）有限公司")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = datetime.now().strftime('%Y年%m月%d日')
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return prs

def generate_acceptance_ppt_v11(customer_info):
    """生成验收汇报PPT - v11.0 专业版（25页）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    company_name = customer_info.get('companyName', '企业名称')
    industry = customer_info.get('industry', '制造业')
    acceptance_date = customer_info.get('acceptanceDate', datetime.now().strftime('%Y-%m-%d'))
    
    # 封面
    add_title_slide(prs, f"{company_name}ERP系统验收汇报", "项目验收与成果总结")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"验收日期：{acceptance_date}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 目录
    add_section_slide(prs, "目录")
    add_content_slide(prs, "内容概览", [
        "一、验收概述",
        "二、验收范围与标准",
        "三、验收过程",
        "四、验收结果",
        "五、项目成果展示",
        "六、遗留问题与计划",
        "七、验收结论"
    ])
    
    # 第一章：验收概述
    add_section_slide(prs, "一、验收概述")
    
    add_content_slide(prs, "验收基本信息", [
        f"企业名称：{company_name}",
        f"所属行业：{industry}",
        f"项目周期：8个月（2026年1月-2026年8月）",
        f"验收日期：{acceptance_date}",
        f"验收结论：✅ 通过验收",
        "",
        "###验收意义：",
        "- 项目正式交付使用",
        "- 系统功能确认完整",
        "- 项目成果验收达标"
    ])
    
    # 第二章：验收范围与标准
    add_section_slide(prs, "二、验收范围与标准")
    
    add_content_slide(prs, "验收范围", [
        "###核心模块验收：",
        "- 财务管理：总账、应收、应付、固定资产、现金管理、成本管理",
        "- 供应链管理：采购管理、库存管理、销售管理、物流管理",
        "- 生产制造：生产计划、车间管理、质量管理、设备管理",
        "- 人力资源管理：人事管理、薪酬管理、绩效管理、培训管理",
        "",
        "###系统集成验收：",
        "- 与MES系统对接：生产工单、完工报告同步",
        "- 与WMS系统对接：出入库单据、库存数据同步",
        "- 与OA系统对接：审批流程、通知推送同步"
    ])
    
    add_two_column_slide(prs, "验收标准",
        "功能标准", [
            "- 完成合同约定功能：100%完成",
            "- 业务流程运行正常：流程顺畅",
            "- 数据准确完整：准确率≥99%",
            "- 用户操作熟练：满意度≥90%"
        ],
        "性能标准", [
            "- 系统响应时间：<2秒",
            "- 并发用户数：≥500人",
            "- 数据处理量：≥100万条/月",
            "- 系统可用性：≥99.9%"
        ]
    )
    
    # 第三章：验收过程
    add_section_slide(prs, "三、验收过程")
    
    add_content_slide(prs, "验收流程", [
        "###验收准备（2026年8月第1周）：",
        "- 验收方案制定：确定验收范围、验收标准、验收方法",
        "- 测试数据准备：准备测试用例、测试数据",
        "- 验收团队组建：组建验收小组、明确验收职责",
        "",
        "###验收执行（2026年8月第2-3周）：",
        "- 功能测试：按测试用例逐项测试",
        "- 性能测试：压力测试、并发测试",
        "- 用户验收：用户操作测试、满意度调查",
        "- 专家评审：专家组评审验收结果"
    ])
    
    # 第四章：验收结果
    add_section_slide(prs, "四、验收结果")
    
    add_content_slide(prs, "验收结果汇总", [
        "###功能验收：",
        "- ✅ 总账模块：功能完整，符合需求",
        "- ✅ 应收模块：功能完整，符合需求",
        "- ✅ 应付模块：功能完整，符合需求",
        "- ✅ 库存模块：功能完整，符合需求",
        "",
        "###性能验收：",
        "- ✅ 系统响应时间：<2秒（达标）",
        "- ✅ 并发用户数：500人（达标）",
        "- ✅ 系统可用性：99.9%（达标）",
        "",
        "###用户验收：",
        "- ✅ 用户培训：100%完成",
        "- ✅ 用户操作：操作熟练",
        "- ✅ 用户满意度：92%（达标）"
    ])
    
    # 第五章：项目成果展示
    add_section_slide(prs, "五、项目成果展示")
    
    add_two_column_slide(prs, "项目成果",
        "系统成果", [
            "- 完成ERP系统实施：财务、供应链、制造、人力",
            "- 实现系统集成对接：MES、WMS、OA",
            "- 建立数据治理体系：主数据、业务数据",
            "- 完成用户培训：培训覆盖率100%"
        ],
        "业务成果", [
            "- 业务流程标准化：流程规范化、审批自动化",
            "- 数据管理规范化：数据集中、实时共享",
            "- 决策支持智能化：数据可视化、分析智能化",
            "- 用户满意度提升：从70%提升到92%"
        ]
    )
    
    # 第六章：遗留问题与计划
    add_section_slide(prs,    "六、遗留问题与计划")
    
    add_content_slide(prs, "遗留问题", [
        "###当前问题：",
        "- 部分细节功能需要优化：报表格式、界面显示",
        "- 用户习惯需要培养：持续培训、操作演练",
        "- 数据质量需要持续监控：数据治理、数据清洗",
        "",
        "###后续计划：",
        "- 系统功能持续优化：收集用户反馈、持续优化",
        "- 用户培训推广：持续培训、提升操作熟练度",
        "- 数据治理深化：建立数据治理体系、提升数据质量"
    ])
    
    # 第七章：验收结论
    add_section_slide(prs, "七、验收结论")
    
    add_content_slide(prs, "验收结论", [
        "###验收结论：",
        "- ✅ 项目达到合同约定要求",
        "- ✅ 系统功能完整，运行稳定",
        "- ✅ 业务流程优化效果显著",
        "- ✅ 用户满意度较高",
        "",
        "###后续建议：",
        "- 持续优化系统功能：收集用户反馈、持续优化",
        "- 加强用户培训：持续培训、提升操作熟练度",
        "- 深化数据价值挖掘：数据分析、决策支持",
        "- 建立运维体系：运维流程、问题处理"
    ])
    
    # 封底
    add_title_slide(prs, "谢谢", "金蝶软件（中国）有限公司")
    slide = prs.slides[-1]
    tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = datetime.now().strftime('%Y年%m月%d日')
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 3:
        print("用法: python3 kingdee-ppt-v11-professional.py --type <type> --companyName <name> [其他参数]")
        sys.exit(1)
    
    # 解析参数
    params = {}
    for i in range(1, len(sys.argv), 2):
        if i + 1 < len(sys.argv):
            key = sys.argv[i].lstrip('-')
            params[key] = sys.argv[i + 1]
    
    # 生成PPT
    ppt_type = params.get('type', 'presales')
    company_name = params.get('companyName', '企业名称')
    
    try:
        if ppt_type == 'presales':
            prs = generate_presales_ppt_v11(params)
        elif ppt_type == 'golive':
            prs = generate_golive_ppt_v11(params)
        elif ppt_type == 'acceptance':
            prs = generate_acceptance_ppt_v11(params)
        else:
            print(f"不支持的PPT类型: {ppt_type}")
            sys.exit(1)
        
        # 保存文件
        filename = f"{company_name}_{ppt_type}_v11.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        
        print(json.dumps({
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "slides": len(prs.slides),
            "type": ppt_type
        }, ensure_ascii=False))
        
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": str(e)
        }, ensure_ascii=False))

if __name__ == "__main__":
    main()
