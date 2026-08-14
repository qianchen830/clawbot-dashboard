#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金蝶售前PPT生成器 v5.0
基于4A架构与价值工程方法
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import sys
import os

# 颜色定义
COLORS = {
    'primary': RGBColor(26, 35, 126),      # 深蓝 #1a237e
    'secondary': RGBColor(13, 71, 161),     # 蓝色 #0d47a1
    'accent': RGBColor(255, 152, 0),        # 橙色 #ff9800
    'text': RGBColor(51, 51, 51),           # 深灰 #333333
    'light': RGBColor(245, 245, 245),       # 浅灰 #f5f5f5
    'white': RGBColor(255, 255, 255),       # 白色
}

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9比例
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景形状
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, number=""):
    """添加章节标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景形状
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['secondary']
    shape.line.fill.background()
    
    # 章节编号
    if number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(2), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(3), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    return slide

def add_content_slide(prs, title, content_items):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    # 标题文本
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 处理不同层级
        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            p.level = item.get('level', 0)
        else:
            p.text = "• " + item
            p.font.size = Pt(18)
        
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.fill.background()
    
    # 标题文本
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 表格
    cols = len(headers)
    row_count = len(rows) + 1  # 包含表头
    
    table = slide.shapes.add_table(
        row_count, cols,
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(0.5 * row_count)
    ).table
    
    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['secondary']
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # 设置数据行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['text']
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_presales_ppt(data):
    """生成售前PPT"""
    prs = create_presentation()
    
    # 转换数据类型
    months = int(data.get('implMonths', 6) or 6)
    employees = data.get('employees', '') or ''
    revenue = data.get('revenue', '') or ''
    budget = data.get('budget', '') or ''
    user_count = data.get('userCount', '') or ''
    
    # 1. 封面
    add_title_slide(
        prs,
        f"{data.get('companyName', '客户')}金蝶云解决方案",
        "金蝶软件（中国）有限公司"
    )
    
    # 2. 目录
    add_content_slide(prs, "目录", [
        "一、企业概况",
        "二、痛点分析",
        "三、解决方案",
        "四、4A架构设计",
        "五、业务架构蓝图",
        "六、价值工程",
        "七、实施路线",
        "八、成功案例"
    ])
    
    # 3. 企业概况
    add_section_slide(prs, "企业概况", "01")
    add_content_slide(prs, "基本信息", [
        f"企业名称：{data.get('companyName', '待确认') or '待确认'}",
        f"所属行业：{data.get('industry', '待确认') or '待确认'}",
        f"企业规模：{data.get('companySize', '待确认') or '待确认'}",
        f"员工人数：{employees}人" if employees else "员工人数：待确认",
        f"年营业额：{revenue}万元" if revenue else "年营业额：待确认",
        f"管控模式：{data.get('controlMode', '待确认') or '待确认'}"
    ])
    
    # 4. 痛点分析
    add_section_slide(prs, "痛点分析", "02")
    pain_points = data.get('painPoints', '企业当前面临的主要痛点').split('\n')
    add_content_slide(prs, "当前痛点", [p.strip() for p in pain_points if p.strip()])
    
    # 5. 解决方案
    add_section_slide(prs, "解决方案", "03")
    modules = data.get('modules', ['finance']) or ['finance']
    if isinstance(modules, str):
        modules = [m.strip() for m in modules.split(',') if m.strip()]
    module_names = {'finance': '财务云', 'supply': '供应链云', 'manufacturing': '制造云', 
                    'hr': '人力云', 'plm': 'PLM云', 'pm': '项目管理'}
    module_list = [module_names.get(m, m) for m in modules]
    
    add_content_slide(prs, "实施范围", [
        f"实施模块：{'、'.join(module_list)}",
        f"用户规模：{user_count}用户" if user_count else "用户规模：待确认",
        f"实施周期：{months}个月",
        f"项目预算：{budget}万元" if budget else "项目预算：待确认"
    ])
    
    # 6. 4A架构设计
    add_section_slide(prs, "4A架构设计", "04")
    
    # 6.1 业务架构
    add_content_slide(prs, "业务架构（BA）", [
        {"text": "业务架构定义", "size": 22, "bold": True},
        "业务架构是业务的结构化表达，描述组织如何运用业务的关键要素来实现其战略意图和目标",
        {"text": "", "size": 12},
        {"text": "核心要素", "size": 22, "bold": True},
        "价值流（Value Stream）：端到端价值创造活动",
        "业务能力（Capability）：企业做什么",
        "业务流程（Process）：如何做",
        "业务对象（Business Object）：信息载体",
        "组织架构（Organization）：谁来执行"
    ])
    
    # 6.2 数据架构
    add_content_slide(prs, "数据架构（DA）", [
        {"text": "数据架构定义", "size": 22, "bold": True},
        "以结构化的方式描述在业务运作和管理决策中所需要的各类信息及其关系的一套整体组件规范",
        {"text": "", "size": 12},
        {"text": "核心要素", "size": 22, "bold": True},
        "数据实体：业务对象的数据定义",
        "数据服务：数据的采集、存储、处理、分析服务",
        "数据治理：数据标准、质量、安全、生命周期管理",
        "主数据管理（MDM）：客户、供应商、物料、会计科目主数据",
        "数据分析平台（BI）：报表、分析、决策支持"
    ])
    
    # 6.3 应用架构
    add_content_slide(prs, "应用架构（AA）", [
        {"text": "应用架构定义", "size": 22, "bold": True},
        "描述了各种用于支持业务架构并对数据架构所定义的各种数据进行处理的应用功能",
        {"text": "", "size": 12},
        {"text": "核心要素", "size": 22, "bold": True},
        "应用系统：核心ERP、周边系统、第三方系统",
        "应用模块：功能模块划分",
        "应用服务：服务接口、API",
        "集成方案：系统间集成",
        {"text": "", "size": 12},
        {"text": f"实施模块：{'、'.join(module_list)}", "size": 18}
    ])
    
    # 6.4 技术架构
    add_content_slide(prs, "技术架构（TA）", [
        {"text": "技术架构定义", "size": 22, "bold": True},
        "代表了各种可以从市场或组织内部获得的软件和硬件组件",
        {"text": "", "size": 12},
        {"text": "核心要素", "size": 22, "bold": True},
        "基础设施：服务器、存储、网络",
        "云平台：公有云、私有云、混合云",
        "中间件：应用服务器、数据库、消息队列",
        "安全技术：身份认证、权限管理、数据加密、审计日志",
        "运维监控：日志管理、性能监控、告警管理"
    ])
    
    # 7. 业务架构蓝图
    add_section_slide(prs, "业务架构蓝图", "05")
    
    add_content_slide(prs, "业务架构蓝图六种类型", [
        {"text": "一级蓝图", "size": 20, "bold": True},
        "1. 业务能力地图（L1-L3层级）",
        "2. 价值流与业务能力映射图",
        "3. 业务流程架构图",
        {"text": "", "size": 12},
        {"text": "二级蓝图", "size": 20, "bold": True},
        "4. 流程串接图",
        "5. 价值流业务场景图",
        "6. 业务流程图（BPMN）"
    ])
    
    add_content_slide(prs, "典型端到端流程", [
        "订单到现金（O2C）：销售订单 → 发货 → 开票 → 收款 → 核销",
        "采购到付款（P2P）：采购申请 → 订单 → 收货 → 发票 → 付款",
        "计划到制造（P2M）：需求计划 → 生产计划 → 生产执行 → 质检入库",
        "招聘到退休（H2R）：招聘 → 入职 → 薪酬 → 绩效 → 离职"
    ])
    
    # 8. 价值工程
    add_section_slide(prs, "价值工程", "06")
    
    add_table_slide(prs, "价值目标模型", 
        ["价值概念", "含义", "典型指标"],
        [
            ["Strategy（战略）", "保障战略有效性", "市场份额、品牌价值"],
            ["Spending（支出回报）", "增收节支", "采购成本、制造成本"],
            ["Situation（情景）", "提升竞争力", "客户满意度、交付周期"],
            ["Structure（结构）", "提升效率", "流程效率、人员效率"]
        ]
    )
    
    # 根据模块选择指标
    value_rows = []
    if 'finance' in modules:
        value_rows.extend([
            ["财务结账时间", "缩短50%", "15天", "7天"],
            ["应收账款天数", "减少20%", "60天", "48天"]
        ])
    if 'supply' in modules:
        value_rows.extend([
            ["库存天数", "减少30%", "60天", "42天"],
            ["采购成本", "降低10%", "100%", "90%"]
        ])
    if 'manufacturing' in modules:
        value_rows.extend([
            ["制造周期", "缩短20%", "30天", "24天"],
            ["准时交付", "提升至95%", "85%", "95%"]
        ])
    
    if value_rows:
        add_table_slide(prs, "价值驱动指标",
            ["指标", "改进目标", "基线值", "目标值"],
            value_rows
        )
    
    # 9. 实施路线
    add_section_slide(prs, "实施路线", "07")
    
    # 转换月份为整数
    months = int(data.get('implMonths', 6) or 6)
    m1 = max(1, months // 4)
    m2 = max(m1 + 1, months // 2)
    m3 = max(m2 + 1, int(months * 0.8))
    
    add_table_slide(prs, "实施计划",
        ["阶段", "时间", "主要工作", "交付物"],
        [
            ["项目启动", "第1月", "项目组建、需求调研", "项目章程、调研报告"],
            ["蓝图设计", f"第2-{m1}月", "业务蓝图设计", "蓝图设计文档"],
            ["系统实施", f"第{m1+1}-{m2}月", "系统配置、数据迁移", "配置文档、迁移方案"],
            ["测试培训", f"第{m2+1}-{m3}月", "系统测试、用户培训", "测试报告、培训材料"],
            ["上线支持", f"第{m3+1}月-结束", "系统上线、运维支持", "上线方案、运维手册"]
        ]
    )
    
    # 10. 成功案例
    add_section_slide(prs, "成功案例", "08")
    
    industry = data.get('industry', '制造业')
    add_content_slide(prs, f"案例：某{industry}企业", [
        {"text": "实施范围", "size": 20, "bold": True},
        f"实施模块：{'、'.join(module_list)}",
        {"text": "", "size": 12},
        {"text": "核心痛点", "size": 20, "bold": True},
        "财务结账周期长、库存周转率低、供应链协同效率低",
        {"text": "", "size": 12},
        {"text": "项目成果", "size": 20, "bold": True},
        "财务结账时间缩短50%",
        "库存周转率提升40%",
        "采购周期缩短30%"
    ])
    
    # 11. 联系方式
    add_title_slide(
        prs,
        "谢谢！",
        "金蝶软件（中国）有限公司"
    )
    
    return prs

def main():
    """主函数"""
    # 从命令行参数读取客户数据
    if len(sys.argv) > 1:
        data = json.loads(sys.argv[1])
    else:
        # 默认测试数据
        data = {
            "companyName": "重庆智造科技有限公司",
            "industry": "制造业",
            "companySize": "中型企业",
            "employees": 800,
            "revenue": 50000,
            "budget": 200,
            "userCount": 150,
            "implMonths": 8,
            "controlMode": "统分结合",
            "modules": ["finance", "supply", "manufacturing"],
            "painPoints": "1. 财务结账周期长，月结需要15天\n2. 库存周转率低，资金占用大\n3. 供应链协同效率低，采购周期长\n4. 生产计划不准确，经常缺料或积压",
            "businessGoals": "1. 财务结账周期缩短50%\n2. 库存周转率提升40%\n3. 采购周期缩短30%\n4. 生产计划准确性提升50%"
        }
    
    # 生成PPT
    prs = generate_presales_ppt(data)
    
    # 保存文件
    output_dir = os.path.expanduser("~/.openclaw/workspace/output")
    os.makedirs(output_dir, exist_ok=True)
    
    company_code = data.get('customerCode', data.get('companyName', '客户')[:4])
    output_path = os.path.join(output_dir, f"{company_code}_售前解决方案.pptx")
    
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
